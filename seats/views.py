from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, HttpResponseRedirect, JsonResponse, StreamingHttpResponse
from django.views.decorators.http import require_POST, require_http_methods
from django.db import transaction, models, IntegrityError, OperationalError, ProgrammingError
from django.utils import timezone
from django.utils.dateparse import parse_datetime
from django.urls import reverse
from django.utils.encoding import escape_uri_path
from django.conf import settings
from django.test import RequestFactory
from pypinyin import lazy_pinyin
import base64
import desktop_runtime
import copy
from .models import (
    Classroom,
    Student,
    Seat,
    SeatCellType,
    SeatGroup,
    LayoutSnapshot,
    SeatConstraint,
    StudentTag,
    StudentTagMembership,
    StudentTagRule,
    FutureModeConfig,
    AIConversation,
    AIConversationMessage,
    FrontendKVStore,
    ONBOARDING_SEEN_STORE_KEY,
    ONBOARDING_SEEN_STORE_VALUE,
    ClassroomHistoryEntry,
    SyncMeta,
    OnboardingState,
)
from io import BytesIO
from http.cookies import CookieError, SimpleCookie
import json
import random
import os
import re
import shlex
import secrets
import threading
import uuid
import time
import html
import openpyxl
import ssl
import urllib.error
import urllib.parse
import urllib.request
import math
import zlib
from collections import defaultdict
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.ciphers.aead import AESGCM
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from openpyxl.styles import Alignment, Border, Side, Font, PatternFill
from openpyxl.utils import get_column_letter
from .constraints import (
    ConstraintServiceError,
    compile_constraint_maps,
    constraint_issue_messages,
    get_constraint_type_definitions,
    get_tag_rule_type_definitions,
    normalize_constraint_payload,
    normalize_tag_rule_payload,
    serialize_constraints,
    serialize_tag_rules,
    tag_rule_issue_messages,
    validate_constraint_candidate,
)
from .plugin_components import plugin_component_library
from .plugin_system import (
    plugin_registry,
    PluginActionMethodNotAllowedError,
    PluginActionNotFoundError,
    PluginNotFoundError,
    PluginUIScriptMethodNotAllowedError,
    PluginUIScriptNotFoundError,
    PluginWorkspaceScriptMethodNotAllowedError,
    PluginWorkspaceScriptNotFoundError,
)
from .cloud import (
    apply_cloud_subscription_payload,
    CloudAPIError,
    build_cloud_login_url,
    clear_cloud_session,
    cloud_api_request,
    cloud_exchange_session_code,
    cloud_public_request,
    get_active_cloud_session,
    get_cloud_server_url,
    refresh_cloud_subscription,
    save_cloud_session_from_payload,
    set_cloud_server_url,
    suspend_sync_version_bump,
)
from .data_sharing import (
    get_data_sharing_config,
    get_data_sharing_enabled,
    set_data_sharing_enabled,
    set_data_sharing_log_retention_days,
    set_data_sharing_prompt_seen_version,
    share_log,
    share_usage_event,
)

APP_MANIFEST_REDIRECT_URL = 'https://apps.577622.xyz/api/user_a6d12cebda652894/7h4sjhx0azr/api.json'
UPDATE_DETAILS_REDIRECT_URL = 'https://apps.577622.xyz/api/user_a6d12cebda652894/7h4sjhx0azr/update.txt'


def _redirect_preserving_query(request, target_url):
    query_string = request.META.get('QUERY_STRING', '')
    if query_string:
        separator = '&' if '?' in target_url else '?'
        target_url = f'{target_url}{separator}{query_string}'
    return HttpResponseRedirect(target_url)


def app_manifest_redirect(request):
    return _redirect_preserving_query(request, APP_MANIFEST_REDIRECT_URL)


def update_details_redirect(request):
    return _redirect_preserving_query(request, UPDATE_DETAILS_REDIRECT_URL)


try:
    import pandas as pd
    _PANDAS_IMPORT_ERROR = None
except Exception as exc:  # pragma: no cover - triggered only in broken runtime env
    pd = None
    _PANDAS_IMPORT_ERROR = exc


def _require_pandas():
    if pd is None:
        detail = str(_PANDAS_IMPORT_ERROR) if _PANDAS_IMPORT_ERROR else 'unknown error'
        raise RuntimeError(f'当前运行环境无法加载 pandas：{detail}')
    return pd


def _is_missing_import_value(value):
    if pd is not None:
        try:
            return bool(pd.isna(value))
        except Exception:
            pass

    if value is None:
        return True
    if isinstance(value, float):
        return math.isnan(value)
    return False

DISABLED_SUGGESTION_TYPES = {'jqj_hzh'}
DEFAULT_AI_CONVERSATION_TITLE = '新对话'
AI_CONVERSATION_FETCH_LIMIT = 50
AI_MESSAGE_FETCH_LIMIT = 120
AI_CONTEXT_MESSAGE_LIMIT = 12
FUTURE_MODE_PENDING_TTL_SECONDS = 1800
FUTURE_MODE_PENDING_CACHE = {}
MAX_LAYOUT_GRID_SIZE = 30
CLASSROOM_HISTORY_LIMIT = 1000
GROUP_MOVE_MODE_FIXED = 'fixed'
GROUP_MOVE_MODE_FOLLOW = 'follow'
EXPORT_FONT_BLACK = '鸿蒙黑体 Medium'
EXPORT_FONT_LIGHT = '鸿蒙黑体 Light'
HISTORY_STUDENT_ID_KEYS = {
    'student_pk',
    'target_student_pk',
    'leader_student_pk',
    'student_id',
    'target_student_id',
    'student_a_id',
    'student_b_id',
    'left_guardian_student_id',
    'right_guardian_student_id',
    'prev_student_id',
    'before_student_id',
    'after_student_id',
}
HISTORY_GROUP_ID_KEYS = {
    'group_pk',
    'group_id',
    'prev_group_id',
    'before_group_id',
    'after_group_id',
    'target_group_id',
}
HISTORY_GROUP_ID_LIST_KEYS = {'source_group_ids'}
HISTORY_SNAPSHOT_ID_KEYS = {'snapshot_id'}
HISTORY_CONSTRAINT_ID_KEYS = {'constraint_id'}
HISTORY_TAG_ID_KEYS = {'tag_id', 'tag_pk'}
HISTORY_TAG_ID_LIST_KEYS = {'tag_ids'}
HISTORY_TAG_RULE_ID_KEYS = {'tag_rule_id'}
FIXED_SEAT_NOTE_MARKER = '系统固定座位'

AI_TOOL_DEFINITIONS = [
    {
        'type': 'function',
        'name': 'get_classroom_overview',
        'description': '获取当前班级的整体概览，包括总人数、已入座人数、未入座人数、小组数量与建议操作。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {},
            'required': [],
            'additionalProperties': False,
        },
    },
    {
        'type': 'function',
        'name': 'get_student_info',
        'description': '按姓名、学号或系统内学生 ID 获取学生信息、座位位置、小组、成绩等。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {
                'student_query': {
                    'type': 'string',
                    'description': '学生姓名、学号或系统内数字 ID。',
                },
            },
            'required': ['student_query'],
            'additionalProperties': False,
        },
    },
    {
        'type': 'function',
        'name': 'get_group_scores',
        'description': '获取小组评分排行，返回每个小组的人数、总分、平均分、组长和成员名单。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {},
            'required': [],
            'additionalProperties': False,
        },
    },
    {
        'type': 'function',
        'name': 'get_student_list',
        'description': '读取学生列表，支持排序、筛选、字段裁剪与分页，便于生成名单或核对数据。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {
                'sort_by': {
                    'type': 'string',
                    'enum': ['id', 'name', 'student_id', 'gender', 'score', 'seat_row', 'seat_col', 'group'],
                },
                'sort_order': {
                    'type': 'string',
                    'enum': ['asc', 'desc'],
                },
                'limit': {
                    'type': 'integer',
                    'minimum': 1,
                    'maximum': 200,
                },
                'offset': {
                    'type': 'integer',
                    'minimum': 0,
                    'maximum': 5000,
                },
                'fields': {
                    'type': 'array',
                    'items': {
                        'type': 'string',
                        'enum': [
                            'id',
                            'name',
                            'student_id',
                            'gender',
                            'score',
                            'score_display',
                            'seat',
                            'group',
                            'is_seated',
                            'is_group_leader',
                            'tags',
                        ],
                    },
                    'uniqueItems': True,
                },
                'filters': {
                    'type': 'object',
                    'properties': {
                        'keyword': {'type': 'string'},
                        'seated': {'type': 'boolean'},
                        'gender': {'type': 'string', 'enum': ['M', 'F']},
                        'min_score': {'type': 'number'},
                        'max_score': {'type': 'number'},
                        'group_query': {'type': 'string'},
                        'tag_id': {'type': 'integer'},
                        'tag_ids': {'type': 'array', 'items': {'type': 'integer'}, 'uniqueItems': True},
                        'tag_name': {'type': 'string'},
                        'tag_names': {'type': 'array', 'items': {'type': 'string'}, 'uniqueItems': True},
                        'tag_match': {'type': 'string', 'enum': ['any', 'all', 'none']},
                        'untagged': {'type': 'boolean'},
                        'row': {'type': 'integer'},
                        'col': {'type': 'integer'},
                    },
                    'required': [],
                    'additionalProperties': False,
                },
            },
            'required': [],
            'additionalProperties': False,
        },
    },
    {
        'type': 'function',
        'name': 'send_card_info',
        'description': '发送结构化卡片信息，支持部分座位图、学生详情图、整体座位图、班级报告图。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {
                'card_type': {
                    'type': 'string',
                    'enum': ['partial_seat_map', 'student_detail', 'full_seat_map', 'class_report'],
                },
                'title': {'type': 'string'},
                'student_query': {'type': 'string'},
                'rows': {
                    'type': 'array',
                    'items': {'type': 'integer'},
                },
                'cols': {
                    'type': 'array',
                    'items': {'type': 'integer'},
                },
                'max_students': {
                    'type': 'integer',
                    'minimum': 1,
                    'maximum': 40,
                },
                'max_groups': {
                    'type': 'integer',
                    'minimum': 1,
                    'maximum': 30,
                },
                'include_fields': {
                    'type': 'array',
                    'items': {
                        'type': 'string',
                        'enum': ['classroom', 'metrics', 'top_students', 'group_ranking', 'suggestions'],
                    },
                    'uniqueItems': True,
                },
            },
            'required': ['card_type'],
            'additionalProperties': False,
        },
    },
    {
        'type': 'function',
        'name': 'swap_students',
        'description': '交换两名已入座学生的座位，并自动校正约束。该工具会真实修改当前班级排座结果。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {
                'student_a': {
                    'type': 'string',
                    'description': '第一名学生的姓名、学号或系统内数字 ID。',
                },
                'student_b': {
                    'type': 'string',
                    'description': '第二名学生的姓名、学号或系统内数字 ID。',
                },
            },
            'required': ['student_a', 'student_b'],
            'additionalProperties': False,
        },
    },
    {
        'type': 'function',
        'name': 'execute_classroom_action',
        'description': '执行班级操作（与人工在页面点击的功能一致），如移动座位、分组、约束、快照、撤销重做、重命名等。该工具会真实修改数据。',
        'strict': True,
        'parameters': {
            'type': 'object',
            'properties': {
                'action': {
                    'type': 'string',
                    'enum': [
                        'move_student',
                        'move_students_batch',
                        'assign_student',
                        'clear_seat',
                        'update_cell_type',
                        'create_group',
                        'rename_group',
                        'delete_group',
                        'assign_group',
                        'assign_group_batch',
                        'set_group_leader',
                        'auto_arrange',
                        'create_constraint',
                        'update_constraint',
                        'toggle_constraint',
                        'delete_constraint',
                        'save_layout_snapshot',
                        'load_layout_snapshot',
                        'delete_layout_snapshot',
                        'undo',
                        'redo',
                        'rename_classroom',
                        'delete_student',
                    ],
                    'description': '要执行的动作名称。',
                },
                'student_query': {'type': 'string'},
                'target_student_query': {'type': 'string'},
                'group_query': {'type': 'string'},
                'snapshot_query': {'type': 'string'},
                'constraint_id': {'type': 'integer'},
                'row': {'type': 'integer'},
                'col': {'type': 'integer'},
                'cell_type': {'type': 'string', 'enum': ['seat', 'aisle', 'podium', 'empty']},
                'name': {'type': 'string'},
                'new_name': {'type': 'string'},
                'mode': {'type': 'string'},
                'constraint_type': {
                    'type': 'string',
                    'enum': ['must_seat', 'forbid_seat', 'must_row', 'forbid_row', 'must_col', 'forbid_col', 'must_together', 'forbid_together'],
                },
                'enabled': {'type': 'boolean'},
                'distance': {'type': 'integer'},
                'note': {'type': 'string'},
                'moves': {
                    'type': 'array',
                    'items': {
                        'type': 'object',
                        'properties': {
                            'student_query': {'type': 'string'},
                            'row': {'type': 'integer'},
                            'col': {'type': 'integer'},
                        },
                        'required': ['student_query', 'row', 'col'],
                        'additionalProperties': False,
                    },
                },
                'seats': {
                    'type': 'array',
                    'items': {
                        'type': 'object',
                        'properties': {
                            'row': {'type': 'integer'},
                            'col': {'type': 'integer'},
                        },
                        'required': ['row', 'col'],
                        'additionalProperties': False,
                    },
                },
            },
            'required': ['action'],
            'additionalProperties': False,
        },
    },
]

AI_TOOL_LABELS = {
    'get_classroom_overview': '读取班级概览',
    'get_student_info': '读取学生信息',
    'get_group_scores': '读取小组评分',
    'get_student_list': '读取学生列表',
    'send_card_info': '发送卡片信息',
    'swap_students': '交换座位',
    'execute_classroom_action': '执行班级操作',
}

CLASSROOM_COMMAND_HELP_ORDER = [
    'help',
    'view',
    'overview',
    'students',
    'seat',
    'group',
    'snapshot',
    'arrange',
    'undo',
    'redo',
]

CLASSROOM_COMMAND_ALIASES = {
    'help': {'help', 'bangzhu', 'bz'},
    'view': {'view', 'shitu', 'goto'},
    'overview': {'overview', 'gaikuang', 'stats', 'state'},
    'students': {'student', 'students', 'xuesheng', 'xs'},
    'seat': {'seat', 'zuowei', 'zw'},
    'group': {'group', 'groups', 'xiaozu', 'xz'},
    'snapshot': {'snapshot', 'kuaizhao', 'kz'},
    'arrange': {'arrange', 'paizuo', 'pz'},
    'undo': {'undo', 'chexiao', 'cx'},
    'redo': {'redo', 'chongzuo', 'cz'},
}

CLASSROOM_VIEW_TARGET_ALIASES = {
    'classroom': {'classroom', 'banji', 'home', 'zhuye'},
    'layout': {'layout', 'buju'},
    'ai': {'ai', 'wendao', 'zhineng'},
}

CLASSROOM_STUDENT_SUBCOMMAND_ALIASES = {
    'info': {'info', 'xinxi', 'detail', 'xiangqing'},
    'top': {'top', 'rank', 'paiming', 'chengji'},
    'unseated': {'unseated', 'weiruzuo', 'weipai', 'noseat'},
    'group': {'group', 'xiaozu'},
    'seat': {'seat', 'zuowei'},
    'search': {'search', 'sousuo', 'find'},
}

CLASSROOM_SEAT_SUBCOMMAND_ALIASES = {
    'assign': {'assign', 'fenpei'},
    'move': {'move', 'yidong'},
    'clear': {'clear', 'qingkong'},
    'swap': {'swap', 'jiaohuan', 'duidiao'},
}

CLASSROOM_GROUP_SUBCOMMAND_ALIASES = {
    'score': {'score', 'defen', 'rank', 'paiming'},
    'create': {'create', 'chuangjian'},
    'rename': {'rename', 'gaiming', 'chongmingming'},
    'delete': {'delete', 'remove', 'shanchu'},
    'leader': {'leader', 'zuzhang'},
}

CLASSROOM_SNAPSHOT_SUBCOMMAND_ALIASES = {
    'list': {'list', 'liebiao'},
    'save': {'save', 'baocun'},
    'load': {'load', 'jiazai'},
    'delete': {'delete', 'remove', 'shanchu'},
}

CLASSROOM_ARRANGE_MODE_ALIASES = {
    'random': {'random', 'suiji'},
    'score_desc': {'scoredesc', 'chengji', 'gaofen'},
    'score_asc': {'scoreasc', 'difen'},
    'good_front': {'goodfront', 'qianpai'},
    'good_back': {'goodback', 'houpai'},
    'score_spread': {'scorespread', 'junheng'},
    'group_balanced': {'groupbalanced', 'xiaozujunheng'},
    'group_mentor': {'groupmentor', 'xiaozuzhidao'},
}

CLASSROOM_COMMAND_HELP = {
    'help': {
        'summary': '查看全部命令或某个命令的详细用法。',
        'examples': ['/help', '/bangzhu zuowei'],
    },
    'view': {
        'summary': '切换到指定页面视图，适合命令面板做页面跳转。',
        'examples': ['/view layout', '/shitu buju', '/view classroom'],
    },
    'overview': {
        'summary': '读取当前班级概览，包括人数、入座情况、小组和建议。',
        'examples': ['/overview', '/gaikuang'],
    },
    'students': {
        'summary': '查询学生详情、排行榜、未入座名单、按组或座位筛选。',
        'examples': ['/students 张三', '/xuesheng top 10', '/xuesheng unseated'],
    },
    'seat': {
        'summary': '执行座位相关操作，例如分配、移动、清空、交换。',
        'examples': ['/seat assign 张三 2 3', '/zuowei move 张三 1 1', '/zuowei swap 张三 李四'],
    },
    'group': {
        'summary': '查看小组排行或维护小组信息。',
        'examples': ['/group score', '/xiaozu create 第一组', '/xiaozu leader 张三'],
    },
    'snapshot': {
        'summary': '列出、保存、加载、删除布局快照。',
        'examples': ['/snapshot list', '/kuaizhao save 期中布局', '/kuaizhao load 期中布局'],
    },
    'arrange': {
        'summary': '执行自动排座，支持随机、成绩、均衡等模式。',
        'examples': ['/arrange random', '/paizuo suiji', '/paizuo junheng'],
    },
    'undo': {
        'summary': '撤销最近一次操作。',
        'examples': ['/undo', '/chexiao'],
    },
    'redo': {
        'summary': '重做最近一次撤销。',
        'examples': ['/redo', '/chongzuo'],
    },
}


ONBOARDING_SAMPLE_NAME = '示例班级（新手引导）'
ONBOARDING_SAMPLE_STUDENTS = [
    ('张明', 'M', 92), ('李华', 'M', 85), ('王芳', 'F', 78),
    ('刘伟', 'M', 88), ('陈静', 'F', 95), ('杨磊', 'M', 70),
    ('赵敏', 'F', 82), ('黄强', 'M', 76), ('周婷', 'F', 90),
    ('吴鹏', 'M', 68), ('徐丽', 'F', 84), ('孙杰', 'M', 73),
]
ONBOARDING_SAMPLE_GROUPS = ('第一组', '第二组')

_onboarding_sample_lock = threading.Lock()


def _frontend_onboarding_seen():
    try:
        value = FrontendKVStore.objects.filter(
            key=ONBOARDING_SEEN_STORE_KEY,
        ).values_list('value', flat=True).first()
        if str(value or '').strip().lower() in {'1', 'true', 'yes', 'seen'}:
            return True
    except Exception:
        pass
    try:
        if OnboardingState.objects.filter(seen=True).exists():
            FrontendKVStore.objects.update_or_create(
                key=ONBOARDING_SEEN_STORE_KEY,
                defaults={'value': ONBOARDING_SEEN_STORE_VALUE},
            )
            return True
    except Exception:
        pass
    return False


def _ensure_onboarding_sample_classroom(request):
    pk = request.session.get('onboarding_sample_pk')
    if pk:
        try:
            classroom = Classroom.objects.get(pk=pk)
            _ensure_onboarding_sample_groups(classroom)
            return classroom
        except Classroom.DoesNotExist:
            pass

    with _onboarding_sample_lock:
        pk = request.session.get('onboarding_sample_pk')
        if pk:
            try:
                classroom = Classroom.objects.get(pk=pk)
                _ensure_onboarding_sample_groups(classroom)
                return classroom
            except Classroom.DoesNotExist:
                pass
        existing = Classroom.objects.filter(name=ONBOARDING_SAMPLE_NAME).order_by('-created_at').first()
        if existing:
            _ensure_onboarding_sample_groups(existing)
            request.session['onboarding_sample_pk'] = existing.pk
            request.session.modified = True
            return existing
        classroom = Classroom.objects.create(name=ONBOARDING_SAMPLE_NAME, rows=6, cols=8)
        Student.objects.bulk_create([
            Student(classroom=classroom, name=name, gender=gender, score=score)
            for name, gender, score in ONBOARDING_SAMPLE_STUDENTS
        ])
        _ensure_onboarding_sample_groups(classroom)
        request.session['onboarding_sample_pk'] = classroom.pk
        request.session.modified = True
        return classroom


def _ensure_onboarding_sample_groups(classroom):
    for index, name in enumerate(ONBOARDING_SAMPLE_GROUPS, start=1):
        SeatGroup.objects.get_or_create(
            classroom=classroom,
            name=name,
            defaults={'order': index},
        )


def _delete_onboarding_sample(request, prefer_pk=0):
    session_pk = _safe_int(request.session.get('onboarding_sample_pk'), 0)
    prefer_pk = _safe_int(prefer_pk, 0)
    classroom = None
    for pk in [prefer_pk, session_pk]:
        if not pk:
            continue
        classroom = Classroom.objects.filter(pk=pk, name=ONBOARDING_SAMPLE_NAME).first()
        if classroom:
            break
    if not classroom:
        classroom = Classroom.objects.filter(name=ONBOARDING_SAMPLE_NAME).order_by('-created_at').first()
    if not classroom:
        if session_pk:
            request.session.pop('onboarding_sample_pk', None)
            request.session.modified = True
        return False
    with suspend_sync_version_bump(), transaction.atomic():
        classroom.left_guardian = None
        classroom.right_guardian = None
        classroom.save(update_fields=['left_guardian', 'right_guardian'])
        classroom.groups.update(leader=None)
        classroom.delete()
    request.session.pop('onboarding_sample_pk', None)
    request.session.modified = True
    return True


_CLEANUP_PENDING_KEY = 'onboarding_cleanup_pending'


def index(request):
    if request.session.get(_CLEANUP_PENDING_KEY):
        request.session.pop(_CLEANUP_PENDING_KEY, None)
        request.session.modified = True
        _delete_onboarding_sample(request)
    classrooms = Classroom.objects.all().order_by('-created_at')
    sample_pk = None
    if _onboarding_should_show_flag(request):
        try:
            sample = _ensure_onboarding_sample_classroom(request)
            sample_pk = sample.pk
            classrooms = Classroom.objects.all().order_by('-created_at')
        except Exception:
            sample_pk = request.session.get('onboarding_sample_pk')
    return render(request, 'seats/index.html', {
        'classrooms': classrooms,
        'onboarding_sample_pk': sample_pk,
    })


def _onboarding_should_show_flag(request):
    """轻量判定：当前会话是否需要展示新手引导（与 context_processor 一致）。"""
    path = getattr(request, 'path', '') or ''
    if path.startswith('/admin') or path.startswith('/static') or path.startswith('/media'):
        return False
    if getattr(request, 'method', 'GET') != 'GET':
        return False
    if _frontend_onboarding_seen():
        return False
    try:
        from seats.models import OnboardingState
        sk = request.session.session_key
        if not sk:
            request.session['ob_init'] = True
            request.session.save()
            sk = request.session.session_key
        if not sk:
            return False
        state = OnboardingState.objects.filter(session_key=sk).first()
        return state is None or not state.seen
    except Exception:
        return False


def settings_page(request):
    return render(request, 'seats/settings.html')


def create_classroom(request):
    if request.method == 'POST':
        name = request.POST.get('name')
        rows = int(request.POST.get('rows', 6))
        cols = int(request.POST.get('cols', 8))
        classroom = Classroom.objects.create(name=name, rows=rows, cols=cols)
        _emit_plugin_hook(
            'classroom_created',
            request=request,
            classroom=classroom,
            payload={'name': name, 'rows': rows, 'cols': cols},
        )
        mode = request.POST.get('mode', 'blank')
        if mode == 'excel_students':
            return redirect('import_students_options_page', pk=classroom.pk)
        if mode == 'excel_layout':
            return redirect('import_layout_excel_options_page', pk=classroom.pk)
        if mode == 'bsce':
            from django.urls import reverse
            return redirect(reverse('classroom_detail', kwargs={'pk': classroom.pk}) + '?open=bsce-import')
        return redirect('classroom_detail', pk=classroom.pk)
    return render(request, 'seats/create_classroom.html')


def _emit_plugin_hook(event_name, *, request=None, classroom=None, payload=None):
    plugin_registry.emit(
        event_name,
        hook=event_name,
        request=request,
        classroom=classroom,
        payload=payload if isinstance(payload, dict) else {},
        timestamp=timezone.now().isoformat(),
    )


def _extract_plugin_payload(request):
    if request.method == 'GET':
        return dict(request.GET.items())

    content_type = str(request.content_type or '').split(';', 1)[0].strip().lower()
    if content_type == 'application/json':
        raw = request.body or b''
        if not raw:
            return {}
        try:
            data = json.loads(raw)
        except json.JSONDecodeError as exc:
            raise ValueError('JSON 数据格式错误') from exc
        if not isinstance(data, dict):
            raise ValueError('JSON 请求体必须为对象')
        return data
    return dict(request.POST.items())


def _resolve_plugin_classroom(payload):
    if not isinstance(payload, dict):
        return None
    classroom_id = payload.get('classroom_id')
    if classroom_id in (None, ''):
        return None
    try:
        return Classroom.objects.get(pk=int(classroom_id))
    except (TypeError, ValueError) as exc:
        raise ValueError('classroom_id 必须为整数') from exc
    except Classroom.DoesNotExist as exc:
        raise ValueError('classroom_id 不存在') from exc


PLUGIN_WORKSPACE_GRANTS_SESSION_KEY = 'plugin_workspace_dom_grants_v1'


def _normalize_plugin_workspace_grants(raw):
    if not isinstance(raw, dict):
        return {}
    normalized = {}
    for classroom_key, plugin_ids in raw.items():
        key = str(classroom_key or '').strip()
        if not key:
            continue
        values = set()
        if isinstance(plugin_ids, (list, tuple, set)):
            for item in plugin_ids:
                plugin_id = str(item or '').strip()
                if plugin_id:
                    values.add(plugin_id)
        normalized[key] = values
    return normalized


def _get_workspace_grants_from_session(request):
    raw = request.session.get(PLUGIN_WORKSPACE_GRANTS_SESSION_KEY, {})
    return _normalize_plugin_workspace_grants(raw)


def _save_workspace_grants_to_session(request, grants):
    payload = {
        classroom_key: sorted(plugin_ids)
        for classroom_key, plugin_ids in grants.items()
        if plugin_ids
    }
    request.session[PLUGIN_WORKSPACE_GRANTS_SESSION_KEY] = payload
    request.session.modified = True


def _resolve_classroom_id_int(raw_value):
    try:
        return int(raw_value)
    except (TypeError, ValueError) as exc:
        raise ValueError('classroom_id 必须为整数') from exc


def _is_workspace_plugin_granted(request, classroom_id, plugin_id):
    if classroom_id in (None, ''):
        return False
    grants = _get_workspace_grants_from_session(request)
    granted = grants.get(str(classroom_id), set())
    return str(plugin_id or '').strip() in granted


def _set_workspace_plugin_grant(request, classroom_id, plugin_id, granted):
    cid = str(classroom_id)
    pid = str(plugin_id or '').strip()
    grants = _get_workspace_grants_from_session(request)
    rows = grants.setdefault(cid, set())
    if granted:
        rows.add(pid)
    else:
        rows.discard(pid)
    if not rows and cid in grants:
        grants.pop(cid, None)
    _save_workspace_grants_to_session(request, grants)
    return _is_workspace_plugin_granted(request, classroom_id, plugin_id)


@require_http_methods(['GET'])
def plugins_overview(request):
    plugin_registry.ensure_loaded()
    return JsonResponse({
        'status': 'success',
        'plugins': plugin_registry.list_plugins(),
        'load_errors': plugin_registry.load_errors,
    })


@require_http_methods(['GET'])
def plugin_components_overview(request):
    return JsonResponse({
        'status': 'success',
        'components': plugin_component_library.names(),
        'count': len(plugin_component_library.names()),
    })


def _get_plugin_row(plugin_id):
    plugin_registry.ensure_loaded()
    key = str(plugin_id or '').strip()
    for row in plugin_registry.list_plugins():
        if row.get('id') == key:
            return row
    return None


def _build_extension_manifest(plugin_row):
    plugin_id = str(plugin_row.get('id') or '').strip()
    actions = plugin_row.get('actions') or []
    ui_scripts = plugin_row.get('ui_scripts') or []
    workspace_scripts = plugin_row.get('workspace_scripts') or []

    popup_url = ''
    if ui_scripts:
        popup_name = ui_scripts[0].get('name')
        if popup_name:
            popup_url = reverse('plugin_ui_page', args=[plugin_id, popup_name])

    workspace_requires_permission = any(bool(item.get('requires_permission')) for item in workspace_scripts)
    permissions = ['plugin.runtime', 'classroom.read']
    if workspace_requires_permission:
        permissions.append('workspace.dom.write')

    manifest = {
        'manifest_version': 3,
        'name': plugin_row.get('name') or plugin_id,
        'short_name': plugin_id,
        'version': plugin_row.get('version') or '0.0.1',
        'description': plugin_row.get('description') or '',
        'author': plugin_row.get('author') or '',
        'homepage_url': plugin_row.get('website') or '',
        'action': {
            'default_title': plugin_row.get('name') or plugin_id,
            **({'default_popup': popup_url} if popup_url else {}),
        },
        'permissions': permissions,
        'host_permissions': ['/plugins/*', '/extensions/*'],
        'commands': {
            item.get('name'): {
                'description': item.get('description') or '',
                'methods': item.get('methods') or [],
            }
            for item in actions
            if item.get('name')
        },
        'plugin_actions': actions,
        'plugin_ui_scripts': [
            {
                **item,
                'data_url': reverse('plugin_ui_dispatch', args=[plugin_id, item.get('name')]),
                'page_url': reverse('plugin_ui_page', args=[plugin_id, item.get('name')]),
            }
            for item in ui_scripts
            if item.get('name')
        ],
        'plugin_workspace_scripts': workspace_scripts,
        'endpoints': {
            'manifest': reverse('extension_manifest', args=[plugin_id]),
            'send_message': reverse('extension_send_message', args=[plugin_id]),
            'permissions': reverse('extension_workspace_permission', args=[plugin_id]),
            'plugin_api_root': reverse('plugins_overview'),
            'components_library': reverse('plugin_components_overview'),
        },
        'externally_connectable': {
            'matches': ['<all_urls>'],
        },
    }
    return manifest


def _request_prefers_json(request):
    format_value = str(request.GET.get('format') or '').strip().lower()
    if format_value == 'json':
        return True

    requested_with = str(request.headers.get('X-Requested-With') or '').strip().lower()
    if requested_with == 'xmlhttprequest':
        return True

    accept = str(request.headers.get('Accept') or '').strip().lower()
    return 'application/json' in accept and 'text/html' not in accept


@require_http_methods(['GET'])
def extensions_overview(request):
    plugin_registry.ensure_loaded()

    classroom_id_value = request.GET.get('classroom_id')
    classroom_id = None
    if classroom_id_value not in (None, ''):
        try:
            classroom_id = _resolve_classroom_id_int(classroom_id_value)
        except ValueError as exc:
            if _request_prefers_json(request):
                return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    rows = []
    for plugin_row in plugin_registry.list_plugins():
        plugin_id = plugin_row.get('id')
        ui_scripts = plugin_row.get('ui_scripts') or []
        actions = plugin_row.get('actions') or []
        workspace_scripts = plugin_row.get('workspace_scripts') or []

        first_ui_page_url = ''
        if ui_scripts:
            first_ui_name = ui_scripts[0].get('name')
            if first_ui_name:
                first_ui_page_url = reverse('plugin_ui_page', args=[plugin_id, first_ui_name])

        workspace_permission_required = any(bool(item.get('requires_permission')) for item in workspace_scripts)
        workspace_permission_granted = bool(
            workspace_permission_required and classroom_id and _is_workspace_plugin_granted(request, classroom_id, plugin_id)
        )

        rows.append({
            'id': plugin_id,
            'name': plugin_row.get('name') or plugin_id,
            'version': plugin_row.get('version') or '0.0.1',
            'description': plugin_row.get('description') or '',
            'manifest_url': reverse('extension_manifest', args=[plugin_id]),
            'send_message_url': reverse('extension_send_message', args=[plugin_id]),
            'permissions_url': reverse('extension_workspace_permission', args=[plugin_id]),
            'first_ui_page_url': first_ui_page_url,
            'ui_scripts': ui_scripts,
            'actions': actions,
            'workspace_scripts': workspace_scripts,
            'workspace_permission_required': workspace_permission_required,
            'workspace_permission_granted': workspace_permission_granted,
        })

    payload = {
        'status': 'success',
        'extensions': rows,
        'count': len(rows),
        'classroom_id': classroom_id,
    }

    if _request_prefers_json(request):
        return JsonResponse(payload)

    return render(request, 'seats/extensions_overview.html', payload)


@require_http_methods(['GET'])
def extension_manifest(request, plugin_id):
    plugin_row = _get_plugin_row(plugin_id)
    if not plugin_row:
        return JsonResponse({'status': 'error', 'message': '扩展不存在'}, status=404)
    return JsonResponse(_build_extension_manifest(plugin_row))


@require_http_methods(['GET', 'POST'])
def extension_workspace_permission(request, plugin_id):
    plugin_row = _get_plugin_row(plugin_id)
    if not plugin_row:
        return JsonResponse({'status': 'error', 'message': '扩展不存在'}, status=404)

    raw_classroom_id = request.GET.get('classroom_id') if request.method == 'GET' else None

    if request.method == 'POST':
        try:
            payload = _extract_plugin_payload(request)
        except ValueError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
        raw_classroom_id = payload.get('classroom_id')
        if raw_classroom_id in (None, ''):
            return JsonResponse({'status': 'error', 'message': '缺少 classroom_id'}, status=400)

        granted_value = payload.get('granted')
        if isinstance(granted_value, str):
            granted = granted_value.strip().lower() in {'1', 'true', 'yes', 'on'}
        else:
            granted = bool(granted_value)
    else:
        granted = None

    if raw_classroom_id in (None, ''):
        return JsonResponse({'status': 'error', 'message': '缺少 classroom_id'}, status=400)

    try:
        classroom_id = _resolve_classroom_id_int(raw_classroom_id)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    if request.method == 'POST':
        current = _set_workspace_plugin_grant(request, classroom_id, plugin_id, granted)
    else:
        current = _is_workspace_plugin_granted(request, classroom_id, plugin_id)

    return JsonResponse({
        'status': 'success',
        'plugin_id': plugin_id,
        'classroom_id': classroom_id,
        'granted': bool(current),
        'workspace_scripts': plugin_row.get('workspace_scripts') or [],
    })


def _normalize_extension_result(result):
    if isinstance(result, HttpResponse):
        content_type = result.get('Content-Type', '')
        body = ''
        try:
            body = result.content.decode('utf-8', errors='replace')
        except Exception:
            body = ''
        return {
            'http_response': True,
            'status_code': result.status_code,
            'content_type': content_type,
            'body': body,
        }
    return result


@require_http_methods(['POST'])
def extension_send_message(request, plugin_id):
    plugin_row = _get_plugin_row(plugin_id)
    if not plugin_row:
        return JsonResponse({'status': 'error', 'message': '扩展不存在'}, status=404)

    try:
        payload = _extract_plugin_payload(request)
        classroom_id_value = payload.get('classroom_id') if isinstance(payload, dict) else None
        classroom = _resolve_plugin_classroom(payload)
        payload.pop('classroom_id', None)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    message = payload.get('message') if isinstance(payload, dict) else None
    if not isinstance(message, dict):
        message = payload if isinstance(payload, dict) else {}

    message_type = str(message.get('type') or message.get('target') or 'action').strip().lower()

    try:
        if message_type in {'action', 'command'}:
            action_name = str(message.get('name') or message.get('action') or '').strip()
            if not action_name:
                return JsonResponse({'status': 'error', 'message': '缺少 action 名称'}, status=400)
            call_method = str(message.get('method') or 'POST').upper()
            action_payload = message.get('payload') if isinstance(message.get('payload'), dict) else {}
            result = plugin_registry.run_action(
                plugin_id,
                action_name,
                method=call_method,
                request=request,
                classroom=classroom,
                payload=action_payload,
                action_name=action_name,
                plugin_name=plugin_id,
                runtime_message=message,
            )
            return JsonResponse({
                'status': 'success',
                'extension': plugin_id,
                'message_type': 'action',
                'name': action_name,
                'result': _normalize_extension_result(result),
            })

        if message_type == 'ui':
            ui_name = str(message.get('name') or message.get('ui_name') or '').strip()
            if not ui_name:
                return JsonResponse({'status': 'error', 'message': '缺少 ui 名称'}, status=400)
            call_method = str(message.get('method') or 'GET').upper()
            ui_payload = message.get('payload') if isinstance(message.get('payload'), dict) else {}
            ui = plugin_registry.run_ui_script(
                plugin_id,
                ui_name,
                method=call_method,
                request=request,
                classroom=classroom,
                payload=ui_payload,
                plugin_name=plugin_id,
                runtime_message=message,
            )
            return JsonResponse({
                'status': 'success',
                'extension': plugin_id,
                'message_type': 'ui',
                'name': ui_name,
                'result': ui,
            })

        if message_type == 'workspace_script':
            script_name = str(message.get('name') or message.get('script_name') or '').strip()
            if not script_name:
                return JsonResponse({'status': 'error', 'message': '缺少 workspace script 名称'}, status=400)
            call_method = str(message.get('method') or 'GET').upper()
            script = plugin_registry.run_workspace_script(plugin_id, script_name, method=call_method)

            if script.get('requires_permission'):
                resolved_classroom_id = classroom.pk if classroom else None
                if resolved_classroom_id is None and classroom_id_value not in (None, ''):
                    try:
                        resolved_classroom_id = _resolve_classroom_id_int(classroom_id_value)
                    except ValueError as exc:
                        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
                if resolved_classroom_id is None:
                    return JsonResponse({'status': 'error', 'message': 'workspace script 需要 classroom_id 以校验授权'}, status=400)
                if not _is_workspace_plugin_granted(request, resolved_classroom_id, plugin_id):
                    return JsonResponse({'status': 'error', 'message': '插件尚未获得页面修改授权'}, status=403)

            return JsonResponse({
                'status': 'success',
                'extension': plugin_id,
                'message_type': 'workspace_script',
                'name': script_name,
                'result': script,
            })

        if message_type == 'manifest':
            return JsonResponse({
                'status': 'success',
                'extension': plugin_id,
                'message_type': 'manifest',
                'result': _build_extension_manifest(plugin_row),
            })

        return JsonResponse({'status': 'error', 'message': f'不支持的消息类型：{message_type}'}, status=400)

    except PluginNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginActionNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginActionMethodNotAllowedError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=405)
    except PluginUIScriptNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginUIScriptMethodNotAllowedError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=405)
    except PluginWorkspaceScriptNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginWorkspaceScriptMethodNotAllowedError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=405)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': f'runtime.sendMessage 执行失败：{exc}'}, status=500)


@require_http_methods(['GET'])
def plugin_ui_page(request, plugin_id, ui_name):
    plugin_row = _get_plugin_row(plugin_id)
    if not plugin_row:
        return HttpResponse('插件不存在', status=404)

    ui_script_names = {item.get('name') for item in (plugin_row.get('ui_scripts') or [])}
    if ui_name not in ui_script_names:
        return HttpResponse('插件 UI 不存在', status=404)

    return render(request, 'seats/plugin_ui_page.html', {
        'plugin_id': plugin_id,
        'plugin_name': plugin_row.get('name') or plugin_id,
        'ui_name': ui_name,
    })


@require_http_methods(['GET', 'POST'])
def plugin_action_dispatch(request, plugin_id, action):
    try:
        payload = _extract_plugin_payload(request)
        classroom = _resolve_plugin_classroom(payload)
        payload.pop('classroom_id', None)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    try:
        result = plugin_registry.run_action(
            plugin_id,
            action,
            method=request.method,
            request=request,
            classroom=classroom,
            payload=payload,
            action_name=action,
            plugin_name=plugin_id,
        )
    except PluginNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginActionNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginActionMethodNotAllowedError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=405)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': f'插件执行失败：{exc}'}, status=500)

    if isinstance(result, HttpResponse):
        return result

    return JsonResponse({
        'status': 'success',
        'plugin': plugin_id,
        'action': action,
        'result': result,
    })


@require_http_methods(['GET', 'POST'])
def plugin_ui_dispatch(request, plugin_id, ui_name):
    try:
        payload = _extract_plugin_payload(request)
        classroom = _resolve_plugin_classroom(payload)
        payload.pop('classroom_id', None)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    try:
        ui = plugin_registry.run_ui_script(
            plugin_id,
            ui_name,
            method=request.method,
            request=request,
            classroom=classroom,
            payload=payload,
            plugin_name=plugin_id,
        )
    except PluginNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginUIScriptNotFoundError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=404)
    except PluginUIScriptMethodNotAllowedError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=405)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': f'插件 UI 生成失败：{exc}'}, status=500)

    return JsonResponse({
        'status': 'success',
        'plugin': plugin_id,
        'ui_name': ui_name,
        'ui': ui,
    })


def _seat_key(row, col):
    return f"{row}-{col}"


def _build_seat_map(seats):
    return {(s.row, s.col): s for s in seats}


def _coerce_score_value(score):
    value = score or 0
    try:
        value = float(value)
    except (TypeError, ValueError):
        value = 0
    return value


def _legacy_podium_guardian_ids(classroom):
    if not classroom:
        return {'left': None, 'right': None}
    return {
        'left': int(classroom.left_guardian_id) if classroom.left_guardian_id else None,
        'right': int(classroom.right_guardian_id) if classroom.right_guardian_id else None,
    }


def _resolve_auto_podium_guardian_ids_from_seats(classroom, seats):
    if not classroom or not seats:
        return {'left': None, 'right': None}

    seat_map = {}
    podium_rows = defaultdict(list)
    max_col = max(int(getattr(classroom, 'cols', 0) or 0), 0)

    for seat in seats:
        row = _safe_int(getattr(seat, 'row', 0), 0)
        col = _safe_int(getattr(seat, 'col', 0), 0)
        if row < 1 or col < 1:
            continue
        seat_map[(row, col)] = seat
        max_col = max(max_col, col)
        if getattr(seat, 'cell_type', None) == SeatCellType.PODIUM:
            podium_rows[row].append(col)

    for row in sorted(podium_rows.keys()):
        podium_cols = sorted(set(podium_rows[row]))
        if not podium_cols:
            continue

        left_student_id = None
        for col in range(podium_cols[0] - 1, 0, -1):
            seat = seat_map.get((row, col))
            if seat and getattr(seat, 'cell_type', None) == SeatCellType.SEAT:
                left_student_id = int(seat.student_id) if getattr(seat, 'student_id', None) else None
                break

        right_student_id = None
        for col in range(podium_cols[-1] + 1, max_col + 1):
            seat = seat_map.get((row, col))
            if seat and getattr(seat, 'cell_type', None) == SeatCellType.SEAT:
                right_student_id = int(seat.student_id) if getattr(seat, 'student_id', None) else None
                break

        if left_student_id or right_student_id:
            return {
                'left': left_student_id,
                'right': right_student_id,
            }

    return {'left': None, 'right': None}


def _get_podium_guardian_ids(classroom, seats=None):
    if not classroom:
        return {'left': None, 'right': None}

    resolved_seats = seats
    if resolved_seats is None:
        resolved_seats = list(classroom.seats.select_related('student').all())

    auto_guardian_ids = _resolve_auto_podium_guardian_ids_from_seats(classroom, resolved_seats)
    if auto_guardian_ids.get('left') or auto_guardian_ids.get('right'):
        return auto_guardian_ids
    return _legacy_podium_guardian_ids(classroom)


def _get_fixed_seat_student_ids(classroom, constraints=None):
    if not classroom:
        return set()
    fixed_seats, *_ = compile_constraint_maps(classroom, constraints=constraints)
    return set(int(student_id) for student_id in fixed_seats.keys())


def _get_podium_guardian_side(classroom, student, *, seats=None, guardian_student_ids=None):
    if not classroom or not student:
        return ''
    resolved_guardian_ids = guardian_student_ids or _get_podium_guardian_ids(classroom, seats=seats)
    if resolved_guardian_ids.get('left') == student.pk:
        return 'left'
    if resolved_guardian_ids.get('right') == student.pk:
        return 'right'
    return ''


def _serialize_podium_guardian_student(classroom, student, side, *, guardian_student_ids=None):
    if not student or student.classroom_id != classroom.pk:
        return None
    if guardian_student_ids and guardian_student_ids.get(side) != student.pk:
        return None
    seat = getattr(student, 'assigned_seat', None)
    return {
        'id': student.pk,
        'name': student.name,
        'student_id': student.student_id or '',
        'score_display': student.display_score if student.score is not None else '',
        'side': side,
        'seat': {
            'row': seat.row,
            'col': seat.col,
        } if seat else None,
    }


def _serialize_podium_guards(classroom, *, seats=None, guardian_student_ids=None):
    resolved_guardian_ids = guardian_student_ids or _get_podium_guardian_ids(classroom, seats=seats)
    student_map = {}
    if seats is not None:
        for seat in seats:
            student = getattr(seat, 'student', None)
            if student:
                student_map[student.pk] = student

    required_ids = [student_id for student_id in resolved_guardian_ids.values() if student_id]
    missing_ids = [student_id for student_id in required_ids if student_id not in student_map]
    if missing_ids:
        for student in classroom.students.filter(pk__in=missing_ids):
            student_map[student.pk] = student

    left_student = student_map.get(resolved_guardian_ids.get('left'))
    right_student = student_map.get(resolved_guardian_ids.get('right'))
    return {
        'left': _serialize_podium_guardian_student(
            classroom,
            left_student,
            'left',
            guardian_student_ids=resolved_guardian_ids,
        ),
        'right': _serialize_podium_guardian_student(
            classroom,
            right_student,
            'right',
            guardian_student_ids=resolved_guardian_ids,
        ),
    }


def _apply_podium_guards(classroom, left_student=None, right_student=None):
    if left_student and left_student.classroom_id != classroom.pk:
        raise ValueError('左护法学生不属于当前班级')
    if right_student and right_student.classroom_id != classroom.pk:
        raise ValueError('右护法学生不属于当前班级')
    if left_student and right_student and left_student.pk == right_student.pk:
        raise ValueError('左右护法不能设置为同一名学生')
    classroom.left_guardian = left_student
    classroom.right_guardian = right_student
    classroom.save(update_fields=['left_guardian', 'right_guardian'])


def _serialize_student_profile(
    student,
    classroom=None,
    *,
    seats=None,
    constraints=None,
    guardian_student_ids=None,
    fixed_student_ids=None,
    tag_map=None,
):
    seat = getattr(student, 'assigned_seat', None)
    group = seat.group if seat else None
    resolved_classroom = classroom or getattr(student, 'classroom', None)
    resolved_guardian_ids = guardian_student_ids or _get_podium_guardian_ids(resolved_classroom, seats=seats)
    resolved_fixed_student_ids = (
        fixed_student_ids
        if fixed_student_ids is not None
        else _get_fixed_seat_student_ids(resolved_classroom, constraints=constraints)
    )
    resolved_tag_map = tag_map
    if resolved_tag_map is None and resolved_classroom and student.pk:
        resolved_tag_map = _build_student_tag_map(resolved_classroom, [student.pk])
    return {
        'id': student.pk,
        'name': student.name,
        'student_id': student.student_id or '',
        'gender': student.get_gender_display() if student.gender else '',
        'score': _coerce_score_value(student.score),
        'score_display': student.display_score if student.score is not None else '',
        'seat': {
            'row': seat.row,
            'col': seat.col,
        } if seat else None,
        'group': {
            'id': group.pk,
            'name': group.name,
            'is_leader': bool(group and group.leader_id == student.pk),
        } if group else None,
        'podium_guardian_side': _get_podium_guardian_side(
            resolved_classroom,
            student,
            seats=seats,
            guardian_student_ids=resolved_guardian_ids,
        ),
        'is_fixed_seat': student.pk in resolved_fixed_student_ids,
        'tags': list((resolved_tag_map or {}).get(student.pk, [])),
    }


DEFAULT_STUDENT_TAG_COLOR = '#0a59f7'


def _request_payload(request):
    content_type = str(request.headers.get('content-type') or '')
    if 'application/json' in content_type:
        try:
            payload = json.loads(request.body or '{}')
        except (json.JSONDecodeError, ValueError):
            raise ValueError('无效的请求数据')
        if not isinstance(payload, dict):
            raise ValueError('请求数据必须是对象')
        return payload

    payload = {}
    for key in request.POST.keys():
        values = request.POST.getlist(key)
        payload[key] = values if len(values) > 1 else request.POST.get(key)
    return payload


def _normalize_tag_name(value):
    return re.sub(r'\s+', ' ', str(value or '').strip())


def _normalize_tag_color(value):
    color = str(value or '').strip()
    if not color:
        return DEFAULT_STUDENT_TAG_COLOR
    if not color.startswith('#'):
        color = f'#{color}'
    if re.fullmatch(r'#[0-9a-fA-F]{3}([0-9a-fA-F]{3})?', color):
        return color.lower()
    return DEFAULT_STUDENT_TAG_COLOR


def _normalize_id_list(value):
    if value in (None, ''):
        return []
    if isinstance(value, (list, tuple, set)):
        raw_items = value
    else:
        raw_items = re.split(r'[,，\s]+', str(value))
    ids = []
    seen = set()
    for raw in raw_items:
        try:
            item = int(raw)
        except (TypeError, ValueError):
            continue
        if item <= 0 or item in seen:
            continue
        ids.append(item)
        seen.add(item)
    return ids


def _normalize_name_list(value):
    if value in (None, ''):
        return []
    if isinstance(value, str):
        raw_items = re.split(r'[,，;；\n]+', value)
    elif isinstance(value, (list, tuple, set)):
        raw_items = value
    else:
        raw_items = [value]
    names = []
    seen = set()
    for raw in raw_items:
        if isinstance(raw, dict):
            raw = raw.get('name')
        name = _normalize_tag_name(raw)
        key = name.lower()
        if name and key not in seen:
            names.append(name)
            seen.add(key)
    return names


def _serialize_student_tag(tag, *, member_count=None, rule_count=None, include_urls=False):
    data = {
        'id': tag.pk,
        'name': tag.name,
        'color': tag.color or DEFAULT_STUDENT_TAG_COLOR,
        'description': tag.description or '',
        'sort_order': tag.sort_order or 0,
    }
    if member_count is not None:
        data['member_count'] = int(member_count or 0)
    if rule_count is not None:
        data['rule_count'] = int(rule_count or 0)
    if include_urls:
        data.update({
            'update_url': reverse('update_student_tag', args=[tag.classroom_id, tag.pk]),
            'delete_url': reverse('delete_student_tag', args=[tag.classroom_id, tag.pk]),
        })
    return data


def _build_student_tag_map(classroom, student_ids=None):
    queryset = StudentTagMembership.objects.filter(classroom=classroom).select_related('tag')
    if student_ids is not None:
        queryset = queryset.filter(student_id__in=list(student_ids))
    tag_map = defaultdict(list)
    for membership in queryset.order_by('tag__sort_order', 'tag__name', 'tag__pk'):
        tag_map[membership.student_id].append(_serialize_student_tag(membership.tag))
    return dict(tag_map)


def _serialize_student_tag_catalog(classroom):
    tags = list(classroom.student_tags.all().order_by('sort_order', 'name', 'pk'))
    member_counts = defaultdict(int)
    for row in (
        StudentTagMembership.objects
        .filter(classroom=classroom)
        .values('tag_id')
        .annotate(count=models.Count('id'))
    ):
        member_counts[row['tag_id']] = row['count']
    rule_counts = defaultdict(int)
    for row in (
        StudentTagRule.objects
        .filter(classroom=classroom)
        .values('tag_id')
        .annotate(count=models.Count('id'))
    ):
        rule_counts[row['tag_id']] = row['count']
    return [
        _serialize_student_tag(
            tag,
            member_count=member_counts.get(tag.pk, 0),
            rule_count=rule_counts.get(tag.pk, 0),
            include_urls=True,
        )
        for tag in tags
    ]


def _resolve_tags_from_payload(classroom, payload, *, allow_create=False):
    tag_ids = _normalize_id_list(payload.get('tag_ids') or payload.get('tags_ids') or payload.get('tag_id'))
    raw_tags = payload.get('tags')
    if isinstance(raw_tags, (list, tuple, set)):
        extra_ids = []
        for item in raw_tags:
            raw_id = item.get('id') if isinstance(item, dict) else item
            try:
                extra_ids.append(int(raw_id))
            except (TypeError, ValueError):
                continue
        tag_ids.extend([tag_id for tag_id in extra_ids if tag_id not in tag_ids])
    tags = []
    seen = set()
    if tag_ids:
        tag_map = classroom.student_tags.in_bulk(tag_ids)
        missing = [tag_id for tag_id in tag_ids if tag_id not in tag_map]
        if missing:
            raise ValueError('存在不属于当前班级的标签')
        for tag_id in tag_ids:
            tag = tag_map[tag_id]
            tags.append(tag)
            seen.add(tag.pk)

    names = _normalize_name_list(payload.get('tag_names') or payload.get('tags'))
    for name in names:
        tag = classroom.student_tags.filter(name=name).first()
        if not tag and allow_create:
            tag = StudentTag.objects.create(classroom=classroom, name=name, color=DEFAULT_STUDENT_TAG_COLOR)
        if not tag:
            raise ValueError(f'标签不存在：{name}')
        if tag.pk not in seen:
            tags.append(tag)
            seen.add(tag.pk)
    return tags


def _set_student_tags(classroom, student, tags, *, mode='set'):
    if student.classroom_id != classroom.pk:
        raise ValueError('学生不属于当前班级')
    tag_ids = {tag.pk for tag in tags}
    existing_ids = set(
        StudentTagMembership.objects
        .filter(classroom=classroom, student=student)
        .values_list('tag_id', flat=True)
    )

    normalized_mode = str(mode or 'set').strip().lower()
    if normalized_mode == 'add':
        target_ids = existing_ids | tag_ids
    elif normalized_mode == 'remove':
        target_ids = existing_ids - tag_ids
    elif normalized_mode == 'toggle':
        target_ids = set(existing_ids)
        for tag_id in tag_ids:
            if tag_id in target_ids:
                target_ids.remove(tag_id)
            else:
                target_ids.add(tag_id)
    else:
        target_ids = set(tag_ids)

    remove_ids = existing_ids - target_ids
    add_ids = target_ids - existing_ids
    if remove_ids:
        StudentTagMembership.objects.filter(classroom=classroom, student=student, tag_id__in=list(remove_ids)).delete()
    for tag in tags:
        if tag.pk in add_ids:
            StudentTagMembership.objects.get_or_create(classroom=classroom, student=student, tag=tag)
    return len(add_ids), len(remove_ids)


def _apply_student_tag_payload(classroom, students, payload, *, default_mode='set'):
    if not any(key in payload for key in ('tag_ids', 'tags_ids', 'tag_id', 'tag_names', 'tags')):
        return {'added': 0, 'removed': 0}
    tags = _resolve_tags_from_payload(classroom, payload, allow_create=True)
    mode = payload.get('tag_mode') or payload.get('mode') or default_mode
    added = 0
    removed = 0
    for student in students:
        item_added, item_removed = _set_student_tags(classroom, student, tags, mode=mode)
        added += item_added
        removed += item_removed
    return {'added': added, 'removed': removed}


def _get_group_score_rows(classroom):
    rows = []
    groups = classroom.groups.all().order_by('order', 'created_at', 'pk')
    for group in groups:
        seats = list(group.seats.select_related('student').filter(student__isnull=False).order_by('row', 'col'))
        members = [seat.student for seat in seats if seat.student]
        total_score = round(sum(_coerce_score_value(member.score) for member in members), 2)
        average_score = round(total_score / len(members), 2) if members else 0
        rows.append({
            'group_id': group.pk,
            'group_name': group.name,
            'member_count': len(members),
            'total_score': total_score,
            'average_score': average_score,
            'leader_name': group.leader.name if group.leader else '',
            'members': [
                {
                    'id': member.pk,
                    'name': member.name,
                    'student_id': member.student_id or '',
                    'score': _coerce_score_value(member.score),
                    'score_display': member.display_score if member.score is not None else '',
                    'is_leader': bool(group.leader_id == member.pk),
                }
                for member in members
            ],
        })
    rows.sort(key=lambda item: (-item['average_score'], -item['total_score'], item['group_name']))
    return rows


def _ensure_session_key(request):
    if request.session.session_key:
        return str(request.session.session_key)
    request.session.save()
    return str(request.session.session_key or '')


def _conversation_owner_queryset(classroom, request):
    owner_key = _ensure_session_key(request)
    return classroom.ai_conversations.filter(session_key=owner_key), owner_key


def _build_conversation_title_from_message(message):
    text = str(message or '').strip()
    if not text:
        return DEFAULT_AI_CONVERSATION_TITLE
    cleaned = re.sub(r'\s+', ' ', text)
    title = cleaned[:24]
    return title or DEFAULT_AI_CONVERSATION_TITLE


def _extract_cards_from_payload(payload):
    if not isinstance(payload, dict):
        return []
    cards = payload.get('cards')
    if isinstance(cards, list):
        return [item for item in cards if isinstance(item, dict)]
    card = payload.get('card')
    if isinstance(card, dict):
        return [card]
    return []


def _serialize_ai_message(message):
    payload = message.payload if isinstance(message.payload, dict) else {}
    return {
        'id': message.pk,
        'role': message.role,
        'content': str(message.content or ''),
        'cards': _extract_cards_from_payload(payload),
        'created_at': message.created_at.isoformat() if message.created_at else '',
    }


def _serialize_ai_conversation(conversation):
    latest_message = conversation.messages.order_by('-created_at', '-pk').first()
    preview = ''
    if latest_message:
        preview = str(latest_message.content or '').strip()
        if not preview and _extract_cards_from_payload(latest_message.payload):
            preview = '已发送卡片信息'
    return {
        'id': conversation.pk,
        'title': conversation.title or DEFAULT_AI_CONVERSATION_TITLE,
        'updated_at': conversation.updated_at.isoformat() if conversation.updated_at else '',
        'preview': preview[:48],
    }


def _list_ai_conversations(classroom, request):
    qs, _ = _conversation_owner_queryset(classroom, request)
    return [_serialize_ai_conversation(item) for item in qs[:AI_CONVERSATION_FETCH_LIMIT]]


def _load_ai_conversation_messages(conversation):
    messages = conversation.messages.all().order_by('created_at', 'pk')[:AI_MESSAGE_FETCH_LIMIT]
    return [_serialize_ai_message(item) for item in messages]


def _touch_ai_conversation(conversation, *, mode=None, response_id=None):
    update_fields = ['updated_at']
    conversation.updated_at = timezone.now()
    if mode is not None:
        conversation.last_mode = str(mode or '').strip()
        update_fields.append('last_mode')
    if response_id is not None:
        conversation.last_response_id = str(response_id or '').strip()
        update_fields.append('last_response_id')
    conversation.save(update_fields=update_fields)


def _create_ai_conversation(classroom, request, title=''):
    _, owner_key = _conversation_owner_queryset(classroom, request)
    resolved_title = str(title or '').strip() or DEFAULT_AI_CONVERSATION_TITLE
    return AIConversation.objects.create(
        classroom=classroom,
        session_key=owner_key,
        title=resolved_title[:120],
    )


def _resolve_ai_conversation(classroom, request, conversation_id=None, create_if_missing=True):
    qs, _ = _conversation_owner_queryset(classroom, request)
    if conversation_id not in (None, ''):
        try:
            conv_id = int(conversation_id)
        except (TypeError, ValueError):
            raise ValueError('对话 ID 格式错误')
        conversation = qs.filter(pk=conv_id).first()
        if not conversation:
            raise ValueError('未找到该对话，或无权限访问')
        return conversation, False

    conversation = qs.order_by('-updated_at', '-pk').first()
    if conversation:
        return conversation, False
    if not create_if_missing:
        raise ValueError('当前没有可用对话')
    return _create_ai_conversation(classroom, request), True


def _append_ai_conversation_message(conversation, role, content='', payload=None):
    normalized_content = str(content or '').strip()
    normalized_payload = payload if isinstance(payload, dict) else {}
    if not normalized_content and not normalized_payload:
        return None
    message = AIConversationMessage.objects.create(
        conversation=conversation,
        role=str(role or AIConversationMessage.MessageRole.USER),
        content=normalized_content[:4000],
        payload=normalized_payload,
    )
    _touch_ai_conversation(conversation)
    return message


def _build_history_from_conversation(conversation, limit=AI_CONTEXT_MESSAGE_LIMIT):
    rows = list(
        conversation.messages
        .filter(role__in=[AIConversationMessage.MessageRole.USER, AIConversationMessage.MessageRole.ASSISTANT])
        .order_by('-created_at', '-pk')[:max(1, int(limit))]
    )
    rows.reverse()
    history = []
    for item in rows:
        content = str(item.content or '').strip()
        if not content:
            continue
        role = 'assistant' if item.role == AIConversationMessage.MessageRole.ASSISTANT else 'user'
        history.append({'role': role, 'content': content[:4000]})
    return history


def _normalize_int_list(values, lower, upper):
    result = []
    if not isinstance(values, list):
        return result
    for value in values:
        try:
            current = int(value)
        except (TypeError, ValueError):
            continue
        if current < lower or current > upper:
            continue
        if current in result:
            continue
        result.append(current)
    result.sort()
    return result


def _safe_int(value, default=0):
    try:
        return int(value)
    except (TypeError, ValueError):
        return int(default)


def _safe_float(value, default=0.0):
    try:
        return float(value)
    except (TypeError, ValueError):
        return float(default)


def _serialize_card_seat_cell(seat):
    student = seat.student
    group = seat.group
    return {
        'row': seat.row,
        'col': seat.col,
        'cell_type': seat.cell_type,
        'student_name': student.name if student else '',
        'student_id': student.student_id if student and student.student_id else '',
        'group_name': group.name if group else '',
    }


def _build_partial_seat_map_card(classroom, arguments):
    rows = _normalize_int_list(arguments.get('rows'), 1, classroom.rows)
    cols = _normalize_int_list(arguments.get('cols'), 1, classroom.cols)
    if not rows:
        rows = list(range(1, classroom.rows + 1))
    if not cols:
        cols = list(range(1, classroom.cols + 1))
    seats = list(
        classroom.seats
        .select_related('student', 'group')
        .filter(row__in=rows, col__in=cols)
        .order_by('row', 'col')
    )
    return {
        'type': 'partial_seat_map',
        'title': str(arguments.get('title') or '部分座位图'),
        'classroom': {
            'id': classroom.pk,
            'name': classroom.name,
        },
        'rows': rows,
        'cols': cols,
        'cells': [_serialize_card_seat_cell(seat) for seat in seats],
    }


def _build_full_seat_map_card(classroom, arguments):
    seats = list(classroom.seats.select_related('student', 'group').all().order_by('row', 'col'))
    return {
        'type': 'full_seat_map',
        'title': str(arguments.get('title') or '整体座位图'),
        'classroom': {
            'id': classroom.pk,
            'name': classroom.name,
            'rows': classroom.rows,
            'cols': classroom.cols,
        },
        'cells': [_serialize_card_seat_cell(seat) for seat in seats],
    }


def _build_student_detail_card(classroom, arguments):
    student = _resolve_student_query(classroom, arguments.get('student_query'))
    return {
        'type': 'student_detail',
        'title': str(arguments.get('title') or f'{student.name} 详情'),
        'student': _serialize_student_profile(student, classroom=classroom),
    }


def _build_class_report_card(classroom, arguments):
    include_fields = arguments.get('include_fields')
    if not isinstance(include_fields, list):
        include_fields = ['classroom', 'metrics', 'top_students', 'group_ranking', 'suggestions']

    include_set = set(str(item or '').strip() for item in include_fields if str(item or '').strip())
    max_students = max(1, min(40, _safe_int(arguments.get('max_students'), 10)))
    max_groups = max(1, min(30, _safe_int(arguments.get('max_groups'), 10)))
    overview = _get_classroom_overview_payload(classroom)

    report = {'title': str(arguments.get('title') or '班级报告图')}
    if 'classroom' in include_set:
        report['classroom'] = overview.get('classroom') or {}
    if 'metrics' in include_set:
        report['metrics'] = overview.get('metrics') or {}
    if 'top_students' in include_set:
        report['top_students'] = (overview.get('top_students') or [])[:max_students]
    if 'group_ranking' in include_set:
        report['group_ranking'] = (overview.get('group_ranking') or [])[:max_groups]
    if 'suggestions' in include_set:
        report['suggestions'] = overview.get('suggestions') or []

    return {
        'type': 'class_report',
        'title': report.pop('title'),
        'report': report,
    }


def _build_student_list_payload(classroom, arguments):
    arguments = arguments or {}
    filters = arguments.get('filters') if isinstance(arguments.get('filters'), dict) else {}
    queryset = classroom.students.select_related('assigned_seat__group').all()

    keyword = str(filters.get('keyword') or '').strip()
    if keyword:
        queryset = queryset.filter(
            models.Q(name__icontains=keyword)
            | models.Q(student_id__icontains=keyword)
            | models.Q(tag_memberships__tag__name__icontains=keyword)
        ).distinct()

    if 'seated' in filters:
        seated_raw = filters.get('seated')
        if isinstance(seated_raw, bool):
            seated = seated_raw
        else:
            seated = str(seated_raw).strip().lower() in {'1', 'true', 'yes', 'y'}
        queryset = queryset.filter(assigned_seat__isnull=not seated)

    gender = str(filters.get('gender') or '').strip()
    if gender in {'M', 'F'}:
        queryset = queryset.filter(gender=gender)

    if filters.get('min_score') not in (None, ''):
        queryset = queryset.filter(score__gte=_safe_float(filters.get('min_score')))
    if filters.get('max_score') not in (None, ''):
        queryset = queryset.filter(score__lte=_safe_float(filters.get('max_score')))

    if filters.get('group_query') not in (None, ''):
        group = _resolve_group_query(classroom, filters.get('group_query'))
        queryset = queryset.filter(assigned_seat__group=group)

    tag_ids = _normalize_id_list(filters.get('tag_ids') or filters.get('tags') or filters.get('tag_id'))
    tag_names = _normalize_name_list(filters.get('tag_names') or filters.get('tag_name'))
    if tag_names:
        tag_ids.extend(
            tag_id
            for tag_id in classroom.student_tags.filter(name__in=tag_names).values_list('pk', flat=True)
            if tag_id not in tag_ids
        )
    tag_match = str(filters.get('tag_match') or filters.get('match') or 'any').strip().lower()
    if tag_ids:
        if tag_match == 'all':
            for tag_id in tag_ids:
                queryset = queryset.filter(tag_memberships__tag_id=tag_id)
        elif tag_match == 'none':
            queryset = queryset.exclude(tag_memberships__tag_id__in=tag_ids)
        else:
            queryset = queryset.filter(tag_memberships__tag_id__in=tag_ids)
        queryset = queryset.distinct()
    if 'untagged' in filters and _parse_bool(filters.get('untagged')):
        queryset = queryset.filter(tag_memberships__isnull=True)

    if filters.get('row') not in (None, ''):
        queryset = queryset.filter(assigned_seat__row=_safe_int(filters.get('row')))
    if filters.get('col') not in (None, ''):
        queryset = queryset.filter(assigned_seat__col=_safe_int(filters.get('col')))

    sort_by = str(arguments.get('sort_by') or 'name').strip()
    sort_order = str(arguments.get('sort_order') or 'asc').strip().lower()
    sort_map = {
        'id': 'pk',
        'name': 'name',
        'student_id': 'student_id',
        'gender': 'gender',
        'score': 'score',
        'seat_row': 'assigned_seat__row',
        'seat_col': 'assigned_seat__col',
        'group': 'assigned_seat__group__name',
    }
    sort_field = sort_map.get(sort_by, 'name')
    ordering = f'-{sort_field}' if sort_order == 'desc' else sort_field
    queryset = queryset.order_by(ordering, 'pk')

    total = queryset.count()
    limit = max(1, min(200, _safe_int(arguments.get('limit'), 30)))
    offset = max(0, _safe_int(arguments.get('offset'), 0))
    students = list(queryset[offset: offset + limit])

    requested_fields = arguments.get('fields') if isinstance(arguments.get('fields'), list) else []
    allowed_fields = {
        'id',
        'name',
        'student_id',
        'gender',
        'score',
        'score_display',
        'seat',
        'group',
        'is_seated',
        'is_group_leader',
        'tags',
    }
    field_set = [item for item in requested_fields if item in allowed_fields]
    if not field_set:
        field_set = ['id', 'name', 'student_id', 'gender', 'score_display', 'seat', 'group', 'is_seated', 'tags']

    items = []
    tag_map = _build_student_tag_map(classroom, [student.pk for student in students])

    for student in students:
        seat = getattr(student, 'assigned_seat', None)
        group = seat.group if seat else None
        row = {
            'id': student.pk,
            'name': student.name,
            'student_id': student.student_id or '',
            'gender': student.get_gender_display() if student.gender else '',
            'score': _coerce_score_value(student.score),
            'score_display': student.display_score if student.score is not None else '',
            'seat': {'row': seat.row, 'col': seat.col} if seat else None,
            'group': {'id': group.pk, 'name': group.name} if group else None,
            'is_seated': bool(seat),
            'is_group_leader': bool(group and group.leader_id == student.pk),
            'tags': tag_map.get(student.pk, []),
        }
        items.append({key: row.get(key) for key in field_set})

    return {
        'message': f'已读取学生列表，共 {total} 人，当前返回 {len(items)} 人。',
        'total': total,
        'offset': offset,
        'limit': limit,
        'sort_by': sort_by,
        'sort_order': sort_order,
        'fields': field_set,
        'items': items,
    }


def _build_card_info_payload(classroom, arguments):
    card_type = str(arguments.get('card_type') or '').strip()
    if card_type == 'partial_seat_map':
        card = _build_partial_seat_map_card(classroom, arguments)
    elif card_type == 'student_detail':
        card = _build_student_detail_card(classroom, arguments)
    elif card_type == 'full_seat_map':
        card = _build_full_seat_map_card(classroom, arguments)
    elif card_type == 'class_report':
        card = _build_class_report_card(classroom, arguments)
    else:
        raise ValueError(f'不支持的卡片类型：{card_type}')
    return {
        'message': f'已生成卡片：{card.get("title") or card_type}',
        'cards': [card],
    }


def _collect_cards_from_tool_events(tool_events):
    cards = []
    for event in tool_events or []:
        if not isinstance(event, dict):
            continue
        result = event.get('result')
        if not isinstance(result, dict) or not result.get('ok'):
            continue
        data = result.get('data')
        if not isinstance(data, dict):
            continue
        for card in _extract_cards_from_payload(data):
            cards.append(card)
    return cards


def _resolve_student_query(classroom, query):
    keyword = str(query or '').strip()
    if not keyword:
        raise ValueError('学生查询不能为空')

    exact_candidates = []
    if keyword.isdigit():
        student = classroom.students.filter(pk=int(keyword)).first()
        if student:
            return student

    direct_matches = list(
        classroom.students.filter(
            models.Q(name=keyword) | models.Q(student_id=keyword)
        ).order_by('name', 'pk')[:6]
    )
    if len(direct_matches) == 1:
        return direct_matches[0]
    if len(direct_matches) > 1:
        exact_candidates = direct_matches

    fuzzy_matches = list(
        classroom.students.filter(
            models.Q(name__icontains=keyword) | models.Q(student_id__icontains=keyword)
        ).order_by('name', 'pk')[:6]
    )
    if len(fuzzy_matches) == 1:
        return fuzzy_matches[0]
    if len(fuzzy_matches) > 1:
        exact_candidates = fuzzy_matches

    if exact_candidates:
        raise ValueError('找到多名匹配学生：' + '、'.join(student.name for student in exact_candidates))
    raise ValueError(f'未找到学生：{keyword}')


def _resolve_group_query(classroom, query):
    keyword = str(query or '').strip()
    if not keyword:
        raise ValueError('小组查询不能为空')
    if keyword.isdigit():
        group = classroom.groups.filter(pk=int(keyword)).first()
        if group:
            return group
    matches = list(classroom.groups.filter(name=keyword).order_by('order', 'pk')[:3])
    if len(matches) == 1:
        return matches[0]
    if len(matches) > 1:
        raise ValueError('找到多个同名小组，请使用小组 ID')
    fuzzy = list(classroom.groups.filter(name__icontains=keyword).order_by('order', 'pk')[:3])
    if len(fuzzy) == 1:
        return fuzzy[0]
    if len(fuzzy) > 1:
        raise ValueError('找到多个匹配小组，请使用更精确名称或小组 ID')
    raise ValueError(f'未找到小组：{keyword}')


def _resolve_snapshot_query(classroom, query):
    keyword = str(query or '').strip()
    if not keyword:
        raise ValueError('快照查询不能为空')
    if keyword.isdigit():
        snapshot = classroom.layout_snapshots.filter(pk=int(keyword)).first()
        if snapshot:
            return snapshot
    exact = list(classroom.layout_snapshots.filter(name=keyword).order_by('-created_at', '-pk')[:3])
    if len(exact) == 1:
        return exact[0]
    if len(exact) > 1:
        raise ValueError('存在同名快照，请使用快照 ID')
    fuzzy = list(classroom.layout_snapshots.filter(name__icontains=keyword).order_by('-created_at', '-pk')[:3])
    if len(fuzzy) == 1:
        return fuzzy[0]
    if len(fuzzy) > 1:
        raise ValueError('找到多个匹配快照，请使用更精确名称或快照 ID')
    raise ValueError(f'未找到快照：{keyword}')


def _decode_json_response(response):
    try:
        return json.loads(response.content.decode('utf-8'))
    except Exception:
        return {}


def _invoke_classroom_action_view(
    classroom,
    request,
    view_func,
    *,
    json_payload=None,
    form_payload=None,
    extra_args=None,
):
    if request is None:
        raise ValueError('当前上下文不支持执行写操作')

    factory = RequestFactory()
    if json_payload is not None:
        fake_request = factory.post(
            '/ai/tool/',
            data=json.dumps(json_payload, ensure_ascii=False),
            content_type='application/json',
            HTTP_X_REQUESTED_WITH='XMLHttpRequest',
        )
    else:
        fake_request = factory.post(
            '/ai/tool/',
            data=form_payload or {},
            HTTP_X_REQUESTED_WITH='XMLHttpRequest',
        )

    fake_request.session = request.session
    fake_request.user = getattr(request, 'user', None)

    args = [fake_request, classroom.pk]
    if extra_args:
        args.extend(extra_args)
    response = view_func(*args)

    if isinstance(response, JsonResponse):
        payload = _decode_json_response(response)
        if response.status_code >= 400 or payload.get('status') == 'error':
            raise ValueError(payload.get('message') or f'操作失败（HTTP {response.status_code}）')
        return payload

    if 300 <= response.status_code < 400:
        return {'status': 'success'}
    if response.status_code >= 400:
        text = response.content.decode('utf-8', errors='ignore') if hasattr(response, 'content') else ''
        raise ValueError(text or f'操作失败（HTTP {response.status_code}）')
    return {'status': 'success'}


def _execute_classroom_action_tool(classroom, arguments, request=None):
    action = str((arguments or {}).get('action') or '').strip()
    if not action:
        raise ValueError('缺少 action 参数')

    if action == 'move_student':
        student = _resolve_student_query(classroom, arguments.get('student_query'))
        row = int(arguments.get('row'))
        col = int(arguments.get('col'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            move_student,
            json_payload={'student_id': student.pk, 'row': row, 'col': col},
        )
        return {'message': f'已移动 {student.name} 到 {row}-{col}', 'payload': payload}

    if action == 'move_students_batch':
        moves = arguments.get('moves') or []
        payload_moves = []
        for item in moves:
            student = _resolve_student_query(classroom, item.get('student_query'))
            payload_moves.append({
                'student_id': student.pk,
                'row': int(item.get('row')),
                'col': int(item.get('col')),
            })
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            move_students_batch,
            json_payload={'moves': payload_moves},
        )
        return {'message': f'已批量移动 {len(payload_moves)} 名学生', 'payload': payload}

    if action == 'assign_student':
        student = _resolve_student_query(classroom, arguments.get('student_query'))
        row = int(arguments.get('row'))
        col = int(arguments.get('col'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            assign_student,
            json_payload={'student_id': student.pk, 'row': row, 'col': col},
        )
        return {'message': f'已指派 {student.name} 到 {row}-{col}', 'payload': payload}

    if action == 'clear_seat':
        row = int(arguments.get('row'))
        col = int(arguments.get('col'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            clear_seat,
            json_payload={'row': row, 'col': col},
        )
        return {'message': f'已清空座位 {row}-{col}', 'payload': payload}

    if action == 'update_cell_type':
        row = int(arguments.get('row'))
        col = int(arguments.get('col'))
        cell_type = str(arguments.get('cell_type') or '').strip()
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            update_cell_type,
            json_payload={'row': row, 'col': col, 'cell_type': cell_type},
        )
        return {'message': f'已将 {row}-{col} 类型改为 {cell_type}', 'payload': payload}

    if action == 'create_group':
        name = str(arguments.get('name') or '').strip()
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            create_group,
            form_payload={'name': name},
        )
        return {'message': f'已创建小组：{name}', 'payload': payload}

    if action == 'rename_group':
        group = _resolve_group_query(classroom, arguments.get('group_query'))
        new_name = str(arguments.get('new_name') or '').strip()
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            rename_group,
            form_payload={'name': new_name},
            extra_args=[group.pk],
        )
        return {'message': f'已将小组“{group.name}”重命名为“{new_name}”', 'payload': payload}

    if action == 'delete_group':
        group = _resolve_group_query(classroom, arguments.get('group_query'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            delete_group,
            extra_args=[group.pk],
        )
        return {'message': f'已删除小组：{group.name}', 'payload': payload}

    if action == 'assign_group':
        row = int(arguments.get('row'))
        col = int(arguments.get('col'))
        group_query = arguments.get('group_query')
        group_id = None
        if group_query:
            group_id = _resolve_group_query(classroom, group_query).pk
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            assign_group,
            json_payload={'row': row, 'col': col, 'group_id': group_id},
        )
        return {'message': f'已更新座位 {row}-{col} 的小组归属', 'payload': payload}

    if action == 'assign_group_batch':
        group_query = arguments.get('group_query')
        group_id = None
        if group_query:
            group_id = _resolve_group_query(classroom, group_query).pk
        seats = arguments.get('seats') or []
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            assign_group_batch,
            json_payload={'group_id': group_id, 'seats': seats},
        )
        return {'message': f'已批量更新 {len(seats)} 个座位的小组归属', 'payload': payload}

    if action == 'set_group_leader':
        student = _resolve_student_query(classroom, arguments.get('student_query'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            set_group_leader,
            json_payload={'student_id': student.pk},
        )
        return {'message': f'已更新组长状态：{student.name}', 'payload': payload}

    if action == 'auto_arrange':
        method = str(arguments.get('mode') or 'random').strip() or 'random'
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            auto_arrange_seats,
            form_payload={'method': method},
        )
        return {'message': f'已执行自动排座（{method}）', 'payload': payload}

    if action == 'create_constraint':
        student = _resolve_student_query(classroom, arguments.get('student_query'))
        target_student_query = arguments.get('target_student_query')
        target_student = _resolve_student_query(classroom, target_student_query) if target_student_query else None
        form_payload = {
            'constraint_type': str(arguments.get('constraint_type') or '').strip(),
            'student_id': student.pk,
            'target_student_id': target_student.pk if target_student else '',
            'row': arguments.get('row') if arguments.get('row') is not None else '',
            'col': arguments.get('col') if arguments.get('col') is not None else '',
            'distance': int(arguments.get('distance') or 1),
            'note': str(arguments.get('note') or ''),
        }
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            create_constraint,
            form_payload=form_payload,
        )
        return {'message': '已创建约束', 'payload': payload}

    if action == 'update_constraint':
        constraint_id = int(arguments.get('constraint_id'))
        constraint = classroom.constraints.filter(pk=constraint_id).first()
        if not constraint:
            raise ValueError(f'未找到约束：{constraint_id}')
        student = _resolve_student_query(classroom, arguments.get('student_query')) if arguments.get('student_query') else constraint.student
        target_student_query = arguments.get('target_student_query')
        target_student = (
            _resolve_student_query(classroom, target_student_query)
            if target_student_query
            else constraint.target_student
        )
        form_payload = {
            'constraint_type': str(arguments.get('constraint_type') or constraint.constraint_type).strip(),
            'student_id': student.pk,
            'target_student_id': target_student.pk if target_student else '',
            'row': arguments.get('row') if arguments.get('row') is not None else (constraint.row if constraint.row is not None else ''),
            'col': arguments.get('col') if arguments.get('col') is not None else (constraint.col if constraint.col is not None else ''),
            'distance': int(arguments.get('distance') or constraint.distance or 1),
            'note': str(arguments.get('note') or constraint.note or ''),
            'enabled': arguments.get('enabled') if arguments.get('enabled') is not None else ('1' if constraint.enabled else '0'),
        }
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            update_constraint,
            form_payload=form_payload,
            extra_args=[constraint_id],
        )
        return {'message': f'已更新约束 #{constraint_id}', 'payload': payload}

    if action == 'toggle_constraint':
        constraint_id = int(arguments.get('constraint_id'))
        form_payload = {}
        if arguments.get('enabled') is not None:
            form_payload['enabled'] = '1' if bool(arguments.get('enabled')) else '0'
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            toggle_constraint,
            form_payload=form_payload,
            extra_args=[constraint_id],
        )
        return {'message': f'已切换约束 #{constraint_id}', 'payload': payload}

    if action == 'delete_constraint':
        constraint_id = int(arguments.get('constraint_id'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            delete_constraint,
            extra_args=[constraint_id],
        )
        return {'message': f'已删除约束 #{constraint_id}', 'payload': payload}

    if action == 'save_layout_snapshot':
        name = str(arguments.get('name') or '').strip()
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            save_layout_snapshot,
            form_payload={'snapshot_name': name},
        )
        return {'message': f'已保存布局快照：{name}', 'payload': payload}

    if action == 'load_layout_snapshot':
        snapshot = _resolve_snapshot_query(classroom, arguments.get('snapshot_query'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            load_layout_snapshot,
            extra_args=[snapshot.pk],
        )
        return {'message': f'已加载布局快照：{snapshot.name}', 'payload': payload}

    if action == 'delete_layout_snapshot':
        snapshot = _resolve_snapshot_query(classroom, arguments.get('snapshot_query'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            delete_layout_snapshot,
            extra_args=[snapshot.pk],
        )
        return {'message': f'已删除布局快照：{snapshot.name}', 'payload': payload}

    if action == 'undo':
        payload = _invoke_classroom_action_view(classroom, request, undo_action)
        return {'message': '已执行撤销', 'payload': payload}

    if action == 'redo':
        payload = _invoke_classroom_action_view(classroom, request, redo_action)
        return {'message': '已执行重做', 'payload': payload}

    if action == 'rename_classroom':
        new_name = str(arguments.get('new_name') or arguments.get('name') or '').strip()
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            rename_classroom,
            json_payload={'name': new_name},
        )
        return {'message': f'已重命名班级为：{new_name}', 'payload': payload}

    if action == 'delete_student':
        student = _resolve_student_query(classroom, arguments.get('student_query'))
        payload = _invoke_classroom_action_view(
            classroom,
            request,
            delete_student,
            extra_args=[student.pk],
        )
        return {'message': f'已删除学生：{student.name}', 'payload': payload}

    raise ValueError(f'不支持的 action：{action}')


def _get_classroom_overview_payload(classroom):
    def _normalize_overview_suggestion(item):
        if isinstance(item, dict):
            title = str(item.get('title') or '').strip()
            message = str(item.get('message') or '').strip()
            return {'title': title, 'message': message}
        message = str(item or '').strip()
        if not message:
            return None
        return {'title': '', 'message': message}

    students = list(classroom.students.select_related('assigned_seat').all())
    seats = list(classroom.seats.select_related('student').all())
    constraints = list(classroom.constraints.select_related('student', 'target_student').all())
    guardian_student_ids = _get_podium_guardian_ids(classroom, seats=seats)
    fixed_student_ids = _get_fixed_seat_student_ids(classroom, constraints=constraints)
    tag_map = _build_student_tag_map(classroom, [student.pk for student in students])
    seated_count = sum(1 for student in students if getattr(student, 'assigned_seat', None))
    suggestions = _evaluate_layout(classroom, None)
    normalized_suggestions = []
    for item in suggestions[:5]:
        normalized = _normalize_overview_suggestion(item)
        if normalized:
            normalized_suggestions.append(normalized)
    group_rows = _get_group_score_rows(classroom)
    top_students = sorted(
        (
            _serialize_student_profile(
                student,
                classroom=classroom,
                seats=seats,
                constraints=constraints,
                guardian_student_ids=guardian_student_ids,
                fixed_student_ids=fixed_student_ids,
                tag_map=tag_map,
            )
            for student in students
        ),
        key=lambda item: (-item['score'], item['name'])
    )[:5]
    podium_guards = _serialize_podium_guards(
        classroom,
        seats=seats,
        guardian_student_ids=guardian_student_ids,
    )
    return {
        'classroom': {
            'id': classroom.pk,
            'name': classroom.name,
            'rows': classroom.rows,
            'cols': classroom.cols,
        },
        'metrics': {
            'student_count': len(students),
            'seated_count': seated_count,
            'unseated_count': max(0, len(students) - seated_count),
            'group_count': classroom.groups.count(),
            'constraint_count': classroom.constraints.count(),
            'podium_guardian_count': int(bool(guardian_student_ids.get('left'))) + int(bool(guardian_student_ids.get('right'))),
        },
        'group_ranking': group_rows[:6],
        'top_students': top_students,
        'podium_guards': podium_guards,
        'suggestions': normalized_suggestions,
    }


def _build_swap_action(student_a, student_b):
    return {
        'type': 'swap',
        'student_a_id': student_a.pk,
        'student_b_id': student_b.pk,
    }


def _get_future_mode_pending_store(request):
    store = request.session.get('future_mode_pending_tools', {})
    request.session['future_mode_pending_tools'] = store
    return store


def _cleanup_future_mode_pending_cache():
    now_ts = time.time()
    expired_tokens = []
    for token, payload in FUTURE_MODE_PENDING_CACHE.items():
        created_at_ts = float(payload.get('_created_at_ts') or 0)
        if created_at_ts <= 0 or (now_ts - created_at_ts) > FUTURE_MODE_PENDING_TTL_SECONDS:
            expired_tokens.append(token)
    for token in expired_tokens:
        FUTURE_MODE_PENDING_CACHE.pop(token, None)


def _store_future_mode_pending(
    request,
    classroom_id,
    conversation_id,
    response_id,
    function_calls,
    mode='responses',
    chat_messages=None,
):
    _cleanup_future_mode_pending_cache()
    owner_session_key = _ensure_session_key(request)
    token = uuid.uuid4().hex
    payload = {
        'classroom_id': classroom_id,
        'conversation_id': conversation_id,
        'response_id': response_id,
        'function_calls': function_calls,
        'mode': mode or 'responses',
        'chat_messages': chat_messages or [],
        'session_key': owner_session_key,
        '_created_at_ts': time.time(),
    }
    store = _get_future_mode_pending_store(request)
    store[token] = payload
    FUTURE_MODE_PENDING_CACHE[token] = payload
    request.session.modified = True
    return token


def _consume_future_mode_pending(request, token, classroom_id, conversation_id=None):
    _cleanup_future_mode_pending_cache()
    store = _get_future_mode_pending_store(request)
    payload = store.get(token) or FUTURE_MODE_PENDING_CACHE.get(token)
    if not payload:
        raise ValueError('授权请求已过期，请重新发送本次指令。')

    owner_session_key = str(payload.get('session_key') or '').strip()
    if owner_session_key:
        current_session_key = _ensure_session_key(request)
        if str(current_session_key or '').strip() != owner_session_key:
            raise ValueError('授权请求与当前会话不匹配。')

    if int(payload.get('classroom_id') or 0) != int(classroom_id):
        raise ValueError('授权请求与当前班级不匹配。')
    if conversation_id not in (None, ''):
        if int(payload.get('conversation_id') or 0) != int(conversation_id):
            raise ValueError('授权请求与当前对话不匹配。')

    if token in store:
        del store[token]
        request.session.modified = True
    FUTURE_MODE_PENDING_CACHE.pop(token, None)
    return payload


def _normalize_ai_client_config(payload):
    payload = payload or {}
    if not isinstance(payload, dict):
        payload = {}
    thinking_mode = str(payload.get('thinking_mode') or '').strip().lower()
    if thinking_mode not in {'enabled', 'disabled'}:
        thinking_mode = ''
    return {
        'api_key': str(payload.get('api_key') or '').strip(),
        'base_url': str(payload.get('base_url') or '').strip(),
        'model': str(payload.get('model') or '').strip(),
        'thinking_mode': thinking_mode,
    }


def _load_persisted_ai_client_config(classroom):
    try:
        config = FutureModeConfig.objects.filter(classroom=classroom).first()
    except (OperationalError, ProgrammingError):
        return {'api_key': '', 'base_url': '', 'model': '', 'thinking_mode': ''}
    if not config:
        return {'api_key': '', 'base_url': '', 'model': '', 'thinking_mode': ''}
    return {
        'api_key': str(config.api_key or '').strip(),
        'base_url': str(config.base_url or '').strip(),
        'model': str(config.model or '').strip(),
        'thinking_mode': str(config.thinking_mode or '').strip().lower() if str(config.thinking_mode or '').strip().lower() in {'enabled', 'disabled'} else '',
    }


def _merge_ai_client_config(classroom, incoming=None):
    incoming_config = _normalize_ai_client_config(incoming)
    persisted = _load_persisted_ai_client_config(classroom)
    merged = {}
    for key in ('api_key', 'base_url', 'model', 'thinking_mode'):
        merged[key] = incoming_config[key] if incoming_config[key] else persisted[key]
    return merged


def _save_ai_client_config(classroom, payload):
    normalized = _normalize_ai_client_config(payload)
    try:
        if not any(normalized.values()):
            FutureModeConfig.objects.filter(classroom=classroom).delete()
            return normalized

        config, _ = FutureModeConfig.objects.get_or_create(classroom=classroom)
        config.api_key = normalized['api_key']
        config.base_url = normalized['base_url']
        config.model = normalized['model']
        config.thinking_mode = normalized['thinking_mode']
        config.save(update_fields=['api_key', 'base_url', 'model', 'thinking_mode', 'updated_at'])
    except (OperationalError, ProgrammingError) as exc:
        raise RuntimeError('数据库未完成迁移，请先执行 `python3 manage.py migrate`。') from exc
    return normalized


def _get_openai_client(client_config=None):
    client_config = _normalize_ai_client_config(client_config)
    try:
        from openai import OpenAI
    except ImportError as exc:
        raise RuntimeError('缺少 `openai` 依赖，请先执行 `pip install -r requirements.txt`。') from exc

    api_key = client_config['api_key'] or os.getenv('OPENAI_API_KEY') or getattr(settings, 'OPENAI_API_KEY', '')
    if not api_key:
        raise RuntimeError('未配置 `OPENAI_API_KEY`，请先在页面右侧填写 API Key 或在服务端配置环境变量。')

    client_kwargs = {'api_key': api_key}
    base_url = client_config['base_url'] or os.getenv('OPENAI_BASE_URL') or getattr(settings, 'OPENAI_BASE_URL', '')
    if base_url:
        client_kwargs['base_url'] = base_url
    return OpenAI(**client_kwargs)


def _get_openai_model(client_config=None):
    client_config = _normalize_ai_client_config(client_config)
    return client_config['model'] or os.getenv('OPENAI_MODEL') or getattr(settings, 'OPENAI_MODEL', 'gpt-4.1-mini')


def _resolve_openai_base_url(client_config=None):
    client_config = _normalize_ai_client_config(client_config)
    return client_config['base_url'] or os.getenv('OPENAI_BASE_URL') or getattr(settings, 'OPENAI_BASE_URL', '')


def _should_use_chat_completions(client_config=None):
    base_url = str(_resolve_openai_base_url(client_config) or '').strip().lower()
    if not base_url:
        return False
    return 'api.openai.com' not in base_url


def _build_chat_completion_extra_body(client_config=None):
    client_config = _normalize_ai_client_config(client_config)
    thinking_mode = client_config.get('thinking_mode')
    if thinking_mode in {'enabled', 'disabled'}:
        return {'thinking': {'type': thinking_mode}}
    return None


def _is_responses_not_supported_error(exc):
    message = str(exc or '').strip().lower()
    if _is_responses_content_unmarshal_error(exc):
        return True
    return (
        'not implemented' in message
        or '/responses' in message
        or 'responses api' in message
    )


def _normalize_future_mode_openai_exception(exc):
    message = str(exc or '').strip()
    lower_message = message.lower()
    if isinstance(exc, NotImplementedError):
        return 400, '当前 Base URL 或模型不支持 Responses API（Window Inteligence ｜ 闻道智能），请改用支持该接口的服务。'
    if 'not implemented' in lower_message and ('status_code=500' in lower_message or 'status code: 500' in lower_message or 'responses' in lower_message):
        return 400, '当前 Base URL 或模型不支持 Responses API（Window Inteligence ｜ 闻道智能），请改用支持该接口的服务。'
    if _is_responses_content_unmarshal_error(exc):
        return 400, '当前 Base URL 的 Responses API 与 Window Inteligence ｜ 闻道智能 消息格式不兼容（content 字段），请更换兼容接口或使用官方接口。'

    try:
        import openai
    except Exception:
        return None

    if isinstance(exc, (openai.AuthenticationError, getattr(openai, 'PermissionDeniedError', tuple()))):
        return 400, 'OpenAI 鉴权失败，请检查 API Key、Base URL 与 Model ID 配置。'

    if isinstance(exc, openai.NotFoundError):
        if 'responses' in lower_message:
            return 400, '当前 Base URL 或模型不支持 Responses API（Window Inteligence ｜ 闻道智能），请改用支持该接口的服务。'
        return 400, f'AI 接口不存在或模型不可用：{message or "Not Found"}'

    if isinstance(exc, openai.BadRequestError):
        return 400, f'AI 请求参数错误：{message or "Bad Request"}'

    if isinstance(exc, openai.RateLimitError):
        return 429, 'AI 请求过于频繁或额度不足，请稍后再试。'

    if isinstance(exc, (openai.APIConnectionError, openai.APITimeoutError)):
        return 502, 'AI 服务连接失败或超时，请检查网络后重试。'

    if isinstance(exc, openai.APIStatusError):
        status_code = getattr(exc, 'status_code', None)
        if isinstance(status_code, int) and 400 <= status_code < 500:
            return 400, f'AI 请求失败（HTTP {status_code}）：{message or "请求未成功"}'
        if isinstance(status_code, int):
            return 502, f'AI 服务异常（HTTP {status_code}），请稍后重试。'
        return 502, 'AI 服务异常，请稍后重试。'

    if isinstance(exc, openai.OpenAIError):
        return 400, f'AI 请求失败：{message or "OpenAI 调用异常"}'

    return None


def _is_tool_output_call_mismatch_error(exc):
    message = str(exc or '').strip().lower()
    return 'no tool call found for function call output' in message


def _is_responses_content_unmarshal_error(exc):
    message = str(exc or '').strip().lower()
    return (
        'cannot unmarshal object into go struct field' in message
        and '.content' in message
        and 'responsesoutputcontent' in message
    )


def _build_tool_events_fallback_reply(tool_events):
    if not tool_events:
        return '工具调用已结束，但 AI 会话同步失败。请重试一次。'

    lines = ['工具已执行，但 AI 会话续答失败。以下是本次执行结果：']
    for event in tool_events:
        name = str(event.get('name') or '').strip()
        label = AI_TOOL_LABELS.get(name, name or '工具')
        approved = bool(event.get('approved'))
        result = event.get('result') if isinstance(event.get('result'), dict) else {}
        if not approved:
            lines.append(f'- {label}：你已拒绝执行。')
            continue

        if result.get('ok'):
            data = result.get('data') if isinstance(result.get('data'), dict) else {}
            message = str(data.get('message') or '').strip()
            if message:
                lines.append(f'- {label}：{message}')
            elif name == 'get_classroom_overview':
                lines.append(f'- {label}：已读取当前班级概览。')
            elif name == 'get_student_info':
                student = data.get('name') if isinstance(data, dict) else ''
                lines.append(f'- {label}：已读取学生信息{f"（{student}）" if student else ""}。')
            elif name == 'get_group_scores':
                lines.append(f'- {label}：已读取小组评分数据。')
            else:
                lines.append(f'- {label}：执行完成。')
            continue

        error_message = str(result.get('message') or '执行失败').strip()
        lines.append(f'- {label}：{error_message}')

    lines.append('可继续发送下一条指令；如果你在使用第三方 Base URL，请确认其完整支持 Responses API 的工具调用链路。')
    return '\n'.join(lines)


def _build_tool_events_success_reply(tool_events):
    if not tool_events:
        return '操作已完成。'
    lines = []
    for event in tool_events:
        name = str(event.get('name') or '').strip()
        label = AI_TOOL_LABELS.get(name, name or '工具')
        approved = bool(event.get('approved'))
        result = event.get('result') if isinstance(event.get('result'), dict) else {}
        if not approved:
            lines.append(f'- {label}：已拒绝执行')
            continue
        if result.get('ok'):
            data = result.get('data') if isinstance(result.get('data'), dict) else {}
            message = str(data.get('message') or '').strip()
            lines.append(f'- {label}：{message or "执行成功"}')
        else:
            lines.append(f'- {label}：{str(result.get("message") or "执行失败").strip()}')
    return '\n'.join(lines) if lines else '操作已完成。'


def _build_ai_tool_error_result(tool_name, exc):
    error_type = exc.__class__.__name__ if exc is not None else 'ToolExecutionError'
    error_message = str(exc or '').strip() or '未知错误'
    return {
        'ok': False,
        'tool': tool_name,
        'message': f'工具执行失败：{error_message}',
        'error': {
            'type': error_type,
            'message': error_message,
        },
    }


def _command_alias_candidates(value):
    text = str(value or '').strip().lower().lstrip('/').strip()
    if not text:
        return set()
    candidates = {text, re.sub(r'[\s_-]+', '', text)}
    for item in list(candidates):
        if item and re.search(r'[\u4e00-\u9fff]', item):
            candidates.add(''.join(lazy_pinyin(item)).lower())
    return {item for item in candidates if item}


def _resolve_command_alias(token, alias_map):
    token_candidates = _command_alias_candidates(token)
    if not token_candidates:
        return None
    for canonical, aliases in alias_map.items():
        alias_candidates = _command_alias_candidates(canonical)
        for alias in aliases or set():
            alias_candidates.update(_command_alias_candidates(alias))
        if token_candidates & alias_candidates:
            return canonical
    return None


def _list_command_aliases(canonical, alias_map):
    aliases = [canonical]
    for alias in sorted(alias_map.get(canonical, set())):
        if alias not in aliases:
            aliases.append(alias)
    return aliases


def _tokenize_classroom_command(command_text):
    raw = str(command_text or '').strip()
    if raw.startswith('／'):
        raw = '/' + raw[1:]
    if not raw.startswith('/'):
        return [], '命令必须以 / 开头，例如 /help'
    body = raw[1:].strip()
    if not body:
        return [], '请输入命令，例如 /help'
    try:
        return shlex.split(body), ''
    except ValueError:
        return [], '命令格式错误，请检查引号是否成对出现'


def _build_classroom_command_manifest(classroom):
    commands = []
    for name in CLASSROOM_COMMAND_HELP_ORDER:
        doc = CLASSROOM_COMMAND_HELP.get(name, {})
        commands.append({
            'name': name,
            'aliases': _list_command_aliases(name, CLASSROOM_COMMAND_ALIASES),
            'summary': doc.get('summary') or '',
            'examples': doc.get('examples') or [],
        })
    return {
        'prefix': '/',
        'endpoint': reverse('classroom_command', args=[classroom.pk]),
        'placeholder': '输入 /help 查看命令',
        'commands': commands,
    }


def _format_command_seat_label(seat):
    if not isinstance(seat, dict) or seat.get('row') in (None, '') or seat.get('col') in (None, ''):
        return '未入座'
    return f"{seat.get('row')}-{seat.get('col')}"


def _format_classroom_overview_command_reply(data):
    classroom_info = data.get('classroom') if isinstance(data, dict) else {}
    metrics = data.get('metrics') if isinstance(data, dict) else {}
    suggestions = data.get('suggestions') if isinstance(data, dict) else []
    lines = [
        f"班级：{classroom_info.get('name') or '未命名班级'}",
        (
            f"1. 学生 {metrics.get('student_count', 0)} 人，"
            f"已入座 {metrics.get('seated_count', 0)} 人，"
            f"未入座 {metrics.get('unseated_count', 0)} 人。"
        ),
        (
            f"2. 小组 {metrics.get('group_count', 0)} 个，"
            f"约束 {metrics.get('constraint_count', 0)} 条。"
        ),
    ]
    if suggestions:
        first_suggestion = suggestions[0]
        message = str((first_suggestion or {}).get('message') or '').strip()
        if message:
            lines.append(f'3. 当前建议：{message}')
    return '\n'.join(lines)


def _format_student_profile_command_reply(data):
    seat_label = _format_command_seat_label(data.get('seat'))
    group = data.get('group') if isinstance(data, dict) else None
    group_name = (group or {}).get('name') or '未分组'
    if group and group.get('is_leader'):
        group_name += '（组长）'
    lines = [
        f"学生：{data.get('name') or '未知'}",
        f"1. 学号：{data.get('student_id') or '未填写'}",
        f"2. 成绩：{data.get('score_display') if data.get('score_display') not in (None, '') else data.get('score', 0)}",
        f"3. 座位：{seat_label}",
        f"4. 小组：{group_name}",
    ]
    if data.get('gender'):
        lines.insert(2, f"2. 性别：{data.get('gender')}")
        lines[3] = f"3. 成绩：{data.get('score_display') if data.get('score_display') not in (None, '') else data.get('score', 0)}"
        lines[4] = f"4. 座位：{seat_label}"
        lines[5] = f"5. 小组：{group_name}"
    return '\n'.join(lines)


def _format_student_list_item(item):
    name = str(item.get('name') or '未知')
    student_id = str(item.get('student_id') or '').strip()
    score = item.get('score_display')
    seat_label = _format_command_seat_label(item.get('seat'))
    group = item.get('group')
    group_name = group.get('name') if isinstance(group, dict) else ''
    parts = [name]
    if student_id:
        parts.append(f'学号 {student_id}')
    if score not in (None, ''):
        parts.append(f'成绩 {score}')
    parts.append(f'座位 {seat_label}')
    if group_name:
        parts.append(f'小组 {group_name}')
    return '，'.join(parts)


def _format_student_list_command_reply(data, title):
    items = data.get('items') if isinstance(data, dict) else []
    total = int(data.get('total') or 0) if isinstance(data, dict) else 0
    if not items:
        return f'{title}\n当前没有匹配结果。'
    lines = [f'{title}（共 {total} 条，本次返回 {len(items)} 条）']
    for index, item in enumerate(items[:8], start=1):
        lines.append(f'{index}. {_format_student_list_item(item)}')
    if len(items) > 8:
        lines.append(f'共返回 {len(items)} 条，已省略后续内容。')
    return '\n'.join(lines)


def _format_group_scores_command_reply(data):
    items = data.get('items') if isinstance(data, dict) else []
    if not items:
        return '当前还没有可统计的小组成绩。'
    lines = [f'小组排行（共 {len(items)} 组）']
    for index, item in enumerate(items[:6], start=1):
        lines.append(
            f"{index}. {item.get('group_name') or '未命名小组'}，"
            f"均分 {item.get('average_score', 0)}，"
            f"总分 {item.get('total_score', 0)}，"
            f"人数 {item.get('member_count', 0)}"
        )
    return '\n'.join(lines)


def _format_snapshot_list_command_reply(items):
    if not items:
        return '当前还没有布局快照。'
    lines = [f'布局快照（共 {len(items)} 个）']
    for index, item in enumerate(items[:10], start=1):
        lines.append(f"{index}. {item.get('name') or '未命名快照'}")
    return '\n'.join(lines)


def _build_classroom_command_result(
    classroom,
    command_text,
    command,
    *,
    subcommand='',
    kind='query',
    reply='',
    ok=True,
    data=None,
    tool_name='',
    tool_arguments=None,
    navigation=None,
    needs_refresh=False,
    extra=None,
):
    result = {
        'ok': bool(ok),
        'raw': str(command_text or ''),
        'command': str(command or ''),
        'subcommand': str(subcommand or ''),
        'kind': str(kind or 'query'),
        'reply': str(reply or '').strip(),
        'classroom': {
            'id': classroom.pk,
            'name': classroom.name,
        },
        'tool_name': str(tool_name or ''),
        'tool_arguments': tool_arguments if isinstance(tool_arguments, dict) else {},
        'needs_refresh': bool(needs_refresh),
        'state_url': reverse('classroom_state', args=[classroom.pk]),
        'data': data if data is not None else {},
    }
    if isinstance(navigation, dict) and navigation:
        result['navigation'] = navigation
    if isinstance(extra, dict) and extra:
        result.update(extra)
    return result


def _build_classroom_command_error(classroom, command_text, reply, *, command='', subcommand='', suggestions=None):
    extra = {}
    if suggestions:
        extra['suggestions'] = list(suggestions)
    return _build_classroom_command_result(
        classroom,
        command_text,
        command,
        subcommand=subcommand,
        kind='error',
        reply=reply,
        ok=False,
        needs_refresh=False,
        extra=extra,
    )


def _run_classroom_command_tool(
    classroom,
    command_text,
    command,
    *,
    subcommand='',
    tool_name,
    arguments=None,
    request=None,
    kind='query',
    needs_refresh=False,
    reply_formatter=None,
):
    normalized_arguments = arguments if isinstance(arguments, dict) else {}
    try:
        tool_result = _execute_ai_tool(classroom, tool_name, normalized_arguments, request=request)
    except Exception as exc:
        return _build_classroom_command_error(
            classroom,
            command_text,
            str(exc),
            command=command,
            subcommand=subcommand,
        )
    data = tool_result.get('data') if isinstance(tool_result, dict) else {}
    reply = reply_formatter(data) if callable(reply_formatter) else str((data or {}).get('message') or '命令已执行。')
    return _build_classroom_command_result(
        classroom,
        command_text,
        command,
        subcommand=subcommand,
        kind=kind,
        reply=reply,
        ok=True,
        data=data,
        tool_name=tool_name,
        tool_arguments=normalized_arguments,
        needs_refresh=needs_refresh,
    )


def _handle_classroom_help_command(classroom, command_text, args):
    manifest = _build_classroom_command_manifest(classroom)
    target_name = _resolve_command_alias(args[0], CLASSROOM_COMMAND_ALIASES) if args else None
    command_detail = None
    reply_prefix = ''
    if args and not target_name:
        reply_prefix = f"未找到命令 “{args[0]}”，下面是全部可用命令。\n"
    if target_name:
        command_detail = next((item for item in manifest['commands'] if item['name'] == target_name), None)
    if command_detail:
        reply_lines = [
            f"{command_detail['name']}：{command_detail.get('summary') or ''}",
            '可用写法：' + '、'.join(f"/{item}" for item in command_detail.get('aliases') or []),
            '示例：',
        ]
        for index, example in enumerate(command_detail.get('examples') or [], start=1):
            reply_lines.append(f'{index}. {example}')
        reply = '\n'.join(reply_lines)
    else:
        reply_lines = ['可用命令：']
        for index, item in enumerate(manifest['commands'], start=1):
            reply_lines.append(f"{index}. /{item['name']}：{item.get('summary') or ''}")
        reply_lines.append('输入 /help 命令名 可以查看某个命令的详细示例。')
        reply = reply_prefix + '\n'.join(reply_lines)
    return _build_classroom_command_result(
        classroom,
        command_text,
        'help',
        kind='help',
        reply=reply,
        data={
            'manifest': manifest,
            'command': command_detail or {},
        },
    )


def _handle_classroom_view_command(classroom, command_text, args):
    if not args:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请指定目标视图，例如 /view layout 或 /shitu buju',
            command='view',
        )
    target = _resolve_command_alias(args[0], CLASSROOM_VIEW_TARGET_ALIASES)
    if not target:
        return _build_classroom_command_error(
            classroom,
            command_text,
            f'不支持的视图：{args[0]}',
            command='view',
            suggestions=['/view classroom', '/view layout', '/view ai'],
        )
    view_targets = {
        'classroom': {
            'target': 'classroom',
            'label': '班级主页',
            'url': reverse('classroom_detail', args=[classroom.pk]),
        },
        'layout': {
            'target': 'layout',
            'label': '布局视图',
            'url': reverse('layout_editor', args=[classroom.pk]),
        },
        'ai': {
            'target': 'ai',
            'label': '闻道智能',
            'url': reverse('ai_workspace', args=[classroom.pk]),
        },
    }
    navigation = view_targets[target]
    return _build_classroom_command_result(
        classroom,
        command_text,
        'view',
        subcommand=target,
        kind='navigate',
        reply=f"已解析到 {navigation['label']}，前端可跳转到 {navigation['url']}",
        navigation=navigation,
        data=navigation,
    )


def _handle_classroom_overview_command(classroom, command_text, request=None):
    return _run_classroom_command_tool(
        classroom,
        command_text,
        'overview',
        tool_name='get_classroom_overview',
        arguments={},
        request=request,
        kind='query',
        reply_formatter=_format_classroom_overview_command_reply,
    )


def _handle_classroom_students_command(classroom, command_text, args, request=None):
    if not args:
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            tool_name='get_student_list',
            arguments={
                'sort_by': 'name',
                'sort_order': 'asc',
                'limit': 20,
                'fields': ['id', 'name', 'student_id', 'score_display', 'seat', 'group', 'is_seated'],
            },
            request=request,
            kind='query',
            reply_formatter=lambda data: _format_student_list_command_reply(data, '学生列表'),
        )

    subcommand = _resolve_command_alias(args[0], CLASSROOM_STUDENT_SUBCOMMAND_ALIASES)
    if subcommand is None:
        student_query = ' '.join(args).strip()
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            subcommand='info',
            tool_name='get_student_info',
            arguments={'student_query': student_query},
            request=request,
            kind='query',
            reply_formatter=_format_student_profile_command_reply,
        )

    if subcommand == 'info':
        student_query = ' '.join(args[1:]).strip()
        if not student_query:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供学生姓名、学号或系统 ID，例如 /students 张三',
                command='students',
                subcommand='info',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            subcommand='info',
            tool_name='get_student_info',
            arguments={'student_query': student_query},
            request=request,
            kind='query',
            reply_formatter=_format_student_profile_command_reply,
        )

    if subcommand == 'top':
        limit = max(1, min(50, _safe_int(args[1], 10))) if len(args) > 1 else 10
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            subcommand='top',
            tool_name='get_student_list',
            arguments={
                'sort_by': 'score',
                'sort_order': 'desc',
                'limit': limit,
                'fields': ['id', 'name', 'student_id', 'score_display', 'seat', 'group', 'is_seated'],
            },
            request=request,
            kind='query',
            reply_formatter=lambda data: _format_student_list_command_reply(data, f'成绩 Top {limit}'),
        )

    if subcommand == 'unseated':
        limit = max(1, min(100, _safe_int(args[1], 50))) if len(args) > 1 else 50
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            subcommand='unseated',
            tool_name='get_student_list',
            arguments={
                'limit': limit,
                'fields': ['id', 'name', 'student_id', 'score_display', 'seat', 'is_seated'],
                'filters': {'seated': False},
            },
            request=request,
            kind='query',
            reply_formatter=lambda data: _format_student_list_command_reply(data, '未入座学生'),
        )

    if subcommand == 'group':
        group_query = ' '.join(args[1:]).strip()
        if not group_query:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供小组名称，例如 /students group 第一组',
                command='students',
                subcommand='group',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            subcommand='group',
            tool_name='get_student_list',
            arguments={
                'limit': 100,
                'fields': ['id', 'name', 'student_id', 'score_display', 'seat', 'group', 'is_seated'],
                'filters': {'group_query': group_query},
            },
            request=request,
            kind='query',
            reply_formatter=lambda data: _format_student_list_command_reply(data, f'小组 {group_query} 的学生'),
        )

    if subcommand == 'seat':
        if len(args) < 3:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供座位行列，例如 /students seat 2 3',
                command='students',
                subcommand='seat',
            )
        row = _safe_int(args[1], 0)
        col = _safe_int(args[2], 0)
        if row <= 0 or col <= 0:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '座位行列必须是正整数，例如 /students seat 2 3',
                command='students',
                subcommand='seat',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'students',
            subcommand='seat',
            tool_name='get_student_list',
            arguments={
                'limit': 10,
                'fields': ['id', 'name', 'student_id', 'score_display', 'seat', 'group', 'is_seated'],
                'filters': {'row': row, 'col': col},
            },
            request=request,
            kind='query',
            reply_formatter=lambda data: _format_student_list_command_reply(data, f'座位 {row}-{col} 的学生'),
        )

    keyword = ' '.join(args[1:]).strip()
    if not keyword:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请提供搜索关键词，例如 /students search 张',
            command='students',
            subcommand='search',
        )
    return _run_classroom_command_tool(
        classroom,
        command_text,
        'students',
        subcommand='search',
        tool_name='get_student_list',
        arguments={
            'limit': 20,
            'fields': ['id', 'name', 'student_id', 'score_display', 'seat', 'group', 'is_seated'],
            'filters': {'keyword': keyword},
        },
        request=request,
        kind='query',
        reply_formatter=lambda data: _format_student_list_command_reply(data, f'搜索 “{keyword}” 的结果'),
    )


def _handle_classroom_seat_command(classroom, command_text, args, request=None):
    if not args:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请提供座位命令，例如 /seat assign 张三 2 3',
            command='seat',
        )

    subcommand = _resolve_command_alias(args[0], CLASSROOM_SEAT_SUBCOMMAND_ALIASES)
    if subcommand is None:
        return _build_classroom_command_error(
            classroom,
            command_text,
            f'不支持的座位命令：{args[0]}',
            command='seat',
            suggestions=['/seat assign 张三 2 3', '/seat move 张三 1 1', '/seat clear 2 3', '/seat swap 张三 李四'],
        )

    if subcommand in {'assign', 'move'}:
        if len(args) < 4:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供学生和目标座位，例如 /seat move 张三 2 3',
                command='seat',
                subcommand=subcommand,
            )
        student_query = ' '.join(args[1:-2]).strip()
        row = _safe_int(args[-2], 0)
        col = _safe_int(args[-1], 0)
        if not student_query or row <= 0 or col <= 0:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '命令格式错误，例如 /seat assign 张三 2 3',
                command='seat',
                subcommand=subcommand,
            )
        action = 'assign_student' if subcommand == 'assign' else 'move_student'
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'seat',
            subcommand=subcommand,
            tool_name='execute_classroom_action',
            arguments={
                'action': action,
                'student_query': student_query,
                'row': row,
                'col': col,
            },
            request=request,
            kind='mutation',
            needs_refresh=True,
        )

    if subcommand == 'clear':
        if len(args) < 3:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供要清空的座位，例如 /seat clear 2 3',
                command='seat',
                subcommand='clear',
            )
        row = _safe_int(args[1], 0)
        col = _safe_int(args[2], 0)
        if row <= 0 or col <= 0:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '座位行列必须是正整数，例如 /seat clear 2 3',
                command='seat',
                subcommand='clear',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'seat',
            subcommand='clear',
            tool_name='execute_classroom_action',
            arguments={'action': 'clear_seat', 'row': row, 'col': col},
            request=request,
            kind='mutation',
            needs_refresh=True,
        )

    if len(args) < 3:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请提供两名学生，例如 /seat swap 张三 李四',
            command='seat',
            subcommand='swap',
        )
    student_a = str(args[1] or '').strip()
    student_b = ' '.join(args[2:]).strip()
    if not student_a or not student_b:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请提供两名学生，例如 /seat swap 张三 李四',
            command='seat',
            subcommand='swap',
        )
    return _run_classroom_command_tool(
        classroom,
        command_text,
        'seat',
        subcommand='swap',
        tool_name='swap_students',
        arguments={'student_a': student_a, 'student_b': student_b},
        request=request,
        kind='mutation',
        needs_refresh=True,
    )


def _handle_classroom_group_command(classroom, command_text, args, request=None):
    if not args:
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'group',
            subcommand='score',
            tool_name='get_group_scores',
            arguments={},
            request=request,
            kind='query',
            reply_formatter=_format_group_scores_command_reply,
        )

    subcommand = _resolve_command_alias(args[0], CLASSROOM_GROUP_SUBCOMMAND_ALIASES)
    if subcommand is None:
        return _build_classroom_command_error(
            classroom,
            command_text,
            f'不支持的小组命令：{args[0]}',
            command='group',
            suggestions=['/group score', '/group create 第一组', '/group rename 第一组 第二组', '/group leader 张三'],
        )

    if subcommand == 'score':
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'group',
            subcommand='score',
            tool_name='get_group_scores',
            arguments={},
            request=request,
            kind='query',
            reply_formatter=_format_group_scores_command_reply,
        )

    if subcommand == 'create':
        name = ' '.join(args[1:]).strip()
        if not name:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供小组名，例如 /group create 第一组',
                command='group',
                subcommand='create',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'group',
            subcommand='create',
            tool_name='execute_classroom_action',
            arguments={'action': 'create_group', 'name': name},
            request=request,
            kind='mutation',
            needs_refresh=True,
        )

    if subcommand == 'rename':
        if len(args) < 3:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供原组名和新组名，例如 /group rename 第一组 第二组',
                command='group',
                subcommand='rename',
            )
        group_query = str(args[1] or '').strip()
        new_name = ' '.join(args[2:]).strip()
        if not group_query or not new_name:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供原组名和新组名，例如 /group rename 第一组 第二组',
                command='group',
                subcommand='rename',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'group',
            subcommand='rename',
            tool_name='execute_classroom_action',
            arguments={'action': 'rename_group', 'group_query': group_query, 'new_name': new_name},
            request=request,
            kind='mutation',
            needs_refresh=True,
        )

    if subcommand == 'delete':
        group_query = ' '.join(args[1:]).strip()
        if not group_query:
            return _build_classroom_command_error(
                classroom,
                command_text,
                '请提供要删除的小组名，例如 /group delete 第一组',
                command='group',
                subcommand='delete',
            )
        return _run_classroom_command_tool(
            classroom,
            command_text,
            'group',
            subcommand='delete',
            tool_name='execute_classroom_action',
            arguments={'action': 'delete_group', 'group_query': group_query},
            request=request,
            kind='mutation',
            needs_refresh=True,
        )

    student_query = ' '.join(args[1:]).strip()
    if not student_query:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请提供学生姓名，例如 /group leader 张三',
            command='group',
            subcommand='leader',
        )
    return _run_classroom_command_tool(
        classroom,
        command_text,
        'group',
        subcommand='leader',
        tool_name='execute_classroom_action',
        arguments={'action': 'set_group_leader', 'student_query': student_query},
        request=request,
        kind='mutation',
        needs_refresh=True,
    )


def _handle_classroom_snapshot_command(classroom, command_text, args, request=None):
    subcommand = _resolve_command_alias(args[0], CLASSROOM_SNAPSHOT_SUBCOMMAND_ALIASES) if args else 'list'
    if subcommand is None:
        return _build_classroom_command_error(
            classroom,
            command_text,
            f'不支持的快照命令：{args[0]}',
            command='snapshot',
            suggestions=['/snapshot list', '/snapshot save 期中布局', '/snapshot load 期中布局', '/snapshot delete 期中布局'],
        )

    if subcommand == 'list':
        items = [
            {
                'id': snapshot.pk,
                'name': snapshot.name,
                'created_at': snapshot.created_at.isoformat() if snapshot.created_at else '',
            }
            for snapshot in classroom.layout_snapshots.all().order_by('-created_at', '-pk')[:20]
        ]
        return _build_classroom_command_result(
            classroom,
            command_text,
            'snapshot',
            subcommand='list',
            kind='query',
            reply=_format_snapshot_list_command_reply(items),
            data={'items': items},
        )

    snapshot_query = ' '.join(args[1:]).strip()
    if not snapshot_query:
        return _build_classroom_command_error(
            classroom,
            command_text,
            '请提供快照名称，例如 /snapshot save 期中布局',
            command='snapshot',
            subcommand=subcommand,
        )

    action_map = {
        'save': 'save_layout_snapshot',
        'load': 'load_layout_snapshot',
        'delete': 'delete_layout_snapshot',
    }
    argument_key = 'name' if subcommand == 'save' else 'snapshot_query'
    return _run_classroom_command_tool(
        classroom,
        command_text,
        'snapshot',
        subcommand=subcommand,
        tool_name='execute_classroom_action',
        arguments={'action': action_map[subcommand], argument_key: snapshot_query},
        request=request,
        kind='mutation',
        needs_refresh=True,
    )


def _handle_classroom_arrange_command(classroom, command_text, args, request=None):
    mode = 'random'
    if args:
        mode = _resolve_command_alias(args[0], CLASSROOM_ARRANGE_MODE_ALIASES) or str(args[0] or '').strip()
    if mode not in CLASSROOM_ARRANGE_MODE_ALIASES:
        mode = _resolve_command_alias(mode, CLASSROOM_ARRANGE_MODE_ALIASES) or mode
    if mode not in CLASSROOM_ARRANGE_MODE_ALIASES:
        return _build_classroom_command_error(
            classroom,
            command_text,
            f'不支持的排座模式：{args[0]}',
            command='arrange',
            suggestions=['/arrange random', '/arrange score_desc', '/arrange score_spread', '/paizuo suiji'],
        )
    return _run_classroom_command_tool(
        classroom,
        command_text,
        'arrange',
        subcommand=mode,
        tool_name='execute_classroom_action',
        arguments={'action': 'auto_arrange', 'mode': mode},
        request=request,
        kind='mutation',
        needs_refresh=True,
    )


def _handle_classroom_undo_redo_command(classroom, command_text, command, request=None):
    action = 'undo' if command == 'undo' else 'redo'
    return _run_classroom_command_tool(
        classroom,
        command_text,
        command,
        tool_name='execute_classroom_action',
        arguments={'action': action},
        request=request,
        kind='mutation',
        needs_refresh=True,
    )


def _execute_classroom_command(classroom, command_text, request=None):
    tokens, error_message = _tokenize_classroom_command(command_text)
    if error_message:
        return _build_classroom_command_error(
            classroom,
            command_text,
            error_message,
            suggestions=['/help', '/view layout', '/students 张三'],
        )

    command = _resolve_command_alias(tokens[0], CLASSROOM_COMMAND_ALIASES)
    if command is None:
        return _build_classroom_command_error(
            classroom,
            command_text,
            f'未知命令：{tokens[0]}。输入 /help 查看全部命令。',
            suggestions=['/help', '/view layout', '/overview', '/students 张三'],
        )

    args = tokens[1:]
    if command == 'help':
        return _handle_classroom_help_command(classroom, command_text, args)
    if command == 'view':
        return _handle_classroom_view_command(classroom, command_text, args)
    if command == 'overview':
        return _handle_classroom_overview_command(classroom, command_text, request=request)
    if command == 'students':
        return _handle_classroom_students_command(classroom, command_text, args, request=request)
    if command == 'seat':
        return _handle_classroom_seat_command(classroom, command_text, args, request=request)
    if command == 'group':
        return _handle_classroom_group_command(classroom, command_text, args, request=request)
    if command == 'snapshot':
        return _handle_classroom_snapshot_command(classroom, command_text, args, request=request)
    if command == 'arrange':
        return _handle_classroom_arrange_command(classroom, command_text, args, request=request)
    if command in {'undo', 'redo'}:
        return _handle_classroom_undo_redo_command(classroom, command_text, command, request=request)
    return _build_classroom_command_error(
        classroom,
        command_text,
        f'命令 {command} 暂未实现。',
        command=command,
    )


def _build_classroom_command_response_payload(classroom, command_result):
    return {
        'status': 'success',
        'reply': str((command_result or {}).get('reply') or '').strip(),
        'command_result': command_result if isinstance(command_result, dict) else {},
        'cards': [],
        'overview': _get_classroom_overview_payload(classroom),
    }


def _extract_direct_swap_call(classroom, message):
    text = str(message or '').strip()
    if not text:
        return None
    lower_text = text.lower()
    keywords = ('交换', '对调', '互换', '换位置', '换座位')
    if not any(word in text for word in keywords):
        return None

    patterns = [
        r'(?:把)?\s*(.+?)\s*(?:和|与|跟)\s*(.+?)\s*(?:交换|对调|互换|换)(?:一下|下)?(?:座位|位置)?',
        r'(?:交换|对调|互换|换)\s*(.+?)\s*(?:和|与|跟)\s*(.+?)(?:的)?(?:座位|位置)?(?:一下|下)?',
    ]

    def clean_name(raw):
        name = str(raw or '').strip()
        for suffix in ['同学', '同桌', '座位', '位置', '的座位', '一下', '下', '吗', '吧', '。', '，', ',', '？', '?']:
            name = name.replace(suffix, '')
        return name.strip()

    for pattern in patterns:
        match = re.search(pattern, text)
        if not match:
            continue
        student_a = clean_name(match.group(1))
        student_b = clean_name(match.group(2))
        if not student_a or not student_b:
            continue
        try:
            resolved_a = _resolve_student_query(classroom, student_a)
            resolved_b = _resolve_student_query(classroom, student_b)
        except ValueError:
            continue
        return {
            'call_id': f'direct_swap_{uuid.uuid4().hex[:12]}',
            'name': 'swap_students',
            'arguments': {
                'student_a': resolved_a.name,
                'student_b': resolved_b.name,
            },
        }
    return None


def _execute_ai_tool(classroom, tool_name, arguments, request=None):
    arguments = arguments or {}

    if tool_name == 'get_classroom_overview':
        overview = _get_classroom_overview_payload(classroom)
        return {
            'ok': True,
            'tool': tool_name,
            'data': {
                'message': '已读取当前班级概览。',
                **overview,
            },
        }

    if tool_name == 'get_student_info':
        student = _resolve_student_query(classroom, arguments.get('student_query'))
        return {
            'ok': True,
            'tool': tool_name,
            'data': {
                'message': f'已读取学生信息：{student.name}',
                **_serialize_student_profile(student, classroom=classroom),
            },
        }

    if tool_name == 'get_group_scores':
        group_rows = _get_group_score_rows(classroom)
        return {
            'ok': True,
            'tool': tool_name,
            'data': {
                'message': f'已读取小组评分，共 {len(group_rows)} 组。',
                'items': group_rows,
            },
        }

    if tool_name == 'get_student_list':
        return {
            'ok': True,
            'tool': tool_name,
            'data': _build_student_list_payload(classroom, arguments),
        }

    if tool_name == 'send_card_info':
        return {
            'ok': True,
            'tool': tool_name,
            'data': _build_card_info_payload(classroom, arguments),
        }

    if tool_name == 'execute_classroom_action':
        return {
            'ok': True,
            'tool': tool_name,
            'data': _execute_classroom_action_tool(classroom, arguments, request=request),
        }

    if tool_name == 'swap_students':
        student_a = _resolve_student_query(classroom, arguments.get('student_a'))
        student_b = _resolve_student_query(classroom, arguments.get('student_b'))
        if student_a.pk == student_b.pk:
            raise ValueError('不能交换同一名学生')
        seat_a = getattr(student_a, 'assigned_seat', None)
        seat_b = getattr(student_b, 'assigned_seat', None)
        if not seat_a or not seat_b:
            raise ValueError('两名学生都需要先入座，才能交换座位')

        before_state = _capture_history_state(classroom)
        with transaction.atomic():
            _swap_seats(seat_a, seat_b)
            violations = _stabilize_layout_with_rules(classroom, request)
            if violations:
                raise ValueError(f'交换失败：{_format_issues_preview(violations)}')
        if request is not None:
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                'swap',
                extra=_build_swap_action(student_a, student_b),
            )
        student_a = classroom.students.get(pk=student_a.pk)
        student_b = classroom.students.get(pk=student_b.pk)
        return {
            'ok': True,
            'tool': tool_name,
            'data': {
                'message': f'已交换 {student_a.name} 和 {student_b.name} 的座位',
                'student_a': _serialize_student_profile(student_a, classroom=classroom),
                'student_b': _serialize_student_profile(student_b, classroom=classroom),
            },
        }

    raise ValueError(f'未知 AI 工具：{tool_name}')


def _extract_function_calls(response):
    def item_get(item, key, default=None):
        if isinstance(item, dict):
            return item.get(key, default)
        return getattr(item, key, default)

    calls = []
    for item in getattr(response, 'output', []) or []:
        if item_get(item, 'type', '') != 'function_call':
            continue
        arguments_raw = item_get(item, 'arguments', {}) or {}
        if isinstance(arguments_raw, dict):
            arguments = arguments_raw
        elif isinstance(arguments_raw, str):
            try:
                arguments = json.loads(arguments_raw)
            except (json.JSONDecodeError, TypeError, ValueError):
                arguments = {}
        else:
            arguments = {}
        call_id = str(item_get(item, 'call_id', '') or item_get(item, 'id', '') or '').strip()
        name = str(item_get(item, 'name', '') or '').strip()
        if not call_id or not name:
            continue
        calls.append({
            'call_id': call_id,
            'name': name,
            'arguments': arguments,
        })
    return calls


def _describe_future_mode_call(tool_name, arguments):
    arguments = arguments or {}
    label = AI_TOOL_LABELS.get(tool_name, tool_name or '未知工具')
    if tool_name == 'get_classroom_overview':
        return f'{label}，用于读取当前班级整体状态。'
    if tool_name == 'get_student_info':
        return f'{label}：{arguments.get("student_query") or "未提供学生"}'
    if tool_name == 'get_group_scores':
        return f'{label}，用于统计当前小组平均分与总分。'
    if tool_name == 'get_student_list':
        return f'{label}，支持排序/筛选并返回名单。'
    if tool_name == 'send_card_info':
        return f'{label}：{arguments.get("card_type") or "未指定卡片类型"}'
    if tool_name == 'execute_classroom_action':
        return f'{label}：{arguments.get("action") or "未指定动作"}'
    if tool_name == 'swap_students':
        return f'{label}：{arguments.get("student_a") or "学生A"} ↔ {arguments.get("student_b") or "学生B"}'
    return f'{label}，参数：{json.dumps(arguments, ensure_ascii=False)}'


def _collect_response_text(response):
    text = getattr(response, 'output_text', '') or ''
    if text:
        return text.strip()

    parts = []
    for item in getattr(response, 'output', []) or []:
        if getattr(item, 'type', '') != 'message':
            continue
        for content in getattr(item, 'content', []) or []:
            if getattr(content, 'type', '') in {'output_text', 'text'}:
                parts.append(getattr(content, 'text', '') or '')
    return '\n'.join(part.strip() for part in parts if str(part).strip()).strip()


def _future_mode_system_prompt():
    return (
        '你是“不想排座位”的 Window Inteligence ｜ 闻道智能 助手，用简体中文回答。'
        'Window Inteligence ｜ 闻道智能 是全球首个面向教师行业的通用型 Agent 应用，当前正处于内部封测阶段，你和整个“不想排座位“项目属于一个小 Demo'
        'Window Inteligence 的开发者是 老三 ，他的个人网站是 www.577622.xyz ，Window Inteligence 的官网是window.577622.xyz'
        '你服务于班主任/老师，回答必须准确、冷静、直接。'
        '涉及当前班级的事实信息时，优先调用工具，不要猜测。'
        '用户要求交换座位时，若信息足够，应直接调用 swap_students。'
        '用户要求名单筛选、学生列表、卡片图表时，应优先调用 get_student_list 或 send_card_info。'
        '一切切实改动都需要调用工具，请勿直接对用户说“已完成”之类的话，记住，一切完成的前提都是调用工具，请在你准备和用户说“完成/搞定”之前，先想想你是否真的调用了工具。'
        '禁止使用 Markdown 语法。'
        '只允许使用换行与数字序号（1. 2. 3.）进行简单排版。'
        '回复尽量简洁，并明确说明你做了什么。'
    )


def _extract_chat_delta_text(content):
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                text = item.get('text')
            else:
                text = getattr(item, 'text', '')
            if text:
                parts.append(str(text))
        return ''.join(parts)
    if hasattr(content, 'text'):
        return str(getattr(content, 'text') or '')
    return ''


def _sse_event(event_name, payload):
    return f"event: {event_name}\ndata: {json.dumps(payload, ensure_ascii=False)}\n\n"


def _create_future_mode_response(client, model, conversation=None, previous_response_id=None, tool_outputs=None):
    if previous_response_id:
        return client.responses.create(
            model=model,
            previous_response_id=previous_response_id,
            input=tool_outputs or [],
            tools=AI_TOOL_DEFINITIONS,
            parallel_tool_calls=False,
        )

    system_prompt = _future_mode_system_prompt()
    def build_input_items(use_content_parts=False):
        input_items = [{
            'role': 'system',
            'content': [{'type': 'input_text', 'text': system_prompt}] if use_content_parts else system_prompt,
        }]
        for message in (conversation or [])[-AI_CONTEXT_MESSAGE_LIMIT:]:
            role = 'assistant' if message.get('role') == 'assistant' else 'user'
            content = str(message.get('content') or '').strip()
            if not content:
                continue
            text = content[:4000]
            input_items.append({
                'role': role,
                'content': [{'type': 'input_text', 'text': text}] if use_content_parts else text,
            })
        return input_items

    try:
        return client.responses.create(
            model=model,
            input=build_input_items(use_content_parts=False),
            tools=AI_TOOL_DEFINITIONS,
            parallel_tool_calls=False,
        )
    except Exception as exc:
        if not _is_responses_content_unmarshal_error(exc):
            raise
        return client.responses.create(
            model=model,
            input=build_input_items(use_content_parts=True),
            tools=AI_TOOL_DEFINITIONS,
            parallel_tool_calls=False,
        )


def _build_chat_tools():
    tools = []
    for item in AI_TOOL_DEFINITIONS:
        tools.append({
            'type': 'function',
            'function': {
                'name': item.get('name'),
                'description': item.get('description') or '',
                'strict': bool(item.get('strict', True)),
                'parameters': item.get('parameters') or {
                    'type': 'object',
                    'properties': {},
                    'required': [],
                    'additionalProperties': False,
                },
            },
        })
    return tools


def _normalize_chat_text_content(content):
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                text = item.get('text')
            else:
                text = getattr(item, 'text', '')
            if text:
                parts.append(str(text).strip())
        return '\n'.join(part for part in parts if part)
    return str(content or '').strip()


def _normalize_chat_reasoning_content(content):
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                text = item.get('text')
            else:
                text = getattr(item, 'text', '')
            if text:
                parts.append(str(text).strip())
        return '\n'.join(part for part in parts if part)
    return str(content or '').strip()


def _extract_chat_function_calls(message):
    tool_calls = getattr(message, 'tool_calls', None)
    if tool_calls is None and isinstance(message, dict):
        tool_calls = message.get('tool_calls')
    calls = []
    for call in tool_calls or []:
        call_id = str(getattr(call, 'id', None) or (call.get('id') if isinstance(call, dict) else '') or '').strip()
        function = getattr(call, 'function', None) if not isinstance(call, dict) else call.get('function')
        name = str(getattr(function, 'name', None) or (function.get('name') if isinstance(function, dict) else '') or '').strip()
        arguments_raw = getattr(function, 'arguments', None) if not isinstance(function, dict) else function.get('arguments')
        if isinstance(arguments_raw, dict):
            arguments = arguments_raw
            arguments_json = json.dumps(arguments_raw, ensure_ascii=False)
        elif isinstance(arguments_raw, str):
            arguments_json = arguments_raw
            try:
                arguments = json.loads(arguments_raw)
            except (TypeError, ValueError, json.JSONDecodeError):
                arguments = {}
        else:
            arguments = {}
            arguments_json = '{}'
        if not call_id or not name:
            continue
        calls.append({
            'call_id': call_id,
            'name': name,
            'arguments': arguments,
            'arguments_json': arguments_json,
        })
    return calls


def _tool_outputs_to_chat_messages(tool_outputs):
    messages = []
    for item in tool_outputs or []:
        if not isinstance(item, dict):
            continue
        call_id = str(item.get('call_id') or '').strip()
        if not call_id:
            continue
        output = item.get('output')
        if not isinstance(output, str):
            output = json.dumps(output, ensure_ascii=False)
        messages.append({
            'role': 'tool',
            'tool_call_id': call_id,
            'content': output or '{}',
        })
    return messages


def _run_future_mode_chat(
    classroom,
    client,
    model,
    conversation,
    request=None,
    client_config=None,
    tool_outputs=None,
    tool_events=None,
    chat_messages=None,
    conversation_id=None,
):
    tools = _build_chat_tools()
    if chat_messages is None:
        system_prompt = _future_mode_system_prompt()
        messages = [{'role': 'system', 'content': system_prompt}]
        for message in (conversation or [])[-AI_CONTEXT_MESSAGE_LIMIT:]:
            role = 'assistant' if message.get('role') == 'assistant' else 'user'
            content = str(message.get('content') or '').strip()
            if not content:
                continue
            messages.append({'role': role, 'content': content[:4000]})
    else:
        messages = list(chat_messages)

    if tool_outputs:
        messages.extend(_tool_outputs_to_chat_messages(tool_outputs))

    completion_kwargs = {
        'model': model,
        'messages': messages,
        'tools': tools,
        'tool_choice': 'auto',
        'parallel_tool_calls': False,
    }
    extra_body = _build_chat_completion_extra_body(client_config=client_config)
    if extra_body:
        completion_kwargs['extra_body'] = extra_body
    completion = client.chat.completions.create(
        **completion_kwargs,
    )
    choice = (getattr(completion, 'choices', None) or [None])[0]
    assistant_message = getattr(choice, 'message', None) if choice is not None else None
    if assistant_message is None:
        return {
            'status': 'completed',
            'reply': 'AI 未返回有效消息，请稍后重试。',
            'tool_events': tool_events or [],
        }

    function_calls = _extract_chat_function_calls(assistant_message)
    if function_calls:
        assistant_tool_calls = []
        for call in function_calls:
            assistant_tool_calls.append({
                'id': call['call_id'],
                'type': 'function',
                'function': {
                    'name': call['name'],
                    'arguments': call.get('arguments_json') or json.dumps(call.get('arguments') or {}, ensure_ascii=False),
                },
            })

        assistant_content_raw = assistant_message.get('content') if isinstance(assistant_message, dict) else getattr(assistant_message, 'content', None)
        assistant_reasoning_raw = assistant_message.get('reasoning_content') if isinstance(assistant_message, dict) else getattr(assistant_message, 'reasoning_content', None)
        assistant_content = _normalize_chat_text_content(assistant_content_raw)
        assistant_reasoning_content = _normalize_chat_reasoning_content(assistant_reasoning_raw)
        pending_messages = list(messages)
        assistant_pending_message = {
            'role': 'assistant',
            'content': assistant_content or '',
            'tool_calls': assistant_tool_calls,
        }
        if assistant_reasoning_content:
            assistant_pending_message['reasoning_content'] = assistant_reasoning_content
        pending_messages.append(assistant_pending_message)
        token = _store_future_mode_pending(
            request,
            classroom.pk,
            conversation_id,
            response_id='',
            function_calls=[{
                'call_id': call['call_id'],
                'name': call['name'],
                'arguments': call['arguments'],
            } for call in function_calls],
            mode='chat',
            chat_messages=pending_messages,
        ) if request is not None else ''
        return {
            'status': 'needs_approval',
            'approval_token': token,
            'pending_calls': [
                {
                    'call_id': call['call_id'],
                    'name': call['name'],
                    'label': AI_TOOL_LABELS.get(call['name'], call['name']),
                    'arguments': call['arguments'],
                    'summary': _describe_future_mode_call(call['name'], call['arguments']),
                }
                for call in function_calls
            ],
            'tool_events': tool_events or [],
        }

    reply = _normalize_chat_text_content(getattr(assistant_message, 'content', None)) or '已完成处理，但没有生成可展示的回复。'
    return {
        'status': 'completed',
        'reply': reply,
        'tool_events': tool_events or [],
    }


def _run_future_mode(
    classroom,
    conversation,
    request=None,
    client_config=None,
    previous_response_id=None,
    tool_outputs=None,
    tool_events=None,
    mode='responses',
    chat_messages=None,
    conversation_id=None,
):
    client = _get_openai_client(client_config=client_config)
    model = _get_openai_model(client_config=client_config)
    selected_mode = mode or 'responses'
    if selected_mode == 'auto':
        selected_mode = 'chat' if _should_use_chat_completions(client_config=client_config) else 'responses'

    if selected_mode == 'chat':
        return _run_future_mode_chat(
            classroom,
            client,
            model,
            conversation=conversation,
            request=request,
            client_config=client_config,
            tool_outputs=tool_outputs,
            tool_events=tool_events,
            chat_messages=chat_messages,
            conversation_id=conversation_id,
        )

    try:
        response = _create_future_mode_response(
            client,
            model,
            conversation=conversation,
            previous_response_id=previous_response_id,
            tool_outputs=tool_outputs,
        )
    except Exception as exc:
        if _is_responses_not_supported_error(exc) and not previous_response_id:
            return _run_future_mode_chat(
                classroom,
                client,
                model,
                conversation=conversation,
                request=request,
                client_config=client_config,
                tool_outputs=tool_outputs,
                tool_events=tool_events,
                chat_messages=chat_messages,
                conversation_id=conversation_id,
            )
        raise

    function_calls = _extract_function_calls(response)
    if function_calls:
        token = _store_future_mode_pending(
            request,
            classroom.pk,
            conversation_id,
            getattr(response, 'id', ''),
            function_calls,
            mode='responses',
            chat_messages=None,
        ) if request is not None else ''
        return {
            'status': 'needs_approval',
            'approval_token': token,
            'pending_calls': [
                {
                    'call_id': call['call_id'],
                    'name': call['name'],
                    'label': AI_TOOL_LABELS.get(call['name'], call['name']),
                    'arguments': call['arguments'],
                    'summary': _describe_future_mode_call(call['name'], call['arguments']),
                }
                for call in function_calls
            ],
            'tool_events': tool_events or [],
        }

    reply = _collect_response_text(response) or '已完成处理，但没有生成可展示的回复。'
    return {
        'status': 'completed',
        'reply': reply,
        'tool_events': tool_events or [],
    }


SVG_EXPORT_THEME_MAP = {
    'classic': {
        'bg': '#f7faff',
        'title': '#0f172a',
        'name': '#111827',
        'sub': '#667085',
        'type': '#475467',
        'podium_fill': '#e7efff',
        'podium_stroke': '#c9dbff',
        'seat_fill_occupied': '#eef4ff',
        'seat_stroke_occupied': '#bfd4ff',
        'seat_fill_empty': '#f8fbff',
        'seat_stroke_empty': '#d3e1ff',
        'nonseat_stroke': '#d0d5dd',
        'nonseat_aisle': '#eff3f8',
        'nonseat_podium': '#fff3e8',
        'nonseat_empty': '#f2f4f7',
        'tag_text': '#ffffff',
        'group_palette': ['#0a59f7', '#00a38c', '#ff8b00', '#e45193', '#6b64ff', '#2ca2ff', '#13a44a', '#c85a0f'],
    },
    'minimal': {
        'bg': '#f8fafc',
        'title': '#1f2937',
        'name': '#111827',
        'sub': '#6b7280',
        'type': '#4b5563',
        'podium_fill': '#edf2f7',
        'podium_stroke': '#d2dae6',
        'seat_fill_occupied': '#f9fafb',
        'seat_stroke_occupied': '#cbd5e1',
        'seat_fill_empty': '#ffffff',
        'seat_stroke_empty': '#d1d5db',
        'nonseat_stroke': '#d1d5db',
        'nonseat_aisle': '#f1f5f9',
        'nonseat_podium': '#f3f4f6',
        'nonseat_empty': '#f8fafc',
        'tag_text': '#ffffff',
        'group_palette': ['#0a59f7', '#64748b', '#0f766e', '#b45309', '#be123c', '#1d4ed8', '#065f46', '#7c3aed'],
    },
    'contrast': {
        'bg': '#0b1220',
        'title': '#e5ecff',
        'name': '#ffffff',
        'sub': '#b7c7e9',
        'type': '#d2dbf5',
        'podium_fill': '#1c2f5d',
        'podium_stroke': '#33509c',
        'seat_fill_occupied': '#172a55',
        'seat_stroke_occupied': '#3b5db7',
        'seat_fill_empty': '#111b34',
        'seat_stroke_empty': '#30477f',
        'nonseat_stroke': '#2e426f',
        'nonseat_aisle': '#1d2a44',
        'nonseat_podium': '#2a2f4d',
        'nonseat_empty': '#202c47',
        'tag_text': '#ffffff',
        'group_palette': ['#0a59f7', '#0fa968', '#f59e0b', '#ef4444', '#06b6d4', '#a855f7', '#f97316', '#14b8a6'],
    },
}


def _name_emphasis_font_size(text):
    length = max(1, len(str(text or '')))
    size = 30 - length * 2
    if size < 16:
        size = 16
    if size > 26:
        size = 26
    return size


def _hex_to_rgb_parts(color):
    raw = str(color or '').strip().lstrip('#')
    if len(raw) == 3:
        raw = ''.join(ch * 2 for ch in raw)
    if len(raw) != 6:
        return 0, 0, 0
    try:
        return int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16)
    except ValueError:
        return 0, 0, 0


def _export_svg_font_style():
    return (
        f'.title{{font-family:"{EXPORT_FONT_BLACK}";font-size:24px;}}'
        f'.cell-name{{font-family:"{EXPORT_FONT_BLACK}";font-size:16px;}}'
        f'.cell-sub{{font-family:"{EXPORT_FONT_LIGHT}";font-size:12px;}}'
        f'.tag{{font-family:"{EXPORT_FONT_BLACK}";font-size:11px;}}'
        f'.cell-type{{font-family:"{EXPORT_FONT_BLACK}";font-size:13px;}}'
    )


def _sync_seats(classroom, rows, cols):
    if classroom.rows != rows or classroom.cols != cols:
        classroom.rows = rows
        classroom.cols = cols
        classroom.save(update_fields=['rows', 'cols'])
    classroom.seats.filter(models.Q(row__gt=rows) | models.Q(col__gt=cols)).delete()
    classroom.generate_seats()


def _empty_layout_cell_payload(row, col):
    return {
        'row': row,
        'col': col,
        'cell_type': SeatCellType.EMPTY,
        'student_pk': None,
        'student_id': None,
        'student_name': None,
        'group_name': None,
    }


def _sort_layout_seats(seats):
    return sorted(
        list(seats),
        key=lambda item: (
            int(item.get('row') or 0),
            int(item.get('col') or 0),
        )
    )


def _is_blank_layout_cell(cell):
    return (
        str(cell.get('cell_type') or SeatCellType.SEAT) == SeatCellType.EMPTY
        and not cell.get('student_pk')
        and not cell.get('student_id')
        and not cell.get('student_name')
        and not cell.get('group_name')
    )


def _normalize_shift_direction(direction):
    normalized = str(direction or '').strip().lower()
    aliases = {
        'up': 'front',
        'forward': 'front',
        'down': 'back',
        'backward': 'back',
    }
    normalized = aliases.get(normalized, normalized)
    if normalized not in {'left', 'right', 'front', 'back'}:
        raise ValueError('移动方向不合法')
    return normalized


def _shift_direction_meta(direction):
    normalized = _normalize_shift_direction(direction)
    if normalized == 'left':
        return {
            'direction': normalized,
            'axis_field': 'col',
            'size_key': 'cols',
            'unit': '列',
            'action_label': '左移',
            'leading_label': '最左侧',
            'expand': False,
        }
    if normalized == 'right':
        return {
            'direction': normalized,
            'axis_field': 'col',
            'size_key': 'cols',
            'unit': '列',
            'action_label': '右移',
            'leading_label': '最左侧',
            'expand': True,
        }
    if normalized == 'front':
        return {
            'direction': normalized,
            'axis_field': 'row',
            'size_key': 'rows',
            'unit': '行',
            'action_label': '前移',
            'leading_label': '最前方',
            'expand': False,
        }
    return {
        'direction': normalized,
        'axis_field': 'row',
        'size_key': 'rows',
        'unit': '行',
        'action_label': '后移',
        'leading_label': '最前方',
        'expand': True,
    }


def _shift_layout_constraints(constraints, axis_field, delta, axis_label, action_label):
    size_map = {'row': 'rows', 'col': 'cols'}
    size_key = size_map.get(axis_field)
    if not size_key:
        raise ValueError(f'{axis_label}约束处理失败')
    shifted = []
    for raw in constraints:
        item = dict(raw)
        axis_value = item.get(axis_field)
        if axis_value in (None, ''):
            shifted.append(item)
            continue
        size = int(item.get('_classroom_size') or 0)
        if size < 1:
            raise ValueError(f'{axis_label}约束处理失败')
        next_value = (int(axis_value) - 1 + int(delta)) % size + 1
        item[axis_field] = next_value
        item.pop('_classroom_size', None)
        shifted.append(item)
    return shifted


def _build_shift_axis_map(size, excluded_values, expand, steps):
    if size < 1:
        return {}

    excluded = {
        _safe_int(value, 0)
        for value in (excluded_values or set())
    }
    eligible_values = [
        value
        for value in range(1, size + 1)
        if value not in excluded
    ]
    if not eligible_values:
        return {value: value for value in range(1, size + 1)}

    normalized_steps = int(steps or 0) % len(eligible_values)
    if normalized_steps == 0:
        normalized_steps = len(eligible_values)

    delta = normalized_steps if expand else -normalized_steps
    axis_map = {value: value for value in range(1, size + 1)}
    for index, value in enumerate(eligible_values):
        axis_map[value] = eligible_values[(index + delta) % len(eligible_values)]
    return axis_map


def _shift_layout_constraints_by_axis_map(constraints, axis_field, axis_map):
    shifted = []
    for raw in constraints:
        item = dict(raw)
        axis_value = item.get(axis_field)
        if axis_value not in (None, ''):
            try:
                normalized_axis_value = int(axis_value)
                item[axis_field] = int(axis_map.get(normalized_axis_value, normalized_axis_value))
            except (TypeError, ValueError):
                pass
        item.pop('_classroom_size', None)
        shifted.append(item)
    return shifted


def _normalize_mirror_axis(axis):
    normalized = str(axis or '').strip().lower()
    aliases = {
        'left_right': 'lr',
        'horizontal': 'lr',
        'flip_lr': 'lr',
        'mirror_lr': 'lr',
    }
    normalized = aliases.get(normalized, normalized)
    if normalized not in {'lr'}:
        raise ValueError('镜像方向不合法')
    return normalized


def _build_mirrored_layout_payload(classroom, axis='lr'):
    normalized_axis = _normalize_mirror_axis(axis)
    payload = copy.deepcopy(
        _snapshot_payload(classroom, include_students=False, include_constraints=True)
    )

    cols = int(payload.get('classroom', {}).get('cols') or classroom.cols)

    mirrored_seats = []
    for seat in payload.get('seats', []):
        item = dict(seat)
        if normalized_axis == 'lr':
            item['col'] = cols - int(item.get('col') or 0) + 1
        mirrored_seats.append(item)

    payload['seats'] = _sort_layout_seats(mirrored_seats)

    mirrored_constraints = []
    for raw in payload.get('constraints', []):
        item = dict(raw)
        if normalized_axis == 'lr':
            current_col = item.get('col')
            if current_col not in (None, ''):
                item['col'] = cols - int(current_col) + 1
        mirrored_constraints.append(item)
    payload['constraints'] = mirrored_constraints

    return payload


def _build_layout_column_payload(payload):
    column_map = defaultdict(list)
    for seat in payload.get('seats', []):
        try:
            col = int(seat.get('col') or 0)
        except Exception:
            continue
        if col < 1:
            continue
        column_map[col].append(dict(seat))
    for cells in column_map.values():
        cells.sort(key=lambda item: int(item.get('row') or 0))
    return column_map


def _classify_horizontal_layout_column(cells, row_count):
    seat_cells = 0
    aisle_cells = 0
    occupied_cells = 0

    for cell in cells:
        cell_type = str(cell.get('cell_type') or SeatCellType.SEAT)
        if cell_type == SeatCellType.SEAT:
            seat_cells += 1
        elif cell_type == SeatCellType.AISLE:
            aisle_cells += 1
        if cell.get('student_pk') or cell.get('student_id') or cell.get('student_name') or cell.get('group_name'):
            occupied_cells += 1

    aisle_threshold = max(1, math.ceil(max(row_count, 1) * 0.5))
    seat_threshold = max(1, math.ceil(max(row_count, 1) * 0.4))

    if occupied_cells == 0 and aisle_cells >= aisle_threshold and aisle_cells > seat_cells:
        return 'aisle'
    if occupied_cells > 0 or seat_cells >= seat_threshold:
        return 'seat'
    return 'other'


def _build_horizontal_template_blocks(payload):
    classroom_data = payload.get('classroom', {})
    cols = int(classroom_data.get('cols') or 0)
    rows = int(classroom_data.get('rows') or 0)
    column_payload = _build_layout_column_payload(payload)

    blocks = []
    current_type = None
    current_cols = []

    for col in range(1, cols + 1):
        cells = column_payload.get(col, [])
        column_type = _classify_horizontal_layout_column(cells, rows)
        if current_type != column_type:
            if current_cols:
                start_col = current_cols[0]
                end_col = current_cols[-1]
                content_cells = []
                for source_col in current_cols:
                    for cell in column_payload.get(source_col, []):
                        item = dict(cell)
                        item['col_offset'] = source_col - start_col
                        content_cells.append(item)
                blocks.append({
                    'block_type': current_type,
                    'start_col': start_col,
                    'end_col': end_col,
                    'width': len(current_cols),
                    'columns': list(current_cols),
                    'content_cells': content_cells,
                    'is_structural': current_type == 'aisle',
                })
            current_type = column_type
            current_cols = [col]
        else:
            current_cols.append(col)

    if current_cols:
        start_col = current_cols[0]
        end_col = current_cols[-1]
        content_cells = []
        for source_col in current_cols:
            for cell in column_payload.get(source_col, []):
                item = dict(cell)
                item['col_offset'] = source_col - start_col
                content_cells.append(item)
        blocks.append({
            'block_type': current_type,
            'start_col': start_col,
            'end_col': end_col,
            'width': len(current_cols),
            'columns': list(current_cols),
            'content_cells': content_cells,
            'is_structural': current_type == 'aisle',
        })

    return blocks


def _format_horizontal_template_signature(blocks):
    return '+'.join(str(int(block.get('width') or 0)) for block in blocks if int(block.get('width') or 0) > 0)


def _analyze_horizontal_layout_template(payload):
    blocks = _build_horizontal_template_blocks(payload)
    if not blocks:
        return {
            'supported': False,
            'reason': '当前布局为空，无法识别横向模板',
            'blocks': [],
            'template_signature': '',
            'seat_block_count': 0,
            'aisle_block_count': 0,
        }

    other_blocks = [block for block in blocks if block.get('block_type') == 'other']
    if other_blocks:
        return {
            'supported': False,
            'reason': '存在无法稳定判定的结构列，请先整理走廊列',
            'blocks': blocks,
            'template_signature': _format_horizontal_template_signature(blocks),
            'seat_block_count': len([block for block in blocks if block.get('block_type') == 'seat']),
            'aisle_block_count': len([block for block in blocks if block.get('block_type') == 'aisle']),
        }

    seat_blocks = [block for block in blocks if block.get('block_type') == 'seat']
    aisle_blocks = [block for block in blocks if block.get('block_type') == 'aisle']
    if len(seat_blocks) < 2:
        return {
            'supported': False,
            'reason': '未识别到至少两个稳定的座位块',
            'blocks': blocks,
            'template_signature': _format_horizontal_template_signature(blocks),
            'seat_block_count': len(seat_blocks),
            'aisle_block_count': len(aisle_blocks),
        }
    if not aisle_blocks:
        return {
            'supported': False,
            'reason': '未识别到结构走廊列，已无法安全执行块级轮换',
            'blocks': blocks,
            'template_signature': _format_horizontal_template_signature(blocks),
            'seat_block_count': len(seat_blocks),
            'aisle_block_count': 0,
        }

    if blocks[0].get('block_type') != 'seat' or blocks[-1].get('block_type') != 'seat':
        return {
            'supported': False,
            'reason': '横向模板必须以座位块开始并以座位块结束',
            'blocks': blocks,
            'template_signature': _format_horizontal_template_signature(blocks),
            'seat_block_count': len(seat_blocks),
            'aisle_block_count': len(aisle_blocks),
        }

    for index, block in enumerate(blocks):
        expected_type = 'seat' if index % 2 == 0 else 'aisle'
        if block.get('block_type') != expected_type:
            return {
                'supported': False,
                'reason': '当前横向结构不是稳定的座位块与走廊块交替模板',
                'blocks': blocks,
                'template_signature': _format_horizontal_template_signature(blocks),
                'seat_block_count': len(seat_blocks),
                'aisle_block_count': len(aisle_blocks),
            }

    if len(aisle_blocks) != len(seat_blocks) - 1:
        return {
            'supported': False,
            'reason': '结构走廊数量与座位块边界不匹配',
            'blocks': blocks,
            'template_signature': _format_horizontal_template_signature(blocks),
            'seat_block_count': len(seat_blocks),
            'aisle_block_count': len(aisle_blocks),
        }

    return {
        'supported': True,
        'reason': '',
        'blocks': blocks,
        'template_signature': _format_horizontal_template_signature(blocks),
        'seat_block_count': len(seat_blocks),
        'aisle_block_count': len(aisle_blocks),
    }


def _rotate_template_seat_blocks(seat_blocks, direction, steps):
    if not seat_blocks:
        return []
    shift = int(steps or 0) % len(seat_blocks)
    if shift == 0:
        return list(seat_blocks)
    if direction == 'left':
        return list(seat_blocks[shift:]) + list(seat_blocks[:shift])
    return list(seat_blocks[-shift:]) + list(seat_blocks[:-shift])


def _should_preserve_horizontal_template_structure(blocks):
    if not blocks:
        return False
    widths = [int(block.get('width') or 0) for block in blocks]
    return widths == list(reversed(widths))


def _build_structure_preserved_template_shift(payload, blocks, direction, steps):
    seat_blocks = [block for block in blocks if block.get('block_type') == 'seat']
    if not seat_blocks:
        return None

    seat_blocks_by_width = defaultdict(list)
    for block in seat_blocks:
        seat_blocks_by_width[int(block.get('width') or 0)].append(block)

    rotated_by_width = {}
    for width, width_blocks in seat_blocks_by_width.items():
        rotated_by_width[width] = _rotate_template_seat_blocks(width_blocks, direction, steps)

    width_offsets = defaultdict(int)
    rebuilt_seats = []
    column_map = {}

    for block in blocks:
        width = int(block.get('width') or 0)
        target_start = int(block.get('start_col') or 0)
        if width < 1 or target_start < 1:
            continue

        if block.get('block_type') == 'seat':
            width_key = int(block.get('width') or 0)
            width_index = width_offsets[width_key]
            source_block = rotated_by_width[width_key][width_index]
            width_offsets[width_key] += 1
        else:
            source_block = block

        source_start = int(source_block.get('start_col') or target_start)
        for offset in range(width):
            column_map[source_start + offset] = target_start + offset

        for cell in source_block.get('content_cells', []):
            item = dict(cell)
            item['col'] = target_start + int(cell.get('col_offset') or 0)
            item.pop('col_offset', None)
            rebuilt_seats.append(item)

    rebuilt_payload = copy.deepcopy(payload)
    rebuilt_payload['seats'] = _sort_layout_seats(rebuilt_seats)

    remapped_constraints = []
    for raw in rebuilt_payload.get('constraints', []):
        item = dict(raw)
        current_col = item.get('col')
        if current_col not in (None, ''):
            try:
                mapped_col = column_map.get(int(current_col))
            except Exception:
                mapped_col = None
            if mapped_col is not None:
                item['col'] = mapped_col
        remapped_constraints.append(item)
    rebuilt_payload['constraints'] = remapped_constraints

    return rebuilt_payload, column_map


def _rotate_single_column_units(seat_columns, direction, steps):
    if not seat_columns:
        return []
    shift = int(steps or 0) % len(seat_columns)
    if shift == 0:
        return list(seat_columns)
    if direction == 'left':
        return list(seat_columns[-shift:]) + list(seat_columns[:-shift])
    return list(seat_columns[shift:]) + list(seat_columns[:shift])


def _build_single_column_horizontal_shift_payload(classroom, direction, steps):
    normalized_steps = _safe_int(steps, 0)
    if normalized_steps < 1:
        raise ValueError('移动模板单位数必须大于 0')

    payload = copy.deepcopy(
        _snapshot_payload(classroom, include_students=False, include_constraints=True)
    )
    classroom_data = payload.get('classroom', {})
    cols = int(classroom_data.get('cols') or classroom.cols)
    rows = int(classroom_data.get('rows') or classroom.rows)
    column_payload = _build_layout_column_payload(payload)
    blocks = _build_horizontal_template_blocks(payload)

    seat_columns = []
    structural_columns = []
    column_map = {}
    for col in range(1, cols + 1):
        column_type = _classify_horizontal_layout_column(column_payload.get(col, []), rows)
        column_map[col] = col
        if column_type == 'seat':
            seat_columns.append(col)
        else:
            structural_columns.append(col)

    rotated_columns = _rotate_single_column_units(seat_columns, direction, normalized_steps)
    for target_col, source_col in zip(seat_columns, rotated_columns):
        column_map[source_col] = target_col

    rebuilt_seats = []
    for source_col in range(1, cols + 1):
        target_col = int(column_map.get(source_col) or source_col)
        for cell in column_payload.get(source_col, []):
            item = dict(cell)
            item['col'] = target_col
            rebuilt_seats.append(item)

    payload['seats'] = _sort_layout_seats(rebuilt_seats)

    remapped_constraints = []
    for raw in payload.get('constraints', []):
        item = dict(raw)
        current_col = item.get('col')
        if current_col not in (None, ''):
            try:
                item['col'] = int(column_map.get(int(current_col), int(current_col)))
            except Exception:
                pass
        remapped_constraints.append(item)
    payload['constraints'] = remapped_constraints

    return payload, {
        'supported': True,
        'reason': '',
        'column_map': column_map,
        'shift_units': normalized_steps,
        'seat_column_count': len(seat_columns),
        'structural_column_count': len(structural_columns),
        'template_signature': _format_horizontal_template_signature(blocks),
        'template_strategy': 'single_column',
    }


def _rebuild_template_blocks_for_shift(blocks, direction, steps):
    seat_blocks = [block for block in blocks if block.get('block_type') == 'seat']
    aisle_blocks = [block for block in blocks if block.get('block_type') == 'aisle']
    rotated_seat_blocks = _rotate_template_seat_blocks(seat_blocks, direction, steps)
    rebuilt = []
    for index, seat_block in enumerate(rotated_seat_blocks):
        rebuilt.append(seat_block)
        if index < len(aisle_blocks):
            rebuilt.append(aisle_blocks[index])
    return rebuilt


def _build_intelligent_horizontal_shift_payload(classroom, direction, steps):
    normalized_steps = _safe_int(steps, 0)
    if normalized_steps < 1:
        raise ValueError('移动模板单位数必须大于 0')

    payload = copy.deepcopy(
        _snapshot_payload(classroom, include_students=False, include_constraints=True)
    )
    analysis = _analyze_horizontal_layout_template(payload)
    if not analysis.get('supported'):
        return None, analysis

    column_map = {}
    if _should_preserve_horizontal_template_structure(analysis.get('blocks', [])):
        rebuilt = _build_structure_preserved_template_shift(
            payload,
            analysis.get('blocks', []),
            direction,
            normalized_steps,
        )
        if rebuilt is not None:
            payload, column_map = rebuilt
            analysis['template_strategy'] = 'preserve_structure'
    if not column_map:
        rebuilt_blocks = _rebuild_template_blocks_for_shift(analysis.get('blocks', []), direction, normalized_steps)
        rebuilt_seats = []
        next_col = 1

        for block in rebuilt_blocks:
            source_start = int(block.get('start_col') or next_col)
            width = int(block.get('width') or 0)
            for offset in range(width):
                column_map[source_start + offset] = next_col + offset
            for cell in block.get('content_cells', []):
                item = dict(cell)
                item['col'] = next_col + int(cell.get('col_offset') or 0)
                item.pop('col_offset', None)
                rebuilt_seats.append(item)
            next_col += width

        payload['seats'] = _sort_layout_seats(rebuilt_seats)

        remapped_constraints = []
        for raw in payload.get('constraints', []):
            item = dict(raw)
            current_col = item.get('col')
            if current_col not in (None, ''):
                try:
                    mapped_col = column_map.get(int(current_col))
                except Exception:
                    mapped_col = None
                if mapped_col is not None:
                    item['col'] = mapped_col
            remapped_constraints.append(item)
        payload['constraints'] = remapped_constraints
        analysis['template_strategy'] = 'rotate_blocks'

    analysis['column_map'] = column_map
    analysis['shift_units'] = normalized_steps
    return payload, analysis


def _build_shifted_layout_payload(classroom, direction, steps):
    meta = _shift_direction_meta(direction)

    requested_steps = _safe_int(steps, 0)
    if requested_steps < 1:
        raise ValueError(f'移动{meta["unit"]}数必须大于 0')

    payload = copy.deepcopy(
        _snapshot_payload(classroom, include_students=False, include_constraints=True)
    )

    rows = int(payload.get('classroom', {}).get('rows') or classroom.rows)
    cols = int(payload.get('classroom', {}).get('cols') or classroom.cols)
    axis_field = meta['axis_field']
    size = cols if meta['size_key'] == 'cols' else rows
    normalized_steps = requested_steps % size if size > 0 else requested_steps
    if normalized_steps == 0:
        normalized_steps = size if size > 0 else normalized_steps

    excluded_axis_values = set()
    if axis_field == 'row':
        for seat in payload.get('seats', []):
            if seat.get('cell_type') != SeatCellType.PODIUM:
                continue
            podium_row = _safe_int(seat.get('row'), 0)
            if podium_row > 0:
                excluded_axis_values.add(podium_row)

    axis_map = None
    if excluded_axis_values:
        axis_map = _build_shift_axis_map(size, excluded_axis_values, meta['expand'], requested_steps)

    delta = normalized_steps if meta['expand'] else -normalized_steps

    shifted_seats = []
    for seat in payload.get('seats', []):
        item = dict(seat)
        current_axis_value = int(item.get(axis_field) or 0)
        if current_axis_value > 0:
            if axis_map is not None:
                item[axis_field] = int(axis_map.get(current_axis_value, current_axis_value))
            else:
                item[axis_field] = (current_axis_value - 1 + delta) % size + 1
        shifted_seats.append(item)

    payload['seats'] = _sort_layout_seats(shifted_seats)

    constraint_items = [
        dict(item, _classroom_size=size)
        for item in payload.get('constraints', [])
    ]
    if axis_map is not None:
        payload['constraints'] = _shift_layout_constraints_by_axis_map(
            constraint_items,
            axis_field,
            axis_map,
        )
    else:
        payload['constraints'] = _shift_layout_constraints(
            constraint_items,
            axis_field,
            delta,
            meta['unit'],
            meta['action_label'],
        )

    return payload


def _build_snapshot_seat_coord_map(payload):
    seat_map = {}
    for seat in payload.get('seats', []):
        try:
            row = int(seat.get('row') or 0)
            col = int(seat.get('col') or 0)
        except (TypeError, ValueError, AttributeError):
            continue
        if row < 1 or col < 1:
            continue
        seat_map[(row, col)] = seat
    return seat_map


def _find_snapshot_student_coord(payload, student_pk):
    if student_pk in (None, ''):
        return None
    try:
        target_student_pk = int(student_pk)
    except (TypeError, ValueError):
        return None

    for seat in payload.get('seats', []):
        try:
            current_student_pk = int(seat.get('student_pk'))
            row = int(seat.get('row') or 0)
            col = int(seat.get('col') or 0)
        except (TypeError, ValueError, AttributeError):
            continue
        if current_student_pk == target_student_pk and row > 0 and col > 0:
            return (row, col)
    return None


def _swap_snapshot_seat_payload(source_seat, target_seat):
    coord_keys = {'row', 'col'}
    source_payload = {key: value for key, value in source_seat.items() if key not in coord_keys}
    target_payload = {key: value for key, value in target_seat.items() if key not in coord_keys}

    for key in list(source_seat.keys()):
        if key not in coord_keys:
            source_seat.pop(key, None)
    for key in list(target_seat.keys()):
        if key not in coord_keys:
            target_seat.pop(key, None)

    source_seat.update(target_payload)
    target_seat.update(source_payload)


def _restore_snapshot_constraints_for_students(before_data, after_data, student_ids):
    if not student_ids:
        return
    before_constraints = before_data.get('constraints') or []
    after_constraints = after_data.get('constraints') or []
    for index, before_item in enumerate(before_constraints):
        if index >= len(after_constraints):
            break
        if before_item.get('student_pk') in student_ids:
            after_constraints[index] = dict(before_item)


def _pin_snapshot_cell_type_positions(before_data, after_data, cell_type):
    before_seat_map = _build_snapshot_seat_coord_map(before_data)
    after_seat_map = _build_snapshot_seat_coord_map(after_data)

    before_coords = {
        coord for coord, item in before_seat_map.items()
        if item.get('cell_type') == cell_type
    }
    after_coords = {
        coord for coord, item in after_seat_map.items()
        if item.get('cell_type') == cell_type
    }

    missing_coords = sorted(before_coords - after_coords)
    extra_coords = sorted(after_coords - before_coords)
    for original_coord, current_coord in zip(missing_coords, extra_coords):
        source_seat = after_seat_map.get(original_coord)
        target_seat = after_seat_map.get(current_coord)
        if not source_seat or not target_seat:
            continue
        _swap_snapshot_seat_payload(source_seat, target_seat)


def _resolve_auto_podium_guardian_ids_from_snapshot(payload):
    if not isinstance(payload, dict):
        return {'left': None, 'right': None}

    classroom_data = payload.get('classroom') or {}
    seat_map = _build_snapshot_seat_coord_map(payload)
    podium_rows = defaultdict(list)
    max_col = _safe_int(classroom_data.get('cols'), 0)

    for seat in payload.get('seats', []):
        try:
            row = int(seat.get('row') or 0)
            col = int(seat.get('col') or 0)
        except (TypeError, ValueError, AttributeError):
            continue
        if row < 1 or col < 1:
            continue
        max_col = max(max_col, col)
        if seat.get('cell_type') == SeatCellType.PODIUM:
            podium_rows[row].append(col)

    for row in sorted(podium_rows.keys()):
        podium_cols = sorted(set(podium_rows[row]))
        if not podium_cols:
            continue

        left_student_id = None
        for col in range(podium_cols[0] - 1, 0, -1):
            seat = seat_map.get((row, col))
            if seat and seat.get('cell_type') == SeatCellType.SEAT:
                left_student_id = _safe_int(seat.get('student_pk'), 0) or None
                break

        right_student_id = None
        for col in range(podium_cols[-1] + 1, max_col + 1):
            seat = seat_map.get((row, col))
            if seat and seat.get('cell_type') == SeatCellType.SEAT:
                right_student_id = _safe_int(seat.get('student_pk'), 0) or None
                break

        if left_student_id or right_student_id:
            return {
                'left': left_student_id,
                'right': right_student_id,
            }

    return {'left': None, 'right': None}


def _get_snapshot_podium_guardian_ids(payload):
    auto_guardian_ids = _resolve_auto_podium_guardian_ids_from_snapshot(payload)
    if auto_guardian_ids.get('left') or auto_guardian_ids.get('right'):
        return auto_guardian_ids

    classroom_data = payload.get('classroom') or {}
    return {
        'left': _safe_int(classroom_data.get('left_guardian_student_pk'), 0) or None,
        'right': _safe_int(classroom_data.get('right_guardian_student_pk'), 0) or None,
    }


def _get_fixed_student_ids_from_snapshot(payload):
    fixed_student_ids = set()
    for item in payload.get('constraints') or []:
        if str(item.get('constraint_type') or '') != SeatConstraint.ConstraintType.MUST_SEAT:
            continue
        if not _parse_bool(item.get('enabled') if isinstance(item, dict) else True):
            continue
        if str(item.get('note') or '').strip() != FIXED_SEAT_NOTE_MARKER:
            continue
        student_pk = _safe_int(item.get('student_pk'), 0) if isinstance(item, dict) else 0
        if student_pk > 0:
            fixed_student_ids.add(student_pk)
    return fixed_student_ids


def _pin_podium_guardians_in_shift_payload(before_data, after_data):
    if not isinstance(before_data, dict) or not isinstance(after_data, dict):
        return after_data

    _pin_snapshot_cell_type_positions(before_data, after_data, SeatCellType.PODIUM)

    pinned_candidates = []
    guardian_ids = _get_snapshot_podium_guardian_ids(before_data)
    for student_pk in guardian_ids.values():
        normalized_pk = _safe_int(student_pk, 0)
        if normalized_pk > 0 and normalized_pk not in pinned_candidates:
            pinned_candidates.append(normalized_pk)

    for student_pk in sorted(_get_fixed_student_ids_from_snapshot(before_data)):
        if student_pk not in pinned_candidates:
            pinned_candidates.append(student_pk)

    if not pinned_candidates:
        return after_data

    pinned_student_ids = set()
    after_seat_map = _build_snapshot_seat_coord_map(after_data)
    for student_pk in pinned_candidates:
        original_coord = _find_snapshot_student_coord(before_data, student_pk)
        current_coord = _find_snapshot_student_coord(after_data, student_pk)
        if not original_coord or not current_coord or original_coord == current_coord:
            continue
        source_seat = after_seat_map.get(original_coord)
        target_seat = after_seat_map.get(current_coord)
        if not source_seat or not target_seat:
            continue
        _swap_snapshot_seat_payload(source_seat, target_seat)
        pinned_student_ids.add(student_pk)

    _restore_snapshot_constraints_for_students(before_data, after_data, pinned_student_ids)
    return after_data


def _snapshot_payload(classroom, include_students=True, include_constraints=True):
    seats = list(classroom.seats.select_related('student', 'group'))
    groups = list(classroom.groups.all())
    students = list(classroom.students.all())
    constraints = list(classroom.constraints.all())
    student_tags = list(classroom.student_tags.all())
    student_tag_memberships = list(StudentTagMembership.objects.filter(classroom=classroom))
    student_tag_rules = list(classroom.student_tag_rules.all())

    data = {
        'meta': {
            'app': '不想排座位',
            'version': '1.0',
            'exported_at': timezone.now().isoformat()
        },
        'classroom': {
            'name': classroom.name,
            'rows': classroom.rows,
            'cols': classroom.cols,
            'left_guardian_student_pk': classroom.left_guardian_id,
            'left_guardian_student_id': classroom.left_guardian.student_id if classroom.left_guardian_id and classroom.left_guardian else None,
            'left_guardian_student_name': classroom.left_guardian.name if classroom.left_guardian_id and classroom.left_guardian else None,
            'right_guardian_student_pk': classroom.right_guardian_id,
            'right_guardian_student_id': classroom.right_guardian.student_id if classroom.right_guardian_id and classroom.right_guardian else None,
            'right_guardian_student_name': classroom.right_guardian.name if classroom.right_guardian_id and classroom.right_guardian else None,
        },
        'seats': [
            {
                'row': seat.row,
                'col': seat.col,
                'cell_type': seat.cell_type,
                'student_pk': seat.student.pk if seat.student else None,
                'student_id': seat.student.student_id if seat.student else None,
                'student_name': seat.student.name if seat.student else None,
                'group_name': seat.group.name if seat.group else None
            }
            for seat in seats
        ],
        'groups': [
            {
                'name': group.name,
                'order': group.order
            }
            for group in groups
        ],
        'student_tags': [
            {
                'tag_pk': tag.pk,
                'name': tag.name,
                'color': tag.color,
                'description': tag.description,
                'sort_order': tag.sort_order,
            }
            for tag in student_tags
        ]
    }

    if include_students:
        data['students'] = [
            {
                'name': student.name,
                'student_id': student.student_id,
                'gender': student.gender,
                'score': student.score
            }
            for student in students
        ]
        data['student_tag_memberships'] = [
            {
                'student_pk': membership.student_id,
                'student_id': membership.student.student_id if membership.student else None,
                'student_name': membership.student.name if membership.student else None,
                'tag_pk': membership.tag_id,
                'tag_name': membership.tag.name if membership.tag else None,
                'note': membership.note,
            }
            for membership in student_tag_memberships
        ]

    if include_constraints:
        data['constraints'] = [
            {
                'constraint_type': c.constraint_type,
                'student_pk': c.student.pk,
                'student_id': c.student.student_id,
                'student_name': c.student.name,
                'target_student_pk': c.target_student.pk if c.target_student else None,
                'target_student_id': c.target_student.student_id if c.target_student else None,
                'target_student_name': c.target_student.name if c.target_student else None,
                'row': c.row,
                'col': c.col,
                'distance': c.distance,
                'enabled': c.enabled,
                'note': c.note
            }
            for c in constraints
        ]
        data['student_tag_rules'] = [
            {
                'tag_rule_pk': rule.pk,
                'tag_pk': rule.tag_id,
                'tag_name': rule.tag.name if rule.tag else None,
                'rule_type': rule.rule_type,
                'row_min': rule.row_min,
                'row_max': rule.row_max,
                'col_min': rule.col_min,
                'col_max': rule.col_max,
                'distance': rule.distance,
                'enabled': rule.enabled,
                'priority': rule.priority,
                'note': rule.note,
            }
            for rule in student_tag_rules
        ]

    return data


def _find_student(classroom, payload):
    if payload.get('student_pk'):
        student = classroom.students.filter(pk=payload['student_pk']).first()
        if student:
            return student
    student_id = payload.get('student_id')
    name = payload.get('student_name') or payload.get('name')
    if student_id:
        student = classroom.students.filter(student_id=student_id).first()
        if student:
            return student
    if name:
        return classroom.students.filter(name=name).first()
    return None


def _find_student_tag(classroom, payload):
    tag_pk = payload.get('tag_pk') or payload.get('tag_id')
    if tag_pk:
        tag = classroom.student_tags.filter(pk=tag_pk).first()
        if tag:
            return tag
    name = payload.get('tag_name') or payload.get('name')
    if name:
        return classroom.student_tags.filter(name=str(name).strip()).first()
    return None


def _apply_layout_data(classroom, data, replace_students=False):
    with transaction.atomic():
        classroom_data = data.get('classroom', {})
        rows = int(classroom_data.get('rows', classroom.rows))
        cols = int(classroom_data.get('cols', classroom.cols))
        _sync_seats(classroom, rows, cols)

        if replace_students:
            SeatConstraint.objects.filter(classroom=classroom).delete()
            StudentTagRule.objects.filter(classroom=classroom).delete()
            StudentTagMembership.objects.filter(classroom=classroom).delete()
            StudentTag.objects.filter(classroom=classroom).delete()
            SeatGroup.objects.filter(classroom=classroom).delete()
            Student.objects.filter(classroom=classroom).delete()

        group_map = {}
        for group_data in data.get('groups', []):
            name = str(group_data.get('name', '')).strip()
            if not name:
                continue
            group, _ = SeatGroup.objects.get_or_create(
                classroom=classroom,
                name=name,
                defaults={'order': int(group_data.get('order', 0))}
            )
            group.order = int(group_data.get('order', group.order))
            group.save(update_fields=['order'])
            group_map[name] = group

        tag_map = {}
        if data.get('student_tags') is not None:
            for tag_data in data.get('student_tags', []):
                name = _normalize_tag_name(tag_data.get('name'))
                if not name:
                    continue
                tag = None
                if not replace_students:
                    tag = classroom.student_tags.filter(name=name).first()
                if not tag:
                    tag = StudentTag(classroom=classroom)
                tag.name = name
                tag.color = _normalize_tag_color(tag_data.get('color'))
                tag.description = str(tag_data.get('description') or '').strip()[:160]
                tag.sort_order = max(0, _safe_int(tag_data.get('sort_order'), 0))
                tag.save()
                tag_map[name] = tag
                if tag_data.get('tag_pk'):
                    tag_map[str(tag_data.get('tag_pk'))] = tag

        if data.get('students') is not None:
            for student_data in data.get('students', []):
                name = str(student_data.get('name', '')).strip()
                if not name:
                    continue
                student_id = str(student_data.get('student_id') or '').strip()
                student = None
                if not replace_students:
                    if student_id:
                        student = classroom.students.filter(student_id=student_id).first()
                    if not student:
                        student = classroom.students.filter(name=name).first()
                if not student:
                    student = Student(classroom=classroom)
                student.name = name
                student.student_id = student_id
                student.gender = student_data.get('gender') or None
                student.score = float(student_data.get('score') or 0)
                student.save()

        if data.get('student_tag_memberships') is not None:
            StudentTagMembership.objects.filter(classroom=classroom).delete()
            for membership_data in data.get('student_tag_memberships', []):
                student = _find_student(classroom, membership_data)
                tag = _find_student_tag(classroom, membership_data)
                if not tag:
                    tag_name = _normalize_tag_name(membership_data.get('tag_name'))
                    tag = tag_map.get(tag_name)
                if not student or not tag:
                    continue
                StudentTagMembership.objects.get_or_create(
                    classroom=classroom,
                    student=student,
                    tag=tag,
                    defaults={'note': str(membership_data.get('note') or '')[:120]},
                )

        left_guardian = _find_student(classroom, {
            'student_pk': classroom_data.get('left_guardian_student_pk'),
            'student_id': classroom_data.get('left_guardian_student_id'),
            'student_name': classroom_data.get('left_guardian_student_name'),
        })
        right_guardian = _find_student(classroom, {
            'student_pk': classroom_data.get('right_guardian_student_pk'),
            'student_id': classroom_data.get('right_guardian_student_id'),
            'student_name': classroom_data.get('right_guardian_student_name'),
        })
        _apply_podium_guards(classroom, left_guardian, right_guardian)

        seats = list(classroom.seats.select_related('student', 'group'))
        seat_map = _build_seat_map(seats)
        for seat in seats:
            seat.student = None
            seat.group = None
            seat.cell_type = seat.cell_type or SeatCellType.SEAT
            seat.save(update_fields=['student', 'group', 'cell_type'])

        for seat_data in data.get('seats', []):
            row = int(seat_data.get('row', 0))
            col = int(seat_data.get('col', 0))
            seat = seat_map.get((row, col))
            if not seat:
                continue
            cell_type = seat_data.get('cell_type') or SeatCellType.SEAT
            seat.cell_type = cell_type
            group_name = seat_data.get('group_name')
            if cell_type == SeatCellType.SEAT and group_name:
                seat.group = group_map.get(group_name)
            else:
                seat.group = None
            seat.student = None
            student_payload = {
                'student_pk': seat_data.get('student_pk'),
                'student_id': seat_data.get('student_id'),
                'student_name': seat_data.get('student_name')
            }
            student = _find_student(classroom, student_payload)
            if student and cell_type == SeatCellType.SEAT:
                seat.student = student
            seat.save()

        if data.get('constraints') is not None:
            SeatConstraint.objects.filter(classroom=classroom).delete()
            for cdata in data.get('constraints', []):
                student = _find_student(classroom, cdata)
                if not student:
                    continue
                target_payload = {
                    'student_pk': cdata.get('target_student_pk'),
                    'student_id': cdata.get('target_student_id'),
                    'student_name': cdata.get('target_student_name')
                }
                target_student = _find_student(classroom, target_payload)
                SeatConstraint.objects.create(
                    classroom=classroom,
                    constraint_type=cdata.get('constraint_type'),
                    student=student,
                    target_student=target_student,
                    row=cdata.get('row') or None,
                    col=cdata.get('col') or None,
                    distance=int(cdata.get('distance') or 1),
                    enabled=bool(cdata.get('enabled', True)),
                    note=str(cdata.get('note') or '')
                )

        if data.get('student_tag_rules') is not None:
            StudentTagRule.objects.filter(classroom=classroom).delete()
            for rule_data in data.get('student_tag_rules', []):
                tag = _find_student_tag(classroom, rule_data)
                if not tag:
                    continue
                try:
                    cleaned = normalize_tag_rule_payload(classroom, {
                        'tag_id': tag.pk,
                        'rule_type': rule_data.get('rule_type'),
                        'row_min': rule_data.get('row_min'),
                        'row_max': rule_data.get('row_max'),
                        'col_min': rule_data.get('col_min'),
                        'col_max': rule_data.get('col_max'),
                        'distance': rule_data.get('distance'),
                        'enabled': rule_data.get('enabled'),
                        'priority': rule_data.get('priority'),
                        'note': rule_data.get('note'),
                    })
                except ConstraintServiceError:
                    continue
                StudentTagRule.objects.create(classroom=classroom, **cleaned)


def _serialize_history_datetime(value):
    if not value:
        return ''
    return value.isoformat()


def _restore_history_datetime(model_cls, pk, value, field_name='created_at'):
    dt = parse_datetime(str(value or '').strip()) if value else None
    if not dt:
        return
    if timezone.is_naive(dt):
        dt = timezone.make_aware(dt, timezone.get_current_timezone())
    model_cls.objects.filter(pk=pk).update(**{field_name: dt})


def _capture_history_state(classroom):
    seats = list(classroom.seats.select_related('student', 'group').all())
    students = list(classroom.students.all().order_by('pk'))
    groups = list(classroom.groups.all().order_by('pk'))
    constraints = list(classroom.constraints.all().order_by('pk'))
    student_tags = list(classroom.student_tags.all().order_by('pk'))
    student_tag_memberships = list(StudentTagMembership.objects.filter(classroom=classroom).order_by('pk'))
    student_tag_rules = list(classroom.student_tag_rules.all().order_by('pk'))
    layout_snapshots = list(classroom.layout_snapshots.all().order_by('pk'))
    return {
        'classroom': {
            'pk': classroom.pk,
            'name': classroom.name,
            'rows': classroom.rows,
            'cols': classroom.cols,
            'left_guardian_student_pk': classroom.left_guardian_id,
            'right_guardian_student_pk': classroom.right_guardian_id,
            'created_at': _serialize_history_datetime(classroom.created_at),
        },
        'students': [
            {
                'pk': student.pk,
                'name': student.name,
                'student_id': student.student_id,
                'gender': student.gender,
                'score': student.score,
            }
            for student in students
        ],
        'groups': [
            {
                'pk': group.pk,
                'name': group.name,
                'order': group.order,
                'leader_student_pk': group.leader_id,
                'created_at': _serialize_history_datetime(group.created_at),
            }
            for group in groups
        ],
        'seats': [
            {
                'row': seat.row,
                'col': seat.col,
                'cell_type': seat.cell_type,
                'student_pk': seat.student_id,
                'group_pk': seat.group_id,
            }
            for seat in seats
        ],
        'constraints': [
            {
                'pk': constraint.pk,
                'constraint_type': constraint.constraint_type,
                'student_pk': constraint.student_id,
                'target_student_pk': constraint.target_student_id,
                'row': constraint.row,
                'col': constraint.col,
                'distance': constraint.distance,
                'enabled': constraint.enabled,
                'note': constraint.note,
                'created_at': _serialize_history_datetime(constraint.created_at),
            }
            for constraint in constraints
        ],
        'student_tags': [
            {
                'pk': tag.pk,
                'name': tag.name,
                'color': tag.color,
                'description': tag.description,
                'sort_order': tag.sort_order,
                'created_at': _serialize_history_datetime(tag.created_at),
                'updated_at': _serialize_history_datetime(tag.updated_at),
            }
            for tag in student_tags
        ],
        'student_tag_memberships': [
            {
                'pk': membership.pk,
                'student_pk': membership.student_id,
                'tag_pk': membership.tag_id,
                'note': membership.note,
                'created_at': _serialize_history_datetime(membership.created_at),
            }
            for membership in student_tag_memberships
        ],
        'student_tag_rules': [
            {
                'pk': rule.pk,
                'tag_pk': rule.tag_id,
                'rule_type': rule.rule_type,
                'row_min': rule.row_min,
                'row_max': rule.row_max,
                'col_min': rule.col_min,
                'col_max': rule.col_max,
                'distance': rule.distance,
                'enabled': rule.enabled,
                'priority': rule.priority,
                'note': rule.note,
                'created_at': _serialize_history_datetime(rule.created_at),
                'updated_at': _serialize_history_datetime(rule.updated_at),
            }
            for rule in student_tag_rules
        ],
        'layout_snapshots': [
            {
                'pk': snapshot.pk,
                'name': snapshot.name,
                'data': copy.deepcopy(snapshot.data),
                'created_at': _serialize_history_datetime(snapshot.created_at),
            }
            for snapshot in layout_snapshots
        ],
    }


def _serialize_future_mode_config(classroom):
    config = FutureModeConfig.objects.filter(classroom=classroom).first()
    if not config:
        return None
    return {
        'api_key': str(config.api_key or ''),
        'base_url': str(config.base_url or ''),
        'model': str(config.model or ''),
        'thinking_mode': str(config.thinking_mode or ''),
        'created_at': _serialize_history_datetime(config.created_at),
        'updated_at': _serialize_history_datetime(config.updated_at),
    }


def _serialize_ai_conversations_for_export(classroom):
    conversations = list(
        classroom.ai_conversations
        .all()
        .prefetch_related('messages')
        .order_by('created_at', 'pk')
    )
    data = []
    for conversation in conversations:
        data.append({
            'session_key': str(conversation.session_key or ''),
            'title': str(conversation.title or ''),
            'last_mode': str(conversation.last_mode or ''),
            'last_response_id': str(conversation.last_response_id or ''),
            'created_at': _serialize_history_datetime(conversation.created_at),
            'updated_at': _serialize_history_datetime(conversation.updated_at),
            'messages': [
                {
                    'role': str(message.role or ''),
                    'content': str(message.content or ''),
                    'payload': copy.deepcopy(message.payload if isinstance(message.payload, dict) else {}),
                    'created_at': _serialize_history_datetime(message.created_at),
                }
                for message in conversation.messages.all()
            ],
        })
    return data


def _serialize_classroom_history_for_export(classroom):
    return [
        {
            'action_type': str(entry.action_type or ''),
            'payload': copy.deepcopy(entry.payload if isinstance(entry.payload, dict) else {}),
            'is_applied': bool(entry.is_applied),
            'created_at': _serialize_history_datetime(entry.created_at),
        }
        for entry in _get_history_queryset(classroom.pk)
    ]


def _serialize_seats_file_bundle(classroom):
    data = _snapshot_payload(classroom, include_students=True, include_constraints=True)
    data['meta'] = {
        **copy.deepcopy(data.get('meta') or {}),
        'version': '2.0',
        'schema': 'full',
    }
    data['current_state'] = _capture_history_state(classroom)
    data['history'] = {
        'entries': _serialize_classroom_history_for_export(classroom),
    }
    data['future_mode_config'] = _serialize_future_mode_config(classroom)
    data['ai_conversations'] = _serialize_ai_conversations_for_export(classroom)
    return data


def _encode_history_blob(data):
    raw = json.dumps(data, ensure_ascii=False, separators=(',', ':')).encode('utf-8')
    return base64.b64encode(zlib.compress(raw, level=6)).decode('ascii')


def _decode_history_blob(blob):
    if not blob:
        return {}
    raw = zlib.decompress(base64.b64decode(str(blob).encode('ascii')))
    return json.loads(raw.decode('utf-8'))


def _collect_state_reference_pks(state, collector=None):
    collector = collector or {
        'student': set(),
        'group': set(),
        'constraint': set(),
        'snapshot': set(),
        'tag': set(),
        'tag_rule': set(),
    }
    if not isinstance(state, dict):
        return collector

    classroom_data = state.get('classroom') or {}
    for key in ('left_guardian_student_pk', 'right_guardian_student_pk'):
        try:
            collector['student'].add(int(classroom_data.get(key)))
        except (TypeError, ValueError):
            pass

    for item in state.get('students', []):
        try:
            collector['student'].add(int(item.get('pk')))
        except (TypeError, ValueError):
            continue

    for item in state.get('groups', []):
        try:
            collector['group'].add(int(item.get('pk')))
        except (TypeError, ValueError):
            pass
        try:
            collector['student'].add(int(item.get('leader_student_pk')))
        except (TypeError, ValueError):
            pass

    for item in state.get('constraints', []):
        try:
            collector['constraint'].add(int(item.get('pk')))
        except (TypeError, ValueError):
            pass
        try:
            collector['student'].add(int(item.get('student_pk')))
        except (TypeError, ValueError):
            pass
        try:
            collector['student'].add(int(item.get('target_student_pk')))
        except (TypeError, ValueError):
            pass

    for item in state.get('student_tags', []):
        try:
            collector['tag'].add(int(item.get('pk')))
        except (TypeError, ValueError):
            pass

    for item in state.get('student_tag_memberships', []):
        try:
            collector['student'].add(int(item.get('student_pk')))
        except (TypeError, ValueError):
            pass
        try:
            collector['tag'].add(int(item.get('tag_pk')))
        except (TypeError, ValueError):
            pass

    for item in state.get('student_tag_rules', []):
        try:
            collector['tag_rule'].add(int(item.get('pk')))
        except (TypeError, ValueError):
            pass
        try:
            collector['tag'].add(int(item.get('tag_pk')))
        except (TypeError, ValueError):
            pass

    for item in state.get('layout_snapshots', []):
        try:
            collector['snapshot'].add(int(item.get('pk')))
        except (TypeError, ValueError):
            continue

    return collector


def _extract_export_history_entries(data):
    history = data.get('history')
    if isinstance(history, dict):
        entries = history.get('entries')
        if isinstance(entries, list):
            return entries
    if isinstance(history, list):
        return history
    return []


def _collect_bundle_reference_pks(data):
    collector = _collect_state_reference_pks(data.get('current_state'))
    for entry in _extract_export_history_entries(data):
        if not isinstance(entry, dict):
            continue
        payload = entry.get('payload')
        if not isinstance(payload, dict):
            continue
        for key in ('before_state', 'after_state'):
            blob = payload.get(key)
            if not blob:
                continue
            try:
                state = _decode_history_blob(blob)
            except Exception:
                continue
            _collect_state_reference_pks(state, collector)
    return collector


def _build_entity_pk_mapping(model_cls, source_pks, classroom=None):
    source_pks = {int(pk) for pk in source_pks if pk not in (None, '')}
    if not source_pks:
        return {}

    queryset = model_cls.objects.all()
    if classroom is not None and any(field.name == 'classroom' for field in model_cls._meta.fields):
        queryset = queryset.exclude(classroom=classroom)

    occupied = set(queryset.filter(pk__in=list(source_pks)).values_list('pk', flat=True))
    current_max = queryset.order_by('-pk').values_list('pk', flat=True).first() or 0
    next_pk = max(current_max, max(source_pks)) + 1
    reserved = set(occupied)
    mapping = {}

    for source_pk in sorted(source_pks):
        if source_pk not in reserved:
            mapping[source_pk] = source_pk
            reserved.add(source_pk)
            continue
        while next_pk in reserved:
            next_pk += 1
        mapping[source_pk] = next_pk
        reserved.add(next_pk)
        next_pk += 1

    return mapping


def _build_seats_file_pk_mappings(classroom, data):
    collector = _collect_bundle_reference_pks(data)
    return {
        'student': _build_entity_pk_mapping(Student, collector['student'], classroom=classroom),
        'group': _build_entity_pk_mapping(SeatGroup, collector['group'], classroom=classroom),
        'constraint': _build_entity_pk_mapping(SeatConstraint, collector['constraint'], classroom=classroom),
        'snapshot': _build_entity_pk_mapping(LayoutSnapshot, collector['snapshot'], classroom=classroom),
        'tag': _build_entity_pk_mapping(StudentTag, collector['tag'], classroom=classroom),
        'tag_rule': _build_entity_pk_mapping(StudentTagRule, collector['tag_rule'], classroom=classroom),
    }


def _remap_scalar_pk(value, mapping):
    if value in (None, ''):
        return value
    try:
        source_pk = int(value)
    except (TypeError, ValueError):
        return value
    return mapping.get(source_pk, source_pk)


def _remap_snapshot_payload_data(data, mappings):
    if not isinstance(data, dict):
        return {}
    payload = copy.deepcopy(data)
    classroom_data = payload.get('classroom') or {}
    classroom_data['left_guardian_student_pk'] = _remap_scalar_pk(classroom_data.get('left_guardian_student_pk'), mappings['student'])
    classroom_data['right_guardian_student_pk'] = _remap_scalar_pk(classroom_data.get('right_guardian_student_pk'), mappings['student'])
    payload['classroom'] = classroom_data
    for item in payload.get('seats', []):
        if not isinstance(item, dict):
            continue
        item['student_pk'] = _remap_scalar_pk(item.get('student_pk'), mappings['student'])
    for item in payload.get('constraints', []):
        if not isinstance(item, dict):
            continue
        item['student_pk'] = _remap_scalar_pk(item.get('student_pk'), mappings['student'])
        item['target_student_pk'] = _remap_scalar_pk(item.get('target_student_pk'), mappings['student'])
    for item in payload.get('student_tags', []):
        if not isinstance(item, dict):
            continue
        item['tag_pk'] = _remap_scalar_pk(item.get('tag_pk'), mappings['tag'])
    for item in payload.get('student_tag_memberships', []):
        if not isinstance(item, dict):
            continue
        item['student_pk'] = _remap_scalar_pk(item.get('student_pk'), mappings['student'])
        item['tag_pk'] = _remap_scalar_pk(item.get('tag_pk'), mappings['tag'])
    for item in payload.get('student_tag_rules', []):
        if not isinstance(item, dict):
            continue
        item['tag_rule_pk'] = _remap_scalar_pk(item.get('tag_rule_pk'), mappings['tag_rule'])
        item['tag_pk'] = _remap_scalar_pk(item.get('tag_pk'), mappings['tag'])
    return payload


def _remap_history_state(state, mappings, classroom_pk):
    if not isinstance(state, dict):
        return {}
    payload = copy.deepcopy(state)
    classroom_data = payload.get('classroom') or {}
    classroom_data['pk'] = classroom_pk
    classroom_data['left_guardian_student_pk'] = _remap_scalar_pk(classroom_data.get('left_guardian_student_pk'), mappings['student'])
    classroom_data['right_guardian_student_pk'] = _remap_scalar_pk(classroom_data.get('right_guardian_student_pk'), mappings['student'])
    payload['classroom'] = classroom_data

    for item in payload.get('students', []):
        if not isinstance(item, dict):
            continue
        item['pk'] = _remap_scalar_pk(item.get('pk'), mappings['student'])

    for item in payload.get('groups', []):
        if not isinstance(item, dict):
            continue
        item['pk'] = _remap_scalar_pk(item.get('pk'), mappings['group'])
        item['leader_student_pk'] = _remap_scalar_pk(item.get('leader_student_pk'), mappings['student'])

    for item in payload.get('seats', []):
        if not isinstance(item, dict):
            continue
        item['student_pk'] = _remap_scalar_pk(item.get('student_pk'), mappings['student'])
        item['group_pk'] = _remap_scalar_pk(item.get('group_pk'), mappings['group'])

    for item in payload.get('constraints', []):
        if not isinstance(item, dict):
            continue
        item['pk'] = _remap_scalar_pk(item.get('pk'), mappings['constraint'])
        item['student_pk'] = _remap_scalar_pk(item.get('student_pk'), mappings['student'])
        item['target_student_pk'] = _remap_scalar_pk(item.get('target_student_pk'), mappings['student'])

    for item in payload.get('student_tags', []):
        if not isinstance(item, dict):
            continue
        item['pk'] = _remap_scalar_pk(item.get('pk'), mappings['tag'])

    for item in payload.get('student_tag_memberships', []):
        if not isinstance(item, dict):
            continue
        item['student_pk'] = _remap_scalar_pk(item.get('student_pk'), mappings['student'])
        item['tag_pk'] = _remap_scalar_pk(item.get('tag_pk'), mappings['tag'])

    for item in payload.get('student_tag_rules', []):
        if not isinstance(item, dict):
            continue
        item['pk'] = _remap_scalar_pk(item.get('pk'), mappings['tag_rule'])
        item['tag_pk'] = _remap_scalar_pk(item.get('tag_pk'), mappings['tag'])

    for item in payload.get('layout_snapshots', []):
        if not isinstance(item, dict):
            continue
        item['pk'] = _remap_scalar_pk(item.get('pk'), mappings['snapshot'])
        item['data'] = _remap_snapshot_payload_data(item.get('data'), mappings)

    return payload


def _remap_history_payload(payload, mappings, classroom_pk):
    def _walk(value, current_key=''):
        if isinstance(value, dict):
            remapped = {}
            for key, item in value.items():
                if key in {'before_state', 'after_state'} and item:
                    try:
                        state = _decode_history_blob(item)
                    except Exception:
                        remapped[key] = item
                    else:
                        remapped[key] = _encode_history_blob(_remap_history_state(state, mappings, classroom_pk))
                    continue
                if key in {'before_data', 'after_data'} and isinstance(item, dict):
                    remapped[key] = _remap_snapshot_payload_data(item, mappings)
                    continue
                if key == 'classroom_id':
                    remapped[key] = classroom_pk
                    continue
                if key in HISTORY_STUDENT_ID_KEYS:
                    remapped[key] = _remap_scalar_pk(item, mappings['student'])
                    continue
                if key in HISTORY_GROUP_ID_KEYS:
                    remapped[key] = _remap_scalar_pk(item, mappings['group'])
                    continue
                if key in HISTORY_SNAPSHOT_ID_KEYS:
                    remapped[key] = _remap_scalar_pk(item, mappings['snapshot'])
                    continue
                if key in HISTORY_CONSTRAINT_ID_KEYS:
                    remapped[key] = _remap_scalar_pk(item, mappings['constraint'])
                    continue
                if key in HISTORY_TAG_ID_KEYS:
                    remapped[key] = _remap_scalar_pk(item, mappings['tag'])
                    continue
                if key in HISTORY_TAG_RULE_ID_KEYS:
                    remapped[key] = _remap_scalar_pk(item, mappings['tag_rule'])
                    continue
                if key in HISTORY_GROUP_ID_LIST_KEYS and isinstance(item, list):
                    remapped[key] = [_remap_scalar_pk(entry, mappings['group']) for entry in item]
                    continue
                if key in HISTORY_TAG_ID_LIST_KEYS and isinstance(item, list):
                    remapped[key] = [_remap_scalar_pk(entry, mappings['tag']) for entry in item]
                    continue
                remapped[key] = _walk(item, key)
            return remapped
        if isinstance(value, list):
            if current_key in HISTORY_GROUP_ID_LIST_KEYS:
                return [_remap_scalar_pk(item, mappings['group']) for item in value]
            if current_key in HISTORY_TAG_ID_LIST_KEYS:
                return [_remap_scalar_pk(item, mappings['tag']) for item in value]
            return [_walk(item) for item in value]
        return value

    if not isinstance(payload, dict):
        return {}
    return _walk(copy.deepcopy(payload))


def _build_history_snapshot_action(before_state, after_state, action_type, extra=None):
    payload = {
        'before_state': _encode_history_blob(before_state),
        'after_state': _encode_history_blob(after_state),
    }
    if extra:
        payload.update(copy.deepcopy(extra))
    payload['type'] = str(action_type or 'history_snapshot')
    payload['history_mode'] = 'snapshot'
    return payload


def _push_snapshot_action(request, classroom, before_state, action_type, extra=None):
    after_state = _capture_history_state(classroom)
    if before_state == after_state:
        return None
    action = _build_history_snapshot_action(before_state, after_state, action_type, extra=extra)
    _push_action(request, classroom.pk, action)
    return action


def _restore_history_state(classroom, state):
    state = state or {}
    classroom_data = state.get('classroom') or {}
    target_name = str(classroom_data.get('name') or classroom.name)
    target_rows = int(classroom_data.get('rows') or classroom.rows)
    target_cols = int(classroom_data.get('cols') or classroom.cols)

    with transaction.atomic():
        if classroom.name != target_name:
            classroom.name = target_name
            classroom.save(update_fields=['name'])

        _sync_seats(classroom, target_rows, target_cols)
        _restore_history_datetime(Classroom, classroom.pk, classroom_data.get('created_at'))

        classroom.seats.update(student=None, group=None, cell_type=SeatCellType.SEAT)
        classroom.constraints.all().delete()
        classroom.student_tag_rules.all().delete()
        StudentTagMembership.objects.filter(classroom=classroom).delete()
        classroom.student_tags.all().delete()
        classroom.layout_snapshots.all().delete()
        classroom.groups.update(leader=None)
        classroom.groups.all().delete()
        classroom.students.all().delete()

        student_map = {}
        for item in state.get('students', []):
            pk = int(item.get('pk'))
            student = Student(
                pk=pk,
                classroom=classroom,
                name=str(item.get('name') or '').strip(),
                student_id=str(item.get('student_id') or '').strip() or None,
                gender=item.get('gender') or None,
                score=float(item.get('score') or 0),
            )
            student.save(force_insert=True)
            student_map[pk] = student

        tag_map = {}
        for item in state.get('student_tags', []):
            pk = int(item.get('pk'))
            tag = StudentTag(
                pk=pk,
                classroom=classroom,
                name=_normalize_tag_name(item.get('name')),
                color=_normalize_tag_color(item.get('color')),
                description=str(item.get('description') or '')[:160],
                sort_order=max(0, _safe_int(item.get('sort_order'), 0)),
            )
            if not tag.name:
                continue
            tag.save(force_insert=True)
            _restore_history_datetime(StudentTag, pk, item.get('created_at'))
            _restore_history_datetime(StudentTag, pk, item.get('updated_at'), field_name='updated_at')
            tag_map[pk] = tag

        group_map = {}
        for item in state.get('groups', []):
            pk = int(item.get('pk'))
            group = SeatGroup(
                pk=pk,
                classroom=classroom,
                name=str(item.get('name') or '').strip(),
                order=int(item.get('order') or 0),
            )
            group.save(force_insert=True)
            _restore_history_datetime(SeatGroup, pk, item.get('created_at'))
            group_map[pk] = group

        _apply_podium_guards(
            classroom,
            student_map.get(classroom_data.get('left_guardian_student_pk')),
            student_map.get(classroom_data.get('right_guardian_student_pk')),
        )

        seat_map = _build_seat_map(classroom.seats.all())
        for item in state.get('seats', []):
            row = int(item.get('row') or 0)
            col = int(item.get('col') or 0)
            seat = seat_map.get((row, col))
            if not seat:
                continue
            cell_type = item.get('cell_type') or SeatCellType.SEAT
            seat.cell_type = cell_type
            if cell_type == SeatCellType.SEAT:
                seat.student = student_map.get(item.get('student_pk'))
                seat.group = group_map.get(item.get('group_pk'))
            else:
                seat.student = None
                seat.group = None
            seat.save(update_fields=['cell_type', 'student', 'group'])

        for item in state.get('groups', []):
            group = group_map.get(item.get('pk'))
            leader = student_map.get(item.get('leader_student_pk'))
            if not group or not leader:
                continue
            in_group = classroom.seats.filter(
                row__gte=1,
                cell_type=SeatCellType.SEAT,
                group_id=group.pk,
                student_id=leader.pk,
            ).exists()
            group.leader = leader if in_group else None
            group.save(update_fields=['leader'])

        for item in state.get('constraints', []):
            pk = int(item.get('pk'))
            constraint = SeatConstraint(
                pk=pk,
                classroom=classroom,
                constraint_type=item.get('constraint_type'),
                student=student_map.get(item.get('student_pk')),
                target_student=student_map.get(item.get('target_student_pk')),
                row=item.get('row') or None,
                col=item.get('col') or None,
                distance=int(item.get('distance') or 1),
                enabled=bool(item.get('enabled', True)),
                note=str(item.get('note') or ''),
            )
            if not constraint.student_id:
                continue
            constraint.save(force_insert=True)
            _restore_history_datetime(SeatConstraint, pk, item.get('created_at'))

        for item in state.get('student_tag_memberships', []):
            student = student_map.get(item.get('student_pk'))
            tag = tag_map.get(item.get('tag_pk'))
            if not student or not tag:
                continue
            membership = StudentTagMembership(
                classroom=classroom,
                student=student,
                tag=tag,
                note=str(item.get('note') or '')[:120],
            )
            membership_pk = _safe_int(item.get('pk'), 0)
            if membership_pk and not StudentTagMembership.objects.filter(pk=membership_pk).exists():
                membership.pk = membership_pk
                membership.save(force_insert=True)
            else:
                membership.save()
            _restore_history_datetime(StudentTagMembership, membership.pk, item.get('created_at'))

        for item in state.get('student_tag_rules', []):
            pk = int(item.get('pk'))
            tag = tag_map.get(item.get('tag_pk'))
            if not tag:
                continue
            rule = StudentTagRule(
                pk=pk,
                classroom=classroom,
                tag=tag,
                rule_type=item.get('rule_type'),
                row_min=item.get('row_min') or None,
                row_max=item.get('row_max') or None,
                col_min=item.get('col_min') or None,
                col_max=item.get('col_max') or None,
                distance=int(item.get('distance') or 1),
                enabled=bool(item.get('enabled', True)),
                priority=max(0, _safe_int(item.get('priority'), 0)),
                note=str(item.get('note') or '')[:120],
            )
            rule.save(force_insert=True)
            _restore_history_datetime(StudentTagRule, pk, item.get('created_at'))
            _restore_history_datetime(StudentTagRule, pk, item.get('updated_at'), field_name='updated_at')

        for item in state.get('layout_snapshots', []):
            pk = int(item.get('pk'))
            snapshot = LayoutSnapshot(
                pk=pk,
                classroom=classroom,
                name=str(item.get('name') or '').strip(),
                data=copy.deepcopy(item.get('data') or {}),
            )
            snapshot.save(force_insert=True)
            _restore_history_datetime(LayoutSnapshot, pk, item.get('created_at'))

        _normalize_group_leaders(classroom)


def _get_history_queryset(classroom_id):
    return ClassroomHistoryEntry.objects.filter(classroom_id=classroom_id).order_by('pk')


def _get_history(request, classroom_id):
    undo = []
    redo = []
    for payload, is_applied in _get_history_queryset(classroom_id).values_list('payload', 'is_applied'):
        if is_applied:
            undo.append(payload or {})
        else:
            redo.append(payload or {})
    return {'undo': undo, 'redo': redo}


def _push_action(request, classroom_id, action):
    if not action:
        return
    _get_history_queryset(classroom_id).filter(is_applied=False).delete()
    ClassroomHistoryEntry.objects.create(
        classroom_id=classroom_id,
        action_type=str(action.get('type') or '')[:40],
        payload=copy.deepcopy(action),
        is_applied=True,
    )
    total = _get_history_queryset(classroom_id).count()
    overflow = total - CLASSROOM_HISTORY_LIMIT
    if overflow > 0:
        stale_ids = list(_get_history_queryset(classroom_id).values_list('pk', flat=True)[:overflow])
        if stale_ids:
            ClassroomHistoryEntry.objects.filter(pk__in=stale_ids).delete()


def _reset_history(request, classroom_id):
    _get_history_queryset(classroom_id).delete()


def _restore_exported_history(classroom, data, mappings):
    _get_history_queryset(classroom.pk).delete()
    for item in _extract_export_history_entries(data):
        if not isinstance(item, dict):
            continue
        payload = _remap_history_payload(item.get('payload'), mappings, classroom.pk)
        entry = ClassroomHistoryEntry.objects.create(
            classroom=classroom,
            action_type=str(item.get('action_type') or payload.get('type') or '')[:40],
            payload=payload,
            is_applied=bool(item.get('is_applied', True)),
        )
        _restore_history_datetime(ClassroomHistoryEntry, entry.pk, item.get('created_at'))


def _restore_future_mode_config(classroom, data):
    FutureModeConfig.objects.filter(classroom=classroom).delete()
    if not isinstance(data, dict):
        return
    payload = {
        'api_key': str(data.get('api_key') or '').strip(),
        'base_url': str(data.get('base_url') or '').strip(),
        'model': str(data.get('model') or '').strip(),
        'thinking_mode': str(data.get('thinking_mode') or '').strip(),
    }
    if not any(payload.values()):
        return
    config = FutureModeConfig.objects.create(classroom=classroom, **payload)
    _restore_history_datetime(FutureModeConfig, config.pk, data.get('created_at'))
    _restore_history_datetime(FutureModeConfig, config.pk, data.get('updated_at'), field_name='updated_at')


def _restore_ai_conversations(classroom, data, request=None):
    classroom.ai_conversations.all().delete()
    if not isinstance(data, list):
        return
    owner_key = _ensure_session_key(request) if request else ''
    for item in data:
        if not isinstance(item, dict):
            continue
        conversation = AIConversation.objects.create(
            classroom=classroom,
            session_key=owner_key or str(item.get('session_key') or ''),
            title=str(item.get('title') or DEFAULT_AI_CONVERSATION_TITLE)[:120],
            last_mode=str(item.get('last_mode') or '')[:16],
            last_response_id=str(item.get('last_response_id') or '')[:120],
        )
        _restore_history_datetime(AIConversation, conversation.pk, item.get('created_at'))
        _restore_history_datetime(AIConversation, conversation.pk, item.get('updated_at'), field_name='updated_at')

        for message_data in item.get('messages', []):
            if not isinstance(message_data, dict):
                continue
            message = AIConversationMessage.objects.create(
                conversation=conversation,
                role=str(message_data.get('role') or AIConversationMessage.MessageRole.USER),
                content=str(message_data.get('content') or '')[:4000],
                payload=copy.deepcopy(message_data.get('payload') if isinstance(message_data.get('payload'), dict) else {}),
            )
            _restore_history_datetime(AIConversationMessage, message.pk, message_data.get('created_at'))


def _import_seats_file_payload(classroom, data, request=None):
    if isinstance(data.get('current_state'), dict):
        with transaction.atomic():
            mappings = _build_seats_file_pk_mappings(classroom, data)
            state = _remap_history_state(data.get('current_state'), mappings, classroom.pk)
            _restore_history_state(classroom, state)
            _restore_exported_history(classroom, data, mappings)
            _restore_future_mode_config(classroom, data.get('future_mode_config'))
            _restore_ai_conversations(classroom, data.get('ai_conversations'), request=request)
        return 'full'

    _apply_layout_data(classroom, data, replace_students=True)
    return 'legacy'


def _is_ajax_request(request):
    return request.headers.get('x-requested-with') == 'XMLHttpRequest'


def _normalize_group_move_mode(value, default=GROUP_MOVE_MODE_FIXED):
    normalized = str(value or '').strip().lower()
    if normalized in {GROUP_MOVE_MODE_FIXED, GROUP_MOVE_MODE_FOLLOW}:
        return normalized
    return default


def _normalize_group_leaders(classroom, group_ids=None):
    groups = classroom.groups.all()
    if group_ids is not None:
        groups = groups.filter(pk__in=list(group_ids))
    for group in groups:
        if not group.leader_id:
            continue
        still_in_group = group.seats.filter(cell_type=SeatCellType.SEAT, student_id=group.leader_id).exists()
        if not still_in_group:
            group.leader = None
            group.save(update_fields=['leader'])


def _invert_move_action(action):
    return {
        'type': 'move',
        'student_id': action.get('student_id'),
        'from_row': action.get('to_row'),
        'from_col': action.get('to_col'),
        'to_row': action.get('from_row'),
        'to_col': action.get('from_col'),
        'target_student_id': action.get('target_student_id')
    }


def _apply_move_action(classroom, action):
    student = classroom.students.filter(pk=action.get('student_id')).first()
    if not student:
        return False
    from_row = action.get('from_row')
    from_col = action.get('from_col')
    to_row = action.get('to_row')
    to_col = action.get('to_col')
    target_student_id = action.get('target_student_id')

    seat_to = None
    if to_row is not None and to_col is not None:
        seat_to = classroom.seats.filter(row=to_row, col=to_col).first()

    if seat_to and seat_to.cell_type != SeatCellType.SEAT:
        return False

    target_student = None
    if target_student_id:
        target_student = classroom.students.filter(pk=target_student_id).first()

    current_seat = getattr(student, 'assigned_seat', None)
    

    seat_from = None
    if from_row is not None and from_col is not None:
        seat_from = classroom.seats.filter(row=from_row, col=from_col).first()

    with transaction.atomic():
        if seat_from:
            seat_from.student = None
            seat_from.save(update_fields=['student'])
        
        if seat_to:
            seat_to.student = None
            seat_to.save(update_fields=['student'])

        if seat_from and target_student:
             seat_from.student = target_student
             seat_from.save(update_fields=['student'])
        
        if seat_to and student:
            seat_to.student = student
            seat_to.save(update_fields=['student'])
    _normalize_group_leaders(classroom)
    return True


def _apply_move_batch_action(classroom, action, forward=True):
    items = action.get('items', [])
    if not isinstance(items, list):
        return False
    if forward:
        sequence = items
    else:
        sequence = [_invert_move_action(item) for item in reversed(items)]
    success = True
    for item in sequence:
        ok = _apply_move_action(classroom, item)
        if not ok:
            success = False
    return success


def _apply_swap_action(classroom, action):
    student_a = classroom.students.filter(pk=action.get('student_a_id')).first()
    student_b = classroom.students.filter(pk=action.get('student_b_id')).first()
    if not student_a or not student_b:
        return False
    seat_a = getattr(student_a, 'assigned_seat', None)
    seat_b = getattr(student_b, 'assigned_seat', None)
    if not seat_a or not seat_b:
        return False
    _swap_seats(seat_a, seat_b)
    return True


def _apply_cell_type_action(classroom, action, forward=True):
    row = action.get('row')
    col = action.get('col')
    seat = classroom.seats.filter(row=row, col=col).first()
    if not seat:
        return False
    target_type = action.get('after') if forward else action.get('before')
    prev_student_id = action.get('prev_student_id')
    prev_group_id = action.get('prev_group_id')

    seat.cell_type = target_type
    if target_type == SeatCellType.SEAT:
        if prev_group_id:
            seat.group = classroom.groups.filter(pk=prev_group_id).first()
        if prev_student_id:
            seat.student = classroom.students.filter(pk=prev_student_id).first()
    else:
        seat.student = None
        seat.group = None
    seat.save(update_fields=['cell_type', 'student', 'group'])
    return True


def _apply_group_action(classroom, action, forward=True):
    row = action.get('row')
    col = action.get('col')
    seat = classroom.seats.filter(row=row, col=col).first()
    if not seat:
        return False
    target_group_id = action.get('after_group_id') if forward else action.get('before_group_id')
    if target_group_id:
        seat.group = classroom.groups.filter(pk=target_group_id).first()
    else:
        seat.group = None
    seat.save(update_fields=['group'])
    affected_group_ids = {gid for gid in [action.get('before_group_id'), action.get('after_group_id')] if gid}
    if affected_group_ids:
        _normalize_group_leaders(classroom, affected_group_ids)
    return True


def _apply_group_batch_action(classroom, action, forward=True):
    items = action.get('items', [])
    affected_group_ids = set()
    for item in items:
        row = item.get('row')
        col = item.get('col')
        seat = classroom.seats.filter(row=row, col=col).first()
        if not seat:
            continue
        target_group_id = item.get('after_group_id') if forward else item.get('before_group_id')
        if target_group_id:
            seat.group = classroom.groups.filter(pk=target_group_id).first()
            affected_group_ids.add(target_group_id)
        else:
            seat.group = None
        if item.get('before_group_id'):
            affected_group_ids.add(item.get('before_group_id'))
        if item.get('after_group_id'):
            affected_group_ids.add(item.get('after_group_id'))
        seat.save(update_fields=['group'])
    if affected_group_ids:
        _normalize_group_leaders(classroom, affected_group_ids)
    return True


def _apply_seat_layout_action(classroom, action, forward=True):
    items = action.get('items', [])
    if not isinstance(items, list):
        return False

    seat_map = {}
    student_ids = set()
    group_ids = set()
    affected_group_ids = set()

    for item in items:
        try:
            row = int(item.get('row'))
            col = int(item.get('col'))
        except Exception:
            continue
        seat = classroom.seats.filter(row=row, col=col, cell_type=SeatCellType.SEAT).first()
        if not seat:
            continue
        key = (row, col)
        seat_map[key] = {
            'seat': seat,
            'item': item,
        }

        before_student_id = item.get('before_student_id')
        after_student_id = item.get('after_student_id')
        if before_student_id:
            student_ids.add(before_student_id)
        if after_student_id:
            student_ids.add(after_student_id)

        before_group_id = item.get('before_group_id')
        after_group_id = item.get('after_group_id')
        if before_group_id:
            group_ids.add(before_group_id)
            affected_group_ids.add(before_group_id)
        if after_group_id:
            group_ids.add(after_group_id)
            affected_group_ids.add(after_group_id)

    if not seat_map:
        return False

    student_map = {s.pk: s for s in classroom.students.filter(pk__in=list(student_ids))}
    group_map = {g.pk: g for g in classroom.groups.filter(pk__in=list(group_ids))}

    with transaction.atomic():
        for payload in seat_map.values():
            seat = payload['seat']
            seat.student = None
            seat.group = None
            seat.save(update_fields=['student', 'group'])

        for payload in seat_map.values():
            seat = payload['seat']
            item = payload['item']
            student_id = item.get('after_student_id') if forward else item.get('before_student_id')
            group_id = item.get('after_group_id') if forward else item.get('before_group_id')
            seat.student = student_map.get(student_id) if student_id else None
            seat.group = group_map.get(group_id) if group_id else None
            seat.save(update_fields=['student', 'group'])

    if affected_group_ids:
        _normalize_group_leaders(classroom, affected_group_ids)
    return True


def _apply_layout_snapshot_action(classroom, action, forward=True):
    data = action.get('after_data') if forward else action.get('before_data')
    if not isinstance(data, dict):
        return False
    _apply_layout_data(classroom, data, replace_students=False)
    _normalize_group_leaders(classroom)
    return True


def _apply_recorded_history_action(classroom, action, forward=True):
    if not isinstance(action, dict):
        return False
    if action.get('history_mode') == 'snapshot':
        blob = action.get('after_state') if forward else action.get('before_state')
        state = _decode_history_blob(blob)
        _restore_history_state(classroom, state)
        return True
    if action.get('type') == 'move':
        inverse = _invert_move_action(action)
        return _apply_move_action(classroom, action if forward else inverse)
    if action.get('type') == 'move_batch':
        return _apply_move_batch_action(classroom, action, forward=forward)
    if action.get('type') == 'swap':
        return _apply_swap_action(classroom, action)
    if action.get('type') == 'cell_type':
        return _apply_cell_type_action(classroom, action, forward=forward)
    if action.get('type') == 'group':
        return _apply_group_action(classroom, action, forward=forward)
    if action.get('type') == 'group_batch':
        return _apply_group_batch_action(classroom, action, forward=forward)
    if action.get('type') == 'seat_layout_batch':
        return _apply_seat_layout_action(classroom, action, forward=forward)
    if action.get('type') == 'layout_snapshot':
        return _apply_layout_snapshot_action(classroom, action, forward=forward)
    return False


def _constraint_issues(classroom):
    constraints = classroom.constraints.select_related('student', 'target_student').all()
    return constraint_issue_messages(classroom, constraints=constraints)


def _format_issues_preview(issues, limit=3):
    preview = '；'.join(issues[:limit])
    if len(issues) > limit:
        preview += '；...'
    return preview


def _layout_hard_issues(classroom):
    issues = []
    unseated_count = classroom.students.filter(assigned_seat__isnull=True).count()
    if unseated_count:
        issues.append(f"当前有 {unseated_count} 名学生未入座")
    issues.extend(_constraint_issues(classroom))
    issues.extend(tag_rule_issue_messages(classroom))
    return issues


def _distance(seat_a, seat_b):
    if not seat_a or not seat_b:
        return 10 ** 9
    return abs(seat_a.row - seat_b.row) + abs(seat_a.col - seat_b.col)


def _current_assignments(classroom):
    return {seat.student_id: seat for seat in classroom.seats.select_related('student').filter(student__isnull=False)}


def _candidate_seats(classroom, predicate=None):
    seats = list(classroom.seats.filter(cell_type=SeatCellType.SEAT).order_by('row', 'col'))
    if predicate is None:
        return seats
    return [s for s in seats if predicate(s)]


def _simulate_move_valid(student, target_seat, assignments, maps):
    sid = student.pk
    current = assignments.get(sid)
    occupant = target_seat.student

    if occupant and occupant.pk == sid:
        return True
    if occupant and not current:
        return False

    simulated = dict(assignments)
    simulated[sid] = target_seat
    if occupant and current:
        simulated[occupant.pk] = current

    others_for_student = {k: v for k, v in simulated.items() if k != sid}
    if not _seat_is_valid(student, target_seat, others_for_student, maps):
        return False

    if occupant and current:
        others_for_occupant = {k: v for k, v in simulated.items() if k != occupant.pk}
        if not _seat_is_valid(occupant, current, others_for_occupant, maps):
            return False

    return True


def _pick_best_target(student, candidates, assignments, maps):
    sid = student.pk
    current = assignments.get(sid)
    best = None
    best_score = None
    for seat in candidates:
        if not _simulate_move_valid(student, seat, assignments, maps):
            continue
        occupied_penalty = 3 if seat.student_id else 0
        score = _distance(current, seat) + occupied_penalty
        if best is None or score < best_score:
            best = seat
            best_score = score
    return best


def _enforce_constraints_by_moves(classroom, max_rounds=6):
    constraints = list(
        classroom.constraints.filter(enabled=True).select_related('student', 'target_student').order_by('created_at', 'pk')
    )
    if not constraints:
        return True

    students = list(classroom.students.all())
    maps = _build_constraint_maps(classroom, students)

    for _ in range(max_rounds):
        if not _constraint_issues(classroom):
            return True
        changed = False

        for c in constraints:
            assignments = _current_assignments(classroom)
            student = c.student
            target_student = c.target_student
            seat = assignments.get(student.pk)
            ctype = c.constraint_type

            if ctype == SeatConstraint.ConstraintType.MUST_SEAT and c.row and c.col:
                target = classroom.seats.filter(row=c.row, col=c.col, cell_type=SeatCellType.SEAT).first()
                if target and (not seat or seat.pk != target.pk):
                    if _simulate_move_valid(student, target, assignments, maps):
                        _perform_move(classroom, student, target)
                        changed = True
                continue

            if ctype == SeatConstraint.ConstraintType.FORBID_SEAT and c.row and c.col:
                if seat and seat.row == c.row and seat.col == c.col:
                    candidates = _candidate_seats(classroom, predicate=lambda s: not (s.row == c.row and s.col == c.col))
                    target = _pick_best_target(student, candidates, assignments, maps)
                    if target:
                        _perform_move(classroom, student, target)
                        changed = True
                continue

            if ctype == SeatConstraint.ConstraintType.MUST_ROW and c.row:
                if not seat or seat.row != c.row:
                    candidates = _candidate_seats(classroom, predicate=lambda s: s.row == c.row)
                    target = _pick_best_target(student, candidates, assignments, maps)
                    if target:
                        _perform_move(classroom, student, target)
                        changed = True
                continue

            if ctype == SeatConstraint.ConstraintType.FORBID_ROW and c.row:
                if seat and seat.row == c.row:
                    candidates = _candidate_seats(classroom, predicate=lambda s: s.row != c.row)
                    target = _pick_best_target(student, candidates, assignments, maps)
                    if target:
                        _perform_move(classroom, student, target)
                        changed = True
                continue

            if ctype == SeatConstraint.ConstraintType.MUST_COL and c.col:
                if not seat or seat.col != c.col:
                    candidates = _candidate_seats(classroom, predicate=lambda s: s.col == c.col)
                    target = _pick_best_target(student, candidates, assignments, maps)
                    if target:
                        _perform_move(classroom, student, target)
                        changed = True
                continue

            if ctype == SeatConstraint.ConstraintType.FORBID_COL and c.col:
                if seat and seat.col == c.col:
                    candidates = _candidate_seats(classroom, predicate=lambda s: s.col != c.col)
                    target = _pick_best_target(student, candidates, assignments, maps)
                    if target:
                        _perform_move(classroom, student, target)
                        changed = True
                continue

            if ctype in [SeatConstraint.ConstraintType.MUST_TOGETHER, SeatConstraint.ConstraintType.FORBID_TOGETHER] and target_student:
                assignments = _current_assignments(classroom)
                seat_a = assignments.get(student.pk)
                seat_b = assignments.get(target_student.pk)
                dist = c.distance or 1
                cur_distance = _distance(seat_a, seat_b)

                if ctype == SeatConstraint.ConstraintType.MUST_TOGETHER:
                    if cur_distance <= dist:
                        continue
                    if seat_b:
                        candidates = _candidate_seats(classroom, predicate=lambda s: _distance(s, seat_b) <= dist)
                        target = _pick_best_target(student, candidates, assignments, maps)
                        if target:
                            _perform_move(classroom, student, target)
                            changed = True
                            continue
                    assignments = _current_assignments(classroom)
                    seat_a = assignments.get(student.pk)
                    if seat_a:
                        candidates = _candidate_seats(classroom, predicate=lambda s: _distance(s, seat_a) <= dist)
                        target = _pick_best_target(target_student, candidates, assignments, maps)
                        if target:
                            _perform_move(classroom, target_student, target)
                            changed = True
                else:
                    if cur_distance > dist:
                        continue
                    if seat_b:
                        candidates = _candidate_seats(classroom, predicate=lambda s: _distance(s, seat_b) > dist)
                        target = _pick_best_target(student, candidates, assignments, maps)
                        if target:
                            _perform_move(classroom, student, target)
                            changed = True
                            continue
                    assignments = _current_assignments(classroom)
                    seat_a = assignments.get(student.pk)
                    if seat_a:
                        candidates = _candidate_seats(classroom, predicate=lambda s: _distance(s, seat_a) > dist)
                        target = _pick_best_target(target_student, candidates, assignments, maps)
                        if target:
                            _perform_move(classroom, target_student, target)
                            changed = True

        if not changed:
            break

    return not _constraint_issues(classroom)


def _stabilize_layout_with_rules(classroom, request=None, trigger_student_id=None):
    _apply_internal_policy(classroom, request, trigger_student_id=trigger_student_id)
    _enforce_constraints_by_moves(classroom)
    _apply_internal_policy(classroom, request, trigger_student_id=trigger_student_id)
    _normalize_group_leaders(classroom)
    issues = _constraint_issues(classroom)
    issues.extend(tag_rule_issue_messages(classroom))
    return issues


def _filter_internal_issues(issues):
    return issues


def _is_internal_policy_student(student):
    return False


def _evaluate_layout(classroom, request=None):
    issues = []
    seats = list(classroom.seats.select_related('student', 'group'))

    unseated_count = classroom.students.filter(assigned_seat__isnull=True).count()
    if unseated_count:
        issues.append(f"当前有 {unseated_count} 名学生未入座")

    issues.extend(_constraint_issues(classroom))
    _apply_internal_policy(classroom, request)

    groups = list(classroom.groups.all())

    ignore_export = request.session.get(f'ignore_export_{classroom.pk}', False) if request else False
    ungrouped_count = sum(1 for s in seats if s.student_id and not s.group_id)
    if unseated_count == 0 and ungrouped_count == 0 and len(groups) > 0 and not ignore_export:
        issues.append({
            'type': 'export_suggestion',
            'message': '所有学生已入座并分组，建议导出小组作业登记表。',
            'action_label': '立即导出',
            'action_url': reverse('export_group_report', args=[classroom.pk]),
            'ignore_label': '不再提示',
            'ignore_url': f'/classroom/{classroom.pk}/suggestion/dismiss/?type=export'
        })

    if groups:
        group_data = []
        for g in groups:
             seats = g.seats.filter(cell_type=SeatCellType.SEAT).select_related('student')
             students = [s.student for s in seats if s.student]
             if not students: continue
             current_sum = sum(s.score or 0 for s in students)
             count = len(students)
             avg = current_sum / count
             group_data.append({
                 'group': g, 
                 'students': students, 
                 'sum': current_sum, 
                 'count': count, 
                 'avg': avg
             })
        
        if len(group_data) > 1:
            group_data.sort(key=lambda x: x['avg'])
            min_g = group_data[0]
            max_g = group_data[-1]
            diff = max_g['avg'] - min_g['avg']
            
            if diff > 5:
                 best_swap = None
                 current_improvement = 0
                 
                 for s_high in max_g['students']:
                     for s_low in min_g['students']:
                         if _is_internal_policy_student(s_high) or _is_internal_policy_student(s_low):
                             continue
                         score_diff = (s_high.score or 0) - (s_low.score or 0)
                         if score_diff > 0:
                             new_max_sum = max_g['sum'] - score_diff
                             new_min_sum = min_g['sum'] + score_diff
                             
                             new_max_avg = new_max_sum / max_g['count']
                             new_min_avg = new_min_sum / min_g['count']
                             
                             new_diff = abs(new_max_avg - new_min_avg)
                             improvement = diff - new_diff
                             
                             if improvement > 1 and improvement > current_improvement:
                                 current_improvement = improvement
                                 best_swap = (s_high, s_low)
                 
                 if best_swap:
                     s1, s2 = best_swap
                     issues.append({
                        'type': 'group_balance',
                        'message': f'建议交换 {s1.name} 和 {s2.name} 以平衡小组均分 (分差 {diff:.1f} → {(diff - current_improvement):.1f})',
                        'action_label': '交换优化',
                        'action_url': reverse('apply_suggestion', args=[classroom.pk]) + f'?type=swap_balance&s1={s1.pk}&s2={s2.pk}',
                        'ignore_label': '忽略此条',
                        'ignore_url': '#'
                     })

    return _filter_internal_issues(issues)


def _serialize_classroom_sync_meta(classroom, meta=None):
    meta = meta or SyncMeta.objects.get_or_create(classroom=classroom)[0]
    local_version = int(meta.local_version or 0)
    cloud_version = int(meta.cloud_version or 0)
    last_operation_at = _sync_meta_operation_time(meta, classroom)
    backed_up = bool(meta.last_sync_at and cloud_version > 0)
    has_local_changes = bool(backed_up and local_version > cloud_version)
    if not backed_up:
        state = 'not_backed_up'
    elif has_local_changes:
        state = 'dirty'
    else:
        state = 'synced'
    return {
        'uuid': str(meta.uuid),
        'local_version': local_version,
        'cloud_version': cloud_version,
        'last_operation_at': last_operation_at.isoformat() if last_operation_at else None,
        'last_sync_at': meta.last_sync_at.isoformat() if meta.last_sync_at else None,
        'last_error': meta.last_error,
        'backed_up': backed_up,
        'has_local_changes': has_local_changes,
        'state': state,
    }


def classroom_detail(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    sync_meta, _ = SyncMeta.objects.get_or_create(classroom=classroom)
    suggestions = _evaluate_layout(classroom, request)
    seats = list(classroom.seats.select_related('student', 'group').all())
    seat_map = _build_seat_map(seats)
    constraints = list(classroom.constraints.select_related('student', 'target_student').all())
    guardian_student_ids = _get_podium_guardian_ids(classroom, seats=seats)
    fixed_student_ids = _get_fixed_seat_student_ids(classroom, constraints=constraints)

    seat_grid = []
    for r in range(1, classroom.rows + 1):
        row_seats = []
        for c in range(1, classroom.cols + 1):
            row_seats.append(seat_map.get((r, c)))
        seat_grid.append(row_seats)

    students = list(classroom.students.all().order_by('name'))
    for student in students:
        student.podium_guardian_side = _get_podium_guardian_side(
            classroom,
            student,
            seats=seats,
            guardian_student_ids=guardian_student_ids,
        )
        student.is_fixed_seat = student.pk in fixed_student_ids

    for seat in seats:
        if seat.student_id and seat.student:
            seat.student.podium_guardian_side = _get_podium_guardian_side(
                classroom,
                seat.student,
                seats=seats,
                guardian_student_ids=guardian_student_ids,
            )
            seat.student.is_fixed_seat = seat.student.pk in fixed_student_ids

    seated_student_ids = {seat.student_id for seat in seats if seat.student_id}
    unseated_students = [student for student in students if student.pk not in seated_student_ids]
    groups = classroom.groups.all()
    snapshots = classroom.layout_snapshots.all()
    constraint_items, constraint_metrics = serialize_constraints(classroom, constraints=constraints)
    tag_rule_items, tag_rule_metrics = serialize_tag_rules(classroom)
    
    return render(request, 'seats/classroom_detail.html', {
        'classroom': classroom,
        'seat_grid': seat_grid,
        'students': students,
        'unseated_students': unseated_students,
        'groups': groups,
        'snapshots': snapshots,
        'constraint_items': constraint_items,
        'constraint_metrics': constraint_metrics,
        'constraint_types': get_constraint_type_definitions(),
        'student_tags': _serialize_student_tag_catalog(classroom),
        'tag_rule_items': tag_rule_items,
        'tag_rule_metrics': tag_rule_metrics,
        'tag_rule_types': get_tag_rule_type_definitions(),
        'suggestions': suggestions,
        'sync_meta': sync_meta,
        'sync_meta_payload': _serialize_classroom_sync_meta(classroom, sync_meta),
    })


@require_http_methods(['GET', 'POST'])
def classroom_command(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'GET':
        return JsonResponse({
            'status': 'success',
            'manifest': _build_classroom_command_manifest(classroom),
        })

    try:
        if request.content_type and 'application/json' in request.content_type:
            payload = json.loads(request.body or '{}')
        else:
            payload = request.POST
    except json.JSONDecodeError:
        return JsonResponse({'status': 'error', 'message': '请求数据格式错误'}, status=400)

    command_text = str(payload.get('command') or payload.get('message') or '').strip()
    if not command_text:
        return JsonResponse({'status': 'error', 'message': '请输入命令内容'}, status=400)

    command_result = _execute_classroom_command(classroom, command_text, request=request)
    return JsonResponse(_build_classroom_command_response_payload(classroom, command_result))


def ai_workspace(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/ai_workspace.html', {
        'classroom': classroom,
        'ai_overview': _get_classroom_overview_payload(classroom),
    })


@require_POST
def ai_chat(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        payload = json.loads(request.body or '{}')
    except json.JSONDecodeError:
        return JsonResponse({'status': 'error', 'message': '请求数据格式错误'}, status=400)

    action = str(payload.get('action') or 'message').strip()
    message = str(payload.get('message') or '').strip()
    conversation_id = payload.get('conversation_id')
    client_config = _normalize_ai_client_config(payload.get('client_config'))

    if action == 'config_get':
        return JsonResponse({
            'status': 'success',
            'client_config': _load_persisted_ai_client_config(classroom),
        })

    if action == 'config_save':
        try:
            saved_config = _save_ai_client_config(classroom, payload.get('client_config'))
        except RuntimeError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
        return JsonResponse({
            'status': 'success',
            'message': '连接设置已保存到数据库。',
            'client_config': saved_config,
        })

    if action == 'conversation_init':
        try:
            conversation, _ = _resolve_ai_conversation(
                classroom,
                request,
                conversation_id=conversation_id,
                create_if_missing=True,
            )
        except ValueError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
        return JsonResponse({
            'status': 'success',
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'messages': _load_ai_conversation_messages(conversation),
        })

    if action == 'conversation_create':
        conversation = _create_ai_conversation(
            classroom,
            request,
            title=str(payload.get('title') or '').strip(),
        )
        return JsonResponse({
            'status': 'success',
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'messages': [],
        })

    if action == 'conversation_switch':
        try:
            conversation, _ = _resolve_ai_conversation(
                classroom,
                request,
                conversation_id=conversation_id,
                create_if_missing=False,
            )
        except ValueError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
        return JsonResponse({
            'status': 'success',
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'messages': _load_ai_conversation_messages(conversation),
        })

    if action == 'conversation_delete':
        try:
            conversation, _ = _resolve_ai_conversation(
                classroom,
                request,
                conversation_id=conversation_id,
                create_if_missing=False,
            )
        except ValueError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

        conversation.delete()
        active_conversation, _ = _resolve_ai_conversation(classroom, request, create_if_missing=True)
        return JsonResponse({
            'status': 'success',
            'conversation_id': active_conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'messages': _load_ai_conversation_messages(active_conversation),
        })

    effective_client_config = _merge_ai_client_config(classroom, client_config)

    if action == 'tool_approval':
        approval_token = str(payload.get('approval_token') or '').strip()
        decisions = payload.get('decisions') or []
        if not approval_token:
            return JsonResponse({'status': 'error', 'message': '缺少授权令牌'}, status=400)
        if not isinstance(decisions, list):
            return JsonResponse({'status': 'error', 'message': '授权数据格式错误'}, status=400)

        tool_events = []
        try:
            pending = _consume_future_mode_pending(
                request,
                approval_token,
                pk,
                conversation_id=conversation_id,
            )
            pending_conversation_id = int(pending.get('conversation_id') or 0)
            if pending_conversation_id > 0:
                conversation, _ = _resolve_ai_conversation(
                    classroom,
                    request,
                    conversation_id=pending_conversation_id,
                    create_if_missing=False,
                )
            else:
                conversation, _ = _resolve_ai_conversation(
                    classroom,
                    request,
                    conversation_id=conversation_id,
                    create_if_missing=True,
                )
            decision_map = {}
            for item in decisions:
                if not isinstance(item, dict):
                    continue
                call_id = str(item.get('call_id') or '').strip()
                if not call_id:
                    continue
                decision_map[call_id] = bool(item.get('approved'))

            tool_outputs = []
            pending_calls = pending.get('function_calls') or []
            missing_decisions = [
                call.get('call_id')
                for call in pending_calls
                if str(call.get('call_id') or '').strip() and str(call.get('call_id') or '').strip() not in decision_map
            ]
            if missing_decisions:
                return JsonResponse({'status': 'error', 'message': '请先完成全部授权选择后再提交'}, status=400)

            for call in pending_calls:
                approved = bool(decision_map.get(call['call_id']))
                if approved:
                    try:
                        result = _execute_ai_tool(classroom, call['name'], call['arguments'], request=request)
                    except Exception as exc:
                        result = _build_ai_tool_error_result(call['name'], exc)
                else:
                    result = {
                        'ok': False,
                        'tool': call['name'],
                        'denied': True,
                        'message': '用户拒绝授权执行该工具',
                    }
                tool_events.append({
                    'name': call['name'],
                    'arguments': call['arguments'],
                    'approved': approved,
                    'result': result,
                })
                tool_outputs.append({
                    'type': 'function_call_output',
                    'call_id': call['call_id'],
                    'output': json.dumps(result, ensure_ascii=False),
                })

            pending_mode = str(pending.get('mode') or 'responses').strip() or 'responses'
            cards = _collect_cards_from_tool_events(tool_events)
            if pending_mode == 'direct':
                reply = _build_tool_events_success_reply(tool_events)
                _append_ai_conversation_message(
                    conversation,
                    AIConversationMessage.MessageRole.ASSISTANT,
                    reply,
                    payload={'cards': cards} if cards else {},
                )
                return JsonResponse({
                    'status': 'success',
                    'reply': reply,
                    'tool_events': tool_events,
                    'cards': cards,
                    'conversation_id': conversation.pk,
                    'conversations': _list_ai_conversations(classroom, request),
                    'overview': _get_classroom_overview_payload(classroom),
                })
            result = _run_future_mode(
                classroom,
                conversation=[],
                request=request,
                client_config=effective_client_config,
                previous_response_id=pending.get('response_id') if pending_mode == 'responses' else None,
                tool_outputs=tool_outputs,
                tool_events=tool_events,
                mode=pending_mode,
                chat_messages=pending.get('chat_messages') or [],
                conversation_id=conversation.pk,
            )
        except RuntimeError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
        except ValueError as exc:
            return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
        except Exception as exc:
            if _is_tool_output_call_mismatch_error(exc) and tool_events:
                return JsonResponse({
                    'status': 'success',
                    'reply': _build_tool_events_fallback_reply(tool_events),
                    'tool_events': tool_events,
                    'overview': _get_classroom_overview_payload(classroom),
                })
            normalized = _normalize_future_mode_openai_exception(exc)
            if normalized is not None:
                status_code, error_message = normalized
                return JsonResponse({'status': 'error', 'message': error_message}, status=status_code)
            return JsonResponse({'status': 'error', 'message': f'工具授权处理失败：{exc}'}, status=500)

        if result['status'] == 'needs_approval':
            cards = _collect_cards_from_tool_events(result.get('tool_events') or [])
            if cards:
                _append_ai_conversation_message(
                    conversation,
                    AIConversationMessage.MessageRole.ASSISTANT,
                    '',
                    payload={'cards': cards},
                )
            return JsonResponse({
                'status': 'needs_approval',
                'approval_token': result.get('approval_token') or '',
                'pending_calls': result.get('pending_calls') or [],
                'tool_events': result.get('tool_events') or [],
                'cards': cards,
                'conversation_id': conversation.pk,
                'conversations': _list_ai_conversations(classroom, request),
                'overview': _get_classroom_overview_payload(classroom),
            })
        cards = _collect_cards_from_tool_events(result.get('tool_events') or [])
        reply = result.get('reply') or ''
        _append_ai_conversation_message(
            conversation,
            AIConversationMessage.MessageRole.ASSISTANT,
            reply,
            payload={'cards': cards} if cards else {},
        )
        return JsonResponse({
            'status': 'success',
            'reply': reply,
            'tool_events': result.get('tool_events') or [],
            'cards': cards,
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'overview': _get_classroom_overview_payload(classroom),
        })

    if not message:
        return JsonResponse({'status': 'error', 'message': '请输入内容后再发送'}, status=400)

    try:
        conversation, _ = _resolve_ai_conversation(
            classroom,
            request,
            conversation_id=conversation_id,
            create_if_missing=True,
        )
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    _append_ai_conversation_message(
        conversation,
        AIConversationMessage.MessageRole.USER,
        message,
    )
    if conversation.title == DEFAULT_AI_CONVERSATION_TITLE:
        user_count = conversation.messages.filter(role=AIConversationMessage.MessageRole.USER).count()
        if user_count <= 1:
            conversation.title = _build_conversation_title_from_message(message)
            conversation.updated_at = timezone.now()
            conversation.save(update_fields=['title', 'updated_at'])

    if str(message).strip().startswith(('/', '／')):
        command_result = _execute_classroom_command(classroom, message, request=request)
        _append_ai_conversation_message(
            conversation,
            AIConversationMessage.MessageRole.ASSISTANT,
            command_result.get('reply') or '命令已处理。',
            payload={
                'command': {
                    'name': command_result.get('command') or '',
                    'subcommand': command_result.get('subcommand') or '',
                    'kind': command_result.get('kind') or '',
                    'ok': bool(command_result.get('ok')),
                }
            },
        )
        response_payload = _build_classroom_command_response_payload(classroom, command_result)
        response_payload.update({
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
        })
        return JsonResponse(response_payload)

    direct_swap_call = _extract_direct_swap_call(classroom, message)
    if direct_swap_call:
        token = _store_future_mode_pending(
            request,
            classroom.pk,
            conversation.pk,
            response_id='',
            function_calls=[direct_swap_call],
            mode='direct',
            chat_messages=[],
        )
        return JsonResponse({
            'status': 'needs_approval',
            'approval_token': token,
            'pending_calls': [
                {
                    'call_id': direct_swap_call['call_id'],
                    'name': direct_swap_call['name'],
                    'label': AI_TOOL_LABELS.get(direct_swap_call['name'], direct_swap_call['name']),
                    'arguments': direct_swap_call['arguments'],
                    'summary': _describe_future_mode_call(direct_swap_call['name'], direct_swap_call['arguments']),
                }
            ],
            'tool_events': [],
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'overview': _get_classroom_overview_payload(classroom),
        })

    normalized_history = _build_history_from_conversation(conversation)

    try:
        result = _run_future_mode(
            classroom,
            normalized_history,
            request=request,
            client_config=effective_client_config,
            mode='auto',
            conversation_id=conversation.pk,
        )
    except RuntimeError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        normalized = _normalize_future_mode_openai_exception(exc)
        if normalized is not None:
            status_code, error_message = normalized
            return JsonResponse({'status': 'error', 'message': error_message}, status=status_code)
        return JsonResponse({'status': 'error', 'message': f'AI 处理失败：{exc}'}, status=500)

    if result['status'] == 'needs_approval':
        return JsonResponse({
            'status': 'needs_approval',
            'approval_token': result.get('approval_token') or '',
            'pending_calls': result.get('pending_calls') or [],
            'tool_events': result.get('tool_events') or [],
            'conversation_id': conversation.pk,
            'conversations': _list_ai_conversations(classroom, request),
            'overview': _get_classroom_overview_payload(classroom),
        })

    cards = _collect_cards_from_tool_events(result.get('tool_events') or [])
    reply = result.get('reply') or ''
    _append_ai_conversation_message(
        conversation,
        AIConversationMessage.MessageRole.ASSISTANT,
        reply,
        payload={'cards': cards} if cards else {},
    )
    return JsonResponse({
        'status': 'success',
        'reply': reply,
        'tool_events': result.get('tool_events') or [],
        'cards': cards,
        'conversation_id': conversation.pk,
        'conversations': _list_ai_conversations(classroom, request),
        'overview': _get_classroom_overview_payload(classroom),
    })


@require_POST
def ai_chat_stream(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        payload = json.loads(request.body or '{}')
    except json.JSONDecodeError:
        return JsonResponse({'status': 'error', 'message': '请求数据格式错误'}, status=400)

    action = str(payload.get('action') or 'message').strip()
    if action != 'message':
        return JsonResponse({'status': 'error', 'message': '流式接口仅支持 message 动作'}, status=400)

    message = str(payload.get('message') or '').strip()
    if not message:
        return JsonResponse({'status': 'error', 'message': '请输入内容后再发送'}, status=400)

    client_config = _normalize_ai_client_config(payload.get('client_config'))
    effective_client_config = _merge_ai_client_config(classroom, client_config)
    conversation_id = payload.get('conversation_id')

    try:
        conversation, _ = _resolve_ai_conversation(
            classroom,
            request,
            conversation_id=conversation_id,
            create_if_missing=True,
        )
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)

    _append_ai_conversation_message(
        conversation,
        AIConversationMessage.MessageRole.USER,
        message,
    )
    if conversation.title == DEFAULT_AI_CONVERSATION_TITLE:
        user_count = conversation.messages.filter(role=AIConversationMessage.MessageRole.USER).count()
        if user_count <= 1:
            conversation.title = _build_conversation_title_from_message(message)
            conversation.updated_at = timezone.now()
            conversation.save(update_fields=['title', 'updated_at'])

    if str(message).strip().startswith(('/', '／')):
        command_result = _execute_classroom_command(classroom, message, request=request)
        _append_ai_conversation_message(
            conversation,
            AIConversationMessage.MessageRole.ASSISTANT,
            command_result.get('reply') or '命令已处理。',
            payload={
                'command': {
                    'name': command_result.get('command') or '',
                    'subcommand': command_result.get('subcommand') or '',
                    'kind': command_result.get('kind') or '',
                    'ok': bool(command_result.get('ok')),
                }
            },
        )

        def command_stream():
            payload_data = _build_classroom_command_response_payload(classroom, command_result)
            payload_data.update({
                'conversation_id': conversation.pk,
                'conversations': _list_ai_conversations(classroom, request),
            })
            yield _sse_event('done', payload_data)

        response = StreamingHttpResponse(command_stream(), content_type='text/event-stream; charset=utf-8')
        response['Cache-Control'] = 'no-cache'
        response['X-Accel-Buffering'] = 'no'
        return response

    direct_swap_call = _extract_direct_swap_call(classroom, message)
    if direct_swap_call:
        token = _store_future_mode_pending(
            request,
            classroom.pk,
            conversation.pk,
            response_id='',
            function_calls=[direct_swap_call],
            mode='direct',
            chat_messages=[],
        )

        def direct_swap_stream():
            payload_data = {
                'status': 'needs_approval',
                'approval_token': token,
                'pending_calls': [
                    {
                        'call_id': direct_swap_call['call_id'],
                        'name': direct_swap_call['name'],
                        'label': AI_TOOL_LABELS.get(direct_swap_call['name'], direct_swap_call['name']),
                        'arguments': direct_swap_call['arguments'],
                        'summary': _describe_future_mode_call(direct_swap_call['name'], direct_swap_call['arguments']),
                    }
                ],
                'conversation_id': conversation.pk,
                'conversations': _list_ai_conversations(classroom, request),
                'cards': [],
            }
            yield _sse_event('done', payload_data)

        response = StreamingHttpResponse(direct_swap_stream(), content_type='text/event-stream; charset=utf-8')
        response['Cache-Control'] = 'no-cache'
        response['X-Accel-Buffering'] = 'no'
        return response

    normalized_history = _build_history_from_conversation(conversation)

    def event_stream():
        try:
            client = _get_openai_client(client_config=effective_client_config)
            model = _get_openai_model(client_config=effective_client_config)
            tools = _build_chat_tools()
            chat_messages = [{'role': 'system', 'content': _future_mode_system_prompt()}]
            for item in normalized_history[-AI_CONTEXT_MESSAGE_LIMIT:]:
                role = 'assistant' if item.get('role') == 'assistant' else 'user'
                content = str(item.get('content') or '').strip()
                if not content:
                    continue
                chat_messages.append({'role': role, 'content': content[:4000]})

            completion_kwargs = {
                'model': model,
                'messages': chat_messages,
                'tools': tools,
                'tool_choice': 'auto',
                'parallel_tool_calls': False,
                'stream': True,
            }
            extra_body = _build_chat_completion_extra_body(client_config=effective_client_config)
            if extra_body:
                completion_kwargs['extra_body'] = extra_body
            stream = client.chat.completions.create(**completion_kwargs)

            reply_parts = []
            reasoning_parts = []
            tool_call_buffer = {}
            for chunk in stream:
                choices = getattr(chunk, 'choices', None) or []
                if not choices:
                    continue
                choice = choices[0]
                delta = getattr(choice, 'delta', None)
                if delta is None and isinstance(choice, dict):
                    delta = choice.get('delta')
                if not delta:
                    continue

                delta_content = delta.get('content') if isinstance(delta, dict) else getattr(delta, 'content', None)
                text_delta = _extract_chat_delta_text(delta_content)
                if text_delta:
                    reply_parts.append(text_delta)
                    yield _sse_event('delta', {'text': text_delta})

                delta_reasoning = delta.get('reasoning_content') if isinstance(delta, dict) else getattr(delta, 'reasoning_content', None)
                reasoning_delta = _extract_chat_delta_text(delta_reasoning)
                if reasoning_delta:
                    reasoning_parts.append(reasoning_delta)

                tool_calls = delta.get('tool_calls') if isinstance(delta, dict) else getattr(delta, 'tool_calls', None)
                for tool_call in tool_calls or []:
                    call_index = tool_call.get('index') if isinstance(tool_call, dict) else getattr(tool_call, 'index', None)
                    if call_index is None:
                        call_index = 0
                    call_index = int(call_index)
                    buffer_item = tool_call_buffer.setdefault(call_index, {
                        'call_id': '',
                        'name': '',
                        'arguments_parts': [],
                    })

                    call_id = tool_call.get('id') if isinstance(tool_call, dict) else getattr(tool_call, 'id', None)
                    if call_id:
                        buffer_item['call_id'] = str(call_id).strip()

                    function_data = tool_call.get('function') if isinstance(tool_call, dict) else getattr(tool_call, 'function', None)
                    if not function_data:
                        continue

                    function_name = function_data.get('name') if isinstance(function_data, dict) else getattr(function_data, 'name', None)
                    if function_name:
                        buffer_item['name'] = str(function_name).strip()

                    function_arguments = function_data.get('arguments') if isinstance(function_data, dict) else getattr(function_data, 'arguments', None)
                    if function_arguments:
                        buffer_item['arguments_parts'].append(str(function_arguments))

            if tool_call_buffer:
                function_calls = []
                assistant_tool_calls = []
                for call_index in sorted(tool_call_buffer.keys()):
                    item = tool_call_buffer[call_index]
                    tool_name = str(item.get('name') or '').strip()
                    if not tool_name:
                        continue
                    call_id = str(item.get('call_id') or f'stream_call_{uuid.uuid4().hex[:12]}').strip()
                    arguments_json = ''.join(item.get('arguments_parts') or []).strip() or '{}'
                    try:
                        arguments = json.loads(arguments_json)
                    except (TypeError, ValueError, json.JSONDecodeError):
                        arguments = {}
                    function_calls.append({
                        'call_id': call_id,
                        'name': tool_name,
                        'arguments': arguments,
                    })
                    assistant_tool_calls.append({
                        'id': call_id,
                        'type': 'function',
                        'function': {
                            'name': tool_name,
                            'arguments': arguments_json,
                        },
                    })

                if function_calls:
                    pending_messages = list(chat_messages)
                    assistant_pending_message = {
                        'role': 'assistant',
                        'content': '',
                        'tool_calls': assistant_tool_calls,
                    }
                    assistant_reasoning_content = ''.join(reasoning_parts).strip()
                    if assistant_reasoning_content:
                        assistant_pending_message['reasoning_content'] = assistant_reasoning_content
                    pending_messages.append(assistant_pending_message)
                    token = _store_future_mode_pending(
                        request,
                        classroom.pk,
                        conversation.pk,
                        response_id='',
                        function_calls=function_calls,
                        mode='chat',
                        chat_messages=pending_messages,
                    )
                    yield _sse_event('done', {
                        'status': 'needs_approval',
                        'approval_token': token,
                        'pending_calls': [
                            {
                                'call_id': call['call_id'],
                                'name': call['name'],
                                'label': AI_TOOL_LABELS.get(call['name'], call['name']),
                                'arguments': call['arguments'],
                                'summary': _describe_future_mode_call(call['name'], call['arguments']),
                            }
                            for call in function_calls
                        ],
                        'conversation_id': conversation.pk,
                        'conversations': _list_ai_conversations(classroom, request),
                        'cards': [],
                    })
                    return

            reply = ''.join(reply_parts).strip() or '已完成处理，但没有生成可展示的回复。'
            _append_ai_conversation_message(
                conversation,
                AIConversationMessage.MessageRole.ASSISTANT,
                reply,
            )
            yield _sse_event('done', {
                'status': 'success',
                'reply': reply,
                'conversation_id': conversation.pk,
                'conversations': _list_ai_conversations(classroom, request),
                'cards': [],
            })
        except RuntimeError as exc:
            yield _sse_event('error', {'message': str(exc)})
        except ValueError as exc:
            yield _sse_event('error', {'message': str(exc)})
        except Exception as exc:
            normalized = _normalize_future_mode_openai_exception(exc)
            if normalized is not None:
                _, error_message = normalized
                yield _sse_event('error', {'message': error_message})
                return
            yield _sse_event('error', {'message': f'AI 流式处理失败：{exc}'})

    response = StreamingHttpResponse(event_stream(), content_type='text/event-stream; charset=utf-8')
    response['Cache-Control'] = 'no-cache'
    response['X-Accel-Buffering'] = 'no'
    return response


def classroom_state(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    suggestions = _evaluate_layout(classroom, request)
    seats = list(classroom.seats.select_related('student', 'group').all())
    unseated_students = list(classroom.students.filter(assigned_seat__isnull=True).order_by('name'))
    constraints = list(classroom.constraints.select_related('student', 'target_student').all())
    guardian_student_ids = _get_podium_guardian_ids(classroom, seats=seats)
    fixed_student_ids = _get_fixed_seat_student_ids(classroom, constraints=constraints)
    constraint_items, constraint_metrics = serialize_constraints(classroom, constraints=constraints)
    all_state_student_ids = {
        seat.student_id
        for seat in seats
        if seat.student_id
    }
    all_state_student_ids.update(student.pk for student in unseated_students)
    tag_map = _build_student_tag_map(classroom, all_state_student_ids)
    tag_rule_items, tag_rule_metrics = serialize_tag_rules(classroom)

    seat_payload = []
    for seat in seats:
        student = seat.student
        group = seat.group
        score_value = student.display_score if student and (student.score or 0) > 0 else None
        guardian_side = _get_podium_guardian_side(
            classroom,
            student,
            seats=seats,
            guardian_student_ids=guardian_student_ids,
        ) if student else ''
        seat_payload.append({
            'row': seat.row,
            'col': seat.col,
            'cell_type': seat.cell_type,
            'cell_type_display': seat.get_cell_type_display(),
            'student': {
                'id': student.pk,
                'name': student.name,
                'student_id': student.student_id or '',
                'score_display': score_value,
                'is_leader': (group and getattr(group, 'leader_id', None) == student.pk),
                'podium_guardian_side': guardian_side,
                'is_fixed_seat': student.pk in fixed_student_ids,
                'tags': tag_map.get(student.pk, []),
            } if student else None,
            'group': {
                'id': group.pk,
                'name': group.name
            } if group else None
        })

    unseated_payload = []
    for student in unseated_students:
        score_value = student.display_score if (student.score or 0) > 0 else None
        unseated_payload.append({
            'id': student.pk,
            'name': student.name,
            'student_id': student.student_id or '',
            'gender': student.gender or '',
            'score': student.score or 0,
            'score_display': score_value,
            'podium_guardian_side': _get_podium_guardian_side(
                classroom,
                student,
                seats=seats,
                guardian_student_ids=guardian_student_ids,
            ),
            'is_fixed_seat': student.pk in fixed_student_ids,
            'tags': tag_map.get(student.pk, []),
            'delete_url': reverse('delete_student', args=[classroom.pk, student.pk]),
            'update_url': reverse('update_student', args=[classroom.pk, student.pk]),
        })

    return JsonResponse({
        'seats': seat_payload,
        'unseated': unseated_payload,
        'podium_guards': _serialize_podium_guards(
            classroom,
            seats=seats,
            guardian_student_ids=guardian_student_ids,
        ),
        'constraints': constraint_items,
        'constraint_metrics': constraint_metrics,
        'tags': _serialize_student_tag_catalog(classroom),
        'tag_rules': tag_rule_items,
        'tag_rule_types': get_tag_rule_type_definitions(),
        'tag_rule_metrics': tag_rule_metrics,
        'suggestions': suggestions,
        'unseated_count': len(unseated_payload),
        'sync_meta': _serialize_classroom_sync_meta(classroom),
    })


def layout_editor(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    seats = list(classroom.seats.select_related('student').all())
    seat_map = _build_seat_map(seats)
    seat_grid = []
    for r in range(1, classroom.rows + 1):
        row_seats = []
        for c in range(1, classroom.cols + 1):
            row_seats.append(seat_map.get((r, c)))
        seat_grid.append(row_seats)
    return render(request, 'seats/layout_editor.html', {
        'classroom': classroom,
        'seat_grid': seat_grid
    })


@require_POST
def update_layout_grid(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    rows = int(request.POST.get('rows', classroom.rows))
    cols = int(request.POST.get('cols', classroom.cols))
    rows = max(1, min(rows, MAX_LAYOUT_GRID_SIZE))
    cols = max(1, min(cols, MAX_LAYOUT_GRID_SIZE))
    before_state = _capture_history_state(classroom)
    _sync_seats(classroom, rows, cols)
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'update_layout_grid',
        extra={'rows': rows, 'cols': cols},
    )
    return redirect('layout_editor', pk=pk)


@require_POST
def shift_layout(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        if request.content_type and 'application/json' in request.content_type:
            payload = json.loads(request.body or '{}')
        else:
            payload = request.POST

        direction = _normalize_shift_direction(payload.get('direction'))
        direction_meta = _shift_direction_meta(direction)
        steps = payload.get('steps')
        use_large_groups_raw = payload.get('use_large_groups') if hasattr(payload, 'get') else None
        use_large_groups = True if use_large_groups_raw in (None, '') else _parse_bool(use_large_groups_raw)
        before_data = _snapshot_payload(classroom, include_students=False, include_constraints=True)
        before_state = _capture_history_state(classroom)
        shift_mode = 'normal'
        shift_strategy = 'normal'
        fallback_reason = ''
        template_signature = ''
        template_meta = {}
        if direction in {'left', 'right'}:
            if use_large_groups:
                intelligent_payload, template_meta = _build_intelligent_horizontal_shift_payload(classroom, direction, steps)
                if intelligent_payload is not None:
                    after_data = intelligent_payload
                    shift_mode = 'template'
                    shift_strategy = str(template_meta.get('template_strategy') or 'large_group')
                    template_signature = str(template_meta.get('template_signature') or '')
                else:
                    after_data = _build_shifted_layout_payload(classroom, direction, steps)
                    fallback_reason = str(template_meta.get('reason') or '当前布局结构不明确')
            else:
                after_data, template_meta = _build_single_column_horizontal_shift_payload(classroom, direction, steps)
                shift_mode = 'column'
                shift_strategy = str(template_meta.get('template_strategy') or 'single_column')
                template_signature = str(template_meta.get('template_signature') or '')
        else:
            after_data = _build_shifted_layout_payload(classroom, direction, steps)
        after_data = _pin_podium_guardians_in_shift_payload(before_data, after_data)
        action = {
            'type': 'layout_snapshot',
            'before_data': before_data,
            'after_data': after_data,
            'direction': direction,
            'steps': _safe_int(steps, 0),
            'shift_mode': shift_mode,
            'shift_strategy': shift_strategy,
            'template_signature': template_signature,
            'use_large_groups': use_large_groups,
        }
        with transaction.atomic():
            if not _apply_layout_snapshot_action(classroom, action, forward=True):
                raise ValueError('整体平移失败')
            _push_snapshot_action(request, classroom, before_state, 'shift_layout', extra=action)
        if shift_mode == 'template':
            seat_block_count = int(template_meta.get('seat_block_count') or 0)
            aisle_block_count = int(template_meta.get('aisle_block_count') or 0)
            message = (
                f'已按 {template_signature} 布局模板完成{direction_meta["action_label"]}'
                f'，识别到 {seat_block_count} 个座位块和 {aisle_block_count} 条结构走廊'
            )
        elif shift_mode == 'column':
            seat_column_count = int(template_meta.get('seat_column_count') or 0)
            structural_column_count = int(template_meta.get('structural_column_count') or 0)
            message = (
                f'已按座位纵列轮换完成{direction_meta["action_label"]}'
                f'，识别到 {seat_column_count} 个座位纵列和 {structural_column_count} 列固定结构'
            )
        elif direction in {'left', 'right'} and fallback_reason:
            message = (
                f'当前布局结构不明确，已回退为普通{direction_meta["action_label"]} '
                f'{action["steps"]} {direction_meta["unit"]}'
            )
        else:
            message = f'已整体{direction_meta["action_label"]} {action["steps"]} {direction_meta["unit"]}'
        return JsonResponse({
            'status': 'success',
            'message': message,
            'shift_mode': shift_mode,
            'shift_strategy': shift_strategy,
            'template_signature': template_signature,
            'fallback_reason': fallback_reason,
            'use_large_groups': use_large_groups,
            'seat_column_count': int(template_meta.get('seat_column_count') or 0),
        })
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': f'整体平移失败：{e}'}, status=400)


@require_POST
def mirror_layout(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        if request.content_type and 'application/json' in request.content_type:
            payload = json.loads(request.body or '{}')
        else:
            payload = request.POST

        axis = _normalize_mirror_axis(payload.get('axis') or payload.get('direction') or 'lr')
        before_data = _snapshot_payload(classroom, include_students=False, include_constraints=True)
        before_state = _capture_history_state(classroom)
        after_data = _build_mirrored_layout_payload(classroom, axis=axis)
        action = {
            'type': 'layout_snapshot',
            'before_data': before_data,
            'after_data': after_data,
            'mirror_axis': axis,
        }
        with transaction.atomic():
            if not _apply_layout_snapshot_action(classroom, action, forward=True):
                raise ValueError('左右镜像失败')
            _push_snapshot_action(request, classroom, before_state, 'mirror_layout', extra=action)
        return JsonResponse({
            'status': 'success',
            'message': '已完成左右镜像',
            'mirror_axis': axis,
        })
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': f'左右镜像失败：{e}'}, status=400)


@require_POST
def insert_delete_row_col(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body or '{}')
        action_type = data.get('action')
        index = _safe_int(data.get('index'), 0)

        if action_type not in ('insert_row', 'delete_row', 'insert_col', 'delete_col'):
            raise ValueError('操作类型不合法')
        if index < 1:
            raise ValueError('索引必须大于 0')

        before_data = _snapshot_payload(classroom, include_students=True, include_constraints=True)
        before_state = _capture_history_state(classroom)

        with transaction.atomic():
            if action_type == 'insert_row':
                if index > classroom.rows + 1:
                    raise ValueError('行索引超出范围')
                affected = classroom.seats.filter(row__gte=index)
                affected.update(row=models.F('row') + 1000000)
                classroom.seats.filter(row__gte=1000000).update(row=models.F('row') - 999999)
                classroom.rows += 1
                classroom.save(update_fields=['rows'])
                for c in range(1, classroom.cols + 1):
                    Seat.objects.create(classroom=classroom, row=index, col=c, cell_type=SeatCellType.SEAT)

            elif action_type == 'delete_row':
                if index > classroom.rows:
                    raise ValueError('行索引超出范围')
                if classroom.rows <= 1:
                    raise ValueError('至少保留一行')
                classroom.seats.filter(row=index).update(student=None, group=None)
                classroom.seats.filter(row=index).delete()
                affected = classroom.seats.filter(row__gt=index)
                affected.update(row=models.F('row') + 1000000)
                classroom.seats.filter(row__gte=1000000).update(row=models.F('row') - 1000001)
                classroom.rows -= 1
                classroom.save(update_fields=['rows'])

            elif action_type == 'insert_col':
                if index > classroom.cols + 1:
                    raise ValueError('列索引超出范围')
                affected = classroom.seats.filter(col__gte=index)
                affected.update(col=models.F('col') + 1000000)
                classroom.seats.filter(col__gte=1000000).update(col=models.F('col') - 999999)
                classroom.cols += 1
                classroom.save(update_fields=['cols'])
                for r in range(1, classroom.rows + 1):
                    Seat.objects.create(classroom=classroom, row=r, col=index, cell_type=SeatCellType.SEAT)

            elif action_type == 'delete_col':
                if index > classroom.cols:
                    raise ValueError('列索引超出范围')
                if classroom.cols <= 1:
                    raise ValueError('至少保留一列')
                classroom.seats.filter(col=index).update(student=None, group=None)
                classroom.seats.filter(col=index).delete()
                affected = classroom.seats.filter(col__gt=index)
                affected.update(col=models.F('col') + 1000000)
                classroom.seats.filter(col__gte=1000000).update(col=models.F('col') - 1000001)
                classroom.cols -= 1
                classroom.save(update_fields=['cols'])

            after_data = _snapshot_payload(classroom, include_students=True, include_constraints=True)
            undo_action = {
                'type': 'layout_snapshot',
                'before_data': before_data,
                'after_data': after_data,
            }
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                action_type,
                extra={'action': action_type, 'index': index, **undo_action},
            )

        labels = {
            'insert_row': f'已在第 {index} 行上方插入新行',
            'delete_row': f'已删除第 {index} 行',
            'insert_col': f'已在第 {index} 列左侧插入新列',
            'delete_col': f'已删除第 {index} 列',
        }
        return JsonResponse({'status': 'success', 'message': labels[action_type]})
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': f'操作失败：{e}'}, status=400)


def shift_layout_options_page(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/shift_layout_options.html', {
        'classroom': classroom
    })


def _ensure_temp_import_dir():
    temp_dir = os.path.join(settings.BASE_DIR, 'temp_imports')
    os.makedirs(temp_dir, exist_ok=True)
    return temp_dir


def _save_uploaded_temp_file(uploaded_file, suffix):
    file_id = str(uuid.uuid4())
    temp_path = os.path.join(_ensure_temp_import_dir(), f'{file_id}{suffix}')
    with open(temp_path, 'wb+') as destination:
        for chunk in uploaded_file.chunks():
            destination.write(chunk)
    return file_id, temp_path


def _parse_bool(value):
    return str(value or '').strip().lower() in {'1', 'true', 'yes', 'on'}


def _normalize_cell_text(value):
    if value is None:
        return ''
    text = str(value).strip()
    if not text:
        return ''
    return re.sub(r'\s+', '', text)


def _parse_manual_terms(raw):
    if not raw:
        return set()
    parts = re.split(r'[\n,，;；\s]+', str(raw))
    return {p.strip() for p in parts if p.strip()}


def _is_name_like_text(text):
    if not text:
        return False
    if not (2 <= len(text) <= 5):
        return False
    if any(ch.isdigit() for ch in text):
        return False
    return bool(re.fullmatch(r'[^\d\s]{2,5}', text))


def _build_merged_value_map(ws, min_row, max_row, min_col, max_col):
    merged_map = {}
    for merged_range in ws.merged_cells.ranges:
        c1, r1, c2, r2 = merged_range.bounds
        if r2 < min_row or r1 > max_row or c2 < min_col or c1 > max_col:
            continue
        master_val = ws.cell(row=r1, column=c1).value
        rr1 = max(r1, min_row)
        rr2 = min(r2, max_row)
        cc1 = max(c1, min_col)
        cc2 = min(c2, max_col)
        for r in range(rr1, rr2 + 1):
            for c in range(cc1, cc2 + 1):
                merged_map[(r, c)] = master_val
    return merged_map


def _detect_layout_bounds(ws):
    min_row = None
    max_row = None
    min_col = None
    max_col = None
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        for cell in row:
            if _normalize_cell_text(cell.value):
                r, c = cell.row, cell.column
                min_row = r if min_row is None else min(min_row, r)
                max_row = r if max_row is None else max(max_row, r)
                min_col = c if min_col is None else min(min_col, c)
                max_col = c if max_col is None else max(max_col, c)
    if min_row is None:
        return {
            'min_row': 1,
            'max_row': 1,
            'min_col': 1,
            'max_col': 1
        }
    return {
        'min_row': min_row,
        'max_row': max_row,
        'min_col': min_col,
        'max_col': max_col
    }


def _calc_col_bounds_for_rows(ws, start_row, end_row, base_min_col, base_max_col):
    min_col = None
    max_col = None
    merged_map = _build_merged_value_map(ws, start_row, end_row, base_min_col, base_max_col)
    for r in range(start_row, end_row + 1):
        for c in range(base_min_col, base_max_col + 1):
            val = merged_map.get((r, c), ws.cell(row=r, column=c).value)
            if _normalize_cell_text(val):
                min_col = c if min_col is None else min(min_col, c)
                max_col = c if max_col is None else max(max_col, c)
    if min_col is None:
        return base_min_col, base_max_col
    return min_col, max_col


def _transform_layout_rows(rows, layout_transform):
    transform = str(layout_transform or 'none').strip().lower()
    if transform == 'flip_ud':
        return list(reversed(rows))
    if transform == 'flip_lr':
        return [list(reversed(row)) for row in rows]
    if transform in {'rotate_180', 'rot180', '180'}:
        return [list(reversed(row)) for row in reversed(rows)]
    return rows


def _detect_layout_import_defaults(temp_path, options):
    wb = openpyxl.load_workbook(temp_path, data_only=True)
    ws = wb.active
    bounds = _detect_layout_bounds(ws)
    min_row = bounds['min_row']
    max_row = bounds['max_row']
    min_col = bounds['min_col']
    max_col = bounds['max_col']
    merged_map = _build_merged_value_map(ws, min_row, max_row, min_col, max_col)

    row_podium_count = defaultdict(int)
    row_name_count = defaultdict(int)

    for r in range(min_row, max_row + 1):
        for c in range(min_col, max_col + 1):
            raw_val = merged_map.get((r, c), ws.cell(row=r, column=c).value)
            text = _normalize_cell_text(raw_val)
            cell_type, student_name, _ = _classify_layout_cell(text, options)
            if cell_type == SeatCellType.PODIUM:
                row_podium_count[r] += 1
            if student_name:
                row_name_count[r] += 1

    wb.close()

    start_row = min_row
    end_row = max_row
    layout_transform = 'none'

    podium_row = None
    if row_podium_count:
        podium_row = max(
            row_podium_count.keys(),
            key=lambda row: (row_podium_count[row], -row)
        )

    if podium_row and row_name_count.get(podium_row, 0) == 0:
        above_names = sum(row_name_count[r] for r in range(min_row, podium_row))
        below_names = sum(row_name_count[r] for r in range(podium_row + 1, max_row + 1))

        if below_names > above_names and below_names > 0:
            start_row = min(max_row, podium_row + 1)
        elif above_names > 0:
            end_row = max(min_row, podium_row - 1)

        if podium_row != min_row:
            layout_transform = 'rotate_180'

    if start_row > end_row:
        start_row = min_row
        end_row = max_row

    return {
        'start_row': start_row,
        'end_row': end_row,
        'layout_transform': layout_transform,
        'podium_row': podium_row
    }


LAYOUT_PODIUM_KEYWORDS = {'讲台', '教师', '老师', '黑板', '主席台'}
LAYOUT_AISLE_KEYWORDS = {'走廊', '过道', '通道'}
LAYOUT_EMPTY_KEYWORDS = {'空位', '留空', '空座', '无人'}


def _classify_layout_cell(text, options):
    manual_name_terms = options.get('manual_name_terms', set())
    manual_podium_terms = options.get('manual_podium_terms', set())
    manual_empty_terms = options.get('manual_empty_terms', set())
    manual_aisle_terms = options.get('manual_aisle_terms', set())
    auto_detect_names = options.get('auto_detect_names', True)

    if not text:
        return SeatCellType.AISLE, None, '空白识别为走廊'

    if text in manual_name_terms:
        return SeatCellType.SEAT, text, '手动姓名'
    if text in manual_podium_terms or any(k in text for k in LAYOUT_PODIUM_KEYWORDS):
        return SeatCellType.PODIUM, None, '讲台关键词'
    if text in manual_empty_terms or any(k in text for k in LAYOUT_EMPTY_KEYWORDS):
        return SeatCellType.EMPTY, None, '空位关键词'
    if text in manual_aisle_terms or any(k in text for k in LAYOUT_AISLE_KEYWORDS):
        return SeatCellType.AISLE, None, '走廊关键词'

    if auto_detect_names and _is_name_like_text(text):
        return SeatCellType.SEAT, text, '自动姓名'

    return SeatCellType.SEAT, None, '默认座位'


def _build_layout_grid_from_excel(temp_path, start_row, end_row, options):
    wb = openpyxl.load_workbook(temp_path, data_only=True)
    ws = wb.active
    bounds = _detect_layout_bounds(ws)

    start_row = max(bounds['min_row'], int(start_row or bounds['min_row']))
    end_row = min(bounds['max_row'], int(end_row or bounds['max_row']))
    if end_row < start_row:
        end_row = start_row

    min_col, max_col = _calc_col_bounds_for_rows(
        ws,
        start_row,
        end_row,
        bounds['min_col'],
        bounds['max_col']
    )
    merged_map = _build_merged_value_map(ws, start_row, end_row, min_col, max_col)

    rows = []
    stats = {
        'seat': 0,
        'aisle': 0,
        'podium': 0,
        'empty': 0,
        'named': 0
    }

    for r in range(start_row, end_row + 1):
        row_items = []
        for c in range(min_col, max_col + 1):
            raw_val = merged_map.get((r, c), ws.cell(row=r, column=c).value)
            text = _normalize_cell_text(raw_val)
            cell_type, student_name, reason = _classify_layout_cell(text, options)
            row_items.append({
                'sheet_row': r,
                'sheet_col': c,
                'raw_text': text,
                'cell_type': cell_type,
                'student_name': student_name,
                'reason': reason
            })
            stats[cell_type] += 1
            if student_name:
                stats['named'] += 1
        rows.append(row_items)

    rows = _transform_layout_rows(rows, options.get('layout_transform', 'none'))

    wb.close()
    return {
        'start_row': start_row,
        'end_row': end_row,
        'min_col': min_col,
        'max_col': max_col,
        'rows': rows,
        'bounds': bounds,
        'stats': stats
    }


def _preview_rows_payload(grid_rows):
    total = len(grid_rows)
    if total == 0:
        return [], []
    front = [(idx, grid_rows[idx]) for idx in range(min(2, total))]
    back_start = max(0, total - 2)
    back = [(idx, grid_rows[idx]) for idx in range(back_start, total)]
    if total <= 2:
        back = []

    def render_row(idx, row):
        return {
            'row_index': idx + 1,
            'cells': [
                {
                    'cell_type': item['cell_type'],
                    'label': item['student_name'] or (
                        '讲台' if item['cell_type'] == SeatCellType.PODIUM else
                        '空位' if item['cell_type'] == SeatCellType.EMPTY else
                        '走廊' if item['cell_type'] == SeatCellType.AISLE else
                        '座位'
                    )
                }
                for item in row
            ]
        }

    return [render_row(i, row) for i, row in front], [render_row(i, row) for i, row in back]


def _build_layout_preview_response(temp_path, start_row, end_row, options):
    grid_data = _build_layout_grid_from_excel(temp_path, start_row, end_row, options)
    front_rows, back_rows = _preview_rows_payload(grid_data['rows'])
    return {
        'layout_transform': options.get('layout_transform', 'none'),
        'start_row': grid_data['start_row'],
        'end_row': grid_data['end_row'],
        'bounds': grid_data['bounds'],
        'grid_rows': len(grid_data['rows']),
        'grid_cols': (grid_data['max_col'] - grid_data['min_col'] + 1),
        'front_preview': front_rows,
        'back_preview': back_rows,
        'stats': grid_data['stats']
    }


def _apply_layout_excel_import(classroom, temp_path, start_row, end_row, options):
    grid_data = _build_layout_grid_from_excel(temp_path, start_row, end_row, options)
    rows = grid_data['rows']
    row_count = len(rows)
    col_count = len(rows[0]) if rows else 0
    if row_count == 0 or col_count == 0:
        return 0, 0

    replace_students = options.get('replace_students', False)

    with transaction.atomic():
        _sync_seats(classroom, row_count, col_count)

        if replace_students:
            SeatConstraint.objects.filter(classroom=classroom).delete()
            Student.objects.filter(classroom=classroom).delete()

        seats = list(classroom.seats.select_related('student').all())
        seat_map = _build_seat_map(seats)

        for seat in seats:
            seat.student = None
            seat.group = None
            seat.save(update_fields=['student', 'group'])

        existing_by_name = defaultdict(list)
        if not replace_students:
            for student in classroom.students.all().order_by('pk'):
                existing_by_name[student.name].append(student)
        consumed_student_ids = set()

        imported_student_count = 0

        for local_r, row in enumerate(rows, start=1):
            for local_c, item in enumerate(row, start=1):
                seat = seat_map.get((local_r, local_c))
                if not seat:
                    continue
                target_student = None
                student_name = item.get('student_name')
                if item['cell_type'] == SeatCellType.SEAT and student_name:
                    candidates = existing_by_name.get(student_name, [])
                    for cand in candidates:
                        if cand.pk not in consumed_student_ids:
                            target_student = cand
                            break
                    if not target_student:
                        target_student = Student.objects.create(
                            classroom=classroom,
                            name=student_name,
                            student_id='',
                            score=0
                        )
                        existing_by_name[student_name].append(target_student)
                        imported_student_count += 1
                    consumed_student_ids.add(target_student.pk)

                seat.student = target_student
                seat.group = None
                seat.cell_type = item['cell_type']
                seat.save(update_fields=['student', 'group', 'cell_type'])

    return row_count * col_count, imported_student_count


def import_layout_excel(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method != 'POST':
        return redirect('classroom_detail', pk=pk)

    action = request.POST.get('action', 'upload')
    auto_detect_names = _parse_bool(request.POST.get('auto_detect_names', '1'))
    options = {
        'manual_name_terms': _parse_manual_terms(request.POST.get('manual_name_terms')),
        'manual_podium_terms': _parse_manual_terms(request.POST.get('manual_podium_terms')),
        'manual_empty_terms': _parse_manual_terms(request.POST.get('manual_empty_terms')),
        'manual_aisle_terms': _parse_manual_terms(request.POST.get('manual_aisle_terms')),
        'auto_detect_names': auto_detect_names,
        'layout_transform': request.POST.get('layout_transform', 'none'),
    }

    if action == 'upload':
        excel_file = request.FILES.get('layout_excel_file')
        if not excel_file:
            return JsonResponse({'status': 'error', 'message': '请先选择 Excel 座位表文件'}, status=400)
        suffix = os.path.splitext(excel_file.name)[1].lower() or '.xlsx'
        file_id, temp_path = _save_uploaded_temp_file(excel_file, suffix)
        try:
            defaults = _detect_layout_import_defaults(temp_path, options)
            preview_options = dict(options)
            preview_options['layout_transform'] = defaults['layout_transform']
            preview = _build_layout_preview_response(
                temp_path,
                defaults['start_row'],
                defaults['end_row'],
                preview_options
            )
            return JsonResponse({
                'status': 'ready',
                'file_id': file_id,
                'message': '文件解析完成，请确认范围后导入',
                'auto_selected': {
                    'podium_row': defaults['podium_row'],
                    'start_row': defaults['start_row'],
                    'end_row': defaults['end_row'],
                    'layout_transform': defaults['layout_transform'],
                },
                **preview
            })
        except Exception as e:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            return JsonResponse({'status': 'error', 'message': f'解析失败：{e}'}, status=400)

    file_id = request.POST.get('file_id', '').strip()
    if not file_id:
        return JsonResponse({'status': 'error', 'message': '缺少文件标识，请重新上传'}, status=400)
    temp_path = None
    for ext in ('.xlsx', '.xlsm', '.xls', ''):
        candidate = os.path.join(_ensure_temp_import_dir(), f'{file_id}{ext}')
        if os.path.exists(candidate):
            temp_path = candidate
            break
    if not temp_path:
        return JsonResponse({'status': 'error', 'message': '临时文件已过期，请重新上传'}, status=400)

    start_row = request.POST.get('start_row')
    end_row = request.POST.get('end_row')

    if action == 'preview':
        try:
            preview = _build_layout_preview_response(temp_path, start_row, end_row, options)
            return JsonResponse({
                'status': 'ready',
                'file_id': file_id,
                **preview
            })
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': f'预览失败：{e}'}, status=400)

    if action == 'confirm':
        options['replace_students'] = _parse_bool(request.POST.get('replace_students'))
        try:
            before_state = _capture_history_state(classroom)
            imported_cells, created_students = _apply_layout_excel_import(
                classroom,
                temp_path,
                start_row,
                end_row,
                options
            )
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                'import_layout_excel',
                extra={
                    'replace_students': options.get('replace_students'),
                    'imported_cells': imported_cells,
                    'created_students': created_students,
                },
            )
            if os.path.exists(temp_path):
                os.remove(temp_path)
            return JsonResponse({
                'status': 'success',
                'message': f'导入完成：共处理 {imported_cells} 个网格，新建学生 {created_students} 人'
            })
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': f'导入失败：{e}'}, status=400)

    return JsonResponse({'status': 'error', 'message': '未知操作'}, status=400)


def import_layout_excel_options_page(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/import_layout_options.html', {
        'classroom': classroom
    })


def import_students(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    
    if request.method == 'POST':
        action = request.POST.get('action', 'upload')
        pd_module = None
        if action in {'upload', 'confirm'}:
            try:
                pd_module = _require_pandas()
            except RuntimeError as e:
                return JsonResponse({'status': 'error', 'message': str(e)}, status=500)
        
        if action == 'upload' and request.FILES.get('excel_file'):
            excel_file = request.FILES['excel_file']
            clear_existing = request.POST.get('clear_existing') == '1'
            import_mode = _resolve_student_import_mode(request.POST.get('import_mode'), clear_existing)
            
            try:
                df = pd_module.read_excel(excel_file)
                columns = list(df.columns)

                def find_column(keys):
                    for key in keys:
                        for col in columns:
                            if key in str(col):
                                return col
                    return None

                def find_exact_column(candidates):
                    normalized_candidates = {str(item).strip().lower() for item in candidates}
                    for col in columns:
                        normalized_col = str(col).strip().lower()
                        if normalized_col in normalized_candidates:
                            return col
                    return None

                name_col = find_column(['姓名', '名字', '学生姓名', '学生'])
                score_col = find_exact_column(['总分', '学生总分'])
                
                if name_col and score_col:
                    student_id_col = find_column(['学号', '学生号', '编号', 'ID'])
                    gender_col = find_column(['性别', '男女性别'])
                    before_state = _capture_history_state(classroom)

                    result = _process_import(
                        classroom,
                        df,
                        name_col,
                        student_id_col,
                        gender_col,
                        score_col,
                        import_mode
                    )
                    _emit_plugin_hook(
                        'students_imported',
                        request=request,
                        classroom=classroom,
                        payload={'import_mode': import_mode, 'result': result},
                    )
                    _push_snapshot_action(
                        request,
                        classroom,
                        before_state,
                        'import_students',
                        extra={'import_mode': import_mode, 'result': result},
                    )
                    return JsonResponse({'status': 'success', 'message': _format_import_result_message(result)})
                
                import os
                import uuid
                from django.conf import settings
                
                file_id = str(uuid.uuid4())
                temp_dir = os.path.join(settings.BASE_DIR, 'temp_imports')
                os.makedirs(temp_dir, exist_ok=True)
                temp_path = os.path.join(temp_dir, f'{file_id}.xlsx')
                
                with open(temp_path, 'wb+') as destination:
                    for chunk in excel_file.chunks():
                        destination.write(chunk)
                
                df_preview = pd_module.read_excel(temp_path, header=None)
                preview_data = df_preview.head(20).fillna('').values.tolist()
                
                return JsonResponse({
                    'status': 'ambiguous',
                    'file_id': file_id,
                    'preview_data': preview_data,
                    'message': '仅当列名精确为“总分”或“学生总分”时才会自动导入，请手动匹配列'
                })

            except Exception as e:
                return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

        elif action == 'confirm':
            file_id = request.POST.get('file_id')
            start_row = int(request.POST.get('start_row', 0))
            name_col_idx = int(request.POST.get('name_col_index'))
            score_col_idx = request.POST.get('score_col_index')
            clear_existing = request.POST.get('clear_existing') == 'true'
            import_mode = _resolve_student_import_mode(request.POST.get('import_mode'), clear_existing)
            
            import os
            from django.conf import settings
            temp_path = os.path.join(settings.BASE_DIR, 'temp_imports', f'{file_id}.xlsx')
            
            if not os.path.exists(temp_path):
                return JsonResponse({'status': 'error', 'message': '临时文件已过期，请重新上传'}, status=400)
                
            try:
                before_state = _capture_history_state(classroom)
                df = pd_module.read_excel(temp_path, header=None)
                
                df_data = df.iloc[start_row + 1:].copy()
                
                df_data.columns = [i for i in range(df_data.shape[1])]
                
                name_col = name_col_idx
                score_col = int(score_col_idx) if score_col_idx and score_col_idx != '' else None
                
                result = _process_import(
                    classroom,
                    df_data,
                    name_col,
                    None,
                    None,
                    score_col,
                    import_mode
                )
                
                os.remove(temp_path)

                _emit_plugin_hook(
                    'students_imported',
                    request=request,
                    classroom=classroom,
                    payload={'import_mode': import_mode, 'result': result},
                )
                _push_snapshot_action(
                    request,
                    classroom,
                    before_state,
                    'import_students',
                    extra={'import_mode': import_mode, 'result': result},
                )
                return JsonResponse({'status': 'success', 'message': _format_import_result_message(result)})
            except Exception as e:
                return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

    return redirect('classroom_detail', pk=pk)


def import_students_options_page(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/import_students_options.html', {
        'classroom': classroom
    })


IMPORT_MODE_REPLACE = 'replace'
IMPORT_MODE_MATCH = 'match'
VALID_IMPORT_MODES = {IMPORT_MODE_REPLACE, IMPORT_MODE_MATCH}


def _resolve_student_import_mode(raw_mode, clear_existing=False):
    mode = str(raw_mode or '').strip().lower()
    if mode in VALID_IMPORT_MODES:
        return mode
    return IMPORT_MODE_REPLACE if clear_existing else IMPORT_MODE_MATCH


def _normalize_import_text(value):
    if _is_missing_import_value(value):
        return ''
    return str(value).strip()


def _parse_import_gender(value):
    text = _normalize_import_text(value).lower()
    if text in {'男', 'm', 'male'}:
        return 'M'
    if text in {'女', 'f', 'female'}:
        return 'F'
    return None


def _parse_import_score(value):
    if _is_missing_import_value(value):
        return 0
    if pd is not None:
        numeric_value = pd.to_numeric(value, errors='coerce')
    else:
        try:
            numeric_value = float(value)
        except (TypeError, ValueError):
            numeric_value = float('nan')
    if _is_missing_import_value(numeric_value):
        return 0
    return float(numeric_value)


def _format_import_result_message(result):
    if result['mode'] == IMPORT_MODE_REPLACE:
        return f"成功导入 {result['created']} 名学生"

    parts = [f"匹配更新 {result['updated']} 人"]
    if result['created'] > 0:
        parts.append(f"新增 {result['created']} 人")
    if result['skipped'] > 0:
        parts.append(f"未匹配 {result['skipped']} 人")
    return "匹配导入完成：" + "，".join(parts)


def _process_import(classroom, df, name_col, student_id_col, gender_col, score_col, import_mode=IMPORT_MODE_MATCH):
    import_mode = _resolve_student_import_mode(import_mode)
    created_count = 0
    updated_count = 0
    skipped_count = 0
    has_score_column = score_col is not None

    with transaction.atomic():
        if import_mode == IMPORT_MODE_REPLACE:
            classroom.students.all().delete()

        existing_students = list(classroom.students.all())
        should_match_existing = import_mode == IMPORT_MODE_MATCH and len(existing_students) > 0

        existing_by_id = {}
        existing_names = defaultdict(list)
        if should_match_existing:
            for student in existing_students:
                sid = _normalize_import_text(student.student_id).lower()
                if sid and sid not in existing_by_id:
                    existing_by_id[sid] = student
                normalized_name = _normalize_import_text(student.name).lower()
                if normalized_name:
                    existing_names[normalized_name].append(student)

        unique_name_map = {
            name_key: students[0]
            for name_key, students in existing_names.items()
            if len(students) == 1
        }

        def index_student(student):
            sid = _normalize_import_text(student.student_id).lower()
            if sid and sid not in existing_by_id:
                existing_by_id[sid] = student
            name_key = _normalize_import_text(student.name).lower()
            if not name_key:
                return
            if student not in existing_names[name_key]:
                existing_names[name_key].append(student)
            if len(existing_names[name_key]) == 1:
                unique_name_map[name_key] = student
            else:
                unique_name_map.pop(name_key, None)

        for _, row in df.iterrows():
            name = _normalize_import_text(row[name_col])
            if not name:
                continue

            if name.lower() in {'姓名', 'name'}:
                continue

            student_id = _normalize_import_text(row.get(student_id_col, '')) if student_id_col is not None else ''
            gender = _parse_import_gender(row.get(gender_col, '')) if gender_col is not None else None
            score_value = _parse_import_score(row.get(score_col, 0)) if has_score_column else 0

            if should_match_existing:
                matched_student = None
                if student_id:
                    matched_student = existing_by_id.get(student_id.lower())
                if not matched_student:
                    matched_student = unique_name_map.get(name.lower())

                if not matched_student:
                    created_student = Student.objects.create(
                        classroom=classroom,
                        name=name,
                        student_id=student_id,
                        gender=gender,
                        score=score_value
                    )
                    index_student(created_student)
                    created_count += 1
                    continue

                update_fields = []
                if matched_student.name != name:
                    matched_student.name = name
                    update_fields.append('name')
                if student_id and _normalize_import_text(matched_student.student_id) != student_id:
                    matched_student.student_id = student_id
                    update_fields.append('student_id')
                    existing_by_id[student_id.lower()] = matched_student
                if gender is not None and matched_student.gender != gender:
                    matched_student.gender = gender
                    update_fields.append('gender')
                if has_score_column and matched_student.score != score_value:
                    matched_student.score = score_value
                    update_fields.append('score')

                if update_fields:
                    matched_student.save(update_fields=update_fields)
                updated_count += 1
                continue

            Student.objects.create(
                classroom=classroom,
                name=name,
                student_id=student_id,
                gender=gender,
                score=score_value
            )
            created_count += 1

    return {
        'mode': import_mode,
        'created': created_count,
        'updated': updated_count,
        'skipped': skipped_count,
    }


def _build_constraint_maps(classroom, students):
    return compile_constraint_maps(classroom)


def _swap_seats(seat_a, seat_b):
    if not seat_a or not seat_b or seat_a.pk == seat_b.pk:
        return
    student_a = seat_a.student
    student_b = seat_b.student
    with transaction.atomic():
        seat_a.student = None
        seat_a.save(update_fields=['student'])
        
        seat_b.student = None
        seat_b.save(update_fields=['student'])

        seat_a.student = student_b
        seat_b.student = student_a
        seat_a.save(update_fields=['student'])
        seat_b.save(update_fields=['student'])
    _normalize_group_leaders(seat_a.classroom)


def _get_adjacent_seats(classroom, seat):
    """返回与给定座位相邻的有效座位对象列表。"""
    if not seat:
        return []
    coords = [
        (seat.row, seat.col - 1),
        (seat.row, seat.col + 1),
        (seat.row - 1, seat.col),
        (seat.row + 1, seat.col),
    ]
    seats = []
    for r, c in coords:
        s = classroom.seats.filter(row=r, col=c, cell_type=SeatCellType.SEAT).first()
        if s:
            seats.append(s)
    return seats


def _apply_internal_policy(classroom, request=None, trigger_student_id=None):
    return False


pass # 此部分代码未被披露至开源版本
    



def _seat_is_valid(student, seat, assignments, maps, required_group_map=None):
    fixed_seats, must_rows, must_cols, forbid_rows, forbid_cols, forbid_seats, must_pairs, forbid_pairs = maps
    sid = student.pk

    if required_group_map and sid in required_group_map:
        if seat.group_id != required_group_map[sid]:
            return False

    if sid in fixed_seats and (seat.row, seat.col) != fixed_seats[sid]:
        return False

    if sid in must_rows and seat.row not in must_rows[sid]:
        return False
    if sid in must_cols and seat.col not in must_cols[sid]:
        return False
    if sid in forbid_rows and seat.row in forbid_rows[sid]:
        return False
    if sid in forbid_cols and seat.col in forbid_cols[sid]:
        return False
    if sid in forbid_seats and (seat.row, seat.col) in forbid_seats[sid]:
        return False

    for other_id, dist in forbid_pairs.get(sid, []):
        if other_id in assignments:
            other_seat = assignments[other_id]
            if abs(seat.row - other_seat.row) + abs(seat.col - other_seat.col) <= dist:
                return False

    for other_id, dist in must_pairs.get(sid, []):
        if other_id in assignments:
            other_seat = assignments[other_id]
            if abs(seat.row - other_seat.row) + abs(seat.col - other_seat.col) > dist:
                return False

    return True


def _assign_pairs(students, seats, seat_map, assignments, maps, required_group_map=None):
    must_pairs = maps[6]
    available = seats[:]
    available_set = set(available)

    for student in students:
        if student.pk in assignments:
            continue
        pairs = must_pairs.get(student.pk, [])
        for other_id, dist in pairs:
            if other_id in assignments:
                other_seat = assignments[other_id]
                for seat in list(available):
                    if abs(seat.row - other_seat.row) + abs(seat.col - other_seat.col) <= dist:
                        if _seat_is_valid(student, seat, assignments, maps, required_group_map):
                            assignments[student.pk] = seat
                            if seat in available_set:
                                available_set.remove(seat)
                                available.remove(seat)
                            break
                continue
            other_student = next((s for s in students if s.pk == other_id), None)
            if not other_student or other_student.pk in assignments:
                continue

            for seat in list(available):
                if not _seat_is_valid(student, seat, assignments, maps, required_group_map):
                    continue
                for r in range(-dist, dist + 1):
                    for c in range(-dist, dist + 1):
                        if abs(r) + abs(c) > dist:
                            continue
                        if r == 0 and c == 0:
                            continue
                        neighbor = seat_map.get((seat.row + r, seat.col + c))
                        if neighbor and neighbor in available_set and neighbor.pk != seat.pk:
                            if _seat_is_valid(other_student, neighbor, assignments, maps, required_group_map):
                                assignments[student.pk] = seat
                                assignments[other_student.pk] = neighbor
                                if seat in available_set:
                                    available_set.remove(seat)
                                    available.remove(seat)
                                if neighbor in available_set:
                                    available_set.remove(neighbor)
                                    available.remove(neighbor)
                                break
                    if student.pk in assignments:
                        break
                if student.pk in assignments:
                    break
    return available


def _arrange_standard(classroom, students, seats, method):
    seats = [s for s in seats if s.cell_type == SeatCellType.SEAT]
    seat_map = _build_seat_map(seats)

    fixed_seats, must_rows, must_cols, forbid_rows, forbid_cols, forbid_seats, must_pairs, forbid_pairs = _build_constraint_maps(classroom, students)
    maps = (fixed_seats, must_rows, must_cols, forbid_rows, forbid_cols, forbid_seats, must_pairs, forbid_pairs)

    assignments = {}
    available = seats.copy()

    for student in students:
        if student.pk in fixed_seats:
            target = seat_map.get(fixed_seats[student.pk])
            if target and _seat_is_valid(student, target, assignments, maps):
                assignments[student.pk] = target
                if target in available:
                    available.remove(target)

    available = _assign_pairs(students, available, seat_map, assignments, maps)

    for student in students:
        if student.pk in assignments:
            continue
        for seat in list(available):
            if _seat_is_valid(student, seat, assignments, maps):
                assignments[student.pk] = seat
                available.remove(seat)
                break

    Seat.objects.filter(classroom=classroom).update(student=None)

    for student in students:
        seat = assignments.get(student.pk)
        if seat:
            seat.student = student
            seat.save(update_fields=['student'])


def _arrange_grouped(classroom, students, method):
    groups = list(classroom.groups.all())
    if not groups:
        return False

    group_seats = {group.pk: list(group.seats.filter(cell_type=SeatCellType.SEAT).order_by('row', 'col')) for group in groups}
    if not any(group_seats.values()):
        return False

    students_sorted = sorted(students, key=lambda s: s.score or 0, reverse=True)
    group_buckets = {group.pk: [] for group in groups}

    if method == 'group_balanced':
        num_groups = len(groups)
        students_per_group = len(students_sorted) // num_groups
        for i, group in enumerate(groups):
            start = i * students_per_group
            end = start + students_per_group if i < num_groups - 1 else len(students_sorted)
            group_buckets[group.pk] = students_sorted[start:end]
    elif method == 'group_mentor':
        pairs = []
        left = 0
        right = len(students_sorted) - 1
        while left <= right:
            if left == right:
                pairs.append([students_sorted[left]])
            else:
                pairs.append([students_sorted[left], students_sorted[right]])
            left += 1
            right -= 1
        
        pairs_with_sum = []
        for p in pairs:
             s = sum(st.score or 0 for st in p)
             pairs_with_sum.append((s, p))
        pairs_with_sum.sort(key=lambda x: x[0], reverse=True)
        

        group_sums = {g.pk: 0.0 for g in groups}
        
        for s_sum, p_students in pairs_with_sum:

            target_group = min(groups, key=lambda g: group_sums[g.pk])
            group_buckets[target_group.pk].extend(p_students)
            group_sums[target_group.pk] += s_sum
    else:
        return False

    fixed_seats, must_rows, must_cols, forbid_rows, forbid_cols, forbid_seats, must_pairs, forbid_pairs = _build_constraint_maps(classroom, students)
    maps = (fixed_seats, must_rows, must_cols, forbid_rows, forbid_cols, forbid_seats, must_pairs, forbid_pairs)

    group_candidate_seats = []
    required_group_map = {}
    for group in groups:
        seats = group_seats.get(group.pk, [])
        group_candidate_seats.extend(seats)
        if seats:
            for student in group_buckets[group.pk]:
                required_group_map[student.pk] = group.pk

    all_seat_cells = list(classroom.seats.filter(cell_type=SeatCellType.SEAT).order_by('row', 'col'))
    all_seat_map = _build_seat_map(all_seat_cells)
    assignments = {}
    available = group_candidate_seats[:]

    students_priority = sorted(students, key=lambda s: (s.pk not in required_group_map, -(s.score or 0), s.pk))

    for student in students_priority:
        if student.pk in fixed_seats:
            target = all_seat_map.get(fixed_seats[student.pk])
            if target and target in available and _seat_is_valid(student, target, assignments, maps, required_group_map):
                assignments[student.pk] = target
                available.remove(target)

    available = _assign_pairs(students_priority, available, all_seat_map, assignments, maps, required_group_map)

    for student in students_priority:
        if student.pk in assignments:
            continue
        for seat in list(available):
            if _seat_is_valid(student, seat, assignments, maps, required_group_map):
                assignments[student.pk] = seat
                available.remove(seat)
                break

    remaining_seats = [seat for seat in all_seat_cells if seat.pk not in {s.pk for s in assignments.values()}]
    remaining_students = [s for s in students_priority if s.pk not in assignments]

    for student in list(remaining_students):
        for seat in list(remaining_seats):
            if _seat_is_valid(student, seat, assignments, maps, required_group_map):
                assignments[student.pk] = seat
                remaining_seats.remove(seat)
                break

    Seat.objects.filter(classroom=classroom).update(student=None)
    for student in students:
        seat = assignments.get(student.pk)
        if seat:
            seat.student = student
            seat.save(update_fields=['student'])

    _normalize_group_leaders(classroom)
    return True


def _run_arrangement(classroom, method):
    students = list(classroom.students.all())
    if method in ['score_desc', 'score_asc', 'good_front', 'good_back']:
        seats = list(classroom.seats.select_related('student').order_by('col', 'row'))
    else:
        seats = list(classroom.seats.select_related('student').order_by('row', 'col'))

    seat_cells = [s for s in seats if s.cell_type == SeatCellType.SEAT]
    if len(seat_cells) < len(students):
        return False

    if method == 'random':
        random.shuffle(students)
    elif method == 'score_desc':
        students.sort(key=lambda s: s.score or 0, reverse=True)
    elif method == 'score_asc':
        students.sort(key=lambda s: s.score or 0)
    elif method == 'good_front':
        students.sort(key=lambda s: s.score or 0, reverse=True)
    elif method == 'good_back':
        students.sort(key=lambda s: s.score or 0, reverse=True)
        seats = list(reversed(seats))
    elif method == 'score_spread':
        students.sort(key=lambda s: s.score or 0)
        spread = []
        while students:
            spread.append(students.pop())
            if students:
                spread.append(students.pop(0))
        students = spread
    elif method in ['group_balanced', 'group_mentor']:
        return _arrange_grouped(classroom, students, method)

    _arrange_standard(classroom, students, seats, method)
    return True


def _attempt_auto_constraint_fix(classroom, preferred_method=None):
    methods = []
    if preferred_method:
        methods.append(preferred_method)
    methods.extend([
        'random',
        'score_spread',
        'score_desc',
        'score_asc',
        'good_front',
        'good_back',
        'group_balanced',
        'group_mentor',
    ])

    seen = set()
    ordered_methods = []
    for m in methods:
        if m not in seen:
            ordered_methods.append(m)
            seen.add(m)

    for method in ordered_methods:
        tries = 16 if method in ['random', 'score_spread'] else 5
        for _ in range(tries):
            try:
                with transaction.atomic():
                    ok = _run_arrangement(classroom, method)
                    if not ok:
                        raise ValueError('arrange_failed')
                    _stabilize_layout_with_rules(classroom)
                    issues = _layout_hard_issues(classroom)
                    if issues:
                        raise ValueError('constraint_failed')
                return True
            except Exception:
                continue
    return False


def auto_arrange_seats(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'POST':
        is_ajax = request.headers.get('x-requested-with') == 'XMLHttpRequest'

        def _arrange_error(message, status=400):
            if is_ajax:
                return JsonResponse({'status': 'error', 'message': message}, status=status)
            return HttpResponse(message, status=status)

        method = request.POST.get('method', 'random')
        students_count = classroom.students.count()
        seat_cells_count = classroom.seats.filter(cell_type=SeatCellType.SEAT).count()
        if seat_cells_count < students_count:
            message = f'可用座位不足(座位:{seat_cells_count} < 学生:{students_count})，无法保证100%入座，请在布局编辑中增加座位。'
            return _arrange_error(message, status=400)

        before_state = _capture_history_state(classroom)
        try:
            with transaction.atomic():
                if not _run_arrangement(classroom, method):
                    raise ValueError('未设置小组或小组没有座位')

                _stabilize_layout_with_rules(classroom, request)
                violations = _layout_hard_issues(classroom)
                if violations:
                    raise ValueError(f'约束未满足，排座已回滚：{_format_issues_preview(violations)}')
        except ValueError as e:
            if _attempt_auto_constraint_fix(classroom, preferred_method=method):
                _push_snapshot_action(
                    request,
                    classroom,
                    before_state,
                    'auto_arrange',
                    extra={'method': method, 'auto_fixed': True},
                )
                _emit_plugin_hook(
                    'seats_arranged',
                    request=request,
                    classroom=classroom,
                    payload={'method': method, 'auto_fixed': True},
                )
                if is_ajax:
                    return JsonResponse({'status': 'success', 'message': '已自动调整并满足约束'})
                return redirect('classroom_detail', pk=pk)
            return _arrange_error(str(e), status=400)

        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'auto_arrange',
            extra={'method': method, 'auto_fixed': False},
        )
        _emit_plugin_hook(
            'seats_arranged',
            request=request,
            classroom=classroom,
            payload={'method': method, 'auto_fixed': False},
        )
        if is_ajax:
            return JsonResponse({'status': 'success'})
        return redirect('classroom_detail', pk=pk)
    if request.headers.get('x-requested-with') == 'XMLHttpRequest':
        return JsonResponse({'status': 'error'}, status=400)
    return redirect('classroom_detail', pk=pk)


def _perform_move_fixed_group(classroom, student, target_seat):
    with transaction.atomic():
        current_seat = getattr(student, 'assigned_seat', None)
        target_student = target_seat.student

        if current_seat:
            current_seat.student = None
            current_seat.save(update_fields=['student'])

        if target_student:
            target_seat.student = None
            target_seat.save(update_fields=['student'])
        
        if current_seat and target_student:
            current_seat.student = target_student
            current_seat.save(update_fields=['student'])

        target_seat.student = student
        target_seat.save(update_fields=['student'])

    def _check_leader_lost(stu):
        if not stu: return
        led_group = getattr(stu, 'led_group', None)
        if led_group:
            current_s = getattr(stu, 'assigned_seat', None)
            if not current_s or current_s.group != led_group:
                led_group.leader = None
                led_group.save(update_fields=['leader'])

    if student: student.refresh_from_db()
    if target_student: target_student.refresh_from_db()
    
    _check_leader_lost(student)
    _check_leader_lost(target_student)

    action = {
        'type': 'move',
        'student_id': student.pk,
        'from_row': current_seat.row if current_seat else None,
        'from_col': current_seat.col if current_seat else None,
        'to_row': target_seat.row,
        'to_col': target_seat.col,
        'target_student_id': target_student.pk if target_student else None,
        'group_move_mode': GROUP_MOVE_MODE_FIXED,
    }
    return action


def _perform_move_follow_group(classroom, student, target_seat):
    current_seat = getattr(student, 'assigned_seat', None)
    target_student = target_seat.student
    source_group = current_seat.group if current_seat else None
    target_group = target_seat.group
    affected_group_ids = {
        group_id
        for group_id in [
            source_group.pk if source_group else None,
            target_group.pk if target_group else None,
        ]
        if group_id
    }

    with transaction.atomic():
        if current_seat:
            current_seat.student = None
            current_seat.group = None
            current_seat.save(update_fields=['student', 'group'])

        target_seat.student = None
        target_seat.group = None
        target_seat.save(update_fields=['student', 'group'])

        if current_seat and target_student:
            current_seat.student = target_student
            current_seat.group = target_group
            current_seat.save(update_fields=['student', 'group'])

        target_seat.student = student
        target_seat.group = source_group if current_seat else target_group
        target_seat.save(update_fields=['student', 'group'])

    if affected_group_ids:
        _normalize_group_leaders(classroom, affected_group_ids)

    action = {
        'type': 'move',
        'student_id': student.pk,
        'from_row': current_seat.row if current_seat else None,
        'from_col': current_seat.col if current_seat else None,
        'to_row': target_seat.row,
        'to_col': target_seat.col,
        'target_student_id': target_student.pk if target_student else None,
        'group_move_mode': GROUP_MOVE_MODE_FOLLOW,
    }
    return action


def _perform_move(classroom, student, target_seat, *, group_move_mode=GROUP_MOVE_MODE_FIXED):
    normalized_group_mode = _normalize_group_move_mode(group_move_mode)
    current_seat = getattr(student, 'assigned_seat', None)
    if normalized_group_mode == GROUP_MOVE_MODE_FOLLOW and current_seat:
        return _perform_move_follow_group(classroom, student, target_seat)
    return _perform_move_fixed_group(classroom, student, target_seat)


def move_student(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            student_id = data.get('student_id')
            target_row = int(data.get('row'))
            target_col = int(data.get('col'))
            group_move_mode = _normalize_group_move_mode(data.get('group_move_mode'))

            student = get_object_or_404(Student, pk=student_id, classroom=classroom)
            target_seat = get_object_or_404(Seat, classroom=classroom, row=target_row, col=target_col)

            if target_seat.cell_type != SeatCellType.SEAT:
                return JsonResponse({'status': 'error', 'message': '目标位置不可入座'}, status=400)

            before_state = _capture_history_state(classroom)
            with transaction.atomic():
                action = _perform_move(classroom, student, target_seat, group_move_mode=group_move_mode)
                violations = _stabilize_layout_with_rules(classroom, request, trigger_student_id=student.pk)
                if violations:
                    raise ValueError(f'移动失败：{_format_issues_preview(violations)}')
            _push_snapshot_action(request, classroom, before_state, 'move', extra=action)
            _emit_plugin_hook(
                'student_moved',
                request=request,
                classroom=classroom,
                payload=action,
            )
            return JsonResponse({'status': 'success'})
        except ValueError as e:
            return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    return JsonResponse({'status': 'error'}, status=400)


@require_POST
def move_students_batch(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body or '{}')
        moves = data.get('moves') or []
        group_move_mode = _normalize_group_move_mode(data.get('group_move_mode'))
        if not isinstance(moves, list) or not moves:
            return JsonResponse({'status': 'error', 'message': '缺少批量移动数据'}, status=400)

        parsed_moves = []
        student_ids = []
        target_coords = []
        seen_students = set()
        seen_targets = set()

        for item in moves:
            sid = int(item.get('student_id'))
            row_value = item.get('row')
            col_value = item.get('col')
            is_clear_move = row_value in (None, '') and col_value in (None, '')

            if not is_clear_move and (row_value in (None, '') or col_value in (None, '')):
                raise ValueError('批量移动参数错误：row/col 需同时提供或同时为空')

            row = None
            col = None
            if not is_clear_move:
                row = int(row_value)
                col = int(col_value)

            if sid in seen_students:
                raise ValueError('同一学生重复出现在批量移动中')
            if not is_clear_move and (row, col) in seen_targets:
                raise ValueError('目标座位存在重复')
            seen_students.add(sid)
            if not is_clear_move:
                seen_targets.add((row, col))
            student_ids.append(sid)
            if not is_clear_move:
                target_coords.append((row, col))
            parsed_moves.append({
                'student_id': sid,
                'row': row,
                'col': col,
                'is_clear_move': is_clear_move,
            })

        students_map = classroom.students.in_bulk(student_ids)
        if len(students_map) != len(student_ids):
            raise ValueError('存在不属于当前班级的学生')
        use_follow_group_mode = (
            group_move_mode == GROUP_MOVE_MODE_FOLLOW
            and all(not item['is_clear_move'] for item in parsed_moves)
            and all(getattr(student, 'assigned_seat', None) for student in students_map.values())
        )

        seat_q = models.Q()
        for row, col in target_coords:
            seat_q |= (models.Q(row=row) & models.Q(col=col))
        seat_map = {}
        if seat_q:
            for seat in classroom.seats.filter(seat_q):
                seat_map[(seat.row, seat.col)] = seat

        for row, col in target_coords:
            seat = seat_map.get((row, col))
            if not seat:
                raise ValueError(f'目标座位不存在: {row}-{col}')
            if seat.cell_type != SeatCellType.SEAT:
                raise ValueError(f'目标位置不可入座: {row}-{col}')

        actions = []
        trigger_student_id = None
        before_state = _capture_history_state(classroom)
        with transaction.atomic():
            for item in parsed_moves:
                sid = item.get('student_id')
                row = item.get('row')
                col = item.get('col')
                is_clear_move = bool(item.get('is_clear_move'))
                student = students_map.get(sid)
                if not student:
                    raise ValueError('存在不属于当前班级的学生')

                if is_clear_move:
                    current_seat = getattr(student, 'assigned_seat', None)
                    from_row = current_seat.row if current_seat else None
                    from_col = current_seat.col if current_seat else None

                    if current_seat:
                        current_seat.student = None
                        current_seat.save(update_fields=['student'])

                    led_group = getattr(student, 'led_group', None)
                    if led_group:
                        led_group.leader = None
                        led_group.save(update_fields=['leader'])

                    action = {
                        'type': 'move',
                        'student_id': student.pk,
                        'from_row': from_row,
                        'from_col': from_col,
                        'to_row': None,
                        'to_col': None,
                        'target_student_id': None,
                    }
                else:
                    target_seat = seat_map.get((row, col))
                    action = _perform_move(
                        classroom,
                        student,
                        target_seat,
                        group_move_mode=GROUP_MOVE_MODE_FOLLOW if use_follow_group_mode else GROUP_MOVE_MODE_FIXED,
                    )
                actions.append(action)
                if trigger_student_id is None:
                    trigger_student_id = sid

            violations = _stabilize_layout_with_rules(classroom, request, trigger_student_id=trigger_student_id)
            if violations:
                raise ValueError(f'批量移动失败：{_format_issues_preview(violations)}')

        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'move_batch',
            extra={'type': 'move_batch', 'items': actions},
        )
        _emit_plugin_hook(
            'students_moved_batch',
            request=request,
            classroom=classroom,
            payload={'moved': len(actions), 'items': actions},
        )
        return JsonResponse({'status': 'success', 'moved': len(actions)})
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


def search_students(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    query = request.GET.get('q', '').strip().lower()
    tag_filter_ids = _normalize_id_list(request.GET.get('tag_ids') or request.GET.get('tags') or request.GET.get('tag_id'))
    tag_match = str(request.GET.get('tag_match') or request.GET.get('match') or 'any').strip().lower()
    if tag_match not in {'any', 'all', 'none'}:
        tag_match = 'any'
    if not query and not tag_filter_ids:
        return JsonResponse({'students': []})

    all_students = list(classroom.students.select_related('assigned_seat').all())
    tag_map = _build_student_tag_map(classroom, [student.pk for student in all_students])
    matches = []

    for student in all_students:
        student_tags = tag_map.get(student.pk, [])
        student_tag_ids = {tag['id'] for tag in student_tags}
        tag_filter_set = set(tag_filter_ids)
        if tag_filter_set:
            if tag_match == 'all' and not tag_filter_set.issubset(student_tag_ids):
                continue
            if tag_match == 'any' and not (tag_filter_set & student_tag_ids):
                continue
            if tag_match == 'none' and (tag_filter_set & student_tag_ids):
                continue

        name_lower = student.name.lower()
        pinyin_parts = [part.lower() for part in lazy_pinyin(student.name) if part]
        pinyin = ''.join(pinyin_parts)
        pinyin_initials = ''.join(part[0] for part in pinyin_parts if part)
        student_number = str(student.student_id or '').lower()
        tag_text = ' '.join(tag['name'].lower() for tag in student_tags)

        if not query or query in name_lower or query in pinyin or query in pinyin_initials or query in student_number or query in tag_text:
            seat_info = None
            assigned_seat = getattr(student, 'assigned_seat', None)
            if assigned_seat:
                seat_info = {'row': assigned_seat.row, 'col': assigned_seat.col}
            matches.append({
                'id': student.pk,
                'name': student.name,
                'student_id': student.student_id or '',
                'seat': seat_info,
                'tags': student_tags,
            })

    return JsonResponse({'students': matches})


@require_POST
def clear_seat(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
        row = int(data.get('row'))
        col = int(data.get('col'))
        seat = get_object_or_404(Seat, classroom=classroom, row=row, col=col)
        if not seat.student:
            return JsonResponse({'status': 'error', 'message': '座位为空'}, status=400)
        action = {
            'type': 'move',
            'student_id': seat.student.pk,
            'from_row': seat.row,
            'from_col': seat.col,
            'to_row': None,
            'to_col': None,
            'target_student_id': None
        }
        before_state = _capture_history_state(classroom)
        with transaction.atomic():
            if seat.student:
                student = seat.student
                led_group = getattr(student, 'led_group', None)
                if led_group:
                    led_group.leader = None
                    led_group.save(update_fields=['leader'])

            seat.student = None
            seat.save(update_fields=['student'])

            violations = _stabilize_layout_with_rules(classroom, request)
            if violations:
                raise ValueError(f'清空失败：{_format_issues_preview(violations)}')
        _push_snapshot_action(request, classroom, before_state, 'clear_seat', extra=action)
        return JsonResponse({'status': 'success'})
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@require_POST
def assign_student(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
        student_id = data.get('student_id')
        row = int(data.get('row'))
        col = int(data.get('col'))
        group_move_mode = _normalize_group_move_mode(data.get('group_move_mode'))
        student = get_object_or_404(Student, pk=student_id, classroom=classroom)
        target_seat = get_object_or_404(Seat, classroom=classroom, row=row, col=col)
        if target_seat.cell_type != SeatCellType.SEAT:
            return JsonResponse({'status': 'error', 'message': '目标位置不可入座'}, status=400)
        before_state = _capture_history_state(classroom)
        with transaction.atomic():
            action = _perform_move(classroom, student, target_seat, group_move_mode=group_move_mode)
            violations = _stabilize_layout_with_rules(classroom, request, trigger_student_id=student.pk)
            if violations:
                raise ValueError(f'指派失败：{_format_issues_preview(violations)}')
        _push_snapshot_action(request, classroom, before_state, 'assign_student', extra=action)
        _emit_plugin_hook(
            'student_assigned',
            request=request,
            classroom=classroom,
            payload=action,
        )
        return JsonResponse({'status': 'success'})
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


def delete_student(request, pk, student_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    student = get_object_or_404(Student, pk=student_id, classroom=classroom)
    is_ajax = request.headers.get('x-requested-with') == 'XMLHttpRequest' or request.headers.get('sec-fetch-mode') == 'cors'
    if classroom.seats.filter(student_id=student.pk).exists():
        if is_ajax:
            return JsonResponse({'status': 'error', 'message': '该学生已入座，请先清空座位后再删除'}, status=400)
        return redirect('classroom_detail', pk=pk)
    before_state = _capture_history_state(classroom)
    deleted_student = {'student_id': student.pk, 'student_name': student.name}
    student.delete()
    _push_snapshot_action(request, classroom, before_state, 'delete_student', extra=deleted_student)
    if is_ajax:
        return JsonResponse({'status': 'success'})
    return redirect('classroom_detail', pk=pk)


@require_POST
def add_student(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
    except (json.JSONDecodeError, ValueError):
        return JsonResponse({'status': 'error', 'message': '无效的请求数据'}, status=400)
    name = (data.get('name') or '').strip()
    if not name:
        return JsonResponse({'status': 'error', 'message': '姓名不能为空'}, status=400)
    student_id_val = (data.get('student_id') or '').strip() or None
    gender = data.get('gender') or None
    if gender not in ('M', 'F'):
        gender = None
    try:
        score = float(data.get('score', 0) or 0)
    except (ValueError, TypeError):
        score = 0
    before_state = _capture_history_state(classroom)
    with transaction.atomic():
        student = Student.objects.create(
            classroom=classroom,
            name=name,
            student_id=student_id_val,
            gender=gender,
            score=score,
        )
        tag_changes = _apply_student_tag_payload(classroom, [student], data, default_mode='set')
    _push_snapshot_action(request, classroom, before_state, 'add_student',
                          extra={'student_id': student.pk, 'student_name': student.name, **tag_changes})
    return JsonResponse({'status': 'success', 'student_pk': student.pk, 'student': _serialize_student_profile(
        student,
        classroom=classroom,
        tag_map=_build_student_tag_map(classroom, [student.pk]),
    )})


@require_POST
def update_student(request, pk, student_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    student = get_object_or_404(Student, pk=student_id, classroom=classroom)
    try:
        data = json.loads(request.body)
    except (json.JSONDecodeError, ValueError):
        return JsonResponse({'status': 'error', 'message': '无效的请求数据'}, status=400)
    name = (data.get('name') or '').strip()
    if not name:
        return JsonResponse({'status': 'error', 'message': '姓名不能为空'}, status=400)
    before_state = _capture_history_state(classroom)
    student.name = name
    student.student_id = (data.get('student_id') or '').strip() or None
    gender = data.get('gender') or None
    student.gender = gender if gender in ('M', 'F') else None
    try:
        student.score = float(data.get('score', 0) or 0)
    except (ValueError, TypeError):
        student.score = 0
    with transaction.atomic():
        student.save()
        tag_changes = _apply_student_tag_payload(classroom, [student], data, default_mode='set')
    _push_snapshot_action(request, classroom, before_state, 'update_student',
                          extra={'student_id': student.pk, 'student_name': student.name, **tag_changes})
    return JsonResponse({'status': 'success', 'student': _serialize_student_profile(
        student,
        classroom=classroom,
        tag_map=_build_student_tag_map(classroom, [student.pk]),
    )})


@require_http_methods(["GET", "POST"])
def student_tags(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'GET':
        tag_rules = list(classroom.student_tag_rules.select_related('tag').prefetch_related('tag__memberships').all())
        tag_rule_items, tag_rule_metrics = serialize_tag_rules(classroom, tag_rules=tag_rules)
        return JsonResponse({
            'status': 'success',
            'tags': _serialize_student_tag_catalog(classroom),
            'tag_rule_types': get_tag_rule_type_definitions(),
            'tag_rules': tag_rule_items,
            'tag_rule_metrics': tag_rule_metrics,
        })

    try:
        data = _request_payload(request)
        name = _normalize_tag_name(data.get('name'))
        if not name:
            return JsonResponse({'status': 'error', 'message': '标签名称不能为空'}, status=400)
        description = str(data.get('description') or '').strip()
        sort_order = _safe_int(data.get('sort_order'), 0)
        color = _normalize_tag_color(data.get('color'))
        before_state = _capture_history_state(classroom)
        tag = StudentTag.objects.create(
            classroom=classroom,
            name=name,
            color=color,
            description=description[:160],
            sort_order=max(0, sort_order),
        )
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'create_student_tag',
            extra={'tag_id': tag.pk, 'tag_name': tag.name},
        )
        return JsonResponse({'status': 'success', 'tag': _serialize_student_tag(tag, member_count=0, rule_count=0, include_urls=True)})
    except IntegrityError:
        return JsonResponse({'status': 'error', 'message': '同名标签已存在'}, status=400)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)


@require_POST
def update_student_tag(request, pk, tag_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    tag = get_object_or_404(StudentTag, pk=tag_id, classroom=classroom)
    try:
        data = _request_payload(request)
        name = _normalize_tag_name(data.get('name')) if 'name' in data else tag.name
        if not name:
            return JsonResponse({'status': 'error', 'message': '标签名称不能为空'}, status=400)
        before_state = _capture_history_state(classroom)
        tag.name = name
        if 'color' in data:
            tag.color = _normalize_tag_color(data.get('color'))
        if 'description' in data:
            tag.description = str(data.get('description') or '').strip()[:160]
        if 'sort_order' in data:
            tag.sort_order = max(0, _safe_int(data.get('sort_order'), 0))
        tag.save()
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'update_student_tag',
            extra={'tag_id': tag.pk, 'tag_name': tag.name},
        )
        return JsonResponse({'status': 'success', 'tag': _serialize_student_tag(tag, include_urls=True)})
    except IntegrityError:
        return JsonResponse({'status': 'error', 'message': '同名标签已存在'}, status=400)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)


@require_POST
def delete_student_tag(request, pk, tag_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    tag = get_object_or_404(StudentTag, pk=tag_id, classroom=classroom)
    before_state = _capture_history_state(classroom)
    tag_name = tag.name
    tag.delete()
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'delete_student_tag',
        extra={'tag_id': tag_id, 'tag_name': tag_name},
    )
    return JsonResponse({'status': 'success'})


@require_POST
def assign_student_tags(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = _request_payload(request)
        student_ids = _normalize_id_list(data.get('student_ids') or data.get('student_id'))
        if not student_ids:
            return JsonResponse({'status': 'error', 'message': '请选择学生'}, status=400)
        students_map = classroom.students.in_bulk(student_ids)
        if len(students_map) != len(student_ids):
            return JsonResponse({'status': 'error', 'message': '存在不属于当前班级的学生'}, status=400)
        mode = str(data.get('mode') or 'add').strip().lower()
        if mode not in {'add', 'remove', 'set', 'toggle'}:
            mode = 'add'
        tags = _resolve_tags_from_payload(classroom, data, allow_create=mode != 'remove')
        if not tags and mode != 'set':
            return JsonResponse({'status': 'error', 'message': '请选择标签'}, status=400)

        before_state = _capture_history_state(classroom)
        added = 0
        removed = 0
        with transaction.atomic():
            for student_id in student_ids:
                item_added, item_removed = _set_student_tags(classroom, students_map[student_id], tags, mode=mode)
                added += item_added
                removed += item_removed
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'assign_student_tags',
            extra={
                'student_ids': student_ids,
                'tag_ids': [tag.pk for tag in tags],
                'mode': mode,
                'added': added,
                'removed': removed,
            },
        )
        tag_map = _build_student_tag_map(classroom, student_ids)
        return JsonResponse({
            'status': 'success',
            'added': added,
            'removed': removed,
            'students': [
                {
                    'id': student_id,
                    'tags': tag_map.get(student_id, []),
                }
                for student_id in student_ids
            ],
            'tags': _serialize_student_tag_catalog(classroom),
        })
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)


def search_students_by_tags(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    query = str(request.GET.get('q') or '').strip().lower()
    match_mode = str(request.GET.get('match') or 'any').strip().lower()
    if match_mode not in {'any', 'all', 'none'}:
        match_mode = 'any'
    include_untagged = _parse_bool(request.GET.get('untagged'))
    limit = max(1, min(300, _safe_int(request.GET.get('limit'), 100)))
    tag_filter_ids = _normalize_id_list(request.GET.get('tag_ids') or request.GET.get('tags') or request.GET.get('tag_id'))
    tag_names = _normalize_name_list(request.GET.get('tag_names') or request.GET.get('tag_name'))
    if tag_names:
        found_tags = list(classroom.student_tags.filter(name__in=tag_names))
        if len(found_tags) != len(tag_names):
            return JsonResponse({'status': 'success', 'students': [], 'total': 0})
        tag_filter_ids.extend([tag.pk for tag in found_tags if tag.pk not in tag_filter_ids])
    tag_filter_set = set(tag_filter_ids)

    students = list(classroom.students.select_related('assigned_seat__group').all().order_by('name', 'pk'))
    tag_map = _build_student_tag_map(classroom, [student.pk for student in students])
    results = []
    for student in students:
        student_tags = tag_map.get(student.pk, [])
        student_tag_ids = {tag['id'] for tag in student_tags}
        if include_untagged and student_tag_ids:
            continue
        if tag_filter_set:
            if match_mode == 'all' and not tag_filter_set.issubset(student_tag_ids):
                continue
            if match_mode == 'any' and not (tag_filter_set & student_tag_ids):
                continue
            if match_mode == 'none' and (tag_filter_set & student_tag_ids):
                continue
        if query:
            name_lower = student.name.lower()
            sid_lower = str(student.student_id or '').lower()
            pinyin_parts = [part.lower() for part in lazy_pinyin(student.name) if part]
            pinyin = ''.join(pinyin_parts)
            pinyin_initials = ''.join(part[0] for part in pinyin_parts if part)
            tag_text = ' '.join(tag['name'].lower() for tag in student_tags)
            if query not in name_lower and query not in sid_lower and query not in pinyin and query not in pinyin_initials and query not in tag_text:
                continue
        seat = getattr(student, 'assigned_seat', None)
        group = seat.group if seat else None
        results.append({
            'id': student.pk,
            'name': student.name,
            'student_id': student.student_id or '',
            'gender': student.gender or '',
            'score': student.score or 0,
            'score_display': student.display_score if student.score is not None else '',
            'seat': {'row': seat.row, 'col': seat.col} if seat else None,
            'group': {'id': group.pk, 'name': group.name} if group else None,
            'tags': student_tags,
        })
        if len(results) >= limit:
            break
    return JsonResponse({'status': 'success', 'students': results, 'total': len(results)})


def _tag_rule_payload(request):
    if str(request.headers.get('content-type') or '').lower().find('application/json') >= 0:
        return _request_payload(request)
    payload = _request_payload(request)
    if 'tag_id' not in payload and request.POST.get('tag'):
        payload['tag_id'] = request.POST.get('tag')
    return payload


@require_POST
def create_tag_rule(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        cleaned = normalize_tag_rule_payload(classroom, _tag_rule_payload(request))
        before_state = _capture_history_state(classroom)
        rule = StudentTagRule.objects.create(classroom=classroom, **cleaned)
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'create_tag_rule',
            extra={'tag_rule_id': rule.pk, 'tag_id': rule.tag_id, 'rule_type': rule.rule_type},
        )
        return JsonResponse({'status': 'success', 'tag_rule': serialize_tag_rules(classroom, [rule])[0][0]})
    except ConstraintServiceError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)


@require_POST
def update_tag_rule(request, pk, rule_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    rule = get_object_or_404(StudentTagRule, pk=rule_id, classroom=classroom)
    try:
        cleaned = normalize_tag_rule_payload(classroom, _tag_rule_payload(request), instance=rule)
        before_state = _capture_history_state(classroom)
        for field, value in cleaned.items():
            setattr(rule, field, value)
        rule.save()
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'update_tag_rule',
            extra={'tag_rule_id': rule.pk, 'tag_id': rule.tag_id, 'rule_type': rule.rule_type},
        )
        return JsonResponse({'status': 'success', 'tag_rule': serialize_tag_rules(classroom, [rule])[0][0]})
    except ConstraintServiceError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)


@require_POST
def toggle_tag_rule(request, pk, rule_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    rule = get_object_or_404(StudentTagRule, pk=rule_id, classroom=classroom)
    try:
        data = _tag_rule_payload(request)
        desired_enabled = data.get('enabled')
        next_enabled = (not rule.enabled) if desired_enabled in (None, '') else _parse_bool(desired_enabled)
        before_state = _capture_history_state(classroom)
        rule.enabled = next_enabled
        rule.save(update_fields=['enabled', 'updated_at'])
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'toggle_tag_rule',
            extra={'tag_rule_id': rule.pk, 'tag_id': rule.tag_id, 'rule_type': rule.rule_type, 'enabled': rule.enabled},
        )
        return JsonResponse({'status': 'success', 'tag_rule': serialize_tag_rules(classroom, [rule])[0][0]})
    except ValueError as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=400)


@require_POST
def delete_tag_rule(request, pk, rule_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    rule = get_object_or_404(StudentTagRule, pk=rule_id, classroom=classroom)
    before_state = _capture_history_state(classroom)
    rule_type = rule.rule_type
    tag_id = rule.tag_id
    rule.delete()
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'delete_tag_rule',
        extra={'tag_rule_id': rule_id, 'tag_id': tag_id, 'rule_type': rule_type},
    )
    return JsonResponse({'status': 'success'})


@require_POST
def update_cell_type(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
        row = int(data.get('row'))
        col = int(data.get('col'))
        cell_type = data.get('cell_type')
        if cell_type not in [c.value for c in SeatCellType]:
            return JsonResponse({'status': 'error', 'message': '类型不合法'}, status=400)
        seat = get_object_or_404(Seat, classroom=classroom, row=row, col=col)
        action = {
            'type': 'cell_type',
            'row': seat.row,
            'col': seat.col,
            'before': seat.cell_type,
            'after': cell_type,
            'prev_student_id': seat.student.pk if seat.student else None,
            'prev_group_id': seat.group.pk if seat.group else None
        }
        before_state = _capture_history_state(classroom)
        seat.cell_type = cell_type
        if cell_type != SeatCellType.SEAT:
            seat.student = None
            seat.group = None
        seat.save(update_fields=['cell_type', 'student', 'group'])
        _push_snapshot_action(request, classroom, before_state, 'cell_type', extra=action)
        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@require_POST
def create_group(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    name = str(request.POST.get('name', '')).strip()
    if not name:
        if _is_ajax_request(request):
            return JsonResponse({'status': 'error', 'message': '小组名称不能为空'}, status=400)
        return redirect('classroom_detail', pk=pk)
    before_state = _capture_history_state(classroom)
    group, created = SeatGroup.objects.get_or_create(classroom=classroom, name=name)
    if not created:
        if _is_ajax_request(request):
            return JsonResponse({'status': 'error', 'message': '小组名称已存在'}, status=400)
        return redirect('classroom_detail', pk=pk)

    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'create_group',
        extra={'group_id': group.pk, 'group_name': group.name},
    )

    _emit_plugin_hook(
        'group_created',
        request=request,
        classroom=classroom,
        payload={'group_id': group.pk, 'group_name': group.name},
    )
    if _is_ajax_request(request):
        return JsonResponse({'status': 'success', 'group': {'id': group.pk, 'name': group.name}})
    return redirect('classroom_detail', pk=pk)


def _next_group_names(reference_name, existing_names, count):
    if count <= 0:
        return []
    existing = set(str(name) for name in existing_names if str(name).strip())
    generated = []
    ref = str(reference_name or '').strip()
    if not ref:
        ref = '小组'

    if ref.isdigit():
        n = int(ref) + 1
        while len(generated) < count:
            candidate = str(n)
            if candidate not in existing:
                generated.append(candidate)
                existing.add(candidate)
            n += 1
        return generated

    prefix = ref
    start = 1
    m = re.match(r'^(.*?)(\d+)$', ref)
    if m and m.group(1):
        prefix = m.group(1)
        start = int(m.group(2)) + 1

    max_used = 0
    pattern = re.compile(rf'^{re.escape(prefix)}(\d+)$')
    for name in existing:
        mm = pattern.match(name)
        if mm:
            max_used = max(max_used, int(mm.group(1)))
    n = max(start, max_used + 1)
    while len(generated) < count:
        candidate = f'{prefix}{n}'
        if candidate not in existing:
            generated.append(candidate)
            existing.add(candidate)
        n += 1
    return generated


def _detect_group_style(reference_group):
    seats = list(
        reference_group.seats
        .filter(cell_type=SeatCellType.SEAT)
        .order_by('row', 'col')
    )
    if len(seats) < 2:
        return 'horizontal'

    rows = [s.row for s in seats]
    cols = [s.col for s in seats]
    unique_rows = len(set(rows))
    unique_cols = len(set(cols))
    if unique_cols == 1 and unique_rows > 1:
        return 'vertical'
    if unique_rows == 1 and unique_cols > 1:
        return 'horizontal'

    min_row, max_row = min(rows), max(rows)
    min_col, max_col = min(cols), max(cols)
    area = (max_row - min_row + 1) * (max_col - min_col + 1)
    density = len(seats) / max(area, 1)
    if unique_rows > 1 and unique_cols > 1 and density >= 0.6:
        return 'nearby'

    horizontal_pairs = 0
    vertical_pairs = 0
    adjacent_pairs = 0
    for i in range(len(seats)):
        for j in range(i + 1, len(seats)):
            a = seats[i]
            b = seats[j]
            if a.row == b.row:
                horizontal_pairs += 1
            if a.col == b.col:
                vertical_pairs += 1
            if abs(a.row - b.row) + abs(a.col - b.col) == 1:
                adjacent_pairs += 1

    if vertical_pairs > horizontal_pairs * 1.3:
        return 'vertical'
    if horizontal_pairs > vertical_pairs * 1.3:
        return 'horizontal'

    if density >= 0.5 or adjacent_pairs > 0:
        return 'nearby'
    return 'horizontal'


def _normalize_shape_points(points):
    if not points:
        return []
    min_row = min(r for r, _ in points)
    min_col = min(c for _, c in points)
    return sorted((r - min_row, c - min_col) for r, c in points)


def _transform_shape_points(points, mode):
    transformed = []
    for r, c in points:
        if mode == 'r90':
            transformed.append((c, -r))
        elif mode == 'r180':
            transformed.append((-r, -c))
        elif mode == 'r270':
            transformed.append((-c, r))
        else:
            transformed.append((r, c))
    return transformed


def _build_nearby_shape_profile(reference_group):
    seats = list(
        reference_group.seats
        .filter(cell_type=SeatCellType.SEAT)
        .order_by('row', 'col')
    )
    if len(seats) < 2:
        return None

    raw_points = [(s.row, s.col) for s in seats]
    normalized_points = _normalize_shape_points(raw_points)
    if not normalized_points:
        return None

    max_r = max(r for r, _ in normalized_points)
    max_c = max(c for _, c in normalized_points)
    height = max_r + 1
    width = max_c + 1
    count = len(normalized_points)
    area = max(1, height * width)
    density = count / area

    if count == area:
        shape_name = f'block_{height}x{width}'
    elif count == 3 and height == 2 and width == 2:
        shape_name = 'corner_2x2'
    elif height == 1 or width == 1:
        shape_name = 'line'
    else:
        shape_name = 'irregular'

    variants = []
    seen = set()
    for mode in ('r0', 'r90', 'r180', 'r270'):
        variant_points = tuple(_normalize_shape_points(_transform_shape_points(normalized_points, mode)))
        if variant_points and variant_points not in seen:
            seen.add(variant_points)
            variants.append(list(variant_points))

    return {
        'shape_name': shape_name,
        'count': count,
        'width': width,
        'height': height,
        'density': density,
        'variants': variants,
    }


def _pick_nearby_cluster_greedy(remaining, target_count):
    if not remaining or target_count <= 0:
        return []
    cluster = [remaining[0]]
    while len(cluster) < target_count and len(cluster) < len(remaining):
        best_idx = None
        best_key = None
        for idx, seat in enumerate(remaining):
            if seat in cluster:
                continue
            min_dist = min(abs(seat.row - s.row) + abs(seat.col - s.col) for s in cluster)
            key = (min_dist, seat.row, seat.col)
            if best_key is None or key < best_key:
                best_key = key
                best_idx = idx
        if best_idx is None:
            break
        cluster.append(remaining[best_idx])
    return cluster


def _pick_nearby_cluster_by_shape(remaining, target_count, shape_profile):
    if not remaining or target_count <= 0:
        return []
    variants = (shape_profile or {}).get('variants') or []
    if not variants:
        return _pick_nearby_cluster_greedy(remaining, target_count)

    remaining = sorted(remaining, key=lambda s: (s.row, s.col))
    anchor = remaining[0]
    anchor_pos = (anchor.row, anchor.col)
    pos_to_seat = {(s.row, s.col): s for s in remaining}
    pos_set = set(pos_to_seat.keys())
    best_translated_positions = None
    best_matched_positions = None
    best_score = None

    for variant_idx, variant in enumerate(variants):
        if not variant:
            continue
        variant_points = [tuple(p) for p in variant]
        for pattern_point in variant_points:
            dr = anchor_pos[0] - pattern_point[0]
            dc = anchor_pos[1] - pattern_point[1]
            translated_positions = [(dr + pr, dc + pc) for pr, pc in variant_points]
            matched = [pos for pos in translated_positions if pos in pos_set]
            if anchor_pos not in matched:
                continue
            if not matched:
                continue

            rows = [r for r, _ in matched]
            cols = [c for _, c in matched]
            bbox_area = (max(rows) - min(rows) + 1) * (max(cols) - min(cols) + 1)
            distance_sum = sum(abs(r - anchor.row) + abs(c - anchor.col) for r, c in matched)
            score = (len(matched), -bbox_area, -variant_idx, -distance_sum)
            if best_score is None or score > best_score:
                best_score = score
                best_translated_positions = translated_positions
                best_matched_positions = matched

    if not best_translated_positions or not best_matched_positions:
        return _pick_nearby_cluster_greedy(remaining, target_count)

    expected_count = min(target_count, len(best_translated_positions))
    min_required = max(2, int(math.ceil(expected_count * 0.6)))
    if len(best_matched_positions) < min_required:
        return _pick_nearby_cluster_greedy(remaining, target_count)

    selected = [pos_to_seat[pos] for pos in best_translated_positions if pos in pos_set][:target_count]
    selected_pos = {(s.row, s.col) for s in selected}
    candidates = [s for s in remaining if (s.row, s.col) not in selected_pos]

    while len(selected) < target_count and candidates:
        best_idx = None
        best_key = None
        for idx, seat in enumerate(candidates):
            min_dist = min(abs(seat.row - s.row) + abs(seat.col - s.col) for s in selected)
            key = (min_dist, seat.row, seat.col)
            if best_key is None or key < best_key:
                best_key = key
                best_idx = idx
        if best_idx is None:
            break
        selected.append(candidates.pop(best_idx))

    return selected


def _ordered_seats_by_style(seats, style, group_size, groups_needed, nearby_shape_profile=None):
    seats = list(seats)
    if not seats:
        return seats
    style = str(style or 'horizontal').strip().lower()

    if style == 'vertical':
        return sorted(seats, key=lambda s: (s.col, s.row))
    if style == 'horizontal':
        return sorted(seats, key=lambda s: (s.row, s.col))
    if style != 'nearby':
        return sorted(seats, key=lambda s: (s.row, s.col))

    remaining = list(sorted(seats, key=lambda s: (s.row, s.col)))
    ordered = []
    if groups_needed <= 0:
        return remaining

    for g_idx in range(groups_needed):
        if not remaining:
            break
        target_count = group_size
        if g_idx == groups_needed - 1:
            target_count = len(remaining)

        if nearby_shape_profile:
            cluster = _pick_nearby_cluster_by_shape(remaining, target_count, nearby_shape_profile)
        else:
            cluster = _pick_nearby_cluster_greedy(remaining, target_count)
        if not cluster:
            break

        selected_ids = {s.pk for s in cluster}
        remaining = [s for s in remaining if s.pk not in selected_ids]
        ordered.extend(cluster)

    ordered.extend(remaining)
    return ordered


def _line_group_key(seat, style):
    if style == 'vertical':
        return seat.col
    return seat.row


@require_POST
def auto_group_from_reference(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)

    try:
        if request.content_type and 'application/json' in request.content_type:
            data = json.loads(request.body or '{}')
            ref_group_id = data.get('reference_group_id')
            remainder_strategy = data.get('remainder_strategy')
            auto_detect_group_style = data.get('auto_detect_group_style')
        else:
            ref_group_id = request.POST.get('reference_group_id')
            remainder_strategy = request.POST.get('remainder_strategy')
            auto_detect_group_style = request.POST.get('auto_detect_group_style')
        ref_group_id = int(ref_group_id)
    except Exception:
        return JsonResponse({'status': 'error', 'message': '请选择参考小组'}, status=400)

    remainder_strategy = str(remainder_strategy or 'new_group').strip().lower()
    if remainder_strategy not in {'merge_prev', 'new_group', 'skip'}:
        remainder_strategy = 'new_group'
    auto_detect_group_style = _parse_bool(auto_detect_group_style if auto_detect_group_style is not None else '1')

    ref_group = get_object_or_404(SeatGroup, classroom=classroom, pk=ref_group_id)
    reference_size = ref_group.seats.filter(cell_type=SeatCellType.SEAT).count()
    if reference_size <= 0:
        reference_size = ref_group.seats.filter(cell_type=SeatCellType.SEAT, student__isnull=False).count()
    if reference_size <= 0:
        return JsonResponse({'status': 'error', 'message': '参考小组没有可用规模，请先给该组分配座位'}, status=400)

    target_seats = list(
        classroom.seats
        .filter(cell_type=SeatCellType.SEAT, student__isnull=False, group__isnull=True)
        .order_by('row', 'col')
    )
    if not target_seats:
        return JsonResponse({'status': 'error', 'message': '没有可继续编组的未分组学生'}, status=400)

    detected_group_style = _detect_group_style(ref_group) if auto_detect_group_style else 'horizontal'
    linear_grouping = auto_detect_group_style and detected_group_style in {'horizontal', 'vertical'}
    nearby_shape_profile = _build_nearby_shape_profile(ref_group) if detected_group_style == 'nearby' else None
    total_target_count = len(target_seats)
    assign_target_seats = []
    groups_needed = 0
    full_groups = 0
    remainder = 0
    ordered_target_seats = []
    line_key_map = {}

    if linear_grouping:
        ordered_target_seats = sorted(
            target_seats,
            key=lambda s: (_line_group_key(s, detected_group_style), s.col if detected_group_style == 'horizontal' else s.row)
        )
        assign_target_seats = ordered_target_seats
        line_keys = sorted({_line_group_key(s, detected_group_style) for s in ordered_target_seats})
        groups_needed = len(line_keys)
        line_key_map = {_line_group_key(s, detected_group_style): None for s in ordered_target_seats}
    else:
        full_groups = total_target_count // reference_size
        remainder = total_target_count % reference_size

        if remainder_strategy == 'skip':
            assignable_count = full_groups * reference_size
            assign_target_seats = target_seats[:assignable_count]
            groups_needed = full_groups
        else:
            assign_target_seats = target_seats
            if remainder_strategy == 'new_group':
                groups_needed = full_groups + (1 if remainder > 0 else 0)
            else:
                if full_groups > 0:
                    groups_needed = full_groups
                else:
                    groups_needed = 1

        ordered_target_seats = _ordered_seats_by_style(
            target_seats,
            detected_group_style,
            reference_size,
            groups_needed,
            nearby_shape_profile=nearby_shape_profile,
        )
        if remainder_strategy == 'skip':
            assign_target_seats = ordered_target_seats[:assignable_count]
        else:
            assign_target_seats = ordered_target_seats

    reusable_groups = list(
        classroom.groups
        .exclude(pk=ref_group.pk)
        .annotate(
            used_seat_count=models.Count(
                'seats',
                filter=models.Q(seats__cell_type=SeatCellType.SEAT),
            )
        )
        .filter(used_seat_count=0)
        .order_by('order', 'pk')
    )
    selected_reusable_groups = reusable_groups[:groups_needed]
    remaining_groups_needed = max(0, groups_needed - len(selected_reusable_groups))

    existing_names = list(classroom.groups.values_list('name', flat=True))
    new_group_names = _next_group_names(ref_group.name, existing_names, remaining_groups_needed)

    created_groups = []
    target_groups = list(selected_reusable_groups)
    action_items = []
    affected_group_ids = set()
    before_state = _capture_history_state(classroom)

    with transaction.atomic():
        current_max_order = classroom.groups.aggregate(m=models.Max('order')).get('m') or 0
        for idx, name in enumerate(new_group_names):
            group = SeatGroup.objects.create(
                classroom=classroom,
                name=name,
                order=current_max_order + idx + 1
            )
            created_groups.append(group)
            target_groups.append(group)

        if assign_target_seats and not target_groups:
            return JsonResponse({'status': 'error', 'message': '没有可用小组可分配'}, status=400)

        if linear_grouping:
            line_keys_in_use = []
            for seat in assign_target_seats:
                key = _line_group_key(seat, detected_group_style)
                if key not in line_key_map:
                    line_key_map[key] = None
                if key not in line_keys_in_use:
                    line_keys_in_use.append(key)
            for idx, key in enumerate(line_keys_in_use):
                line_key_map[key] = target_groups[idx]

        for idx, seat in enumerate(assign_target_seats):
            if linear_grouping:
                group = line_key_map[_line_group_key(seat, detected_group_style)]
            else:
                group = target_groups[min(idx // reference_size, len(target_groups) - 1)]
            before_group_id = seat.group_id
            seat.group = group
            seat.save(update_fields=['group'])
            if before_group_id:
                affected_group_ids.add(before_group_id)
            affected_group_ids.add(group.pk)
            action_items.append({
                'row': seat.row,
                'col': seat.col,
                'before_group_id': before_group_id,
                'after_group_id': group.pk
            })

        if affected_group_ids:
            _normalize_group_leaders(classroom, affected_group_ids)
        if action_items:
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                'auto_group_from_reference',
                extra={
                    'type': 'group_batch',
                    'items': action_items,
                    'reference_group_id': ref_group.pk,
                    'remainder_strategy': remainder_strategy,
                    'group_style': detected_group_style,
                    'linear_grouping': linear_grouping,
                },
            )

    unassigned_count = total_target_count - len(assign_target_seats)
    strategy_label = {
        'merge_prev': '并入上一组',
        'new_group': '剩余单独成组',
        'skip': '不编组余数'
    }.get(remainder_strategy, remainder_strategy)

    return JsonResponse({
        'status': 'success',
        'assigned_count': len(assign_target_seats),
        'unassigned_count': unassigned_count,
        'group_size': reference_size,
        'linear_grouping': linear_grouping,
        'remainder_strategy': remainder_strategy,
        'group_style': detected_group_style,
        'group_shape': (nearby_shape_profile or {}).get('shape_name'),
        'auto_detect_group_style': auto_detect_group_style,
        'reused_groups': [{'id': g.pk, 'name': g.name} for g in selected_reusable_groups],
        'created_groups': [{'id': g.pk, 'name': g.name} for g in created_groups],
        'message': (
            f'自动编组完成（{strategy_label}，样式：{detected_group_style}'
            + ('/整行列直编' if linear_grouping else '')
            + (f'/{nearby_shape_profile["shape_name"]}' if nearby_shape_profile else '')
            + f'）：复用 {len(selected_reusable_groups)} 组，'
            f'新增 {len(created_groups)} 组，分配 {len(assign_target_seats)} 人'
            + (f'，未编组 {unassigned_count} 人' if unassigned_count > 0 else '')
        )
    })


@require_POST
def merge_groups(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)

    try:
        if request.content_type and 'application/json' in request.content_type:
            data = json.loads(request.body or '{}')
            target_group_id = int(data.get('target_group_id'))
            source_group_ids = data.get('source_group_ids') or []
        else:
            target_group_id = int(request.POST.get('target_group_id'))
            source_group_ids = request.POST.getlist('source_group_ids')
        source_group_ids = [int(gid) for gid in source_group_ids]
    except Exception:
        return JsonResponse({'status': 'error', 'message': '参数错误'}, status=400)

    source_group_ids = list({gid for gid in source_group_ids if gid != target_group_id})
    if not source_group_ids:
        return JsonResponse({'status': 'error', 'message': '请至少选择一个来源组'}, status=400)

    target_group = get_object_or_404(SeatGroup, classroom=classroom, pk=target_group_id)
    source_groups = list(classroom.groups.filter(pk__in=source_group_ids))
    if not source_groups:
        return JsonResponse({'status': 'error', 'message': '来源组不存在'}, status=400)

    source_ids = [g.pk for g in source_groups]
    source_names = [g.name for g in source_groups]
    before_state = _capture_history_state(classroom)

    with transaction.atomic():
        affected_rows = list(
            classroom.seats
            .filter(cell_type=SeatCellType.SEAT, group_id__in=source_ids)
            .values('row', 'col', 'group_id')
        )
        moved_count = len(affected_rows)

        if moved_count:
            classroom.seats.filter(cell_type=SeatCellType.SEAT, group_id__in=source_ids).update(group=target_group)

        if not target_group.leader_id:
            for gid in source_ids:
                g = next((item for item in source_groups if item.pk == gid), None)
                if not g or not g.leader_id:
                    continue
                in_target = classroom.seats.filter(
                    cell_type=SeatCellType.SEAT,
                    group=target_group,
                    student_id=g.leader_id
                ).exists()
                if in_target:
                    SeatGroup.objects.filter(pk__in=source_ids, leader_id=g.leader_id).update(leader=None)
                    target_group.leader_id = g.leader_id
                    target_group.save(update_fields=['leader'])
                    break

        action_items = [
            {
                'row': item['row'],
                'col': item['col'],
                'before_group_id': item['group_id'],
                'after_group_id': target_group.pk
            }
            for item in affected_rows
        ]

        SeatGroup.objects.filter(classroom=classroom, pk__in=source_ids).delete()
        _normalize_group_leaders(classroom, [target_group.pk])

    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'merge_groups',
        extra={
            'type': 'group_batch',
            'items': action_items,
            'target_group_id': target_group.pk,
            'source_group_ids': source_ids,
        },
    )

    return JsonResponse({
        'status': 'success',
        'target_group': {'id': target_group.pk, 'name': target_group.name},
        'deleted_groups': [{'id': gid, 'name': gname} for gid, gname in zip(source_ids, source_names)],
        'moved_count': moved_count,
        'message': f'已合并 {len(source_ids)} 个来源组到 {target_group.name}'
    })


@require_POST
def rotate_groups(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)

    groups = list(classroom.groups.all())
    if len(groups) < 2:
        return JsonResponse({'status': 'error', 'message': '至少需要 2 个小组才能轮换'}, status=400)

    ordered_groups = []
    expected_size = None
    for group in groups:
        group_seats = list(
            group.seats
            .filter(cell_type=SeatCellType.SEAT)
            .select_related('student', 'group')
            .order_by('row', 'col')
        )
        if not group_seats:
            return JsonResponse({'status': 'error', 'message': f'小组【{group.name}】没有可轮换的座位'}, status=400)
        if expected_size is None:
            expected_size = len(group_seats)
        elif len(group_seats) != expected_size:
            return JsonResponse({'status': 'error', 'message': '小组座位数量不一致，无法执行平移轮换'}, status=400)

        avg_row = sum(seat.row for seat in group_seats) / len(group_seats)
        avg_col = sum(seat.col for seat in group_seats) / len(group_seats)
        ordered_groups.append({
            'group': group,
            'seats': group_seats,
            'avg_row': avg_row,
            'avg_col': avg_col,
        })

    ordered_groups.sort(
        key=lambda item: (
            round(item['avg_row'], 6),
            round(item['avg_col'], 6),
            item['group'].order,
            item['group'].pk
        )
    )

    action_items = []
    for idx, source in enumerate(ordered_groups):
        target = ordered_groups[(idx + 1) % len(ordered_groups)]
        source_group = source['group']
        source_seats = source['seats']
        target_seats = target['seats']

        for source_seat, target_seat in zip(source_seats, target_seats):
            action_items.append({
                'row': target_seat.row,
                'col': target_seat.col,
                'before_student_id': target_seat.student_id,
                'after_student_id': source_seat.student_id,
                'before_group_id': target_seat.group_id,
                'after_group_id': source_group.pk
            })

    if not action_items:
        return JsonResponse({'status': 'error', 'message': '没有可轮换的数据'}, status=400)

    action = {'type': 'seat_layout_batch', 'items': action_items}

    before_state = _capture_history_state(classroom)
    try:
        with transaction.atomic():
            if not _apply_seat_layout_action(classroom, action, forward=True):
                raise ValueError('轮换失败：无法应用座位布局')
            violations = _stabilize_layout_with_rules(classroom, request)
            if violations:
                raise ValueError(f'轮换失败：{_format_issues_preview(violations)}')
            _push_snapshot_action(request, classroom, before_state, 'rotate_groups', extra=action)
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': f'轮换失败：{e}'}, status=400)

    order_preview = ' -> '.join(item['group'].name for item in ordered_groups)
    return JsonResponse({
        'status': 'success',
        'message': f'已完成小组平移轮换：{order_preview}'
    })


@require_POST
def rename_group(request, pk, group_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    group = get_object_or_404(SeatGroup, classroom=classroom, pk=group_id)
    new_name = str(request.POST.get('name') or '').strip()
    if not new_name:
        if _is_ajax_request(request):
            return JsonResponse({'status': 'error', 'message': '小组名称不能为空'}, status=400)
        return redirect('classroom_detail', pk=pk)
    if new_name == group.name:
        if _is_ajax_request(request):
            return JsonResponse({'status': 'success'})
        return redirect('classroom_detail', pk=pk)
    if classroom.groups.exclude(pk=group.pk).filter(name=new_name).exists():
        if _is_ajax_request(request):
            return JsonResponse({'status': 'error', 'message': '小组名称已存在'}, status=400)
        return redirect('classroom_detail', pk=pk)
    before_state = _capture_history_state(classroom)
    try:
        group.name = new_name
        group.save(update_fields=['name'])
    except IntegrityError:
        if _is_ajax_request(request):
            return JsonResponse({'status': 'error', 'message': '小组名称已存在'}, status=400)
        return redirect('classroom_detail', pk=pk)
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'rename_group',
        extra={'group_id': group.pk, 'group_name': group.name},
    )
    if _is_ajax_request(request):
        return JsonResponse({'status': 'success', 'group': {'id': group.pk, 'name': group.name}})
    return redirect('classroom_detail', pk=pk)


@require_POST
def delete_group(request, pk, group_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    group = get_object_or_404(SeatGroup, pk=group_id, classroom=classroom)
    before_state = _capture_history_state(classroom)
    deleted_group_id = group.pk
    deleted_group_name = group.name
    group.delete()
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'delete_group',
        extra={'group_id': deleted_group_id, 'group_name': deleted_group_name},
    )
    if _is_ajax_request(request):
        return JsonResponse({'status': 'success', 'deleted_group_id': deleted_group_id})
    return redirect('classroom_detail', pk=pk)


@require_POST
def assign_group(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
        row = int(data.get('row'))
        col = int(data.get('col'))
        group_id = data.get('group_id')
        seat = get_object_or_404(Seat, classroom=classroom, row=row, col=col)
        if seat.cell_type != SeatCellType.SEAT:
            return JsonResponse({'status': 'error', 'message': '当前单元不可分组'}, status=400)
        before_group_id = seat.group.pk if seat.group else None
        affected_group_ids = set()
        if before_group_id:
            affected_group_ids.add(before_group_id)
        if group_id:
            group = get_object_or_404(SeatGroup, pk=group_id, classroom=classroom)
            seat.group = group
            affected_group_ids.add(group.pk)
        else:
            seat.group = None
        before_state = _capture_history_state(classroom)
        seat.save(update_fields=['group'])
        _normalize_group_leaders(classroom, affected_group_ids)
        action = {
            'type': 'group',
            'row': seat.row,
            'col': seat.col,
            'before_group_id': before_group_id,
            'after_group_id': seat.group.pk if seat.group else None
        }
        _push_snapshot_action(request, classroom, before_state, 'assign_group', extra=action)
        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@require_POST
def assign_group_batch(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
        seats_payload = data.get('seats', [])
        group_id = data.get('group_id') or None

        group = None
        if group_id:
            group = get_object_or_404(SeatGroup, pk=group_id, classroom=classroom)

        items = []
        affected_group_ids = set()
        if group:
            affected_group_ids.add(group.pk)
        before_state = _capture_history_state(classroom)
        for seat_data in seats_payload:
            row = int(seat_data.get('row'))
            col = int(seat_data.get('col'))
            seat = classroom.seats.filter(row=row, col=col).first()
            if not seat or seat.cell_type != SeatCellType.SEAT:
                continue
            before_group_id = seat.group.pk if seat.group else None
            if before_group_id:
                affected_group_ids.add(before_group_id)
            seat.group = group
            seat.save(update_fields=['group'])
            items.append({
                'row': row,
                'col': col,
                'before_group_id': before_group_id,
                'after_group_id': group.pk if group else None
            })

        if items:
            _normalize_group_leaders(classroom, affected_group_ids)
            action = {'type': 'group_batch', 'items': items}
            _push_snapshot_action(request, classroom, before_state, 'assign_group_batch', extra=action)
        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


def _constraint_error_response(request, action_label, error_message, *, status=400):
    if _is_ajax_request(request):
        return JsonResponse({'status': 'error', 'message': error_message}, status=status)
    return HttpResponse(f'{action_label}失败: {error_message}', status=status)


def _constraint_success_response(request, constraint, *, metrics_message='操作成功'):
    if _is_ajax_request(request):
        return JsonResponse({
            'status': 'success',
            'constraint_id': constraint.pk,
            'enabled': constraint.enabled,
            'message': metrics_message,
        })
    return redirect('classroom_detail', pk=constraint.classroom_id)


def _constraint_form_payload(request):
    return {
        'constraint_type': request.POST.get('constraint_type'),
        'student_id': request.POST.get('student_id'),
        'target_student_id': request.POST.get('target_student_id'),
        'row': request.POST.get('row'),
        'col': request.POST.get('col'),
        'distance': request.POST.get('distance'),
        'note': request.POST.get('note'),
        'enabled': request.POST.get('enabled'),
    }


@require_POST
def create_constraint(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        cleaned = normalize_constraint_payload(classroom, _constraint_form_payload(request))
        validate_constraint_candidate(classroom, cleaned)

        before_state = _capture_history_state(classroom)
        constraint = SeatConstraint.objects.create(
            classroom=classroom,
            constraint_type=cleaned['constraint_type'],
            student=cleaned['student'],
            target_student=cleaned['target_student'],
            row=cleaned['row'],
            col=cleaned['col'],
            distance=cleaned['distance'],
            enabled=cleaned['enabled'],
            note=cleaned['note'],
        )
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'create_constraint',
            extra={'constraint_id': constraint.pk, 'constraint_type': cleaned['constraint_type']},
        )
    except ConstraintServiceError as exc:
        return _constraint_error_response(request, '创建约束', str(exc))
    except Exception as exc:
        return _constraint_error_response(request, '创建约束', str(exc))
    return _constraint_success_response(request, constraint, metrics_message='约束已创建')


@require_POST
def update_constraint(request, pk, constraint_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    constraint = get_object_or_404(SeatConstraint, pk=constraint_id, classroom=classroom)
    try:
        cleaned = normalize_constraint_payload(classroom, _constraint_form_payload(request), instance=constraint)
        validate_constraint_candidate(classroom, cleaned, instance=constraint)

        before_state = _capture_history_state(classroom)
        constraint.constraint_type = cleaned['constraint_type']
        constraint.student = cleaned['student']
        constraint.target_student = cleaned['target_student']
        constraint.row = cleaned['row']
        constraint.col = cleaned['col']
        constraint.distance = cleaned['distance']
        constraint.enabled = cleaned['enabled']
        constraint.note = cleaned['note']
        constraint.save()
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'update_constraint',
            extra={'constraint_id': constraint.pk, 'constraint_type': cleaned['constraint_type']},
        )
    except ConstraintServiceError as exc:
        return _constraint_error_response(request, '更新约束', str(exc))
    except Exception as exc:
        return _constraint_error_response(request, '更新约束', str(exc))
    return _constraint_success_response(request, constraint, metrics_message='约束已更新')


@require_POST
def toggle_constraint(request, pk, constraint_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    constraint = get_object_or_404(SeatConstraint, pk=constraint_id, classroom=classroom)
    try:
        desired_enabled = request.POST.get('enabled')
        next_enabled = (not constraint.enabled) if desired_enabled in (None, '') else str(desired_enabled).strip().lower() in {'1', 'true', 'yes', 'on'}

        if next_enabled:
            cleaned = normalize_constraint_payload(
                classroom,
                {
                    'constraint_type': constraint.constraint_type,
                    'student': constraint.student,
                    'target_student': constraint.target_student,
                    'row': constraint.row,
                    'col': constraint.col,
                    'distance': constraint.distance,
                    'note': constraint.note,
                    'enabled': True,
                },
                instance=constraint,
            )
            validate_constraint_candidate(classroom, cleaned, instance=constraint)

        before_state = _capture_history_state(classroom)
        constraint.enabled = next_enabled
        constraint.save(update_fields=['enabled'])
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'toggle_constraint',
            extra={'constraint_id': constraint.pk, 'constraint_type': constraint.constraint_type, 'enabled': constraint.enabled},
        )
    except ConstraintServiceError as exc:
        return _constraint_error_response(request, '切换约束状态', str(exc))
    except Exception as exc:
        return _constraint_error_response(request, '切换约束状态', str(exc))
    return _constraint_success_response(
        request,
        constraint,
        metrics_message='约束已启用' if constraint.enabled else '约束已停用',
    )


@require_POST
def delete_constraint(request, pk, constraint_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    constraint = get_object_or_404(SeatConstraint, pk=constraint_id, classroom=classroom)
    before_state = _capture_history_state(classroom)
    constraint_type = constraint.constraint_type
    constraint.delete()
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'delete_constraint',
        extra={'constraint_id': constraint_id, 'constraint_type': constraint_type},
    )
    if _is_ajax_request(request):
        return JsonResponse({'status': 'success'})
    return redirect('classroom_detail', pk=pk)


def _csis_gender_value(student):
    if student.gender == 'M':
        return 'male'
    if student.gender == 'F':
        return 'female'
    return 'unknown'


def _csis_student_number(student):
    raw_student_id = str(student.student_id or '').strip()
    if raw_student_id.isdigit():
        number = int(raw_student_id)
        if number > 0:
            return number
    return None


def _csis_seat_extra(classroom, seat):
    if not seat:
        return {
            'assigned': False,
            'row': None,
            'col': None,
            'coordinate': '',
            'position': None,
            'cell_type': None,
            'cell_type_display': '',
            'classroom_rows': classroom.rows,
            'classroom_cols': classroom.cols,
        }

    return {
        'assigned': True,
        'row': seat.row,
        'col': seat.col,
        'coordinate': f'{seat.row}-{seat.col}',
        'position': [seat.col - 1, seat.row - 1],
        'cell_type': seat.cell_type,
        'cell_type_display': seat.get_cell_type_display(),
        'classroom_rows': classroom.rows,
        'classroom_cols': classroom.cols,
    }


def _build_csis_csls_payload(classroom):
    students = list(
        classroom.students
        .prefetch_related('tag_memberships__tag')
        .order_by('pk')
    )
    seat_by_student_id = {
        seat.student_id: seat
        for seat in classroom.seats.select_related('student', 'group').filter(student__isnull=False)
    }
    group_names = {group.name for group in classroom.groups.all()}
    csis_students = []

    for student in students:
        seat = seat_by_student_id.get(student.pk)
        group_name = seat.group.name if seat and seat.group_id and seat.group else 'unknown'
        group_names.add(group_name)
        number = _csis_student_number(student)
        tag_names = [
            membership.tag.name
            for membership in student.tag_memberships.all()
            if membership.tag_id and membership.tag
        ]
        student_payload = {
            'id': student.pk,
            'name': student.name,
            'group': group_name,
            'tags': tag_names,
            'gender': _csis_gender_value(student),
            'extra': {
                'source': 'fuckseats',
                'student_pk': student.pk,
                'student_id': student.student_id or '',
                'score': student.score,
                'seat': _csis_seat_extra(classroom, seat),
            },
        }
        if number is not None:
            student_payload['number'] = number
        csis_students.append(student_payload)

    groups = [
        {
            'name': name,
            'tags': [],
            'extra': {},
        }
        for name in sorted(group_names)
    ]

    return {
        'version': 1,
        'classes': [
            {
                'name': classroom.name,
                'class': 0,
                'grade': 0,
                'groups': groups,
                'students': csis_students,
            }
        ],
    }


def export_students(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)

    layout_transform = str(request.GET.get('layout_transform', 'none')).strip().lower()
    rotate_180 = layout_transform in {'rotate_180', 'rot180', '180'}
    if not rotate_180:
        rotate_flag = str(request.GET.get('rotate_180', '')).strip().lower()
        rotate_180 = rotate_flag in {'1', 'true', 'yes', 'on'}

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = classroom.name

    thin_border = Border(left=Side(style='thin'),
                         right=Side(style='thin'),
                         top=Side(style='thin'),
                         bottom=Side(style='thin'))
    center_align = Alignment(horizontal='center', vertical='center', wrap_text=True)

    header_font = Font(name=EXPORT_FONT_BLACK, bold=False, size=20)
    podium_font = Font(name=EXPORT_FONT_BLACK, bold=False, size=14)
    seat_font = Font(name=EXPORT_FONT_LIGHT, size=12, bold=False)

    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=classroom.cols)
    title_suffix = "（180°翻转）" if rotate_180 else ""
    cell = ws.cell(row=1, column=1, value=f"{classroom.name} 座位表{title_suffix}")
    cell.font = header_font
    cell.alignment = center_align
    ws.row_dimensions[1].height = 40

    seat_start_row = 2 if rotate_180 else 3
    podium_row = seat_start_row + classroom.rows if rotate_180 else 2

    ws.merge_cells(start_row=podium_row, start_column=1, end_row=podium_row, end_column=classroom.cols)
    podium_cell = ws.cell(row=podium_row, column=1, value="讲台")
    podium_cell.font = podium_font
    podium_cell.alignment = center_align
    ws.row_dimensions[podium_row].height = 30

    seats = classroom.seats.select_related('student').all()
    seat_map = _build_seat_map(seats)

    for c in range(1, classroom.cols + 1):
        ws.column_dimensions[get_column_letter(c)].width = 14

    for visual_row in range(1, classroom.rows + 1):
        row_index = seat_start_row + visual_row - 1
        ws.row_dimensions[row_index].height = 50
        for c in range(1, classroom.cols + 1):
            cell = ws.cell(row=row_index, column=c)
            source_row = classroom.rows - visual_row + 1 if rotate_180 else visual_row
            source_col = classroom.cols - c + 1 if rotate_180 else c
            seat = seat_map.get((source_row, source_col))

            value = ""
            is_seat = False
            if seat:
                if seat.cell_type == SeatCellType.SEAT:
                    is_seat = True
                    if seat.student:
                        value = seat.student.name
                    else:
                        value = ""
                elif seat.cell_type == SeatCellType.AISLE or seat.cell_type == SeatCellType.EMPTY:
                    value = ""
                else:
                    value = seat.get_cell_type_display()

            cell.value = value
            cell.alignment = center_align
            cell.font = seat_font

            if is_seat and seat.student:
                cell.border = thin_border

    from openpyxl.worksheet.page import PageMargins
    ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.margins = PageMargins(left=0.25, right=0.25, top=0.25, bottom=0.25, header=0, footer=0)
    ws.print_options.horizontalCentered = True
    ws.print_options.verticalCentered = True

    ws.page_setup.fitToPage = True
    ws.page_setup.fitToHeight = 1
    ws.page_setup.fitToWidth = 1

    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    filename_suffix = "_座次图_180度翻转.xlsx" if rotate_180 else "_座次图.xlsx"
    response['Content-Disposition'] = f'attachment; filename="{classroom.name}{filename_suffix}"'
    wb.save(response)

    return response


def export_students_csis(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    payload = _build_csis_csls_payload(classroom)
    content = json.dumps(payload, ensure_ascii=False, indent=2)
    response = HttpResponse(content, content_type='application/json; charset=utf-8')
    filename = escape_uri_path(f'{classroom.name}_CSIS.csls')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


def export_students_options_page(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/export_excel_options.html', {
        'classroom': classroom
    })


def export_students_svg(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    seats = list(classroom.seats.select_related('student', 'group').all())
    seat_map = _build_seat_map(seats)

    def _qbool(key, default=True):
        raw = request.GET.get(key)
        if raw is None or raw == '':
            return default
        return str(raw).strip().lower() in {'1', 'true', 'yes', 'on'}

    show_title = _qbool('show_title', True)
    show_podium = _qbool('show_podium', True)
    show_coords = _qbool('show_coords', True)
    show_name = _qbool('show_name', True)
    show_score = _qbool('show_score', True)
    show_group = _qbool('show_group', True)
    show_empty_label = _qbool('show_empty_label', True)
    show_seat_type = _qbool('show_seat_type', True)
    name_emphasis_mode = show_name and (not show_coords) and (not show_score)

    theme = str(request.GET.get('theme', 'classic')).strip().lower()
    if theme not in SVG_EXPORT_THEME_MAP:
        theme = 'classic'
    style = SVG_EXPORT_THEME_MAP[theme]

    cell_w = 120
    cell_h = 86
    gap = 10
    padding_x = 24
    padding_y = 24
    if show_title and show_podium:
        header_h = 90
    elif show_title or show_podium:
        header_h = 64
    else:
        header_h = 16

    grid_w = classroom.cols * cell_w + max(0, classroom.cols - 1) * gap
    grid_h = classroom.rows * cell_h + max(0, classroom.rows - 1) * gap

    width = padding_x * 2 + grid_w
    height = padding_y * 2 + header_h + grid_h
    grid_top = padding_y + header_h

    podium_w = min(340, max(180, int(grid_w * 0.42)))
    podium_h = 34
    podium_x = padding_x + (grid_w - podium_w) // 2

    def group_color(group_id):
        if not group_id:
            return '#9aa6c2'
        group_palette = style['group_palette']
        return group_palette[(int(group_id) - 1) % len(group_palette)]

    chunks = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        '<defs>',
        '<style><![CDATA['
        + _export_svg_font_style() +
        ']]></style>',
        '</defs>',
        f'<rect x="0" y="0" width="{width}" height="{height}" fill="{style["bg"]}"/>',
    ]

    if show_title:
        title_y = padding_y + 28 if show_podium else padding_y + 30
        chunks.append(
            f'<text x="{padding_x}" y="{title_y}" class="title" fill="{style["title"]}">{html.escape(classroom.name)} 座次图</text>'
        )

    if show_podium:
        podium_y = padding_y + 32 if show_title else padding_y + 14
        chunks.append(
            f'<rect x="{podium_x}" y="{podium_y}" width="{podium_w}" height="{podium_h}" rx="12" fill="{style["podium_fill"]}" stroke="{style["podium_stroke"]}"/>'
        )
        chunks.append(
            f'<text x="{podium_x + podium_w / 2}" y="{podium_y + 22}" text-anchor="middle" class="cell-type" fill="{style["type"]}">讲台</text>'
        )

    for r in range(1, classroom.rows + 1):
        for c in range(1, classroom.cols + 1):
            seat = seat_map.get((r, c))
            if not seat:
                continue

            x = padding_x + (c - 1) * (cell_w + gap)
            y = grid_top + (r - 1) * (cell_h + gap)

            if seat.cell_type == SeatCellType.SEAT:
                if seat.student_id:
                    fill = style['seat_fill_occupied']
                    stroke = style['seat_stroke_occupied']
                else:
                    fill = style['seat_fill_empty']
                    stroke = style['seat_stroke_empty']

                chunks.append(
                    f'<rect x="{x}" y="{y}" width="{cell_w}" height="{cell_h}" rx="16" fill="{fill}" stroke="{stroke}"/>'
                )
                if show_coords:
                    chunks.append(
                        f'<text x="{x + 8}" y="{y + 16}" class="cell-sub" fill="{style["sub"]}">({r}-{c})</text>'
                    )

                if show_group and seat.group_id and seat.group:
                    tag_w = max(36, min(66, 18 + len(seat.group.name) * 12))
                    tag_color = group_color(seat.group_id)
                    chunks.append(
                        f'<rect x="{x + cell_w - tag_w - 8}" y="{y + 8}" width="{tag_w}" height="20" rx="10" fill="{tag_color}"/>'
                    )
                    chunks.append(
                        f'<text x="{x + cell_w - tag_w / 2 - 8}" y="{y + 22}" text-anchor="middle" class="tag" fill="{style["tag_text"]}">{html.escape(seat.group.name)}</text>'
                    )

                if seat.student:
                    base_name_y = y + (48 if show_coords else 42)
                    if show_name:
                        if name_emphasis_mode:
                            name_size = _name_emphasis_font_size(seat.student.name)
                            center_y = y + cell_h / 2 + (6 if (show_group and seat.group_id) else 0)
                            chunks.append(
                                f'<text x="{x + cell_w / 2}" y="{center_y}" text-anchor="middle" dominant-baseline="middle" class="cell-name" font-size="{name_size}" fill="{style["name"]}">{html.escape(seat.student.name)}</text>'
                            )
                        else:
                            chunks.append(
                                f'<text x="{x + 12}" y="{base_name_y}" class="cell-name" fill="{style["name"]}">{html.escape(seat.student.name)}</text>'
                            )
                    if show_score and (seat.student.score or 0) > 0:
                        score_y = base_name_y + 20 if show_name else y + (56 if show_coords else 50)
                        chunks.append(
                            f'<text x="{x + 12}" y="{score_y}" class="cell-sub" fill="{style["sub"]}">{seat.student.display_score}分</text>'
                        )
                elif show_empty_label:
                    empty_y = y + (56 if show_coords else 50)
                    chunks.append(
                        f'<text x="{x + 12}" y="{empty_y}" class="cell-sub" fill="{style["sub"]}">空座位</text>'
                    )
                continue

            if seat.cell_type == SeatCellType.AISLE:
                fill = style['nonseat_aisle']
            elif seat.cell_type == SeatCellType.PODIUM:
                fill = style['nonseat_podium']
            else:
                fill = style['nonseat_empty']

            chunks.append(
                f'<rect x="{x}" y="{y}" width="{cell_w}" height="{cell_h}" rx="16" fill="{fill}" stroke="{style["nonseat_stroke"]}"/>'
            )
            if show_seat_type:
                chunks.append(
                    f'<text x="{x + cell_w / 2}" y="{y + 50}" text-anchor="middle" class="cell-type" fill="{style["type"]}">{html.escape(seat.get_cell_type_display())}</text>'
                )

    chunks.append('</svg>')
    svg_content = ''.join(chunks)

    response = HttpResponse(svg_content, content_type='image/svg+xml; charset=utf-8')
    filename = escape_uri_path(f'{classroom.name}_座次图.svg')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


def export_students_svg_preview_student(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    student_ids = list(classroom.students.values_list('pk', flat=True))
    if not student_ids:
        return JsonResponse({'status': 'empty', 'message': '当前班级暂无学生'})

    random_student_id = random.choice(student_ids)
    student = classroom.students.filter(pk=random_student_id).first()
    if not student:
        return JsonResponse({'status': 'empty', 'message': '当前班级暂无学生'})

    seat = getattr(student, 'assigned_seat', None)
    group_name = ''
    group_index = 0
    coord = ''
    if seat:
        coord = f'{seat.row}-{seat.col}'
        if seat.group_id and seat.group:
            group_name = seat.group.name
            group_index = int(seat.group_id)

    score_display = student.display_score if (student.score or 0) > 0 else ''

    return JsonResponse({
        'status': 'success',
        'sample': {
            'classroom': classroom.name,
            'name': student.name,
            'score': score_display,
            'group': group_name,
            'group_index': group_index,
            'coord': coord
        }
    })


def export_students_svg_options_page(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/export_svg_options.html', {
        'classroom': classroom
    })


def export_students_pptx(request, pk):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls, qn
        from pptx.util import Inches, Pt
    except ImportError:
        return HttpResponse('缺少 python-pptx 依赖，请先安装 requirements.txt', status=500)

    classroom = get_object_or_404(Classroom, pk=pk)
    seats = list(classroom.seats.select_related('student', 'group').all())
    seat_map = _build_seat_map(seats)

    def _qbool(key, default=True):
        raw = request.GET.get(key)
        if raw is None or raw == '':
            return default
        return str(raw).strip().lower() in {'1', 'true', 'yes', 'on'}

    show_title = _qbool('show_title', True)
    show_podium = _qbool('show_podium', True)
    show_coords = _qbool('show_coords', True)
    show_name = _qbool('show_name', True)
    show_score = _qbool('show_score', True)
    show_group = _qbool('show_group', True)
    show_empty_label = _qbool('show_empty_label', True)
    show_seat_type = _qbool('show_seat_type', True)
    name_emphasis_mode = show_name and (not show_coords) and (not show_score)

    theme = str(request.GET.get('theme', 'classic')).strip().lower()
    if theme not in SVG_EXPORT_THEME_MAP:
        theme = 'classic'
    style = SVG_EXPORT_THEME_MAP[theme]

    cell_w = 120
    cell_h = 86
    gap = 10
    padding_x = 24
    padding_y = 24
    if show_title and show_podium:
        header_h = 90
    elif show_title or show_podium:
        header_h = 64
    else:
        header_h = 16

    grid_w = classroom.cols * cell_w + max(0, classroom.cols - 1) * gap
    grid_h = classroom.rows * cell_h + max(0, classroom.rows - 1) * gap
    grid_top = padding_y + header_h

    content_w = padding_x * 2 + grid_w
    content_h = padding_y * 2 + header_h + grid_h

    podium_w = min(340, max(180, int(grid_w * 0.42)))
    podium_h = 34
    podium_x = padding_x + (grid_w - podium_w) // 2

    slide_w = 13.333
    slide_h = 7.5
    margin = 0.3
    usable_w = max(0.1, slide_w - margin * 2)
    usable_h = max(0.1, slide_h - margin * 2)
    scale = min(usable_w / max(1, content_w), usable_h / max(1, content_h))
    offset_x = (slide_w - content_w * scale) / 2
    offset_y = (slide_h - content_h * scale) / 2

    def sx(value):
        return offset_x + value * scale

    def sy(value):
        return offset_y + value * scale

    def sw(value):
        return value * scale

    def sh(value):
        return value * scale

    def font_pt(base_px):
        return max(8, base_px * scale * 72)

    def rgb(hex_color):
        return RGBColor(*_hex_to_rgb_parts(hex_color))

    def group_color(group_id):
        if not group_id:
            return '#9aa6c2'
        group_palette = style['group_palette']
        return group_palette[(int(group_id) - 1) % len(group_palette)]

    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(slide_w),
        Inches(slide_h),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(style['bg'])
    bg.line.fill.background()

    def add_round_rect(x, y, w, h, fill_color, stroke_color=None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(sx(x)),
            Inches(sy(y)),
            Inches(sw(w)),
            Inches(sh(h)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
        if stroke_color:
            shape.line.color.rgb = rgb(stroke_color)
            shape.line.width = Pt(max(0.75, scale * 72))
        else:
            shape.line.fill.background()
        sp_pr = shape._element.spPr
        for child in list(sp_pr):
            if child.tag == qn('a:effectLst'):
                sp_pr.remove(child)
        effect_xml = parse_xml(
            f'<a:effectLst {nsdecls("a")}>'
            '<a:outerShdw blurRad="38100" dist="19050" dir="5400000" algn="ctr" rotWithShape="0">'
            '<a:srgbClr val="000000"><a:alpha val="12000"/></a:srgbClr>'
            '</a:outerShdw>'
            '</a:effectLst>'
        )
        sp_pr.append(effect_xml)
        return shape

    def add_text(x, y, w, h, text, color, size_px, bold=False, center=False, middle=True):
        if text is None:
            return
        font_name = EXPORT_FONT_BLACK if bold else EXPORT_FONT_LIGHT
        shape = slide.shapes.add_textbox(
            Inches(sx(x)),
            Inches(sy(y)),
            Inches(sw(w)),
            Inches(sh(h)),
        )
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE if middle else MSO_ANCHOR.TOP
        paragraph = tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = str(text)
        run.font.name = font_name
        run.font.bold = False
        run.font.size = Pt(font_pt(size_px))
        run.font.color.rgb = rgb(color)
        r_pr = run._r.get_or_add_rPr()
        for tag in ('latin', 'ea', 'cs'):
            node = r_pr.find(qn(f'a:{tag}'))
            if node is None:
                node = parse_xml(f'<a:{tag} {nsdecls("a")} typeface="{font_name}"/>')
                r_pr.append(node)
            else:
                node.set('typeface', font_name)

    if show_title:
        title_y = padding_y + 28 if show_podium else padding_y + 30
        add_text(
            padding_x,
            title_y - 24,
            grid_w,
            32,
            f'{classroom.name} 座次图',
            style['title'],
            24,
            bold=True,
            center=False,
            middle=True,
        )

    if show_podium:
        podium_y = padding_y + 32 if show_title else padding_y + 14
        add_round_rect(podium_x, podium_y, podium_w, podium_h, style['podium_fill'], style['podium_stroke'])
        add_text(
            podium_x,
            podium_y + 4,
            podium_w,
            24,
            '讲台',
            style['type'],
            13,
            bold=True,
            center=True,
            middle=True,
        )

    for r in range(1, classroom.rows + 1):
        for c in range(1, classroom.cols + 1):
            seat = seat_map.get((r, c))
            if not seat:
                continue

            x = padding_x + (c - 1) * (cell_w + gap)
            y = grid_top + (r - 1) * (cell_h + gap)

            if seat.cell_type == SeatCellType.SEAT:
                if seat.student_id:
                    fill = style['seat_fill_occupied']
                    stroke = style['seat_stroke_occupied']
                else:
                    fill = style['seat_fill_empty']
                    stroke = style['seat_stroke_empty']

                add_round_rect(x, y, cell_w, cell_h, fill, stroke)

                if show_coords:
                    add_text(
                        x + 8,
                        y + 5,
                        cell_w - 16,
                        16,
                        f'({r}-{c})',
                        style['sub'],
                        12,
                        bold=False,
                        center=False,
                        middle=False,
                    )

                if show_group and seat.group_id and seat.group:
                    tag_w = max(36, min(66, 18 + len(seat.group.name) * 12))
                    tag_color = group_color(seat.group_id)
                    add_round_rect(x + cell_w - tag_w - 8, y + 8, tag_w, 20, tag_color, None)
                    add_text(
                        x + cell_w - tag_w - 8,
                        y + 8,
                        tag_w,
                        20,
                        seat.group.name,
                        style['tag_text'],
                        11,
                        bold=True,
                        center=True,
                        middle=True,
                    )

                if seat.student:
                    base_name_y = y + (48 if show_coords else 42)
                    if show_name:
                        if name_emphasis_mode:
                            name_size = _name_emphasis_font_size(seat.student.name)
                            center_y = y + cell_h / 2 + (6 if (show_group and seat.group_id) else 0)
                            add_text(
                                x + 10,
                                center_y - 18,
                                cell_w - 20,
                                36,
                                seat.student.name,
                                style['name'],
                                name_size,
                                bold=True,
                                center=True,
                                middle=True,
                            )
                        else:
                            add_text(
                                x + 12,
                                base_name_y - 16,
                                cell_w - 24,
                                22,
                                seat.student.name,
                                style['name'],
                                16,
                                bold=True,
                                center=False,
                                middle=False,
                            )
                    if show_score and (seat.student.score or 0) > 0:
                        score_y = base_name_y + 20 if show_name else y + (56 if show_coords else 50)
                        add_text(
                            x + 12,
                            score_y - 14,
                            cell_w - 24,
                            20,
                            f'{seat.student.display_score}分',
                            style['sub'],
                            12,
                            bold=False,
                            center=False,
                            middle=False,
                        )
                elif show_empty_label:
                    empty_y = y + (56 if show_coords else 50)
                    add_text(
                        x + 12,
                        empty_y - 14,
                        cell_w - 24,
                        20,
                        '空座位',
                        style['sub'],
                        12,
                        bold=False,
                        center=False,
                        middle=False,
                    )
                continue

            if seat.cell_type == SeatCellType.AISLE:
                fill = style['nonseat_aisle']
            elif seat.cell_type == SeatCellType.PODIUM:
                fill = style['nonseat_podium']
            else:
                fill = style['nonseat_empty']

            add_round_rect(x, y, cell_w, cell_h, fill, style['nonseat_stroke'])
            if show_seat_type:
                add_text(
                    x + 8,
                    y + 33,
                    cell_w - 16,
                    28,
                    seat.get_cell_type_display(),
                    style['type'],
                    13,
                    bold=True,
                    center=True,
                    middle=True,
                )

    buffer = BytesIO()
    prs.save(buffer)
    response = HttpResponse(
        buffer.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )
    filename = escape_uri_path(f'{classroom.name}_座次图.pptx')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


def export_students_pptx_options_page(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    return render(request, 'seats/export_pptx_options.html', {
        'classroom': classroom
    })


def export_group_report(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    groups = list(classroom.groups.all())

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = '小组作业登记表'
    
    from openpyxl.worksheet.page import PageMargins

    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin'),
    )
    group_font = Font(name=EXPORT_FONT_BLACK, bold=False, size=13)
    name_font = Font(name=EXPORT_FONT_LIGHT, bold=False, size=11, color="000000")
    leader_name_font = Font(name=EXPORT_FONT_LIGHT, bold=False, size=11, color="FF0000")
    center     = Alignment(horizontal='center', vertical='center')

    header_text = f"&\"{EXPORT_FONT_BLACK}\"&20 {classroom.name} (          ) 登记表"
    ws.oddHeader.center.text = header_text
    ws.evenHeader.center.text = header_text

    
    flat_entries = []
    
    WEIGHT_HEADER = 26
    WEIGHT_MEMBER = 24
    WEIGHT_GAP = 10
    
    total_weight = 0
    
    for i, group in enumerate(groups):
        flat_entries.append({
            'type': 'header', 
            'text': group.name, 
            'group_id': group.pk,
            'weight': WEIGHT_HEADER
        })
        total_weight += WEIGHT_HEADER
        
        seats = group.seats.select_related('student').filter(student__isnull=False)
        members = []
        for s in seats:
             is_ldr = (group.leader_id == s.student_id)
             members.append({'name': s.student.name, 'is_leader': is_ldr})
        
        members.sort(key=lambda x: not x['is_leader'])

        for m in members:
            flat_entries.append({
                'type': 'member', 
                'text': m['name'], 
                'is_leader': m['is_leader'],
                'group_id': group.pk,
                'group_name': group.name,
                'weight': WEIGHT_MEMBER
            })
            total_weight += WEIGHT_MEMBER
            
        if i < len(groups) - 1:
            flat_entries.append({'type': 'gap', 'weight': WEIGHT_GAP})
            total_weight += WEIGHT_GAP

    target_weight = total_weight / 2
    current_weight = 0
    split_index = 0
    
    for i, entry in enumerate(flat_entries):
        current_weight += entry.get('weight', 0)
        if current_weight >= target_weight:
            split_index = i + 1
            break
            
    left_entries = flat_entries[:split_index]
    right_entries = flat_entries[split_index:]
    
    if right_entries:
        first = right_entries[0]
        if first['type'] == 'gap':
            right_entries.pop(0)
            if right_entries:
                first = right_entries[0]
                
        if right_entries and first['type'] == 'member':
            continuation_header = {
                'type': 'header',
                'text': f"{first['group_name']} (续)",
                'group_id': first['group_id'],
                'weight': WEIGHT_HEADER
            }
            right_entries.insert(0, continuation_header)

    PAGE_H_MM = 297
    MARGIN_V_MM = 12.7 * 2
    HEADER_RES_MM = 15
    AVAILABLE_H_MM = PAGE_H_MM - MARGIN_V_MM - HEADER_RES_MM
    AVAILABLE_H_PTS = AVAILABLE_H_MM * 2.835
    
    max_rows = max(len(left_entries), len(right_entries))
    if max_rows == 0:
        max_rows = 1
        
    W_HEADER = 1.0
    W_MEMBER = 1.0
    W_GAP    = 0.4
    
    total_weight = 0
    row_weights = []
    
    for i in range(max_rows):
        l = left_entries[i] if i < len(left_entries) else None
        r = right_entries[i] if i < len(right_entries) else None
        
        w_l = 0
        if l:
            if l['type'] == 'header': w_l = W_HEADER
            elif l['type'] == 'member': w_l = W_MEMBER
            elif l['type'] == 'gap': w_l = W_GAP
            
        w_r = 0
        if r:
            if r['type'] == 'header': w_r = W_HEADER
            elif r['type'] == 'member': w_r = W_MEMBER
            elif r['type'] == 'gap': w_r = W_GAP
            
        cur_w = max(w_l, w_r)
        if cur_w == 0: cur_w = W_GAP
        
        row_weights.append(cur_w)
        total_weight += cur_w
        
    unit_h = AVAILABLE_H_PTS / total_weight
    
    MAX_UNIT_H = 45 
    MIN_UNIT_H = 18
    
    if unit_h > MAX_UNIT_H:
        unit_h = MAX_UNIT_H
    if unit_h < MIN_UNIT_H:
        unit_h = MIN_UNIT_H
        
    TOTAL_COL_WIDTH = 98
    
    boxes_count = 5
    box_width = 4.5
    gap_width = 2
    
    fixed_used = (2 * boxes_count * box_width) + gap_width
    remain_for_names = TOTAL_COL_WIDTH - fixed_used
    name_col_width = remain_for_names / 2
    
    if name_col_width < 12: name_col_width = 12

    left_col_idx  = 1
    gap_col_idx   = 1 + boxes_count + 1
    right_col_idx = gap_col_idx + 1

    ws.column_dimensions[get_column_letter(left_col_idx)].width = name_col_width
    for b in range(1, boxes_count + 1):
        ws.column_dimensions[get_column_letter(left_col_idx + b)].width = box_width
    ws.column_dimensions[get_column_letter(gap_col_idx)].width = gap_width
    ws.column_dimensions[get_column_letter(right_col_idx)].width = name_col_width
    for b in range(1, boxes_count + 1):
        ws.column_dimensions[get_column_letter(right_col_idx + b)].width = box_width

    def _write_entry(ws, row, start_col, entry):
        kind = entry['type']
        
        if kind == 'header':
            ws.merge_cells(
                start_row=row, start_column=start_col,
                end_row=row, end_column=start_col + boxes_count
            )
            cell = ws.cell(row=row, column=start_col, value=entry['text'])
            cell.font = group_font
            cell.alignment = center
            cell.fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
            for c in range(start_col, start_col + boxes_count + 1):
                ws.cell(row=row, column=c).border = thin_border
            
        elif kind == 'member':
            cell_name = ws.cell(row=row, column=start_col, value=entry['text'])
            cell_name.font = leader_name_font if entry.get('is_leader') else name_font
            cell_name.alignment = center
            cell_name.border = thin_border
            for b in range(1, boxes_count + 1):
                ws.cell(row=row, column=start_col + b).border = thin_border

    start_row = 1
    
    for i in range(max_rows):
        r = start_row + i
        
        l_entry = left_entries[i] if i < len(left_entries) else None
        r_entry = right_entries[i] if i < len(right_entries) else None
        
        if l_entry: _write_entry(ws, r, left_col_idx, l_entry)
        if r_entry: _write_entry(ws, r, right_col_idx, r_entry)
            
        h_pts = row_weights[i] * unit_h
        ws.row_dimensions[r].height = h_pts
            
    last_col_letter = get_column_letter(right_col_idx + boxes_count)
    ws.print_area = f"A1:{last_col_letter}{start_row + max_rows - 1}"

    ws.page_setup.orientation = ws.ORIENTATION_PORTRAIT
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.margins = PageMargins(
        left=0.25, right=0.25,
        top=0.5, bottom=0.25,
        header=0.3, footer=0.2
    )
    ws.print_options.horizontalCentered = True
    
    ws.page_setup.fitToPage = True
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1

    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = f'attachment; filename="{classroom.name}_小组作业表.xlsx"'
    wb.save(response)

    return response


def save_layout_snapshot(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'POST':
        name = str(request.POST.get('snapshot_name', '')).strip()
        if not name:
            return redirect('classroom_detail', pk=pk)
        before_state = _capture_history_state(classroom)
        data = _snapshot_payload(classroom, include_students=False)
        snapshot, _ = LayoutSnapshot.objects.update_or_create(
            classroom=classroom,
            name=name,
            defaults={'data': data}
        )
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'save_layout_snapshot',
            extra={'snapshot_id': snapshot.pk, 'snapshot_name': snapshot.name},
        )
        if _is_ajax_request(request):
            return JsonResponse({'status': 'success', 'snapshot': {'id': snapshot.pk, 'name': snapshot.name}})
    return redirect('classroom_detail', pk=pk)


def load_layout_snapshot(request, pk, snapshot_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    snapshot = get_object_or_404(LayoutSnapshot, pk=snapshot_id, classroom=classroom)
    before_state = _capture_history_state(classroom)
    _apply_layout_data(classroom, snapshot.data, replace_students=False)
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'load_layout_snapshot',
        extra={'snapshot_id': snapshot.pk, 'snapshot_name': snapshot.name},
    )
    return redirect('classroom_detail', pk=pk)


def delete_layout_snapshot(request, pk, snapshot_id):
    classroom = get_object_or_404(Classroom, pk=pk)
    snapshot = get_object_or_404(LayoutSnapshot, pk=snapshot_id, classroom=classroom)
    before_state = _capture_history_state(classroom)
    snapshot_name = snapshot.name
    snapshot.delete()
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'delete_layout_snapshot',
        extra={'snapshot_id': snapshot_id, 'snapshot_name': snapshot_name},
    )
    if _is_ajax_request(request):
        return JsonResponse({'status': 'success'})
    return redirect('classroom_detail', pk=pk)


def export_seats_file(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    data = _serialize_seats_file_bundle(classroom)
    payload = json.dumps(data, ensure_ascii=False, indent=2)
    response = HttpResponse(payload, content_type='application/octet-stream')
    filename = escape_uri_path(f'{classroom.name}.seats')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


def import_seats_file(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'POST' and request.FILES.get('seats_file'):
        seats_file = request.FILES['seats_file']
        try:
            raw = seats_file.read().decode('utf-8')
            data = json.loads(raw)
            import_mode = 'full'
            before_state = None
            if not isinstance(data.get('current_state'), dict):
                before_state = _capture_history_state(classroom)
                import_mode = 'legacy'
            resolved_mode = _import_seats_file_payload(classroom, data, request=request)
            if import_mode == 'legacy' and resolved_mode == 'legacy' and before_state is not None:
                _push_snapshot_action(request, classroom, before_state, 'import_seats_file')
        except Exception:
            pass
    return redirect('classroom_detail', pk=pk)


BSCE_CLOUD_AUTH_URL = 'https://sce.jbyc.cc/api/auth.php'
BSCE_CLOUD_WORKSPACE_URL = 'https://sce.jbyc.cc/api/workspace.php'
BSCE_CLOUD_TIMEOUT = 15
BSCE_TRANSPORT_SALT = b'sce-transport-salt-v1'
BSCE_PBKDF2_ITERATIONS = 100_000
BSCE_CLOUD_HEADERS = {
    'Accept': 'application/json, text/plain, */*',
    'Content-Type': 'application/json',
    'Origin': 'https://sce.jbyc.cc',
    'Priority': 'u=1, i',
    'Referer': 'https://sce.jbyc.cc/',
    'Sec-Fetch-Dest': 'empty',
    'Sec-Fetch-Mode': 'cors',
    'Sec-Fetch-Site': 'same-origin',
}
BSCE_CLOUD_BROWSER_PROFILES = [
    {
        'User-Agent': (
            'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) '
            'AppleWebKit/537.36 (KHTML, like Gecko) '
            'Chrome/147.0.0.0 Safari/537.36 Edg/147.0.0.0'
        ),
        'Sec-CH-UA': '"Microsoft Edge";v="147", "Not.A/Brand";v="8", "Chromium";v="147"',
        'Sec-CH-UA-Mobile': '?0',
        'Sec-CH-UA-Platform': '"macOS"',
    },
    {
        'User-Agent': (
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
            'AppleWebKit/537.36 (KHTML, like Gecko) '
            'Chrome/147.0.0.0 Safari/537.36 Edg/147.0.0.0'
        ),
        'Sec-CH-UA': '"Microsoft Edge";v="147", "Not.A/Brand";v="8", "Chromium";v="147"',
        'Sec-CH-UA-Mobile': '?0',
        'Sec-CH-UA-Platform': '"Windows"',
    },
    {
        'User-Agent': (
            'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) '
            'AppleWebKit/537.36 (KHTML, like Gecko) '
            'Chrome/147.0.0.0 Safari/537.36'
        ),
        'Sec-CH-UA': '"Google Chrome";v="147", "Not.A/Brand";v="8", "Chromium";v="147"',
        'Sec-CH-UA-Mobile': '?0',
        'Sec-CH-UA-Platform': '"macOS"',
    },
    {
        'User-Agent': (
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
            'AppleWebKit/537.36 (KHTML, like Gecko) '
            'Chrome/147.0.0.0 Safari/537.36'
        ),
        'Sec-CH-UA': '"Google Chrome";v="147", "Not.A/Brand";v="8", "Chromium";v="147"',
        'Sec-CH-UA-Mobile': '?0',
        'Sec-CH-UA-Platform': '"Windows"',
    },
]
BSCE_CLOUD_ACCEPT_LANGUAGES = [
    'zh-CN,zh;q=0.9,en;q=0.8',
    'zh-CN,zh;q=0.9',
    'zh-CN,zh-Hans;q=0.9,en-US;q=0.8,en;q=0.7',
]


def _bsce_cloud_cookie_header(csrf_token):
    cookies = [
        f'rth-uid={uuid.uuid4()}',
        f'sce_csrf={csrf_token}',
    ]
    last_workspace = str(
        getattr(settings, 'BSCE_CLOUD_LAST_WORKSPACE_COOKIE', os.environ.get('BSCE_CLOUD_LAST_WORKSPACE_COOKIE', '')) or ''
    ).strip()
    clearance = str(
        getattr(settings, 'BSCE_CLOUD_CLEARANCE', os.environ.get('BSCE_CLOUD_CLEARANCE', '')) or ''
    ).strip()
    if last_workspace:
        cookies.append(f'sce_last_workspace={last_workspace}')
    if clearance:
        cookies.append(f'rth-clearance={clearance}')
    return '; '.join(cookies)


def _bsce_generate_csrf():
    return secrets.token_hex(20)


def _bsce_parse_cookie_header(cookie_header):
    cookies = {}
    for part in str(cookie_header or '').split(';'):
        if '=' not in part:
            continue
        name, value = part.split('=', 1)
        name = name.strip()
        if not name:
            continue
        cookies[name] = value.strip()
    return cookies


def _bsce_format_cookie_header(cookies):
    return '; '.join(
        f'{name}={value}'
        for name, value in cookies.items()
        if name and value is not None
    )


def _bsce_extract_set_cookie_headers(response):
    headers = getattr(response, 'headers', None)
    if headers is None and callable(getattr(response, 'info', None)):
        headers = response.info()

    cookie_headers = []
    if headers is not None:
        get_all = getattr(headers, 'get_all', None)
        if callable(get_all):
            cookie_headers.extend(get_all('Set-Cookie') or [])
        elif hasattr(headers, 'get'):
            header_value = headers.get('Set-Cookie')
            if header_value:
                cookie_headers.append(header_value)

    getheaders = getattr(response, 'getheaders', None)
    if callable(getheaders):
        for name, value in getheaders():
            if str(name).lower() == 'set-cookie' and value:
                cookie_headers.append(value)

    return cookie_headers


def _bsce_update_browser_session_cookies(browser_session, response):
    cookie_headers = _bsce_extract_set_cookie_headers(response)
    if not cookie_headers:
        return

    cookie_jar = browser_session.setdefault(
        'cookies',
        _bsce_parse_cookie_header(browser_session.get('headers', {}).get('Cookie')),
    )
    for header in cookie_headers:
        parsed = SimpleCookie()
        try:
            parsed.load(header)
        except CookieError:
            continue
        for name, morsel in parsed.items():
            cookie_jar[name] = morsel.value

    browser_session.setdefault('headers', {})['Cookie'] = _bsce_format_cookie_header(cookie_jar)


def _bsce_derive_transport_key(username):
    key_material = f'sce-auth-{username}'.encode('utf-8')
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=BSCE_TRANSPORT_SALT,
        iterations=BSCE_PBKDF2_ITERATIONS,
    )
    return kdf.derive(key_material)


def _bsce_encrypt_login_password(password, username):
    key = _bsce_derive_transport_key(username)
    iv = secrets.token_bytes(12)
    ciphertext = AESGCM(key).encrypt(iv, str(password or '').encode('utf-8'), None)
    return base64.b64encode(iv + ciphertext).decode('ascii')


def _bsce_cloud_browser_session():
    csrf_token = _bsce_generate_csrf()
    cookie_header = _bsce_cloud_cookie_header(csrf_token)
    headers = {
        **BSCE_CLOUD_HEADERS,
        **secrets.choice(BSCE_CLOUD_BROWSER_PROFILES),
        'Accept-Language': secrets.choice(BSCE_CLOUD_ACCEPT_LANGUAGES),
        'Cookie': cookie_header,
        'X-CSRF-Token': csrf_token,
    }
    return {
        'csrf': csrf_token,
        'headers': headers,
        'cookies': _bsce_parse_cookie_header(cookie_header),
    }


def _extract_json_or_form_payload(request):
    content_type = str(request.content_type or '').split(';', 1)[0].strip().lower()
    if content_type == 'application/json':
        try:
            payload = json.loads(request.body or b'{}')
        except json.JSONDecodeError as exc:
            raise ValueError('JSON 请求体格式错误') from exc
        if not isinstance(payload, dict):
            raise ValueError('JSON 请求体必须是对象')
        return payload
    return dict(request.POST.items())


def _bsce_cloud_request_payload(action, username, password=None, token=None, file_id=None):
    payload = {
        'action': action,
        'username': username,
    }
    if action == 'login':
        payload['encryptedPassword'] = _bsce_encrypt_login_password(password, username)
    if token:
        payload['token'] = token
    if file_id:
        payload['fileId'] = file_id
    return payload


def _bsce_json_post(url, payload, timeout=BSCE_CLOUD_TIMEOUT, browser_session=None):
    browser_session = browser_session or _bsce_cloud_browser_session()
    request_payload = dict(payload)
    request_payload['_csrf'] = browser_session['csrf']

    def execute(context):
        body = json.dumps(request_payload, ensure_ascii=False).encode('utf-8')
        req = urllib.request.Request(
            url,
            data=body,
            headers=browser_session['headers'],
            method='POST',
        )
        with urllib.request.urlopen(req, timeout=timeout, context=context) as response:
            _bsce_update_browser_session_cookies(browser_session, response)
            raw = response.read().decode('utf-8', errors='replace')
        try:
            return json.loads(raw)
        except json.JSONDecodeError as exc:
            raise ValueError('云端返回的数据不是有效 JSON') from exc

    contexts = [ssl.create_default_context()]
    used_cert_fallback = False
    while contexts:
        context = contexts.pop(0)
        try:
            return execute(context)
        except urllib.error.HTTPError as exc:
            detail = ''
            try:
                detail = exc.read().decode('utf-8', errors='replace')
            except Exception:
                detail = str(exc)
            raise ValueError(f'云端请求失败：HTTP {exc.code} {detail[:180]}') from exc
        except urllib.error.URLError as exc:
            reason = getattr(exc, 'reason', exc)
            cert_failed = isinstance(reason, ssl.SSLCertVerificationError) or 'CERTIFICATE_VERIFY_FAILED' in str(reason)
            if cert_failed and not used_cert_fallback:
                used_cert_fallback = True
                contexts.append(ssl._create_unverified_context())
                continue
            raise ValueError(f'无法连接 BSCE 云端：{reason}') from exc

    raise ValueError('无法连接 BSCE 云端')


def _bsce_require_success(response, fallback_message):
    if not isinstance(response, dict):
        raise ValueError(fallback_message)
    if not response.get('success'):
        raise ValueError(str(response.get('message') or fallback_message))
    return response.get('data')


def _bsce_cloud_login(username, password, browser_session=None):
    browser_session = browser_session or _bsce_cloud_browser_session()
    data = _bsce_json_post(
        BSCE_CLOUD_AUTH_URL,
        _bsce_cloud_request_payload('login', username, password=password),
        browser_session=browser_session,
    )
    _bsce_require_success(data, 'BSCE 云端登录失败')
    token = str(browser_session.get('cookies', {}).get('sce_token') or '').strip()
    if not token:
        raise ValueError('BSCE 云端登录成功但响应 Cookie 没有返回 sce_token')
    settings_data = _bsce_json_post(
        BSCE_CLOUD_AUTH_URL,
        _bsce_cloud_request_payload('get_settings', username, token=token),
        browser_session=browser_session,
    )
    _bsce_require_success(settings_data, '获取 BSCE 云端设置失败')
    return token


def _bsce_cloud_list_workspaces(username, password):
    browser_session = _bsce_cloud_browser_session()
    token = _bsce_cloud_login(username, password, browser_session=browser_session)
    data = _bsce_json_post(
        BSCE_CLOUD_WORKSPACE_URL,
        _bsce_cloud_request_payload('list', username, token=token),
        browser_session=browser_session,
    )
    rows = _bsce_require_success(data, '获取 BSCE 云端工作区失败')
    if not isinstance(rows, list):
        raise ValueError('BSCE 云端工作区列表格式错误')

    normalized = []
    for item in rows:
        if not isinstance(item, dict):
            continue
        file_id = str(item.get('fileId') or '').strip()
        if not file_id:
            continue
        metadata = item.get('metadata') if isinstance(item.get('metadata'), dict) else {}
        normalized.append({
            'fileId': file_id,
            'metadata': {
                'author': str(metadata.get('author') or ''),
                'name': str(metadata.get('name') or ''),
                'time': str(metadata.get('time') or ''),
                'size': metadata.get('size') or 0,
            },
        })
    return normalized


def _bsce_cloud_load_workspace(file_id, username, password):
    file_id = str(file_id or '').strip()
    if not file_id:
        raise ValueError('缺少 BSCE 云端工作区 fileId')

    browser_session = _bsce_cloud_browser_session()
    token = _bsce_cloud_login(username, password, browser_session=browser_session)
    data = _bsce_json_post(
        BSCE_CLOUD_WORKSPACE_URL,
        _bsce_cloud_request_payload('load', username, token=token, file_id=file_id),
        browser_session=browser_session,
    )
    workspace = _bsce_require_success(data, '加载 BSCE 云端工作区失败')
    if not isinstance(workspace, dict):
        raise ValueError('BSCE 云端工作区内容格式错误')
    content = workspace.get('content')
    if not isinstance(content, dict):
        raise ValueError('BSCE 云端工作区没有可导入内容')
    metadata = workspace.get('metadata') if isinstance(workspace.get('metadata'), dict) else {}
    return {
        'fileId': file_id,
        'metadata': metadata,
        'content': content,
    }


def _parse_bsce_number(value, default=0):
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _parse_bsce_score(value):
    text = str(value or '').strip()
    if not text:
        return 0
    try:
        return float(text)
    except (TypeError, ValueError):
        return 0


def _is_bsce_payload(data):
    if not isinstance(data, dict):
        return False
    meta = data.get('meta') if isinstance(data.get('meta'), dict) else {}
    if str(meta.get('app') or '').strip() == 'SeatingChartEditor':
        return True
    return isinstance(data.get('students'), list) and isinstance(data.get('layout'), dict) and isinstance(data.get('tags'), list)


def _extract_bsce_student_score(student_data, tag_map):
    for tag_id in student_data.get('tags') or []:
        score = _parse_bsce_score(tag_map.get(tag_id))
        if score > 100:
            return score
    return 0


def _build_bsce_group_columns(config):
    group_count = _parse_bsce_number(config.get('groupCount'), 0)
    default_columns = _parse_bsce_number(config.get('columnsPerGroup'), 2)
    default_rows = _parse_bsce_number(config.get('seatsPerColumn'), 1)
    raw_groups = config.get('groups') if isinstance(config.get('groups'), list) else []

    group_columns = []
    group_rows = []
    total_groups = max(group_count, len(raw_groups), 1)
    for index in range(total_groups):
        group_data = raw_groups[index] if index < len(raw_groups) and isinstance(raw_groups[index], dict) else {}
        group_columns.append(max(1, _parse_bsce_number(group_data.get('columns'), default_columns)))
        group_rows.append(max(1, _parse_bsce_number(group_data.get('rows'), default_rows)))
    return group_columns, group_rows


def _apply_bsce_payload(classroom, data):
    if not _is_bsce_payload(data):
        raise ValueError('不是有效的 BSCE 座位表文件')

    students_payload = data.get('students') if isinstance(data.get('students'), list) else []
    tags_payload = data.get('tags') if isinstance(data.get('tags'), list) else []
    layout = data.get('layout') if isinstance(data.get('layout'), dict) else {}
    config = layout.get('config') if isinstance(layout.get('config'), dict) else {}
    seats_payload = layout.get('seats') if isinstance(layout.get('seats'), list) else []

    if not students_payload:
        raise ValueError('BSCE 文件中没有学生数据')
    if not seats_payload:
        raise ValueError('BSCE 文件中没有座位布局')

    tag_map = {}
    for tag in tags_payload:
        if not isinstance(tag, dict):
            continue
        tag_map[tag.get('id')] = str(tag.get('name') or '').strip()

    group_columns, group_rows = _build_bsce_group_columns(config)
    group_count = len(group_columns)
    group_start_cols = []
    cursor_col = 1
    for columns in group_columns:
        group_start_cols.append(cursor_col)
        cursor_col += columns + 1

    max_source_row = 0
    for seat_data in seats_payload:
        if not isinstance(seat_data, dict):
            continue
        max_source_row = max(max_source_row, _parse_bsce_number(seat_data.get('row'), 0) + 1)

    row_count = max([max_source_row, *group_rows, 1])
    col_count = sum(group_columns) + max(0, group_count - 1)

    with transaction.atomic():
        before_students = len(students_payload)
        _sync_seats(classroom, row_count, col_count)
        _apply_podium_guards(classroom, None, None)
        classroom.seats.update(student=None, group=None, cell_type=SeatCellType.EMPTY)
        SeatConstraint.objects.filter(classroom=classroom).delete()
        SeatGroup.objects.filter(classroom=classroom).delete()
        Student.objects.filter(classroom=classroom).delete()

        student_map = {}
        for student_data in students_payload:
            if not isinstance(student_data, dict):
                continue
            name = str(student_data.get('name') or '').strip()
            if not name:
                continue
            source_id = student_data.get('id')
            student = Student.objects.create(
                classroom=classroom,
                name=name,
                student_id=str(student_data.get('studentNumber') or '').strip() or None,
                score=_extract_bsce_student_score(student_data, tag_map),
            )
            student_map[source_id] = student

        group_map = {}
        for group_index in range(group_count):
            group_map[group_index] = SeatGroup.objects.create(
                classroom=classroom,
                name=f'第{group_index + 1}组',
                order=group_index,
            )

        seat_map = _build_seat_map(classroom.seats.all())
        for aisle_col in group_start_cols[1:]:
            gap_col = aisle_col - 1
            for row in range(1, row_count + 1):
                seat = seat_map.get((row, gap_col))
                if seat:
                    seat.cell_type = SeatCellType.AISLE
                    seat.save(update_fields=['cell_type'])

        assigned_count = 0
        source_seat_count = 0
        for seat_data in seats_payload:
            if not isinstance(seat_data, dict):
                continue
            group_index = _parse_bsce_number(seat_data.get('group'), -1)
            if group_index < 0 or group_index >= group_count:
                continue
            local_col = _parse_bsce_number(seat_data.get('col'), 0)
            row = _parse_bsce_number(seat_data.get('row'), 0) + 1
            col = group_start_cols[group_index] + local_col
            seat = seat_map.get((row, col))
            if not seat:
                continue

            source_seat_count += 1
            is_empty_cell = bool(seat_data.get('empty'))
            if is_empty_cell:
                seat.cell_type = SeatCellType.EMPTY
                seat.student = None
                seat.group = None
            else:
                seat.cell_type = SeatCellType.SEAT
                seat.group = group_map.get(group_index)
                seat.student = student_map.get(seat_data.get('studentId'))
                if seat.student_id:
                    assigned_count += 1
            seat.save(update_fields=['cell_type', 'student', 'group'])

    return {
        'students': before_students,
        'created_students': len(student_map),
        'rows': row_count,
        'cols': col_count,
        'groups': group_count,
        'source_seats': source_seat_count,
        'assigned': assigned_count,
    }


def import_bsce_file(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    if request.method == 'POST' and request.FILES.get('bsce_file'):
        bsce_file = request.FILES['bsce_file']
        try:
            raw = bsce_file.read().decode('utf-8-sig')
            data = json.loads(raw)
            before_state = _capture_history_state(classroom)
            result = _apply_bsce_payload(classroom, data)
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                'import_bsce_file',
                extra=result,
            )
            if _is_ajax_request(request):
                return JsonResponse({'status': 'success', 'message': 'BSCE 导入完成', **result})
        except Exception as e:
            if _is_ajax_request(request):
                return JsonResponse({'status': 'error', 'message': f'BSCE 导入失败：{e}'}, status=400)
    return redirect('classroom_detail', pk=pk)


@require_POST
def import_bsce_cloud(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        payload = _extract_json_or_form_payload(request)
        action = str(payload.get('action') or 'list').strip().lower()
        username = str(payload.get('username') or '').strip()
        password = str(payload.get('password') or '').strip()
        if not username or not password:
            return JsonResponse({'status': 'error', 'message': '请输入 BSCE 云端账号和密码'}, status=400)

        if action in {'list', 'workspaces'}:
            workspaces = _bsce_cloud_list_workspaces(username, password)
            return JsonResponse({
                'status': 'success',
                'workspaces': workspaces,
                'count': len(workspaces),
            })

        if action in {'load', 'import'}:
            file_id = payload.get('fileId') or payload.get('file_id')
            workspace = _bsce_cloud_load_workspace(file_id, username, password)
            before_state = _capture_history_state(classroom)
            result = _apply_bsce_payload(classroom, workspace['content'])
            metadata = workspace.get('metadata') or {}
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                'import_bsce_cloud',
                extra={
                    **result,
                    'file_id': workspace['fileId'],
                    'workspace_name': str(metadata.get('name') or ''),
                    'workspace_time': str(metadata.get('time') or ''),
                },
            )
            return JsonResponse({
                'status': 'success',
                'message': 'BSCE 云导入完成',
                'fileId': workspace['fileId'],
                'metadata': metadata,
                **result,
            })

        return JsonResponse({'status': 'error', 'message': '未知 BSCE 云导入操作'}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


def undo_action(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    entry = _get_history_queryset(pk).filter(is_applied=True).order_by('-pk').first()
    if not entry:
        return JsonResponse({'status': 'error', 'message': '没有可撤销操作'}, status=400)
    action = entry.payload or {}
    with transaction.atomic():
        if not _apply_recorded_history_action(classroom, action, forward=False):
            return JsonResponse({'status': 'error', 'message': '撤销失败：历史记录不可用'}, status=400)
        entry.is_applied = False
        entry.save(update_fields=['is_applied'])
    return JsonResponse({'status': 'success'})


def redo_action(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    entry = _get_history_queryset(pk).filter(is_applied=False).order_by('pk').first()
    if not entry:
        return JsonResponse({'status': 'error', 'message': '没有可重做操作'}, status=400)
    action = entry.payload or {}
    with transaction.atomic():
        if not _apply_recorded_history_action(classroom, action, forward=True):
            return JsonResponse({'status': 'error', 'message': '重做失败：历史记录不可用'}, status=400)
        entry.is_applied = True
        entry.save(update_fields=['is_applied'])
    return JsonResponse({'status': 'success'})


def _cloud_delete_backed_up_classroom(meta):
    if not meta:
        return None

    backed_up = bool(meta.last_sync_at and int(meta.cloud_version or 0) > 0)
    if not backed_up:
        return None

    session = get_active_cloud_session()
    if not session:
        raise CloudAPIError('该班级已同步到云端，删除前请先登录云服务以同步删除操作', status_code=401)

    try:
        return cloud_api_request(session, 'DELETE', f'/api/sync/{meta.uuid}', {
            'base_version': int(meta.cloud_version or 0),
            'device_id': 'local-delete',
        })
    except CloudAPIError as exc:
        payload = exc.payload or {}
        message = str(payload.get('message') or exc)
        if exc.status_code == 404 or '班级不存在' in message:
            return {'ok': True, 'status': 'success', 'already_deleted': True, 'uuid': str(meta.uuid)}
        raise


def delete_classroom(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    sync_meta = SyncMeta.objects.filter(classroom=classroom).first()
    try:
        _cloud_delete_backed_up_classroom(sync_meta)
    except Exception as exc:
        return _cloud_error_response(exc)

    with suspend_sync_version_bump(), transaction.atomic():
        classroom.left_guardian = None
        classroom.right_guardian = None
        classroom.save(update_fields=['left_guardian', 'right_guardian'])
        classroom.groups.update(leader=None)
        classroom.delete()
    return redirect('index')


@require_POST
def rename_classroom(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    is_json = bool(request.content_type and 'application/json' in request.content_type)
    try:
        if is_json:
            payload = json.loads(request.body or '{}')
            new_name = str(payload.get('name') or '').strip()
        else:
            new_name = str(request.POST.get('name') or '').strip()
    except Exception:
        return JsonResponse({'status': 'error', 'message': '请求数据格式错误'}, status=400)

    if not new_name:
        return JsonResponse({'status': 'error', 'message': '班级名称不能为空'}, status=400)
    if len(new_name) > 100:
        return JsonResponse({'status': 'error', 'message': '班级名称不能超过 100 个字符'}, status=400)

    before_state = _capture_history_state(classroom)
    classroom.name = new_name
    classroom.save(update_fields=['name'])
    _push_snapshot_action(
        request,
        classroom,
        before_state,
        'rename_classroom',
        extra={'name': classroom.name},
    )

    if is_json or _is_ajax_request(request):
        return JsonResponse({'status': 'success', 'name': classroom.name})
    return redirect('classroom_detail', pk=pk)


@require_POST
def apply_suggestion(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    suggestion_type = (request.GET.get('type') or '').strip()

    if suggestion_type in DISABLED_SUGGESTION_TYPES:
        return JsonResponse({'status': 'success', 'message': '该建议已停用'})
    
    if suggestion_type == 'swap_balance':
        s1_id = request.GET.get('s1')
        s2_id = request.GET.get('s2')
        try:
            if not s1_id or not s2_id:
                return JsonResponse({'status': 'error', 'message': '缺少学生参数'}, status=400)
            if str(s1_id) == str(s2_id):
                return JsonResponse({'status': 'error', 'message': '不能交换同一名学生'}, status=400)
            s1 = classroom.students.filter(pk=s1_id).first()
            s2 = classroom.students.filter(pk=s2_id).first()
            if not s1 or not s2:
                return JsonResponse({'status': 'error', 'message': '学生不属于当前班级'}, status=400)
            seat1 = getattr(s1, 'assigned_seat', None)
            seat2 = getattr(s2, 'assigned_seat', None)
            if not seat1 or not seat2:
                return JsonResponse({'status': 'error', 'message': '学生未入座，无法交换'}, status=400)
            if seat1.classroom_id != classroom.pk or seat2.classroom_id != classroom.pk:
                return JsonResponse({'status': 'error', 'message': '座位不属于当前班级'}, status=400)
            
            before_state = _capture_history_state(classroom)
            with transaction.atomic():
                _swap_seats(seat1, seat2)
                violations = _stabilize_layout_with_rules(classroom, request)
                if violations:
                    raise ValueError(f'交换失败：{_format_issues_preview(violations)}')
            _push_snapshot_action(
                request,
                classroom,
                before_state,
                'swap',
                extra=_build_swap_action(s1, s2),
            )
            return JsonResponse({'status': 'success', 'message': f'已执行交换并自动校正约束：{s1.name} / {s2.name}'})
        except ValueError as e:
            return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    
    return JsonResponse({'status': 'error', 'message': '未知建议'}, status=400)


@require_POST
def dismiss_suggestion(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    suggestion_type = request.GET.get('type')
    
    if suggestion_type == 'export':
        request.session[f'ignore_export_{pk}'] = True
        return JsonResponse({'status': 'success'})
        
    return JsonResponse({'status': 'success'})

@require_POST
def set_group_leader(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        data = json.loads(request.body)
        student_id = data.get('student_id')
        
        student = get_object_or_404(Student, pk=student_id, classroom=classroom)
        seat = getattr(student, 'assigned_seat', None)
        if not seat or not seat.group:
            return JsonResponse({'status': 'error', 'message': '该学生未分配或未在小组中'}, status=400)
            
        group = seat.group
        before_state = _capture_history_state(classroom)
        
        if group.leader == student:
            group.leader = None
        else:
            group.leader = student
            
        group.save()
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'set_group_leader',
            extra={'group_id': group.pk, 'student_id': student.pk},
        )
        
        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@require_POST
def set_podium_guards(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        if request.content_type and 'application/json' in request.content_type:
            data = json.loads(request.body or '{}')
        else:
            data = request.POST

        marker = object()
        left_raw = data.get('left_student_id', marker)
        right_raw = data.get('right_student_id', marker)

        def _resolve_student(raw_value, label):
            if raw_value in (None, ''):
                return None
            try:
                student_id = int(raw_value)
            except (TypeError, ValueError) as exc:
                raise ValueError(f'{label}学生 ID 不合法') from exc
            student = classroom.students.filter(pk=student_id).first()
            if not student:
                raise ValueError(f'{label}学生不存在')
            return student

        left_student = classroom.left_guardian if left_raw is marker else _resolve_student(left_raw, '左护法')
        right_student = classroom.right_guardian if right_raw is marker else _resolve_student(right_raw, '右护法')

        before_state = _capture_history_state(classroom)
        _apply_podium_guards(classroom, left_student, right_student)
        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'set_podium_guards',
            extra={
                'left_guardian_student_id': classroom.left_guardian_id,
                'right_guardian_student_id': classroom.right_guardian_id,
            },
        )
        return JsonResponse({
            'status': 'success',
            'podium_guards': _serialize_podium_guards(classroom),
        })
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@require_POST
def toggle_fixed_seat(request, pk):
    classroom = get_object_or_404(Classroom, pk=pk)
    try:
        if request.content_type and 'application/json' in request.content_type:
            data = json.loads(request.body or '{}')
        else:
            data = request.POST

        row = _safe_int(data.get('row'), 0)
        col = _safe_int(data.get('col'), 0)
        if row < 1 or col < 1:
            raise ValueError('请选择有效座位后再固定')

        seat = classroom.seats.select_related('student').filter(
            row=row,
            col=col,
            cell_type=SeatCellType.SEAT,
        ).first()
        if not seat:
            raise ValueError('目标位置不是可用座位')
        if not seat.student_id or not seat.student:
            raise ValueError('当前座位没有学生，无法固定')

        constraint = classroom.constraints.filter(
            student=seat.student,
            constraint_type=SeatConstraint.ConstraintType.MUST_SEAT,
        ).order_by('-enabled', 'pk').first()

        desired_enabled_raw = data.get('enabled') if hasattr(data, 'get') else None
        currently_enabled = bool(constraint and constraint.enabled)
        next_enabled = (not currently_enabled) if desired_enabled_raw in (None, '') else _parse_bool(desired_enabled_raw)

        before_state = _capture_history_state(classroom)

        if next_enabled:
            payload = {
                'constraint_type': SeatConstraint.ConstraintType.MUST_SEAT,
                'student': seat.student,
                'row': row,
                'col': col,
                'distance': 1,
                'enabled': True,
                'note': constraint.note if constraint else '',
            }
            cleaned = normalize_constraint_payload(classroom, payload, instance=constraint)
            validate_constraint_candidate(classroom, cleaned, instance=constraint)

            if constraint:
                constraint.constraint_type = cleaned['constraint_type']
                constraint.student = cleaned['student']
                constraint.target_student = cleaned['target_student']
                constraint.row = cleaned['row']
                constraint.col = cleaned['col']
                constraint.distance = cleaned['distance']
                constraint.enabled = cleaned['enabled']
                constraint.note = FIXED_SEAT_NOTE_MARKER
                constraint.save()
            else:
                constraint = SeatConstraint.objects.create(
                    classroom=classroom,
                    constraint_type=cleaned['constraint_type'],
                    student=cleaned['student'],
                    target_student=cleaned['target_student'],
                    row=cleaned['row'],
                    col=cleaned['col'],
                    distance=cleaned['distance'],
                    enabled=cleaned['enabled'],
                    note=FIXED_SEAT_NOTE_MARKER,
                )
        elif constraint:
            constraint.delete()

        _push_snapshot_action(
            request,
            classroom,
            before_state,
            'toggle_fixed_seat',
            extra={
                'student_id': seat.student_id,
                'row': row,
                'col': col,
                'enabled': next_enabled,
            },
        )
        return JsonResponse({
            'status': 'success',
            'enabled': next_enabled,
            'message': '已固定此座位' if next_enabled else '已取消固定座位',
        })
    except ConstraintServiceError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except ValueError as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

from django.views.decorators.csrf import csrf_exempt


def _cloud_json_body(request):
    try:
        raw = request.body.decode('utf-8') if request.body else '{}'
        return json.loads(raw or '{}')
    except Exception:
        raise ValueError('请求数据格式错误')


def _cloud_error_response(exc, default_status=400):
    if isinstance(exc, CloudAPIError):
        payload = dict(exc.payload or {})
        payload.setdefault('status', 'error')
        payload.setdefault('message', str(exc))
        return JsonResponse(payload, status=exc.status_code)
    return JsonResponse({'status': 'error', 'message': str(exc)}, status=default_status)


def _cloud_callback_url(request):
    return request.build_absolute_uri(reverse('cloud_callback'))


def _cloud_local_session_or_401():
    session = get_active_cloud_session()
    if not session:
        raise CloudAPIError('尚未登录云服务', status_code=401, payload={
            'status': 'error',
            'error': 'not_logged_in',
            'message': '尚未登录云服务',
        })
    return session


def _refresh_cloud_subscription_if_logged_in(session, *, strict=False):
    if not session:
        return None
    return refresh_cloud_subscription(session, strict=strict)


def _cloud_remote_versions(session):
    payload = _cloud_remote_status(session)
    return _cloud_versions_from_status(payload)


def _cloud_remote_status(session):
    payload = cloud_api_request(session, 'GET', '/api/sync/status')
    return payload if isinstance(payload, dict) else {}


def _cloud_versions_from_status(payload):
    versions = payload.get('versions') if isinstance(payload, dict) else {}
    return versions if isinstance(versions, dict) else {}


def _cloud_classrooms_from_status(payload):
    rows = payload.get('classrooms') if isinstance(payload, dict) else []
    return rows if isinstance(rows, list) else []


def _cloud_remote_version_for(versions, classroom_uuid):
    raw = versions.get(str(classroom_uuid))
    if raw in (None, ''):
        return None
    try:
        return int(raw)
    except (TypeError, ValueError):
        return None


def _parse_cloud_operation_time(value):
    if not value:
        return None
    parsed = parse_datetime(str(value))
    if parsed and timezone.is_naive(parsed):
        parsed = timezone.make_aware(parsed, timezone.get_current_timezone())
    return parsed


def _sync_meta_operation_time(meta, classroom=None):
    if meta is None:
        return None
    return meta.last_operation_at or meta.updated_at or getattr(classroom, 'created_at', None)


def _operation_time_newer(left, right):
    if not left:
        return False
    if not right:
        return True
    return left > right


def _cloud_operation_times_from_status(payload):
    if not isinstance(payload, dict):
        return {}

    operation_times = {}
    raw_map = payload.get('operation_times')
    if isinstance(raw_map, dict):
        for key, value in raw_map.items():
            if key and value:
                operation_times[str(key)] = value

    for row in _cloud_classrooms_from_status(payload):
        if not isinstance(row, dict):
            continue
        classroom_uuid = str(row.get('uuid') or '').strip()
        operation_at = row.get('last_operation_at') or row.get('last_modified_at')
        if classroom_uuid and operation_at:
            operation_times[classroom_uuid] = operation_at
    return operation_times


def _cloud_remote_operation_time_for(operation_times, classroom_uuid):
    return _parse_cloud_operation_time(operation_times.get(str(classroom_uuid)))


def _cloud_session_payload(session, request=None):
    callback_url = _cloud_callback_url(request) if request else None
    return {
        'logged_in': True,
        'uid': session.uid,
        'nickname': session.nickname,
        'avatar_url': session.avatar_url,
        'email': session.email,
        'tier': session.subscription_tier,
        'tier_display': session.subscription_display_name,
        'expires_at': session.subscription_expires_at.isoformat() if session.subscription_expires_at else None,
        'token_expires_at': session.token_expires_at.isoformat(),
        'limits': session.limits if isinstance(session.limits, dict) else {},
        'cloud_server_url': get_cloud_server_url(),
        'cloud_server_locked': True,
        'official_cloud_server_url': get_cloud_server_url(),
        'data_sharing': get_data_sharing_config(),
        'login_url': build_cloud_login_url(callback_url) if callback_url else None,
    }


def _cloud_export_payload(classroom, session):
    limits = session.limits if isinstance(session.limits, dict) else {}
    max_history_steps = int(limits.get('max_history_steps', 0) or 0)
    sync_ai_conversations = bool(limits.get('sync_ai_conversations', False))

    payload = _serialize_seats_file_bundle(classroom)
    payload.pop('future_mode_config', None)
    meta = payload.get('meta') if isinstance(payload.get('meta'), dict) else {}
    sync_meta = SyncMeta.objects.filter(classroom=classroom).first()
    operation_at = _sync_meta_operation_time(sync_meta, classroom)
    if operation_at:
        meta['last_operation_at'] = operation_at.isoformat()
        payload['meta'] = meta

    history = payload.get('history')
    if isinstance(history, dict):
        entries = history.get('entries')
        if isinstance(entries, list):
            if max_history_steps == 0:
                history['entries'] = []
            elif max_history_steps > 0:
                history['entries'] = entries[-max_history_steps:]
        payload['history'] = history

    legacy_history = payload.get('history_entries')
    if isinstance(legacy_history, list):
        if max_history_steps == 0:
            payload['history_entries'] = []
        elif max_history_steps > 0:
            payload['history_entries'] = legacy_history[-max_history_steps:]

    if not sync_ai_conversations:
        payload.pop('ai_conversations', None)

    return payload


def _cloud_restore_classroom_data(request, classroom_uuid, data, version=None, operation_time=None):
    if not isinstance(data, dict):
        raise ValueError('云端班级数据格式错误')

    meta = SyncMeta.objects.select_related('classroom').filter(uuid=classroom_uuid).first()
    classroom = meta.classroom if meta else None
    classroom_data = data.get('classroom') if isinstance(data.get('classroom'), dict) else {}
    name = str(classroom_data.get('name') or data.get('name') or '云端班级').strip() or '云端班级'
    rows = int(classroom_data.get('rows') or data.get('rows') or 6)
    cols = int(classroom_data.get('cols') or data.get('cols') or 8)
    data_meta = data.get('meta') if isinstance(data.get('meta'), dict) else {}
    operation_at = (
        _parse_cloud_operation_time(operation_time)
        or _parse_cloud_operation_time(data_meta.get('last_operation_at'))
        or _parse_cloud_operation_time(data_meta.get('operation_time'))
    )

    with suspend_sync_version_bump():
        if classroom is None:
            classroom = Classroom.objects.create(name=name, rows=rows, cols=cols)

        _import_seats_file_payload(classroom, data, request=request)
        meta, _ = SyncMeta.objects.get_or_create(classroom=classroom)
        meta.uuid = classroom_uuid
        if version is not None:
            meta.cloud_version = int(version or 0)
            meta.local_version = int(version or 0)
        else:
            meta.local_version = max(int(meta.local_version or 0), int(meta.cloud_version or 0)) + 1
        if operation_at:
            meta.last_operation_at = operation_at
        meta.last_sync_at = timezone.now()
        meta.last_error = ''
        meta.save(update_fields=['uuid', 'cloud_version', 'local_version', 'last_operation_at', 'last_sync_at', 'last_error', 'updated_at'])

    return classroom, meta


def _cloud_pull_and_restore_classroom(request, session, classroom_uuid, fallback_version=None, fallback_operation_time=None):
    payload = cloud_api_request(session, 'GET', f'/api/sync/pull/{classroom_uuid}')
    data = payload.get('data') or payload.get('data_snapshot')
    version = payload.get('version') if payload.get('version') is not None else fallback_version
    operation_time = payload.get('last_operation_at') or payload.get('last_modified_at') or fallback_operation_time
    return _cloud_restore_classroom_data(request, classroom_uuid, data, version=version, operation_time=operation_time)


@require_http_methods(['GET', 'POST'])
@csrf_exempt
def cloud_config(request):
    try:
        if request.method == 'POST':
            data = _cloud_json_body(request)
            server_url = set_cloud_server_url(data.get('cloud_server_url') or data.get('url'))
            if 'data_sharing_enabled' in data:
                set_data_sharing_enabled(data.get('data_sharing_enabled'))
            data_sharing_payload = data.get('data_sharing') if isinstance(data.get('data_sharing'), dict) else {}
            if data_sharing_payload:
                if 'enabled' in data_sharing_payload:
                    set_data_sharing_enabled(data_sharing_payload.get('enabled'))
                if 'local_log_retention_days' in data_sharing_payload:
                    set_data_sharing_log_retention_days(data_sharing_payload.get('local_log_retention_days'))
                if 'log_retention_days' in data_sharing_payload:
                    set_data_sharing_log_retention_days(data_sharing_payload.get('log_retention_days'))
            if 'data_sharing_local_log_retention_days' in data:
                set_data_sharing_log_retention_days(data.get('data_sharing_local_log_retention_days'))
            if 'data_sharing_log_retention_days' in data:
                set_data_sharing_log_retention_days(data.get('data_sharing_log_retention_days'))
            if 'data_sharing_prompt_seen_version' in data:
                set_data_sharing_prompt_seen_version(data.get('data_sharing_prompt_seen_version'))
        else:
            server_url = get_cloud_server_url()
        callback_url = _cloud_callback_url(request)
        return JsonResponse({
            'status': 'success',
            'cloud_server_url': server_url,
            'cloud_server_locked': True,
            'official_cloud_server_url': server_url,
            'callback_url': callback_url,
            'login_url': build_cloud_login_url(callback_url),
            'data_sharing': get_data_sharing_config(),
        })
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_POST
def cloud_open_external(request):
    import webbrowser
    try:
        data = json.loads(request.body or b'{}')
    except Exception:
        return JsonResponse({'status': 'error', 'message': '请求格式错误'}, status=400)
    url = str(data.get('url') or '').strip()
    if not url:
        return JsonResponse({'status': 'error', 'message': '缺少 url'}, status=400)
    try:
        webbrowser.open(url)
        return JsonResponse({'status': 'success'})
    except Exception as exc:
        return JsonResponse({'status': 'error', 'message': str(exc)}, status=500)


@require_http_methods(['GET'])
def cloud_login(request):
    callback_url = request.GET.get('callback') or _cloud_callback_url(request)
    login_url = build_cloud_login_url(callback_url)
    if request.GET.get('format') == 'json':
        return JsonResponse({
            'status': 'success',
            'cloud_server_url': get_cloud_server_url(),
            'cloud_server_locked': True,
            'official_cloud_server_url': get_cloud_server_url(),
            'data_sharing': get_data_sharing_config(),
            'callback_url': callback_url,
            'login_url': login_url,
        })
    return redirect(login_url)


def _cloud_callback_done_page(success, error_detail=None):
    from django.http import HttpResponse
    if success:
        title = '登录成功'
        icon_svg = '<svg width="56" height="56" viewBox="0 0 56 56" fill="none"><circle cx="28" cy="28" r="28" fill="#0a59f7" opacity="0.1"/><circle cx="28" cy="28" r="20" fill="#0a59f7" opacity="0.15"/><path d="M20 28.5L25.5 34L36 22" stroke="#0a59f7" stroke-width="3" stroke-linecap="round" stroke-linejoin="round"/></svg>'
        message = '登录完毕，请关闭此浏览器窗口'
        sub_message = '你可以返回应用继续使用'
    else:
        title = '登录失败'
        icon_svg = '<svg width="56" height="56" viewBox="0 0 56 56" fill="none"><circle cx="28" cy="28" r="28" fill="#ff3b30" opacity="0.1"/><circle cx="28" cy="28" r="20" fill="#ff3b30" opacity="0.15"/><path d="M22 22L34 34M34 22L22 34" stroke="#ff3b30" stroke-width="3" stroke-linecap="round"/></svg>'
        message = '登录失败，请关闭此窗口后重试'
        sub_message = str(error_detail or '')[:120] if error_detail else ''
    html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title} - FuckSeats</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:-apple-system,BlinkMacSystemFont,"SF Pro Display","Helvetica Neue",sans-serif;
background:#f5f5f7;display:flex;align-items:center;justify-content:center;min-height:100vh;
color:#1d1d1f}}
.card{{text-align:center;padding:48px 40px;max-width:380px}}
.icon{{margin-bottom:20px}}
.title{{font-size:22px;font-weight:600;margin-bottom:8px}}
.message{{font-size:15px;color:#86868b;line-height:1.5}}
.sub{{font-size:13px;color:#aeaeb2;margin-top:12px}}
</style>
</head>
<body>
<div class="card">
<div class="icon">{icon_svg}</div>
<div class="title">{title}</div>
<div class="message">{message}</div>
{"<div class='sub'>" + sub_message + "</div>" if sub_message else ""}
</div>
</body>
</html>'''
    return HttpResponse(html, content_type='text/html; charset=utf-8')


@require_http_methods(['GET'])
def cloud_callback(request):
    code = str(request.GET.get('code') or '').strip()
    error = str(request.GET.get('error') or '').strip()
    is_desktop = getattr(settings, 'APP_SHELL', 'browser') != 'browser'

    if error:
        if is_desktop:
            return _cloud_callback_done_page(False, error)
        return redirect(f'/?cloud_login=failed&error={error}')
    if not code:
        if is_desktop:
            return _cloud_callback_done_page(False, 'missing_code')
        return redirect('/?cloud_login=missing_code')

    try:
        payload = cloud_exchange_session_code(code)
        session = save_cloud_session_from_payload(payload)
        try:
            _cloud_sync_classrooms(request, session, {
                'auto': True,
                'device_id': 'login-auto-sync',
            })
        except Exception:
            pass
        if is_desktop:
            return _cloud_callback_done_page(True)
        return redirect('/?cloud_login=success')
    except Exception as exc:
        if is_desktop:
            return _cloud_callback_done_page(False, str(exc))
        return redirect('/?cloud_login=failed')


@require_http_methods(['GET'])
def cloud_userinfo(request):
    session = get_active_cloud_session()
    callback_url = _cloud_callback_url(request)
    if not session:
        return JsonResponse({
            'logged_in': False,
            'cloud_server_url': get_cloud_server_url(),
            'cloud_server_locked': True,
            'official_cloud_server_url': get_cloud_server_url(),
            'data_sharing': get_data_sharing_config(),
            'callback_url': callback_url,
            'login_url': build_cloud_login_url(callback_url),
        })
    _refresh_cloud_subscription_if_logged_in(session)
    return JsonResponse(_cloud_session_payload(session, request=request))


@csrf_exempt
@require_POST
def cloud_refresh_subscription(request):
    try:
        session = _cloud_local_session_or_401()
        refresh_cloud_subscription(session, strict=True)
        return JsonResponse({'status': 'success', **_cloud_session_payload(session, request=request)})
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_POST
def cloud_logout(request):
    session = get_active_cloud_session()
    if session:
        try:
            cloud_api_request(session, 'POST', '/auth/logout')
        except CloudAPIError:
            pass
    clear_cloud_session()
    return JsonResponse({'status': 'success', 'logged_in': False})


@require_http_methods(['GET'])
def cloud_sync_status(request):
    try:
        session = _cloud_local_session_or_401()
        remote = cloud_api_request(session, 'GET', '/api/sync/status')
        classroom_id = request.GET.get('classroom_id')
        local = []
        queryset = Classroom.objects.all().order_by('pk')
        if classroom_id not in (None, ''):
            queryset = queryset.filter(pk=classroom_id)
        for classroom in queryset:
            meta, _ = SyncMeta.objects.get_or_create(classroom=classroom)
            local.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **_serialize_classroom_sync_meta(classroom, meta),
            })
        return JsonResponse({
            'status': 'success',
            'local': local,
            'remote': remote,
        })
    except Exception as exc:
        return _cloud_error_response(exc)


def _cloud_sync_classrooms(request, session, data=None):
    data = data if isinstance(data, dict) else {}
    _refresh_cloud_subscription_if_logged_in(session)
    classroom_ids = data.get('classroom_ids')
    sync_all_classrooms = not classroom_ids
    force = bool(data.get('force'))
    auto_sync = bool(data.get('auto'))
    queryset = Classroom.objects.all().order_by('pk')
    if classroom_ids:
        if not isinstance(classroom_ids, list):
            classroom_ids = [classroom_ids]
        queryset = queryset.filter(pk__in=classroom_ids)

    limits = session.limits if isinstance(session.limits, dict) else {}
    max_classrooms = int(limits.get('max_classrooms', 3) or 3)
    remote_status = _cloud_remote_status(session)
    remote_versions = _cloud_versions_from_status(remote_status)
    remote_operation_times = _cloud_operation_times_from_status(remote_status)
    remote_classrooms = _cloud_classrooms_from_status(remote_status)
    results = []

    for index, classroom in enumerate(queryset):
        meta, _ = SyncMeta.objects.get_or_create(classroom=classroom)
        sync_payload = _serialize_classroom_sync_meta(classroom, meta)
        remote_version = _cloud_remote_version_for(remote_versions, meta.uuid)
        remote_operation_at = _cloud_remote_operation_time_for(remote_operation_times, meta.uuid)
        local_operation_at = _sync_meta_operation_time(meta, classroom)
        cloud_operation_newer = bool(auto_sync and _operation_time_newer(remote_operation_at, local_operation_at))
        upload_by_operation_time = bool(auto_sync and remote_operation_at and not cloud_operation_newer)
        local_has_changes = int(sync_payload['local_version'] or 0) > int(sync_payload['cloud_version'] or 0)
        if max_classrooms != -1 and index >= max_classrooms:
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **sync_payload,
                'status': 'skipped',
                'message': f'当前订阅最多同步 {max_classrooms} 个班级',
            })
            continue

        if cloud_operation_newer and not force:
            try:
                classroom, meta = _cloud_pull_and_restore_classroom(
                    request,
                    session,
                    str(meta.uuid),
                    fallback_version=remote_version,
                    fallback_operation_time=remote_operation_at,
                )
            except Exception as exc:
                meta.last_error = str(exc)
                meta.save(update_fields=['last_error', 'updated_at'])
                results.append({
                    'classroom_id': classroom.pk,
                    'name': classroom.name,
                    **_serialize_classroom_sync_meta(classroom, meta),
                    'status': 'error',
                    'message': str(exc),
                })
            else:
                results.append({
                    'classroom_id': classroom.pk,
                    'name': classroom.name,
                    **_serialize_classroom_sync_meta(classroom, meta),
                    'status': 'pulled',
                    'version': int(meta.cloud_version or remote_version or 0),
                    'message': '检测到云端操作时间更新，已从云端拉取',
                })
            continue

        if (
            remote_version is not None
            and remote_version > int(meta.cloud_version or 0)
            and not force
            and not upload_by_operation_time
        ):
            if not local_has_changes:
                try:
                    classroom, meta = _cloud_pull_and_restore_classroom(
                        request,
                        session,
                        str(meta.uuid),
                        fallback_version=remote_version,
                        fallback_operation_time=remote_operation_at,
                    )
                except Exception as exc:
                    meta.last_error = str(exc)
                    meta.save(update_fields=['last_error', 'updated_at'])
                    results.append({
                        'classroom_id': classroom.pk,
                        'name': classroom.name,
                        **_serialize_classroom_sync_meta(classroom, meta),
                        'status': 'error',
                        'message': str(exc),
                    })
                else:
                    results.append({
                        'classroom_id': classroom.pk,
                        'name': classroom.name,
                        **_serialize_classroom_sync_meta(classroom, meta),
                        'status': 'pulled',
                        'version': int(meta.cloud_version or remote_version or 0),
                        'message': '已从云端更新',
                    })
                continue

            meta.last_error = 'conflict'
            meta.save(update_fields=['last_error', 'updated_at'])
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **_serialize_classroom_sync_meta(classroom, meta),
                'status': 'conflict',
                'cloud_version': remote_version,
                'message': '云端版本更新，请先从云恢复或处理冲突',
            })
            continue

        if (
            not force
            and remote_version is not None
            and sync_payload['backed_up']
            and sync_payload['local_version'] <= sync_payload['cloud_version']
            and remote_version <= sync_payload['cloud_version']
            and not (remote_operation_at and _operation_time_newer(local_operation_at, remote_operation_at))
        ):
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **sync_payload,
                'status': 'up_to_date',
                'message': '云端已是最新',
            })
            continue

        push_force = bool(
            force
            or (
                upload_by_operation_time
                and remote_version is not None
                and remote_version > int(meta.cloud_version or 0)
            )
        )
        body = {
            'uuid': str(meta.uuid),
            'base_version': int(meta.cloud_version or 0),
            'cloud_version': int(meta.cloud_version or 0),
            'local_version': int(meta.local_version or 0),
            'force': push_force,
            'device_id': data.get('device_id') or 'local-desktop',
            'last_operation_at': local_operation_at.isoformat() if local_operation_at else None,
            'data': _cloud_export_payload(classroom, session),
        }

        try:
            payload = cloud_api_request(session, 'POST', '/api/sync/push', body)
        except CloudAPIError as exc:
            payload = exc.payload or {}
            if payload.get('conflict'):
                meta.last_error = 'conflict'
                meta.save(update_fields=['last_error', 'updated_at'])
                results.append({
                    'classroom_id': classroom.pk,
                    'name': classroom.name,
                    **_serialize_classroom_sync_meta(classroom, meta),
                    'status': 'conflict',
                    'cloud_version': payload.get('version'),
                    'message': payload.get('message') or '云端版本更新，请先拉取或让用户选择保留版本',
                })
                continue
            meta.last_error = payload.get('message') or str(exc)
            meta.save(update_fields=['last_error', 'updated_at'])
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **_serialize_classroom_sync_meta(classroom, meta),
                'status': 'error',
                'message': meta.last_error,
            })
            continue

        if payload.get('ok') or payload.get('status') == 'success':
            version = int(payload.get('version') or meta.cloud_version or 0)
            operation_at = (
                _parse_cloud_operation_time(payload.get('last_operation_at') or payload.get('last_modified_at'))
                or local_operation_at
            )
            meta.cloud_version = version
            meta.local_version = version
            if operation_at:
                meta.last_operation_at = operation_at
            meta.last_sync_at = timezone.now()
            meta.last_error = ''
            meta.save(update_fields=['cloud_version', 'local_version', 'last_operation_at', 'last_sync_at', 'last_error', 'updated_at'])
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **_serialize_classroom_sync_meta(classroom, meta),
                'status': 'ok',
                'version': version,
            })
        else:
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **_serialize_classroom_sync_meta(classroom, meta),
                'status': 'error',
                'message': payload.get('message') or payload.get('error') or '云端未接受同步',
            })

    if sync_all_classrooms:
        local_uuids = {
            str(value)
            for value in SyncMeta.objects.values_list('uuid', flat=True)
            if value
        }
        for remote_row in remote_classrooms:
            if not isinstance(remote_row, dict):
                continue
            remote_uuid = str(remote_row.get('uuid') or '').strip()
            if not remote_uuid or remote_uuid in local_uuids:
                continue
            remote_name = str(remote_row.get('name') or '云端班级')
            remote_version = _cloud_remote_version_for(remote_versions, remote_uuid)
            remote_operation_at = _parse_cloud_operation_time(
                remote_row.get('last_operation_at') or remote_row.get('last_modified_at')
            )
            try:
                classroom, meta = _cloud_pull_and_restore_classroom(
                    request,
                    session,
                    remote_uuid,
                    fallback_version=remote_version,
                    fallback_operation_time=remote_operation_at,
                )
            except Exception as exc:
                results.append({
                    'uuid': remote_uuid,
                    'name': remote_name,
                    'status': 'error',
                    'message': str(exc),
                })
                continue

            local_uuids.add(remote_uuid)
            results.append({
                'classroom_id': classroom.pk,
                'name': classroom.name,
                **_serialize_classroom_sync_meta(classroom, meta),
                'status': 'pulled',
                'remote_only': True,
                'version': int(meta.cloud_version or remote_version or 0),
                'message': '已从云端恢复到本地',
            })

    return results


@csrf_exempt
@require_POST
def cloud_sync(request):
    try:
        session = _cloud_local_session_or_401()
        try:
            data = _cloud_json_body(request)
        except ValueError:
            data = {}
        results = _cloud_sync_classrooms(request, session, data)
        return JsonResponse({'status': 'success', 'results': results})
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_POST
def cloud_sync_pull(request, classroom_uuid):
    try:
        session = _cloud_local_session_or_401()
        payload = cloud_api_request(session, 'GET', f'/api/sync/pull/{classroom_uuid}')
        data = payload.get('data') or payload.get('data_snapshot')
        classroom, meta = _cloud_restore_classroom_data(
            request,
            classroom_uuid,
            data,
            version=payload.get('version'),
            operation_time=payload.get('last_operation_at') or payload.get('last_modified_at'),
        )
        return JsonResponse({
            'status': 'success',
            'classroom_id': classroom.pk,
            'name': classroom.name,
            'version': meta.cloud_version,
            **_serialize_classroom_sync_meta(classroom, meta),
        })
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_http_methods(['POST', 'DELETE'])
def cloud_sync_delete(request, classroom_uuid):
    try:
        session = _cloud_local_session_or_401()
        payload = cloud_api_request(session, 'DELETE', f'/api/sync/{classroom_uuid}')
        return JsonResponse({'status': 'success', **payload})
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_http_methods(['GET', 'POST'])
def cloud_snapshots(request):
    try:
        session = _cloud_local_session_or_401()
        if request.method == 'GET':
            classroom_uuid = str(request.GET.get('classroom_uuid') or '').strip()
            if not classroom_uuid:
                raise ValueError('缺少 classroom_uuid')
            payload = cloud_api_request(session, 'GET', f'/api/snapshots/{classroom_uuid}')
            return JsonResponse({'status': 'success', **payload})

        data = _cloud_json_body(request)
        classroom_uuid = str(data.get('classroom_uuid') or '').strip()
        snapshot_data = data.get('data') if isinstance(data.get('data'), dict) else None
        classroom_id = data.get('classroom_id')
        if classroom_id:
            classroom = get_object_or_404(Classroom, pk=classroom_id)
            meta, _ = SyncMeta.objects.get_or_create(classroom=classroom)
            classroom_uuid = str(meta.uuid)
            snapshot_data = _cloud_export_payload(classroom, session)
        if not classroom_uuid or not snapshot_data:
            raise ValueError('缺少快照班级或快照数据')

        payload = cloud_api_request(session, 'POST', '/api/snapshots', {
            'classroom_uuid': classroom_uuid,
            'name': data.get('name') or '手动快照',
            'data': snapshot_data,
        })
        return JsonResponse({'status': 'success', **payload})
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_POST
def cloud_snapshot_restore(request, snapshot_id):
    try:
        session = _cloud_local_session_or_401()
        payload = cloud_api_request(session, 'GET', f'/api/snapshots/{snapshot_id}/download')
        classroom_uuid = payload.get('classroom_uuid')
        data = payload.get('data')
        if not classroom_uuid:
            raise ValueError('云端快照缺少 classroom_uuid')
        classroom, meta = _cloud_restore_classroom_data(request, classroom_uuid, data)
        return JsonResponse({
            'status': 'success',
            'classroom_id': classroom.pk,
            'uuid': str(meta.uuid),
            'name': classroom.name,
        })
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_http_methods(['POST', 'DELETE'])
def cloud_snapshot_delete(request, snapshot_id):
    try:
        session = _cloud_local_session_or_401()
        payload = cloud_api_request(session, 'DELETE', f'/api/snapshots/{snapshot_id}')
        return JsonResponse({'status': 'success', **payload})
    except Exception as exc:
        return _cloud_error_response(exc)


@require_http_methods(['GET'])
def cloud_subscription_plans(request):
    current_tier = None
    active_session = get_active_cloud_session()
    if active_session:
        _refresh_cloud_subscription_if_logged_in(active_session)
        current_tier = active_session.subscription_tier
    try:
        payload = cloud_public_request('GET', '/api/subscription/plans')
        plans = payload.get('plans', [])
        for p in plans:
            p.setdefault('key', p.get('tier', ''))
        return JsonResponse({'status': 'success', 'plans': plans, 'current_tier': current_tier})
    except Exception as exc:
        return _cloud_error_response(exc)


@csrf_exempt
@require_POST
def cloud_subscription_redeem(request):
    try:
        session = _cloud_local_session_or_401()
        data = _cloud_json_body(request)
        payload = cloud_api_request(session, 'POST', '/api/subscription/redeem', {
            'code': data.get('code'),
        })
        if payload.get('uid') or payload.get('subscription'):
            session_payload = {
                'uid': payload.get('uid') or session.uid,
                'nickname': payload.get('nickname') or session.nickname,
                'email': payload.get('email') or session.email,
                'avatar_url': payload.get('avatar_url') or session.avatar_url,
                'session_token': session.session_token,
                'client_key_id': session.client_key_id,
                'client_public_key': session.client_public_key_pem,
                'client_private_key': session.client_private_key_pem,
                'server_key_id': session.server_key_id,
                'server_public_key': session.server_public_key_pem,
                'token_expires_at': session.token_expires_at.isoformat(),
                'subscription': payload.get('subscription') or {},
            }
            save_cloud_session_from_payload(session_payload)
        elif isinstance(payload, dict):
            apply_cloud_subscription_payload(session, payload)
        return JsonResponse({'status': 'success', **payload})
    except Exception as exc:
        return _cloud_error_response(exc)


@require_http_methods(['GET'])
def cloud_purchase_url(request):
    try:
        tier_key = request.GET.get('tier') or request.GET.get('plan') or ''
        payload = cloud_public_request('GET', f'/api/subscription/purchase-url?tier={tier_key}')
        return JsonResponse({'status': 'success', 'ok': payload.get('ok', False), 'tier': tier_key, 'url': payload.get('url', '')})
    except Exception as exc:
        return _cloud_error_response(exc)


def frontend_store_js(request):
    kvs = FrontendKVStore.objects.all()
    store_dict = {kv.key: kv.value for kv in kvs}
    js_code = f"""
    window.BACKEND_STORE = {json.dumps(store_dict)};
    try {{
        const oldSetItem = Storage.prototype.setItem;
        const oldGetItem = Storage.prototype.getItem;
        const oldRemoveItem = Storage.prototype.removeItem;
        const hasBackendValue = (key) => Object.prototype.hasOwnProperty.call(window.BACKEND_STORE, key);
        const normalizeStoreKey = (key) => {{
            if (key === undefined || key === null) return '';
            return String(key);
        }};
        const syncBackendSet = (rawKey, rawValue) => {{
            const key = normalizeStoreKey(rawKey);
            if (!key || key.length > 255) return;
            const value = String(rawValue);
            if (hasBackendValue(key) && window.BACKEND_STORE[key] === value) return;
            window.BACKEND_STORE[key] = value;
            fetch('/api/store/set/', {{
                method: 'POST',
                headers: {{'Content-Type': 'application/json'}},
                body: JSON.stringify({{key: key, value: value}})
            }}).catch(() => {{}});
        }};
        const syncBackendDelete = (rawKey) => {{
            const key = normalizeStoreKey(rawKey);
            if (!key || key.length > 255 || !hasBackendValue(key)) return;
            delete window.BACKEND_STORE[key];
            fetch('/api/store/delete/', {{
                method: 'POST',
                headers: {{'Content-Type': 'application/json'}},
                body: JSON.stringify({{key: key}})
            }}).catch(() => {{}});
        }};

        Storage.prototype.setItem = function(key, val) {{
            try {{ oldSetItem.call(this, key, val); }} catch(e) {{}}
            if (this === window.localStorage) {{
                syncBackendSet(key, val);
            }}
        }};
        
        Storage.prototype.getItem = function(key) {{
            const storeKey = normalizeStoreKey(key);
            if (this === window.localStorage && hasBackendValue(storeKey)) {{
                return window.BACKEND_STORE[storeKey];
            }}
            try {{ return oldGetItem.call(this, key); }} catch(e) {{ return null; }}
        }};
        
        Storage.prototype.removeItem = function(key) {{
            try {{ oldRemoveItem.call(this, key); }} catch(e) {{}}
            if (this === window.localStorage) {{
                syncBackendDelete(key);
            }}
        }};
    }} catch(e) {{
        console.error("Storage prototype override failed:", e);
    }}

    window.addEventListener('DOMContentLoaded', () => {{
        if (navigator.userAgent.includes('Mac OS X') || navigator.platform.toUpperCase().indexOf('MAC') >= 0) {{
            if (!localStorage.getItem('mac_os_warning_seen')) {{
                localStorage.setItem('mac_os_warning_seen', 'true');
                const message = "当前系统可能部分动效、功能不稳定，若您有相关需要，可以前往 Windows/使用 Chrome 打开http://127.0.0.1:23948";
                if (typeof window.showToast === 'function') {{
                    window.showToast(message);
                }} else {{
                    console.warn(message);
                }}
            }}
        }}
    }});
    """
    return HttpResponse(js_code, content_type="application/javascript")

@csrf_exempt
@require_POST
def frontend_store_set(request):
    try:
        data = json.loads(request.body)
        raw_key = data.get('key', '')
        raw_value = data.get('value', '')
        key = '' if raw_key is None else str(raw_key)
        value = '' if raw_value is None else str(raw_value)
        if key and len(key) <= 255:
            FrontendKVStore.objects.update_or_create(key=key, defaults={'value': value})
        return JsonResponse({'status': 'success', 'persisted': bool(key and len(key) <= 255)})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

@csrf_exempt
@require_POST
def frontend_store_delete(request):
    try:
        data = json.loads(request.body)
        raw_key = data.get('key', '')
        key = '' if raw_key is None else str(raw_key)
        if key and len(key) <= 255:
            FrontendKVStore.objects.filter(key=key).delete()
        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@require_http_methods(['GET'])
def desktop_update_check(request):
    try:
        payload = desktop_runtime.check_for_updates()
        return JsonResponse({'status': 'success', **payload})
    except Exception as exc:
        return JsonResponse({
            'status': 'error',
            'state': 'error',
            'platform': desktop_runtime.get_platform_name(),
            'supported': desktop_runtime.is_update_api_supported(),
            'message': str(exc),
        }, status=502)


@csrf_exempt
@require_POST
def desktop_update_start(request):
    try:
        payload = json.loads(request.body or b'{}')
    except Exception:
        return JsonResponse({
            'status': 'error',
            'state': 'invalid_json',
            'message': '请求数据格式错误',
        }, status=400)

    try:
        result = desktop_runtime.start_manual_update(payload.get('target_version'))
        return JsonResponse({'status': 'success', **result})
    except ValueError as exc:
        return JsonResponse({
            'status': 'error',
            'state': 'version_mismatch',
            'message': str(exc),
        }, status=409)
    except RuntimeError as exc:
        return JsonResponse({
            'status': 'error',
            'state': 'unsupported',
            'platform': desktop_runtime.get_platform_name(),
            'supported': desktop_runtime.is_update_api_supported(),
            'message': str(exc),
        }, status=400)
    except Exception as exc:
        return JsonResponse({
            'status': 'error',
            'state': 'error',
            'platform': desktop_runtime.get_platform_name(),
            'supported': desktop_runtime.is_update_api_supported(),
            'message': str(exc),
        }, status=500)


@csrf_exempt
@require_POST
def desktop_update_install(request):
    try:
        result = desktop_runtime.launch_prepared_update()
        return JsonResponse({'status': 'success', **result})
    except RuntimeError as exc:
        return JsonResponse({
            'status': 'error',
            'state': 'not_ready',
            'platform': desktop_runtime.get_platform_name(),
            'supported': desktop_runtime.is_update_api_supported(),
            'message': str(exc),
        }, status=409)
    except Exception as exc:
        return JsonResponse({
            'status': 'error',
            'state': 'error',
            'platform': desktop_runtime.get_platform_name(),
            'supported': desktop_runtime.is_update_api_supported(),
            'message': str(exc),
        }, status=500)


@require_http_methods(['GET'])
def desktop_update_status(request):
    return JsonResponse({'status': 'success', **desktop_runtime.get_update_status()})


@require_POST
def mark_onboarding_seen(request):
    sk = request.session.session_key
    if not sk:
        request.session['ob_init'] = True
        request.session.save()
        sk = request.session.session_key
    completed = ''
    body = {}
    try:
        body = json.loads(request.body or b'{}')
        if not isinstance(body, dict):
            body = {}
        if isinstance(body, dict) and isinstance(body.get('completed_steps'), str):
            completed = body['completed_steps'][:120]
    except Exception:
        completed = ''
        body = {}
    if sk:
        OnboardingState.objects.update_or_create(
            session_key=sk,
            defaults={'seen': True, 'completed_steps': completed},
        )
    FrontendKVStore.objects.update_or_create(
        key=ONBOARDING_SEEN_STORE_KEY,
        defaults={'value': ONBOARDING_SEEN_STORE_VALUE},
    )
    stage = str((body or {}).get('completed_steps') or '').strip()
    if stage in {'detail_done', 'tour_done'}:
        request.session[_CLEANUP_PENDING_KEY] = True
        request.session.modified = True
    payload = {'ok': True, 'sample_deleted': False}
    return JsonResponse(payload)


@csrf_exempt
@require_POST
def frontend_improve_events(request):
    if not get_data_sharing_enabled():
        return JsonResponse({'ok': False, 'status': 'disabled'}, status=503)
    try:
        data = json.loads(request.body or b'{}')
    except Exception:
        return JsonResponse({'ok': False, 'status': 'error', 'message': '请求格式错误'}, status=400)
    events = data.get('events') if isinstance(data.get('events'), list) else []
    accepted = 0
    for item in events[:200]:
        if not isinstance(item, dict):
            continue
        share_usage_event(
            item.get('feature', 'unknown'),
            item.get('action', 'use'),
            success=item.get('success', True),
            duration_ms=int(item.get('duration_ms') or 0),
            count=int(item.get('count') or 1),
            metadata=item.get('metadata') if isinstance(item.get('metadata'), dict) else None,
        )
        accepted += 1
    return JsonResponse({'ok': True, 'status': 'success', 'accepted': accepted})


@csrf_exempt
@require_POST
def frontend_improve_logs(request):
    if not get_data_sharing_enabled():
        return JsonResponse({'ok': False, 'status': 'disabled'}, status=503)
    try:
        data = json.loads(request.body or b'{}')
    except Exception:
        return JsonResponse({'ok': False, 'status': 'error', 'message': '请求格式错误'}, status=400)
    logs = data.get('logs') if isinstance(data.get('logs'), list) else []
    accepted = 0
    for item in logs[:200]:
        if not isinstance(item, dict):
            continue
        share_log(
            item.get('level', 'INFO'),
            item.get('source', 'frontend'),
            item.get('code', ''),
            message=str(item.get('message') or '')[:240],
            context=item.get('context') if isinstance(item.get('context'), dict) else None,
        )
        accepted += 1
    return JsonResponse({'ok': True, 'status': 'success', 'accepted': accepted})


@csrf_exempt
@require_http_methods(['GET'])
def ai_session_status(request):
    from seats.open_api import ai_session, realtime
    payload = ai_session.status()
    return JsonResponse({
        'status': 'success',
        'session': payload,
        'realtime': realtime.snapshot(),
    })


@csrf_exempt
@require_POST
def ai_session_end(request):
    from seats.open_api import ai_session
    ai_session.end()
    return JsonResponse({'status': 'success', 'session': ai_session.status()})


@require_http_methods(['GET'])
def ai_session_stream(request):
    from seats.open_api import ai_session, realtime

    def event_stream():
        last_global_seq = None
        while True:
            rt = realtime.wait_for_change(last_global_seq, timeout=25.0)
            last_global_seq = rt['global_seq']
            payload = json.dumps({
                'status': 'success',
                'session': ai_session.status(),
                'realtime': rt,
            }, ensure_ascii=False)
            yield 'event: session\ndata: ' + payload + '\n\n'
            yield ': ping\n\n'

    response = StreamingHttpResponse(event_stream(), content_type='text/event-stream; charset=utf-8')
    response['Cache-Control'] = 'no-cache'
    response['X-Accel-Buffering'] = 'no'
    return response
